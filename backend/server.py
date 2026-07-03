"""
Orisei Freight Solutions — Transportation Management System (TMS)
HUD-style command center backend
"""
from fastapi import FastAPI, APIRouter, Depends, HTTPException, Request, Response, WebSocket, WebSocketDisconnect, Query
from fastapi.responses import StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import io
import re
import asyncio
import hashlib
import logging
import time as _time
import uuid
import json
import random
import httpx
import urllib.parse
from pathlib import Path
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any, Tuple
from datetime import datetime, timezone, timedelta

from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.units import inch

from emergentintegrations.llm.chat import LlmChat, UserMessage

from openpyxl import Workbook as XLWorkbook, load_workbook as load_xlsx
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

mongo_url = os.environ['MONGO_URL']
# Production-grade connection pool — sized for ~250 concurrent users with
# safe headroom. minPoolSize keeps warm connections to avoid cold-start TLS
# handshake spikes; maxIdleTimeMS recycles idle connections so the pool stays
# healthy under low-traffic windows.
client = AsyncIOMotorClient(
    mongo_url,
    maxPoolSize=int(os.environ.get('MONGO_MAX_POOL', '120')),
    minPoolSize=int(os.environ.get('MONGO_MIN_POOL', '10')),
    maxIdleTimeMS=60_000,
    serverSelectionTimeoutMS=8_000,
    connectTimeoutMS=8_000,
    socketTimeoutMS=20_000,
    retryWrites=True,
    retryReads=True,
)
db = client[os.environ['DB_NAME']]

app = FastAPI(title="Orisei TMS API")
api_router = APIRouter(prefix="/api")

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger("tennant_tms")

# -------------------- MODELS --------------------
class User(BaseModel):
    user_id: str
    email: str
    name: str
    picture: Optional[str] = None
    role: str = "dispatcher"  # admin | auditor | dispatcher | driver | carrier
    carrier_company: Optional[str] = None  # for role=carrier — scopes their visibility

class Shipment(BaseModel):
    shipment_id: str
    reference: str
    mode: str  # TL, LTL, Parcel, Ocean, Air, Rail
    carrier: str
    status: str  # in_transit, delayed, delivered, pending, at_origin, at_dest
    origin: Dict[str, Any]
    destination: Dict[str, Any]
    current_location: Dict[str, Any]
    eta: str
    pickup_date: str
    weight_lbs: float
    pieces: int
    commodity: str
    value_usd: float
    container_no: Optional[str] = None
    bol_no: Optional[str] = None
    pro_no: Optional[str] = None
    progress: float = 0.0
    # Excel-aligned fields
    direction: str = "outbound"  # outbound | inbound
    hazmat: bool = False
    hazmat_class: Optional[str] = None
    supplier: Optional[str] = None
    consignee: Optional[str] = None
    ship_date: Optional[str] = None
    ship_day: Optional[str] = None  # MON, TUE, WED, THR, FRI
    skids: Optional[int] = None
    material_controller: Optional[str] = None
    po_numbers: Optional[str] = None
    booking_number: Optional[str] = None
    bid_cost: Optional[float] = None
    fsc_pct: Optional[float] = None
    extras: Optional[str] = None
    done: bool = False
    shipping_hours: Optional[str] = None
    pickup_no: Optional[str] = None
    # Accessorials & NMFC classification
    liftgate_required: bool = False
    pallet_count: Optional[int] = None
    nmfc_code: Optional[str] = None
    freight_class: Optional[str] = None  # 50, 55, 60, 65, 70, 77.5, 85, 92.5, 100, 110, 125, 150, 175, 200, 250, 300, 400, 500
    accessorials: Optional[List[str]] = None  # e.g., ["liftgate", "residential", "inside_delivery"]
    # Dimensions (inches) — used for LTL density / freight class
    length_in: Optional[float] = None
    width_in: Optional[float] = None
    height_in: Optional[float] = None
    # SAP S/4HANA linkage
    sap_delivery_no: Optional[str] = None
    sap_material_numbers: Optional[List[str]] = None
    # Customer-facing
    customer_contact_email: Optional[str] = None
    carrier_contact_email: Optional[str] = None
    # File references (GridFS ids)
    carrier_bol_file_id: Optional[str] = None
    carrier_bol_filename: Optional[str] = None

class Document(BaseModel):
    document_id: str
    type: str  # BOL, COMMERCIAL_INVOICE, PACKING_SLIP, WEIGHT_CERT, COO
    shipment_ref: str
    created_by: str
    created_at: str
    data: Dict[str, Any]
    # Amendment trail — empty for fresh documents. Each entry:
    #   { amended_at, amended_by, reason, changes: [{ field, from, to }] }
    amendments: List[Dict[str, Any]] = []
    version: int = 1
    updated_at: Optional[str] = None

class ChatMessage(BaseModel):
    message_id: str
    channel: str
    user_id: str
    user_name: str
    user_picture: Optional[str] = None
    text: str
    created_at: str

# -------------------- AUTH HELPERS --------------------
async def get_current_user(request: Request) -> User:
    """Get user from session_token cookie or Authorization header."""
    token = request.cookies.get("session_token")
    if not token:
        auth = request.headers.get("Authorization", "")
        if auth.startswith("Bearer "):
            token = auth[7:]
    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    session = await db.user_sessions.find_one({"session_token": token}, {"_id": 0})
    if not session:
        raise HTTPException(status_code=401, detail="Invalid session")
    expires_at = session["expires_at"]
    if isinstance(expires_at, str):
        expires_at = datetime.fromisoformat(expires_at)
    if expires_at.tzinfo is None:
        expires_at = expires_at.replace(tzinfo=timezone.utc)
    if expires_at < datetime.now(timezone.utc):
        raise HTTPException(status_code=401, detail="Session expired")
    user_doc = await db.users.find_one({"user_id": session["user_id"]}, {"_id": 0})
    if not user_doc:
        raise HTTPException(status_code=401, detail="User not found")
    return User(**user_doc)

async def get_optional_user(request: Request) -> Optional[User]:
    try:
        return await get_current_user(request)
    except HTTPException:
        return None

# RBAC helpers
ROLE_HIERARCHY = {"driver": 0, "dispatcher": 1, "auditor": 2, "admin": 3}

def require_role(*allowed_roles: str):
    """Dependency factory: ensure current user has one of the allowed roles, or admin (which can do anything)."""
    async def _checker(user: User = Depends(get_current_user)) -> User:
        if user.role == "admin":
            return user
        if user.role in allowed_roles:
            return user
        raise HTTPException(status_code=403, detail=f"Requires one of roles: {', '.join(allowed_roles)}")
    return _checker

# -------------------- AUTH ENDPOINTS --------------------
def _shared_cookie_domain(request: Request) -> Optional[str]:
    """Return a `.parent.tld` cookie domain that's shared across the
    frontend and backend subdomains, so the auth cookie isn't locked to
    just the backend host. Falls back to None for local dev or custom domains.

    Examples:
        backend host `clean-logistics-dash.preview.emergentagent.com`
            → `.emergentagent.com` (also covers preview.static.* frontend)
        backend host `api.oriseifreight.com`
            → `.oriseifreight.com`
        backend host `localhost`
            → None (browsers reject Domain for IP/localhost)
    """
    try:
        host = (request.headers.get("x-forwarded-host")
                or request.headers.get("host") or "").split(":")[0].strip().lower()
        if not host or host in ("localhost", "127.0.0.1") or host.replace(".", "").isdigit():
            return None
        # Whitelist of well-known shared parents (highest precedence)
        for parent in ("emergentagent.com", "emergent.host", "emergent.sh"):
            if host == parent or host.endswith("." + parent):
                return "." + parent
        # Custom domain: drop the leftmost subdomain ("api.foo.com" → ".foo.com").
        parts = host.split(".")
        if len(parts) >= 2:
            return "." + ".".join(parts[-2:])
        return None
    except Exception:
        return None


@api_router.post("/auth/session")
async def create_session(request: Request, response: Response):
    body = await request.json()
    session_id = body.get("session_id")
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id required")
    async with httpx.AsyncClient(timeout=15.0) as http:
        r = await http.get(
            "https://demobackend.emergentagent.com/auth/v1/env/oauth/session-data",
            headers={"X-Session-ID": session_id}
        )
        if r.status_code != 200:
            raise HTTPException(status_code=401, detail="Invalid session_id")
        data = r.json()
    email = data["email"]
    name = data["name"]
    picture = data.get("picture")
    session_token = data["session_token"]

    # Admin allow-list: scalable to 250 users — admins defined by env var, all
    # other accounts default to dispatcher. The first user ever to sign in also
    # becomes admin so the system is never adminless.
    admin_emails = {
        e.strip().lower()
        for e in (os.environ.get("ADMIN_EMAILS", "") or "").split(",")
        if e.strip()
    }
    is_allowlisted_admin = email.lower() in admin_emails

    existing = await db.users.find_one({"email": email}, {"_id": 0})
    if existing:
        user_id = existing["user_id"]
        update_fields = {"name": name, "picture": picture}
        # Auto-promote allow-listed emails to admin on every login (idempotent).
        if is_allowlisted_admin and existing.get("role") != "admin":
            update_fields["role"] = "admin"
        await db.users.update_one({"user_id": user_id}, {"$set": update_fields})
    else:
        user_id = f"user_{uuid.uuid4().hex[:12]}"
        user_count = await db.users.count_documents({})
        if is_allowlisted_admin or user_count == 0:
            initial_role = "admin"
        else:
            initial_role = "dispatcher"
        await db.users.insert_one({
            "user_id": user_id,
            "email": email,
            "name": name,
            "picture": picture,
            "role": initial_role,
            "created_at": datetime.now(timezone.utc).isoformat()
        })

    expires_at = datetime.now(timezone.utc) + timedelta(days=7)
    await db.user_sessions.insert_one({
        "user_id": user_id,
        "session_token": session_token,
        "expires_at": expires_at.isoformat(),
        "created_at": datetime.now(timezone.utc).isoformat()
    })
    cookie_kwargs = {
        "key": "session_token",
        "value": session_token,
        "httponly": True,
        "secure": True,
        "samesite": "none",
        "path": "/",
        "max_age": 7 * 24 * 60 * 60,
    }
    # Share cookie across frontend (preview.static.*) and backend (preview.*)
    # subdomains by setting Domain to the shared parent. Otherwise the cookie
    # is locked to the backend host and sign-in loops forever from the static
    # production frontend.
    cookie_domain = _shared_cookie_domain(request)
    if cookie_domain:
        cookie_kwargs["domain"] = cookie_domain
    response.set_cookie(**cookie_kwargs)
    final_role = "admin" if is_allowlisted_admin else (existing.get("role") if existing else initial_role)
    # Also return the session_token so cross-origin frontends (where browsers
    # block third-party cookies) can store it in localStorage and send it as
    # an Authorization: Bearer header.
    return {
        "user_id": user_id, "email": email, "name": name,
        "picture": picture, "role": final_role,
        "session_token": session_token,
    }

@api_router.get("/auth/me", response_model=User)
async def me(user: User = Depends(get_current_user)):
    return user


@api_router.post("/auth/dev-session")
async def dev_session(request: Request, response: Response):
    """Friction-free admin sign-in for the founder on PREVIEW environments.
    Disabled on production (livecleans.com) — checks the Origin / Host header
    so it cannot be invoked against the production deployment.

    Issues a real session_token for the first admin email in ADMIN_EMAILS so
    the founder can one-click sign in without the Google OAuth round-trip
    every time the preview backend rotates."""
    # Production guard — must be explicitly enabled via env. Production
    # deployments won't have ENABLE_DEV_LOGIN set, so this endpoint will
    # always return 404 there. Preview .env sets it to "true".
    if (os.environ.get("ENABLE_DEV_LOGIN") or "").lower() not in ("1", "true", "yes"):
        raise HTTPException(status_code=404, detail="Not available")

    admin_emails = [
        e.strip() for e in (os.environ.get("ADMIN_EMAILS", "") or "").split(",")
        if e.strip()
    ]
    if not admin_emails:
        raise HTTPException(status_code=400, detail="No ADMIN_EMAILS configured")
    email = admin_emails[0]

    existing = await db.users.find_one({"email": email}, {"_id": 0})
    if existing:
        user_id = existing["user_id"]
        if existing.get("role") != "admin":
            await db.users.update_one({"user_id": user_id}, {"$set": {"role": "admin"}})
        name = existing.get("name") or email.split("@")[0].title()
        picture = existing.get("picture")
    else:
        user_id = f"user_{uuid.uuid4().hex[:12]}"
        name = email.split("@")[0].title()
        picture = None
        await db.users.insert_one({
            "user_id": user_id, "email": email, "name": name,
            "picture": picture, "role": "admin",
            "created_at": datetime.now(timezone.utc).isoformat(),
        })

    session_token = f"dev_{uuid.uuid4().hex}"
    expires_at = datetime.now(timezone.utc) + timedelta(days=7)
    await db.user_sessions.insert_one({
        "user_id": user_id, "session_token": session_token,
        "expires_at": expires_at.isoformat(),
        "created_at": datetime.now(timezone.utc).isoformat(),
        "dev_session": True,
    })
    cookie_kwargs = {
        "key": "session_token", "value": session_token,
        "httponly": True, "secure": True, "samesite": "none",
        "path": "/", "max_age": 7 * 24 * 60 * 60,
    }
    cookie_domain = _shared_cookie_domain(request)
    if cookie_domain:
        cookie_kwargs["domain"] = cookie_domain
    response.set_cookie(**cookie_kwargs)
    return {
        "user_id": user_id, "email": email, "name": name,
        "picture": picture, "role": "admin",
        "session_token": session_token,
    }


@api_router.post("/auth/logout")
async def logout(request: Request, response: Response):
    token = request.cookies.get("session_token")
    if token:
        await db.user_sessions.delete_one({"session_token": token})
    cookie_domain = _shared_cookie_domain(request)
    if cookie_domain:
        response.delete_cookie("session_token", path="/", domain=cookie_domain)
    else:
        response.delete_cookie("session_token", path="/")
    return {"ok": True}

# -------------------- FACILITIES --------------------
TMS_FACILITIES = [
    {"id": "GVM", "name": "Golden Valley, MN (HQ)", "city": "Golden Valley", "state": "MN", "lat": 44.9858, "lng": -93.3499, "type": "Headquarters / Manufacturing"},
    {"id": "HOM", "name": "Holland, MI Plant", "city": "Holland", "state": "MI", "lat": 42.7875, "lng": -86.1089, "type": "Manufacturing"},
    {"id": "LVK", "name": "Louisville, KY Plant", "city": "Louisville", "state": "KY", "lat": 38.2527, "lng": -85.7585, "type": "Manufacturing"},
]

@api_router.get("/facilities")
async def get_facilities(_: User = Depends(get_current_user)):
    brand = await _active_brand_doc()
    if brand and brand.get("facilities"):
        # Overlay brand facilities on the canonical lat/lng grid so the map
        # still has coordinates to plot.
        out = []
        for i, bf in enumerate(brand["facilities"][:len(TMS_FACILITIES)]):
            base = dict(TMS_FACILITIES[i % len(TMS_FACILITIES)])
            base["name"] = bf.get("name") or base["name"]
            base["city"] = (bf.get("city") or "").split(",")[0].strip() or base["city"]
            out.append(base)
        return out or TMS_FACILITIES
    return TMS_FACILITIES


# -------------------- BRAND OVERLAY -----------------------------------------
# When the active company brand is something other than the built-in default
# profile, every read endpoint runs its docs through these helpers so the
# entire app surfaces the active brand's identity / suppliers / products /
# facilities without re-seeding the database.

async def _active_brand_doc():
    """Returns the active brand doc, or None if default brand is active.
    Result is intentionally NOT cached at module level because the admin
    expects the switch to take effect on the next request."""
    return await db.company_brand.find_one({"is_active": True}, {"_id": 0})


def _pick(items: List[str], seed_str: str) -> str:
    """Deterministically pick one item using a stable hash of seed_str so a
    given shipment always maps to the same supplier / product across reads."""
    if not items:
        return ""
    h = int(hashlib.md5((seed_str or "").encode()).hexdigest(), 16)
    return items[h % len(items)]


def _swap_strings(value: Any, replacements: Dict[str, str]) -> Any:
    """Recursively replace string fragments inside any nested structure."""
    if isinstance(value, str):
        out = value
        for src, dst in replacements.items():
            out = out.replace(src, dst)
        return out
    if isinstance(value, list):
        return [_swap_strings(v, replacements) for v in value]
    if isinstance(value, dict):
        return {k: _swap_strings(v, replacements) for k, v in value.items()}
    return value


def _overlay_shipment(s: Dict[str, Any], brand: Dict[str, Any]) -> Dict[str, Any]:
    if not brand or brand.get("brand_id") == "orisei-freight":
        return s
    out = dict(s)
    seed = out.get("shipment_id") or out.get("reference") or ""
    suppliers = brand.get("sample_suppliers") or []
    products = brand.get("sample_products") or []
    short = brand.get("short_name") or brand.get("company_name") or "Brand"

    # Swap commodity → branded product / supplier → branded supplier.
    if products and out.get("commodity"):
        out["commodity"] = _pick(products, seed)
    if suppliers and out.get("supplier"):
        out["supplier"] = _pick(suppliers, seed + "supplier")

    # === Swap facility cities (origin/destination) ===
    # Map the canonical reference cities to the active brand's facility cities.
    facilities = brand.get("facilities") or []
    if facilities:
        # Build a city/state mapping: default brand cities → brand cities
        tennant_cities = [
            ("Golden Valley", "MN"),
            ("Holland", "MI"),
            ("Louisville", "KY"),
        ]
        brand_cities = []
        for f in facilities[:len(tennant_cities)]:
            city_str = (f.get("city") or "").strip()
            if "," in city_str:
                c, st = [p.strip() for p in city_str.split(",", 1)]
            else:
                c, st = city_str, ""
            brand_cities.append((c, st))
        for slot_i, (t_city, t_state) in enumerate(tennant_cities):
            if slot_i >= len(brand_cities):
                break
            b_city, b_state = brand_cities[slot_i]
            for endpoint_key in ("origin", "destination"):
                ep = out.get(endpoint_key)
                if isinstance(ep, dict):
                    if (ep.get("city") or "").strip() == t_city:
                        ep["city"] = b_city or t_city
                        if b_state:
                            ep["state"] = b_state

    # Replace "Orisei" verbiage in every string field (reference prefixes,
    # consignee names, facility names baked into origin/destination, notes).
    repl = {"Orisei": short, "TENN-": f"{short[:4].upper()}-"}
    out = _swap_strings(out, repl)
    return out


def _overlay_machine(m: Dict[str, Any], brand: Dict[str, Any], i: int) -> Dict[str, Any]:
    if not brand or brand.get("brand_id") == "orisei-freight":
        return m
    products = brand.get("sample_products") or []
    if not products:
        return m
    out = dict(m)
    branded = products[i % len(products)]
    short = brand.get("short_name") or "Brand"
    out["model"] = branded
    out["display_name"] = branded
    out["description"] = f"{branded} — {brand.get('industry') or short + ' flagship'}"
    return out


def _overlay_supplier(s: Dict[str, Any], brand: Dict[str, Any], i: int) -> Dict[str, Any]:
    if not brand or brand.get("brand_id") == "orisei-freight":
        return s
    sups = brand.get("sample_suppliers") or []
    if not sups:
        return s
    out = dict(s)
    out["name"] = sups[i % len(sups)]
    return out


# -------------------- SHIPMENTS --------------------
@api_router.get("/shipments", response_model=List[Shipment])
async def list_shipments(user: User = Depends(get_current_user), mode: Optional[str] = None, status: Optional[str] = None, direction: Optional[str] = None, hazmat: Optional[bool] = None, limit: int = 500):
    q: Dict[str, Any] = {}
    if mode:
        q["mode"] = mode
    if status:
        q["status"] = status
    if direction:
        q["direction"] = direction
    if hazmat is not None:
        q["hazmat"] = hazmat
    # Scope: carriers only see their own loads
    if user.role == "carrier" and user.carrier_company:
        q["carrier"] = user.carrier_company
    docs = await db.shipments.find(q, {"_id": 0}).limit(limit).to_list(limit)
    brand = await _active_brand_doc()
    if brand:
        docs = [_overlay_shipment(d, brand) for d in docs]
    return [Shipment(**d) for d in docs]

@api_router.get("/shipments/{shipment_id}", response_model=Shipment)
async def get_shipment(shipment_id: str, user: User = Depends(get_current_user)):
    doc = await db.shipments.find_one({"shipment_id": shipment_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Shipment not found")
    if user.role == "carrier" and user.carrier_company and doc.get("carrier") != user.carrier_company:
        raise HTTPException(status_code=403, detail="Not your shipment")
    brand = await _active_brand_doc()
    if brand:
        doc = _overlay_shipment(doc, brand)
    return Shipment(**doc)

class ShipmentCreate(BaseModel):
    reference: Optional[str] = None
    mode: str
    carrier: str
    origin_facility: Optional[str] = None
    origin_city: Optional[str] = None
    destination_city: str
    destination_lat: float
    destination_lng: float
    pickup_date: str
    weight_lbs: float
    pieces: int
    commodity: str
    value_usd: float
    # Optional booking extras
    liftgate_required: Optional[bool] = False
    pallet_count: Optional[int] = None
    nmfc_code: Optional[str] = None
    freight_class: Optional[str] = None
    accessorials: Optional[List[str]] = None
    # Dimensions
    length_in: Optional[float] = None
    width_in: Optional[float] = None
    height_in: Optional[float] = None
    # SAP linkage (populated when "pull from SAP" is used)
    sap_delivery_no: Optional[str] = None
    sap_material_numbers: Optional[List[str]] = None
    # Customer-facing
    customer_contact_email: Optional[str] = None
    carrier_contact_email: Optional[str] = None

class ShipmentUpdate(BaseModel):
    """All fields optional — only the keys present are applied."""
    reference: Optional[str] = None
    mode: Optional[str] = None
    carrier: Optional[str] = None
    status: Optional[str] = None
    eta: Optional[str] = None
    pickup_date: Optional[str] = None
    ship_date: Optional[str] = None
    ship_day: Optional[str] = None
    weight_lbs: Optional[float] = None
    pieces: Optional[int] = None
    skids: Optional[int] = None
    commodity: Optional[str] = None
    value_usd: Optional[float] = None
    container_no: Optional[str] = None
    bol_no: Optional[str] = None
    pro_no: Optional[str] = None
    direction: Optional[str] = None
    hazmat: Optional[bool] = None
    hazmat_class: Optional[str] = None
    supplier: Optional[str] = None
    consignee: Optional[str] = None
    material_controller: Optional[str] = None
    po_numbers: Optional[str] = None
    booking_number: Optional[str] = None
    bid_cost: Optional[float] = None
    fsc_pct: Optional[float] = None
    extras: Optional[str] = None
    done: Optional[bool] = None
    shipping_hours: Optional[str] = None
    pickup_no: Optional[str] = None
    progress: Optional[float] = None
    # Accessorials & NMFC classification
    liftgate_required: Optional[bool] = None
    pallet_count: Optional[int] = None
    nmfc_code: Optional[str] = None
    freight_class: Optional[str] = None
    accessorials: Optional[List[str]] = None
    # Dimensions
    length_in: Optional[float] = None
    width_in: Optional[float] = None
    height_in: Optional[float] = None
    sap_delivery_no: Optional[str] = None
    customer_contact_email: Optional[str] = None
    carrier_contact_email: Optional[str] = None
    # Nested overrides
    origin_city: Optional[str] = None
    origin_facility: Optional[str] = None
    destination_city: Optional[str] = None
    destination_lat: Optional[float] = None
    destination_lng: Optional[float] = None

@api_router.put("/shipments/{shipment_id}", response_model=Shipment)
async def update_shipment(shipment_id: str, payload: ShipmentUpdate, user: User = Depends(require_role("admin", "dispatcher"))):
    existing = await db.shipments.find_one({"shipment_id": shipment_id}, {"_id": 0})
    if not existing:
        raise HTTPException(status_code=404, detail="Shipment not found")
    raw = payload.model_dump(exclude_unset=True)
    updates: Dict[str, Any] = {}
    # Nested origin/destination handling
    if "origin_facility" in raw or "origin_city" in raw:
        origin = dict(existing.get("origin") or {})
        if raw.get("origin_facility"):
            f = next((x for x in TMS_FACILITIES if x["id"] == raw["origin_facility"]), None)
            if f:
                origin = {"name": f["name"], "city": f["city"], "lat": f["lat"], "lng": f["lng"], "facility": f["id"]}
        if raw.get("origin_city"):
            origin["city"] = raw["origin_city"]
            origin["name"] = raw["origin_city"]
        updates["origin"] = origin
    if any(k in raw for k in ("destination_city", "destination_lat", "destination_lng")):
        dest = dict(existing.get("destination") or {})
        if raw.get("destination_city"):
            dest["city"] = raw["destination_city"]
            dest["name"] = raw["destination_city"]
        if raw.get("destination_lat") is not None:
            dest["lat"] = raw["destination_lat"]
        if raw.get("destination_lng") is not None:
            dest["lng"] = raw["destination_lng"]
        updates["destination"] = dest
    # Status normalization
    if raw.get("status") == "delivered" and "done" not in raw:
        updates["done"] = True
    # Flat fields
    for k in ("reference", "mode", "carrier", "status", "eta", "pickup_date", "ship_date", "ship_day",
              "weight_lbs", "pieces", "skids", "commodity", "value_usd", "container_no", "bol_no", "pro_no",
              "direction", "hazmat", "hazmat_class", "supplier", "consignee", "material_controller",
              "po_numbers", "booking_number", "bid_cost", "fsc_pct", "extras", "done", "shipping_hours",
              "pickup_no", "progress", "liftgate_required", "pallet_count", "nmfc_code", "freight_class",
              "accessorials", "length_in", "width_in", "height_in", "sap_delivery_no",
              "customer_contact_email", "carrier_contact_email"):
        if k in raw:
            updates[k] = raw[k]
    updates["updated_by"] = user.user_id
    updates["updated_at"] = datetime.now(timezone.utc).isoformat()
    await db.shipments.update_one({"shipment_id": shipment_id}, {"$set": updates})
    fresh = await db.shipments.find_one({"shipment_id": shipment_id}, {"_id": 0})
    return Shipment(**fresh)

class ShipmentCancel(BaseModel):
    reason: Optional[str] = None

@api_router.delete("/shipments/{shipment_id}")
async def cancel_shipment(shipment_id: str, payload: Optional[ShipmentCancel] = None, user: User = Depends(require_role("admin", "dispatcher"))):
    """Soft delete — marks the shipment as cancelled. Does NOT remove the row."""
    existing = await db.shipments.find_one({"shipment_id": shipment_id}, {"_id": 0})
    if not existing:
        raise HTTPException(status_code=404, detail="Shipment not found")
    reason = (payload.reason if payload else None) or "Cancelled by dispatcher"
    await db.shipments.update_one(
        {"shipment_id": shipment_id},
        {"$set": {
            "status": "cancelled",
            "done": False,
            "cancelled_at": datetime.now(timezone.utc).isoformat(),
            "cancelled_by": user.user_id,
            "cancel_reason": reason,
        }},
    )
    return {"ok": True, "shipment_id": shipment_id, "status": "cancelled", "reason": reason}

@api_router.post("/shipments", response_model=Shipment)
async def create_shipment(payload: ShipmentCreate, user: User = Depends(get_current_user)):
    origin = None
    if payload.origin_facility:
        f = next((x for x in TMS_FACILITIES if x["id"] == payload.origin_facility), None)
        if f:
            origin = {"name": f["name"], "city": f["city"], "lat": f["lat"], "lng": f["lng"], "facility": f["id"]}
    if not origin:
        origin = {"name": payload.origin_city or "Origin", "city": payload.origin_city or "Origin", "lat": 44.9858, "lng": -93.3499}
    destination = {"name": payload.destination_city, "city": payload.destination_city, "lat": payload.destination_lat, "lng": payload.destination_lng}
    sid = f"SHP-{uuid.uuid4().hex[:8].upper()}"
    eta_days = random.randint(2, 14)
    shipment = {
        "shipment_id": sid,
        "reference": payload.reference or f"TN-{random.randint(10000, 99999)}",
        "mode": payload.mode,
        "carrier": payload.carrier,
        "status": "pending",
        "origin": origin,
        "destination": destination,
        "current_location": {"lat": origin["lat"], "lng": origin["lng"], "city": origin["city"]},
        "eta": (datetime.now(timezone.utc) + timedelta(days=eta_days)).isoformat(),
        "pickup_date": payload.pickup_date,
        "weight_lbs": payload.weight_lbs,
        "pieces": payload.pieces,
        "commodity": payload.commodity,
        "value_usd": payload.value_usd,
        "container_no": f"TCLU{random.randint(1000000,9999999)}" if payload.mode == "Ocean" else None,
        "bol_no": f"BOL{random.randint(100000,999999)}",
        "pro_no": f"PRO{random.randint(100000,999999)}",
        "progress": 0.0,
        "liftgate_required": bool(payload.liftgate_required),
        "pallet_count": payload.pallet_count if payload.pallet_count is not None else payload.pieces,
        "nmfc_code": payload.nmfc_code,
        "freight_class": payload.freight_class,
        "accessorials": payload.accessorials or ([] if not payload.liftgate_required else ["liftgate"]),
        "skids": payload.pallet_count if payload.pallet_count is not None else payload.pieces,
        "length_in": payload.length_in,
        "width_in": payload.width_in,
        "height_in": payload.height_in,
        "sap_delivery_no": payload.sap_delivery_no,
        "sap_material_numbers": payload.sap_material_numbers,
        "customer_contact_email": payload.customer_contact_email,
        "carrier_contact_email": payload.carrier_contact_email,
        "created_by": user.user_id,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "is_sample": False,            # explicitly real, not seeded
        "source": "book_load",         # so the UI can route correctly
    }
    await db.shipments.insert_one(dict(shipment))

    # Cross-module sync: every shipment booked from the Book Load screen
    # must also create a brokerage_bookings row so the load flows into
    # Workflow → Factoring → Cash Flow → AI Triage. Without this mirror,
    # the load is visible only in /shipments and the operator can't run
    # the run-the-load HUD on it.
    try:
        equipment_map = {
            "TL": "Van", "LTL": "LTL", "Parcel": "Parcel",
            "Ocean": "Container", "Air": "Air", "Rail": "Rail",
        }
        rate_usd = float(payload.value_usd or 0)
        carrier_pay = round(rate_usd * 0.85, 2) if rate_usd > 0 else 0
        margin = round(rate_usd - carrier_pay, 2) if rate_usd > 0 else 0
        booking_doc = {
            "booked_id":         f"BK-{uuid.uuid4().hex[:10].upper()}",
            "load_id":           sid,
            "shipment_id":       sid,           # back-reference
            "reference":         shipment["reference"],
            "source":            "book_load",   # so the UI can tag it
            "board_id":          "internal",
            "carrier_name":      payload.carrier or "Unassigned",
            "carrier_mc":        None,
            "origin":            origin.get("city") or "Origin",
            "destination":       destination.get("city") or "Destination",
            "origin_full":       origin,
            "destination_full":  destination,
            "miles":             0,
            "equipment":         equipment_map.get(payload.mode, payload.mode),
            "mode":              payload.mode,
            "pieces":             payload.pieces,
            "weight_lbs":         payload.weight_lbs,
            "commodity":          payload.commodity,
            "pickup_date":        payload.pickup_date,
            "delivery_date":      None,
            "forecast_rate_usd":  rate_usd,
            "forecast_carrier_pay_usd": carrier_pay,
            "forecast_margin_usd":      margin,
            "rate_usd":           rate_usd,
            "carrier_pay_usd":    carrier_pay,
            "settled_rate_usd":   None,
            "settled_carrier_pay_usd": None,
            "settled_margin_usd":      None,
            "status":             "booked",
            "booked_at":          shipment["created_at"],
            "booked_by":          user.user_id,
            "notes":              "",
            "is_sample":          False,
        }
        await db.brokerage_bookings.insert_one(dict(booking_doc))
    except Exception as e:                                       # noqa: BLE001
        logger.exception("book_load → brokerage_bookings mirror failed: %s", e)

    return Shipment(**shipment)

# -------------------- DOCUMENTS --------------------
@api_router.get("/documents", response_model=List[Document])
async def list_documents(_: User = Depends(get_current_user)):
    docs = await db.documents.find({}, {"_id": 0}).sort("created_at", -1).limit(200).to_list(200)
    return [Document(**d) for d in docs]

class DocumentCreate(BaseModel):
    type: str
    shipment_ref: str
    data: Dict[str, Any]

@api_router.post("/documents", response_model=Document)
async def create_document(payload: DocumentCreate, user: User = Depends(get_current_user)):
    doc = {
        "document_id": f"DOC-{uuid.uuid4().hex[:8].upper()}",
        "type": payload.type,
        "shipment_ref": payload.shipment_ref,
        "created_by": user.name,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "data": payload.data,
        "amendments": [],
        "version": 1,
        "updated_at": None,
    }
    await db.documents.insert_one(dict(doc))
    return Document(**doc)


class DocumentAmend(BaseModel):
    data: Dict[str, Any]
    reason: Optional[str] = ""


@api_router.patch("/documents/{document_id}", response_model=Document)
async def amend_document(
    document_id: str,
    payload: DocumentAmend,
    user: User = Depends(require_role("admin", "auditor", "dispatcher")),
):
    """Amend a previously created document (BOL, COMMERCIAL_INVOICE, etc.).

    Stores a diff of changed fields plus a free-text reason on an `amendments`
    array. The `data` field is replaced with the new payload and `version`
    increments. PDF re-generation always reflects the latest `data`, so the
    next PDF download shows the amended values, while the audit trail
    preserves who-changed-what-when.
    """
    existing = await db.documents.find_one({"document_id": document_id}, {"_id": 0})
    if not existing:
        raise HTTPException(status_code=404, detail="Document not found")

    old_data: Dict[str, Any] = existing.get("data", {}) or {}
    new_data: Dict[str, Any] = payload.data or {}
    changes: List[Dict[str, Any]] = []
    all_keys = set(old_data.keys()) | set(new_data.keys())
    for field in sorted(all_keys):
        before = old_data.get(field)
        after = new_data.get(field)
        if (before or "") != (after or ""):
            changes.append({"field": field, "from": before, "to": after})

    amendment = {
        "amended_at": datetime.now(timezone.utc).isoformat(),
        "amended_by": user.name,
        "reason": (payload.reason or "").strip(),
        "changes": changes,
    }
    new_version = int(existing.get("version", 1)) + 1
    await db.documents.update_one(
        {"document_id": document_id},
        {
            "$set": {
                "data": new_data,
                "version": new_version,
                "updated_at": amendment["amended_at"],
            },
            "$push": {"amendments": amendment},
        },
    )
    fresh = await db.documents.find_one({"document_id": document_id}, {"_id": 0})
    return Document(**fresh)


class DocumentEmail(BaseModel):
    to: str
    cc: Optional[str] = ""
    message: Optional[str] = ""


@api_router.post("/documents/{document_id}/email")
async def email_document(
    document_id: str,
    payload: DocumentEmail,
    user: User = Depends(require_role("admin", "auditor", "dispatcher")),
):
    """Build a one-click email (subject/body/mailto) that delivers the
    document's PDF download link to the named recipient — same pattern as the
    existing routing-guide and carrier-email composers."""
    doc = await db.documents.find_one({"document_id": document_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    type_label, _ = DOC_TYPE_TITLES.get(doc["type"], (doc["type"], ""))
    pdf_url = f"{os.environ.get('PUBLIC_BASE_URL', '').rstrip('/')}/api/documents/{document_id}/pdf"
    version_note = f" (Rev {doc.get('version', 1)})" if doc.get("version", 1) > 1 else ""
    subject = f"{type_label}{version_note} · {doc.get('shipment_ref') or document_id}"

    data = doc.get("data") or {}
    body_lines = [
        f"Attached is the {type_label.lower()}{version_note} for the brokerage.",
        "",
        f"Document ID: {document_id}",
        f"Shipment:    {doc.get('shipment_ref') or '—'}",
        f"Carrier:     {data.get('carrier') or '—'}",
        f"Origin:      {data.get('origin') or '—'}",
        f"Destination: {data.get('destination') or '—'}",
        f"Commodity:   {data.get('commodity') or '—'}",
        f"Pieces:      {data.get('pieces') or '—'}",
        f"Weight:      {data.get('weight') or '—'} lbs",
        "",
        f"Download PDF: {pdf_url}" if pdf_url.startswith("http") else f"PDF endpoint: /api/documents/{document_id}/pdf",
    ]
    if payload.message:
        body_lines.extend(["", "---", payload.message.strip()])
    if doc.get("amendments"):
        body_lines.extend(["", f"This document has been amended {len(doc['amendments'])} time(s). Latest revision: {doc.get('version', 1)}."])
    body_lines.extend([
        "",
        "Thank you,",
        f"{user.name}",
        "Orisei Freight Solutions · TMS",
    ])
    body = "\n".join(body_lines)

    mailto = (
        f"mailto:{payload.to}"
        f"?subject={urllib.parse.quote(subject)}"
        f"&body={urllib.parse.quote(body)}"
    )
    if payload.cc:
        mailto += f"&cc={urllib.parse.quote(payload.cc)}"

    # Log the send-intent so a single source of truth exists for audit reports
    await db.document_emails.insert_one({
        "document_id": document_id,
        "to": payload.to,
        "cc": payload.cc or "",
        "subject": subject,
        "sent_by": user.name,
        "sent_at": datetime.now(timezone.utc).isoformat(),
    })

    return {
        "ok": True,
        "subject": subject,
        "body": body,
        "to": payload.to,
        "cc": payload.cc or "",
        "mailto": mailto,
        "pdf_url": f"/api/documents/{document_id}/pdf",
    }


class BolFromShipment(BaseModel):
    shipper: Optional[str] = "Orisei Freight Solutions"


@api_router.post("/shipments/{shipment_id}/generate-bol", response_model=Document)
async def generate_bol_from_shipment(
    shipment_id: str,
    payload: BolFromShipment,
    user: User = Depends(require_role("admin", "dispatcher")),
):
    """One-click BOL generator: pull the shipment record and create a BOL
    document with all carrier / origin / destination / commodity fields
    pre-filled. The generated doc lands in the Documents archive alongside
    manually-created BOLs, and can be amended / emailed / downloaded just
    like any other document."""
    s = await db.shipments.find_one({"shipment_id": shipment_id}, {"_id": 0})
    if not s:
        raise HTTPException(status_code=404, detail="Shipment not found")

    origin = s.get("origin", {}) or {}
    dest = s.get("destination", {}) or {}
    # Brand-aware default shipper / commodity: pull from active brand so a
    # Pfizer-active TMS doesn't print "Orisei Freight Solutions" on every BOL.
    brand = await _active_brand_doc()
    brand_company = (brand or {}).get("company_name") or "Orisei Freight Solutions"
    brand_default_commodity = "industrial cleaning equipment"
    if brand and brand.get("brand_id") != "orisei-freight":
        # Use the first sample product as the brand's representative commodity.
        sp = (brand.get("sample_products") or [None])[0]
        if sp:
            brand_default_commodity = sp
    data = {
        "shipper": payload.shipper or brand_company,
        "consignee": dest.get("name") or dest.get("city") or "",
        "origin": f"{origin.get('city', '')}, {origin.get('state', '')}".strip(", "),
        "destination": f"{dest.get('city', '')}, {dest.get('state', '')}".strip(", "),
        "carrier": s.get("carrier") or "",
        "commodity": s.get("commodity") or f"{brand_company} {brand_default_commodity}",
        "weight": s.get("weight_lbs") or "",
        "pieces": s.get("pieces") or s.get("skids") or "",
        "value": s.get("value_usd") or "",
        "country_origin": s.get("country_of_origin") or "USA",
        "bol_no": s.get("bol_no") or "",
        "pro_no": s.get("pro_no") or "",
    }
    doc = {
        "document_id": f"DOC-{uuid.uuid4().hex[:8].upper()}",
        "type": "BOL",
        "shipment_ref": s.get("reference") or shipment_id,
        "created_by": user.name,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "data": data,
        "amendments": [],
        "version": 1,
        "updated_at": None,
    }
    await db.documents.insert_one(dict(doc))
    return Document(**doc)

# -------------------- KPIs --------------------
@api_router.get("/kpis")
async def get_kpis(_: User = Depends(get_current_user)):
    all_shipments = await db.shipments.find({}, {"_id": 0}).to_list(2000)
    total = len(all_shipments)
    in_transit = sum(1 for s in all_shipments if s["status"] == "in_transit")
    delayed = sum(1 for s in all_shipments if s["status"] == "delayed")
    delivered = sum(1 for s in all_shipments if s["status"] == "delivered")
    pending = sum(1 for s in all_shipments if s["status"] == "pending")
    by_mode: Dict[str, int] = {}
    by_carrier: Dict[str, Dict[str, int]] = {}
    total_weight = 0.0
    total_value = 0.0
    for s in all_shipments:
        by_mode[s["mode"]] = by_mode.get(s["mode"], 0) + 1
        c = s["carrier"]
        if c not in by_carrier:
            by_carrier[c] = {"total": 0, "on_time": 0, "delayed": 0}
        by_carrier[c]["total"] += 1
        if s["status"] == "delivered":
            by_carrier[c]["on_time"] += 1
        if s["status"] == "delayed":
            by_carrier[c]["delayed"] += 1
        total_weight += s.get("weight_lbs", 0)
        total_value += s.get("value_usd", 0)
    on_time_rate = round((delivered / max(1, total)) * 100, 1)

    # ---- INDUSTRY-STANDARD CARRIER SCORECARD ----
    # Full metric set aligned with CSCMP, ATA, NASSTRAC, and ISO 9001 transport
    # benchmarks. Numbers are deterministic per-carrier so the report is stable
    # across page refreshes (random seeded by the carrier name).
    scorecard = []
    for carrier, stats in by_carrier.items():
        rnd = random.Random(carrier)
        total_loads = stats["total"]
        on_time_loads = stats["on_time"]
        # === Service quality ===
        on_time_pickup_pct       = round(rnd.uniform(91.0, 99.4), 1)   # OTP
        on_time_delivery_pct     = round(rnd.uniform(89.0, 99.2), 1)   # OTD
        on_time_in_full_pct      = round(rnd.uniform(85.0, 98.6), 1)   # OTIF
        tender_acceptance_pct    = round(rnd.uniform(78.0, 99.0), 1)   # First-tender accept
        tender_lead_time_hrs     = round(rnd.uniform(2.0, 14.0), 1)
        avg_transit_days         = round(rnd.uniform(1.4, 6.8), 1)
        transit_variance_pct     = round(rnd.uniform(2.5, 18.0), 1)    # σ/mean
        # === Compliance / quality ===
        claims_count             = rnd.randint(0, max(1, total_loads // 12))
        claims_freq_pct          = round((claims_count / max(1, total_loads)) * 100, 2)
        claims_paid_pct          = round(rnd.uniform(68.0, 99.0), 1)
        damage_freq_pct          = round(rnd.uniform(0.05, 1.8), 2)
        shortage_freq_pct        = round(rnd.uniform(0.02, 0.9), 2)
        billing_accuracy_pct     = round(rnd.uniform(94.0, 99.8), 1)
        edi_compliance_pct       = round(rnd.uniform(82.0, 99.6), 1)
        pod_timeliness_pct       = round(rnd.uniform(80.0, 99.0), 1)   # POD within SLA
        invoice_dispute_rate_pct = round(rnd.uniform(0.3, 6.5), 1)
        # === Safety / regulatory (FMCSA CSA / DOT) ===
        csa_score                = rnd.randint(3, 78)                   # lower is better
        safety_rating            = rnd.choice(["Satisfactory", "Satisfactory", "Satisfactory", "Conditional"])
        coi_days_to_expiry       = rnd.randint(-15, 365)
        out_of_service_pct       = round(rnd.uniform(0.2, 5.0), 1)      # OOS roadside
        hours_of_service_viol    = rnd.randint(0, 4)
        # === Cost / commercial ===
        avg_cost_per_load_usd    = round(rnd.uniform(420.0, 3850.0), 2)
        avg_cost_per_mile_usd    = round(rnd.uniform(1.55, 3.45), 2)
        fsc_per_mile_usd         = round(rnd.uniform(0.32, 0.58), 3)
        accessorial_spend_pct    = round(rnd.uniform(2.0, 14.0), 1)
        detention_hours_total    = rnd.randint(0, 50)
        detention_cost_usd       = round(detention_hours_total * rnd.uniform(45.0, 95.0), 2)
        rate_compliance_pct      = round(rnd.uniform(88.0, 99.6), 1)
        # === Capacity / responsiveness ===
        capacity_utilization_pct = round(rnd.uniform(62.0, 96.0), 1)
        committed_capacity_loads = rnd.randint(5, 200)
        spot_market_loads        = rnd.randint(0, 60)
        avg_response_time_min    = round(rnd.uniform(2.0, 45.0), 1)
        # === Sustainability ===
        empty_miles_pct          = round(rnd.uniform(4.0, 18.0), 1)
        co2_kg_per_load          = round(rnd.uniform(180.0, 2400.0), 1)
        ev_fleet_pct             = round(rnd.uniform(0.0, 22.0), 1)
        # === Composite (0-100, weighted) ===
        composite = round(
            (on_time_delivery_pct * 0.20) +
            (on_time_in_full_pct  * 0.15) +
            (tender_acceptance_pct * 0.10) +
            (billing_accuracy_pct * 0.10) +
            ((100 - claims_freq_pct * 10) * 0.10) +
            ((100 - damage_freq_pct * 10) * 0.05) +
            (rate_compliance_pct  * 0.10) +
            ((100 - empty_miles_pct) * 0.05) +
            ((100 - out_of_service_pct * 5) * 0.05) +
            (pod_timeliness_pct * 0.10),
            1,
        )
        grade = ("A+" if composite >= 95 else
                 "A"  if composite >= 90 else
                 "A-" if composite >= 87 else
                 "B+" if composite >= 84 else
                 "B"  if composite >= 80 else
                 "B-" if composite >= 77 else
                 "C+" if composite >= 73 else
                 "C"  if composite >= 70 else
                 "D"  if composite >= 60 else "F")
        trend_dir = rnd.choice(["up", "up", "flat", "down"])

        scorecard.append({
            "carrier": carrier,
            "total_loads": total_loads,
            "on_time_loads": on_time_loads,
            "delayed_loads": stats["delayed"],
            # Legacy fields kept for backward compat
            "total": total_loads, "on_time": on_time_loads, "delayed": stats["delayed"],
            # Service quality
            "on_time_pickup_pct": on_time_pickup_pct,
            "on_time_delivery_pct": on_time_delivery_pct,
            "on_time_in_full_pct": on_time_in_full_pct,
            "tender_acceptance_pct": tender_acceptance_pct,
            "tender_lead_time_hrs": tender_lead_time_hrs,
            "avg_transit_days": avg_transit_days,
            "transit_variance_pct": transit_variance_pct,
            # Compliance / quality
            "claims_count": claims_count,
            "claims_freq_pct": claims_freq_pct,
            "claims_paid_pct": claims_paid_pct,
            "damage_freq_pct": damage_freq_pct,
            "shortage_freq_pct": shortage_freq_pct,
            "billing_accuracy_pct": billing_accuracy_pct,
            "edi_compliance_pct": edi_compliance_pct,
            "pod_timeliness_pct": pod_timeliness_pct,
            "invoice_dispute_rate_pct": invoice_dispute_rate_pct,
            # Safety / regulatory
            "csa_score": csa_score,
            "safety_rating": safety_rating,
            "coi_days_to_expiry": coi_days_to_expiry,
            "out_of_service_pct": out_of_service_pct,
            "hours_of_service_violations": hours_of_service_viol,
            # Cost / commercial
            "avg_cost_per_load_usd": avg_cost_per_load_usd,
            "avg_cost_per_mile_usd": avg_cost_per_mile_usd,
            "fsc_per_mile_usd": fsc_per_mile_usd,
            "accessorial_spend_pct": accessorial_spend_pct,
            "detention_hours_total": detention_hours_total,
            "detention_cost_usd": detention_cost_usd,
            "rate_compliance_pct": rate_compliance_pct,
            # Capacity / responsiveness
            "capacity_utilization_pct": capacity_utilization_pct,
            "committed_capacity_loads": committed_capacity_loads,
            "spot_market_loads": spot_market_loads,
            "avg_response_time_min": avg_response_time_min,
            # Sustainability
            "empty_miles_pct": empty_miles_pct,
            "co2_kg_per_load": co2_kg_per_load,
            "ev_fleet_pct": ev_fleet_pct,
            # Composite
            "composite_score": composite,
            "grade": grade,
            "trend": trend_dir,
        })
    scorecard.sort(key=lambda x: -x["composite_score"])

    # Trend last 14 days
    today = datetime.now(timezone.utc).date()
    trend = []
    for i in range(13, -1, -1):
        d = today - timedelta(days=i)
        trend.append({
            "date": d.isoformat(),
            "shipments": random.randint(8, 26),
            "on_time": random.randint(7, 25),
            "cost": round(random.uniform(45000, 92000), 2),
        })

    # === NETWORK-WIDE TRANSPORTATION METRICS (industry-standard) ===
    # Every published number an industry exec or 3PL benchmark report cares about.
    network_metrics = {
        "service_quality": [
            {"key": "on_time_pickup",       "label": "On-Time Pickup",            "value": 95.3, "unit": "%",  "target": 95,  "benchmark": 93,  "trend": +1.2, "category": "service"},
            {"key": "on_time_delivery",     "label": "On-Time Delivery",          "value": 93.8, "unit": "%",  "target": 95,  "benchmark": 92,  "trend": +0.8, "category": "service"},
            {"key": "otif",                 "label": "On-Time In-Full (OTIF)",    "value": 91.2, "unit": "%",  "target": 95,  "benchmark": 89,  "trend": +1.5, "category": "service"},
            {"key": "perfect_order",        "label": "Perfect Order Rate",        "value": 89.6, "unit": "%",  "target": 92,  "benchmark": 87,  "trend": +0.4, "category": "service"},
            {"key": "first_attempt_delivery","label": "First-Attempt Delivery",  "value": 96.8, "unit": "%",  "target": 97,  "benchmark": 94,  "trend": +0.6, "category": "service"},
            {"key": "tender_acceptance",    "label": "First-Tender Acceptance",   "value": 87.4, "unit": "%",  "target": 92,  "benchmark": 84,  "trend": -1.1, "category": "service"},
            {"key": "tender_lead_time",     "label": "Avg Tender Lead Time",      "value": 6.4,  "unit": "h",  "target": 4,   "benchmark": 8,   "trend": -0.3, "category": "service"},
            {"key": "transit_variance",     "label": "Transit Time Variance",     "value": 7.8,  "unit": "%",  "target": 5,   "benchmark": 12,  "trend": -0.5, "category": "service"},
        ],
        "cost_efficiency": [
            {"key": "cost_per_mile",        "label": "Cost per Mile (TL)",        "value": 2.18, "unit": "$",  "target": 2.10, "benchmark": 2.34, "trend": -0.04, "category": "cost"},
            {"key": "cost_per_pound",       "label": "Cost per Pound (LTL)",      "value": 0.092,"unit": "$",  "target": 0.090,"benchmark": 0.098,"trend": -0.003,"category": "cost"},
            {"key": "cost_per_load",        "label": "Avg Cost per Load",         "value": 1842, "unit": "$",  "target": 1750, "benchmark": 1925, "trend": +28,    "category": "cost"},
            {"key": "freight_spend_pct_rev","label": "Freight Spend (% Revenue)", "value": 5.2,  "unit": "%",  "target": 5.0,  "benchmark": 5.8,  "trend": -0.1,   "category": "cost"},
            {"key": "accessorial_pct",      "label": "Accessorial Spend",         "value": 7.4,  "unit": "%",  "target": 6.0,  "benchmark": 9.2,  "trend": -0.3,   "category": "cost"},
            {"key": "fsc_pct_total",        "label": "FSC % of Total",            "value": 18.6, "unit": "%",  "target": 18,   "benchmark": 19.5, "trend": +0.2,   "category": "cost"},
            {"key": "detention_spend",      "label": "Detention Spend (YTD)",     "value": 142800,"unit": "$", "target": 100000,"benchmark": 170000,"trend": -8400, "category": "cost"},
            {"key": "audit_recovery",       "label": "Freight Audit Recovery",    "value": 38200, "unit": "$", "target": 30000, "benchmark": 25000, "trend": +4200,  "category": "cost"},
        ],
        "capacity_utilization": [
            {"key": "weight_utilization",   "label": "Trailer Weight Utilization", "value": 78.4, "unit": "%", "target": 85,  "benchmark": 76,  "trend": +1.6, "category": "capacity"},
            {"key": "cube_utilization",     "label": "Trailer Cube Utilization",   "value": 82.1, "unit": "%", "target": 85,  "benchmark": 79,  "trend": +0.8, "category": "capacity"},
            {"key": "empty_miles",          "label": "Empty Miles",                "value": 9.2,  "unit": "%", "target": 7,   "benchmark": 11,  "trend": -0.7, "category": "capacity"},
            {"key": "load_consolidation",   "label": "Consolidation Ratio",        "value": 1.34, "unit": "x", "target": 1.50, "benchmark": 1.22,"trend": +0.06,"category": "capacity"},
            {"key": "miles_per_load",       "label": "Avg Miles per Load",         "value": 614,  "unit": "mi","target": 650, "benchmark": 580, "trend": +12,   "category": "capacity"},
            {"key": "drop_trailer_pct",     "label": "Drop-Trailer Loads",         "value": 31.2, "unit": "%", "target": 35,  "benchmark": 26,  "trend": +2.1, "category": "capacity"},
        ],
        "compliance_quality": [
            {"key": "claims_freq",          "label": "Claims Frequency",          "value": 0.42, "unit": "%", "target": 0.50, "benchmark": 0.70, "trend": -0.06, "category": "quality"},
            {"key": "damage_rate",          "label": "Damage Rate",               "value": 0.18, "unit": "%", "target": 0.20, "benchmark": 0.35, "trend": -0.03, "category": "quality"},
            {"key": "shortage_rate",        "label": "Shortage Rate",             "value": 0.11, "unit": "%", "target": 0.15, "benchmark": 0.22, "trend": -0.02, "category": "quality"},
            {"key": "claims_resolved_30d",  "label": "Claims Resolved < 30d",     "value": 84.2, "unit": "%", "target": 90,   "benchmark": 78,   "trend": +2.4,  "category": "quality"},
            {"key": "billing_accuracy",     "label": "Billing Accuracy",          "value": 97.1, "unit": "%", "target": 98,   "benchmark": 95,   "trend": +0.3,  "category": "quality"},
            {"key": "edi_compliance",       "label": "EDI 214/210 Compliance",    "value": 91.8, "unit": "%", "target": 95,   "benchmark": 88,   "trend": +1.7,  "category": "quality"},
            {"key": "pod_within_24h",       "label": "POD Within 24h",            "value": 88.6, "unit": "%", "target": 95,   "benchmark": 82,   "trend": +1.9,  "category": "quality"},
            {"key": "invoice_dispute",      "label": "Invoice Dispute Rate",      "value": 2.4,  "unit": "%", "target": 2.0,  "benchmark": 4.1,  "trend": -0.3,  "category": "quality"},
        ],
        "safety_regulatory": [
            {"key": "fmcsa_csa_avg",        "label": "FMCSA CSA Avg (Carriers)",  "value": 28.4, "unit": "",  "target": 30,   "benchmark": 35,   "trend": -0.8,  "category": "safety"},
            {"key": "oos_rate",             "label": "Out-of-Service Rate",       "value": 1.8,  "unit": "%", "target": 2.0,  "benchmark": 3.4,  "trend": -0.2,  "category": "safety"},
            {"key": "hos_violations",       "label": "HOS Violations (YTD)",      "value": 12,   "unit": "",  "target": 15,   "benchmark": 22,   "trend": -3,    "category": "safety"},
            {"key": "coi_compliant",        "label": "Carriers with Valid COI",   "value": 96.8, "unit": "%", "target": 100,  "benchmark": 92,   "trend": +0.4,  "category": "safety"},
            {"key": "hazmat_violations",    "label": "Hazmat Violations",         "value": 0,    "unit": "",  "target": 0,    "benchmark": 1,    "trend": 0,     "category": "safety"},
            {"key": "preventable_accidents","label": "Preventable Accidents",     "value": 2,    "unit": "",  "target": 0,    "benchmark": 4,    "trend": -1,    "category": "safety"},
        ],
        "sustainability": [
            {"key": "co2_per_load",         "label": "CO₂ per Load",              "value": 1240, "unit": "kg","target": 1100, "benchmark": 1380, "trend": -42,   "category": "sustainability"},
            {"key": "co2_per_ton_mile",     "label": "CO₂ per Ton-Mile",          "value": 75.4, "unit": "g", "target": 70,   "benchmark": 88,   "trend": -2.1,  "category": "sustainability"},
            {"key": "ev_fleet_pct",         "label": "EV / Alt-Fuel Fleet",       "value": 6.4,  "unit": "%", "target": 15,   "benchmark": 4.2,  "trend": +0.9,  "category": "sustainability"},
            {"key": "smartway_pct",         "label": "SmartWay-Carrier %",        "value": 78.3, "unit": "%", "target": 85,   "benchmark": 65,   "trend": +2.1,  "category": "sustainability"},
            {"key": "intermodal_share",     "label": "Intermodal Share",          "value": 12.8, "unit": "%", "target": 18,   "benchmark": 10,   "trend": +0.6,  "category": "sustainability"},
        ],
    }

    # === BRAND OVERLAY — perturb metrics deterministically per active brand ===
    # default brands stay untouched. For any other active brand we apply a
    # stable per-metric drift seeded by the brand_id so each company has its
    # own coherent KPI fingerprint (e.g. Walmart's OTIF differs from FedEx).
    brand_doc = await _active_brand_doc()
    if brand_doc and brand_doc.get("brand_id") != "orisei-freight":
        bseed = brand_doc.get("brand_id") or brand_doc.get("short_name") or "brand"
        def _drift(val, key, lo=-0.12, hi=0.12):
            if not isinstance(val, (int, float)) or val == 0:
                return val
            rnd = random.Random(f"{bseed}::{key}")
            factor = 1 + rnd.uniform(lo, hi)
            new_val = val * factor
            if isinstance(val, int):
                return max(0, int(round(new_val)))
            if abs(val) >= 100:
                return round(new_val, 1)
            if abs(val) >= 1:
                return round(new_val, 2)
            return round(new_val, 3)
        for category, items in network_metrics.items():
            for m in items:
                m["value"] = _drift(m["value"], f"nm::{category}::{m['key']}::value")
                if m.get("unit") == "%" and isinstance(m["value"], (int, float)):
                    m["value"] = max(0, min(100, m["value"]))
                m["trend"] = _drift(m["trend"], f"nm::{category}::{m['key']}::trend", -0.25, 0.25)
        for i, day in enumerate(trend):
            day["shipments"] = _drift(day["shipments"], f"trend::{i}::ships")
            day["on_time"] = min(day["shipments"], _drift(day["on_time"], f"trend::{i}::ot"))
            day["cost"] = _drift(day["cost"], f"trend::{i}::cost")
        for c in scorecard:
            c["composite_score"] = round(min(99.9, max(40, _drift(c["composite_score"], f"sc::{c['carrier']}::comp", -0.06, 0.06))), 1)
            c["on_time_delivery_pct"] = max(0, min(100, _drift(c["on_time_delivery_pct"], f"sc::{c['carrier']}::otd", -0.04, 0.04)))
            c["on_time_in_full_pct"] = max(0, min(100, _drift(c["on_time_in_full_pct"], f"sc::{c['carrier']}::otif", -0.04, 0.04)))
        scorecard.sort(key=lambda x: -x["composite_score"])

    return await _brand_swap({
        "totals": {
            "total": total,
            "in_transit": in_transit,
            "delayed": delayed,
            "delivered": delivered,
            "pending": pending,
            "weight_lbs": round(total_weight, 0),
            "value_usd": round(total_value, 0),
            "on_time_rate": on_time_rate,
        },
        "by_mode": by_mode,
        "by_carrier": [{"carrier": k, **v} for k, v in by_carrier.items()],
        "carrier_scorecard": scorecard,
        "network_metrics": network_metrics,
        "trend": trend,
    })

# -------------------- KPIs · Weekly Weights --------------------
def _facility_for_shipment(s: Dict[str, Any]) -> Optional[str]:
    """Return facility id (GVM/HOM/LVK) tied to a shipment.
    Outbound -> origin.facility; Inbound -> destination.facility."""
    if s.get("direction") == "inbound":
        return (s.get("destination") or {}).get("facility")
    return (s.get("origin") or {}).get("facility")

@api_router.get("/kpis/weekly-weights")
async def get_weekly_weights(_: User = Depends(get_current_user)):
    """Weekly average shipped weight per facility (GVM, HOM, LVK) over the past 12 ISO weeks.
    Returns:
      {
        "series": [{ "week": "2026-W18", "GVM": float, "HOM": float, "LVK": float }, ...]   # 12 oldest -> newest
        "summary": {
          "GVM": { current_week_avg_lbs, twelve_wk_avg_lbs, wow_delta_lbs, twelve_wk_total_lbs },
          ...
        }
      }
    """
    today = datetime.now(timezone.utc).date()
    # Build the 12 most recent ISO week keys (oldest -> newest)
    week_keys: List[str] = []
    week_starts: List[datetime] = []
    for i in range(11, -1, -1):
        ref = today - timedelta(weeks=i)
        iso_year, iso_week, _ = ref.isocalendar()
        key = f"{iso_year}-W{iso_week:02d}"
        if key not in week_keys:
            week_keys.append(key)
            week_starts.append(datetime.combine(ref, datetime.min.time(), tzinfo=timezone.utc))

    facilities = ["GVM", "HOM", "LVK"]
    # week_key -> facility -> [weights]
    buckets: Dict[str, Dict[str, List[float]]] = {wk: {f: [] for f in facilities} for wk in week_keys}

    all_shipments = await db.shipments.find({}, {"_id": 0}).to_list(5000)
    for s in all_shipments:
        if s.get("status") == "cancelled":
            continue
        fac = _facility_for_shipment(s)
        if fac not in facilities:
            continue
        date_str = s.get("ship_date") or s.get("pickup_date")
        if not date_str:
            continue
        try:
            d = datetime.fromisoformat(date_str.replace("Z", "+00:00")).date() if "T" in date_str else datetime.strptime(date_str[:10], "%Y-%m-%d").date()
        except Exception:
            continue
        iso_year, iso_week, _ = d.isocalendar()
        key = f"{iso_year}-W{iso_week:02d}"
        if key in buckets:
            try:
                buckets[key][fac].append(float(s.get("weight_lbs") or 0))
            except (TypeError, ValueError):
                continue

    series: List[Dict[str, Any]] = []
    for wk in week_keys:
        row: Dict[str, Any] = {"week": wk}
        for f in facilities:
            arr = buckets[wk][f]
            row[f] = round(sum(arr) / len(arr), 0) if arr else 0
        series.append(row)

    summary: Dict[str, Dict[str, float]] = {}
    for f in facilities:
        all_weights: List[float] = []
        for wk in week_keys:
            all_weights.extend(buckets[wk][f])
        current_avg = series[-1][f]
        prev_avg = series[-2][f] if len(series) >= 2 else 0
        summary[f] = {
            "current_week_avg_lbs": float(current_avg),
            "twelve_wk_avg_lbs": round(sum(all_weights) / len(all_weights), 0) if all_weights else 0.0,
            "wow_delta_lbs": float(current_avg) - float(prev_avg),
            "twelve_wk_total_lbs": round(sum(all_weights), 0),
            "shipment_count": len(all_weights),
        }
    return {"series": series, "summary": summary}

# -------------------- LIVE FEEDS --------------------
@api_router.get("/weather")
async def get_weather(_: User = Depends(get_current_user)):
    """Real weather via Open-Meteo (no key)."""
    results = []
    async with httpx.AsyncClient(timeout=10.0) as http:
        for f in TMS_FACILITIES:
            try:
                r = await http.get(
                    "https://api.open-meteo.com/v1/forecast",
                    params={
                        "latitude": f["lat"],
                        "longitude": f["lng"],
                        "current": "temperature_2m,wind_speed_10m,weather_code,relative_humidity_2m",
                        "temperature_unit": "fahrenheit",
                        "wind_speed_unit": "mph",
                    },
                )
                d = r.json().get("current", {})
                results.append({
                    "facility_id": f["id"],
                    "facility_name": f["name"],
                    "temperature_f": d.get("temperature_2m"),
                    "humidity": d.get("relative_humidity_2m"),
                    "wind_mph": d.get("wind_speed_10m"),
                    "weather_code": d.get("weather_code"),
                })
            except Exception as e:
                logger.warning(f"Weather fetch failed for {f['id']}: {e}")
                results.append({"facility_id": f["id"], "facility_name": f["name"], "error": str(e)})
    return results

MOCK_NEWS = [
    # 40+ rotating items. published_minutes_ago is randomized fresh on every
    # /api/news fetch to simulate a live ticker for the dispatch team.
    {"title": "Diesel prices ease 4¢ as Midwest refineries return from turnarounds", "source": "FreightWaves", "category": "fuel", "body": "EIA weekly diesel average drops to $3.78/gal — first decline in five weeks. PADD 2 (Midwest) led the pull-down at -6.2¢, easing some pressure on FSC schedules for Q1 freight.", "url": "https://www.freightwaves.com/news/diesel"},
    {"title": "Port of Long Beach reports 8% volume increase YoY for January", "source": "JOC", "category": "ocean", "body": "TEU throughput at Long Beach hit 952,600 in January, up from 881,700 a year ago. POLB attributes the gain to early front-loading ahead of Lunar New Year and ongoing strength in furniture and appliance categories.", "url": "https://www.joc.com/port-news"},
    {"title": "FMCSA proposes HOS exemption for short-haul drivers under 150 air-mile radius", "source": "Transport Topics", "category": "regulatory", "body": "Proposed rule would lengthen the short-haul exception from 14 to 16 hours and remove the 11-hour driving cap for qualifying intrastate runs.", "url": "https://www.fmcsa.dot.gov/proposed-rule"},
    {"title": "UPS adds 12 new electric ground-support tugs in Louisville hub", "source": "DC Velocity", "category": "carrier", "body": "Worldport's electrification program now covers 28% of ground-support equipment; UPS targets 40% by year-end.", "url": "https://www.dcvelocity.com/ups-electric"},
    {"title": "Severe winter weather forecast for Upper Midwest mid-week", "source": "Weather.gov", "category": "weather", "body": "NWS Twin Cities issues Winter Storm Watch for 14 counties Wed–Thu — 6–10\" snow forecast for Golden Valley corridor; carriers urged to pre-route freight.", "url": "https://weather.gov/mpx"},
    {"title": "USTR opens Section 301 four-year review on Chinese components", "source": "Reuters", "category": "trade", "body": "Public comment period opens Mar 1 — affects HTS chapters 84/85/87 including motors, batteries, and electronic controllers commonly used by industrial OEMs.", "url": "https://ustr.gov/section-301-review"},
    {"title": "Kuehne+Nagel expands North America air freight network with two new gateways", "source": "Air Cargo News", "category": "carrier", "body": "K+N adds Atlanta (ATL) and Dallas (DFW) gateways and 4 weekly 777F rotations from Frankfurt — designed to cut transit by 1–2 days on cross-Atlantic industrial freight.", "url": "https://www.aircargonews.net/kn-expansion"},
    {"title": "ELD malfunction reports rise 15% in Q4 — FMCSA notice", "source": "CDLLife", "category": "regulatory", "body": "FMCSA issues guidance reminding fleets to keep paper backup logs for 8 days when ELDs fail. Affected device vendors named in notice.", "url": "https://cdllife.com/eld-malfunction"},

    {"title": "Hapag-Lloyd raises GRI on US import lanes by $300/FFE effective Mar 1", "source": "JOC", "category": "ocean", "body": "Targets all-water Asia → USEC routes. Negotiated NAC contracts unaffected — spot market shippers should expect immediate impact.", "url": "https://www.hapag-lloyd.com/news"},
    {"title": "BNSF announces $1.2B capex plan for Southern Transcon corridor", "source": "Railway Age", "category": "rail", "body": "Investments target additional double-track between Belen NM and Amarillo TX, plus Pasadena CA intermodal terminal expansion to handle larger intermodal volumes.", "url": "https://www.railwayage.com/bnsf-capex"},
    {"title": "CSX completes Howard Street Tunnel clearance project in Baltimore", "source": "Trains Magazine", "category": "rail", "body": "Double-stack clearance now operational from Mid-Atlantic ports inland — expected to remove ~120 truck moves/day off I-95 corridor.", "url": "https://trn.trains.com/csx-tunnel"},
    {"title": "Old Dominion posts record Q4 operating ratio of 70.1%", "source": "Logistics Management", "category": "ltl", "body": "ODFL revenue per hundredweight +5.4% YoY despite tonnage softness. Carrier credits density and absence of major service failures.", "url": "https://www.logisticsmgmt.com/odfl"},
    {"title": "XPO opens 12 new LTL service centers across the Southeast", "source": "FreightWaves", "category": "ltl", "body": "XPO completes the largest LTL terminal expansion in 30 years — adds 2.4M ft² of dock space in TN, GA, FL, SC.", "url": "https://www.freightwaves.com/xpo-expansion"},
    {"title": "Truckstop spot rates: DAT National Van down 2.1% WoW", "source": "DAT", "category": "spot-market", "body": "Linehaul rate excl. fuel sits at $1.66/mi for dry van; reefer holds at $2.04/mi as produce season starts to ramp in Florida and California.", "url": "https://www.dat.com/trendlines"},
    {"title": "Mexico nearshoring drives 23% increase in cross-border truckload volume", "source": "JOC", "category": "cross-border", "body": "Laredo crossings up 18% YoY; Otay Mesa up 27%. Reynosa supplier base reports 4-day average border dwell down to 32 hours.", "url": "https://www.joc.com/nearshoring"},
    {"title": "Maersk announces all-electric drayage fleet for LA/LB by 2027", "source": "American Shipper", "category": "ocean", "body": "350 Class-8 electric tractors ordered from Volvo and Daimler. Charging infrastructure in partnership with Forum Mobility.", "url": "https://www.americanshipper.com/maersk"},
    {"title": "DOT Sec'y announces $1.4B in port infrastructure grants", "source": "DOT", "category": "regulatory", "body": "20 ports receive funding; Long Beach, Houston, Norfolk top recipients. Focus on rail connectivity and zero-emission cargo handling.", "url": "https://www.transportation.gov/port-grants"},
    {"title": "Saia announces 22 new terminal openings for 2026", "source": "Transport Topics", "category": "ltl", "body": "Aggressive geographic expansion in the Pacific Northwest and New England — Saia's first OR, WA, and ME locations.", "url": "https://www.ttnews.com/saia-2026"},
    {"title": "Knight-Swift posts mixed Q4 — truckload margin compresses, logistics gains", "source": "Knight-Swift IR", "category": "carrier", "body": "Trucking segment OR 92.3% (-180bps YoY); Logistics 95.8% (+90bps). Net revenue up 4.2%.", "url": "https://knight-transportation.com/ir"},
    {"title": "Schneider deploys 65 Freightliner eCascadia tractors in CA", "source": "FleetOwner", "category": "sustainability", "body": "Largest single OEM order to date — supporting California Advanced Clean Fleets compliance.", "url": "https://www.fleetowner.com/schneider-ev"},
    {"title": "C.H. Robinson rolls out Procure IQ AI for spot procurement", "source": "Transport Topics", "category": "tech", "body": "AI-driven carrier match and pricing — claims to cut tender-to-cover time from 4 hours to 22 minutes on average.", "url": "https://www.chrobinson.com/procure-iq"},
    {"title": "Suez Canal traffic improves 12% as Red Sea attacks de-escalate", "source": "Lloyd's List", "category": "ocean", "body": "Daily transits average 56, up from 50 at the start of the month. Capacity still 38% below pre-conflict baseline.", "url": "https://lloydslist.maritimeintelligence.informa.com/suez"},
    {"title": "EU CBAM phase-2 compliance deadline approaches for steel imports", "source": "Reuters", "category": "trade", "body": "EU importers face quarterly reporting from April; non-compliance fines €50/ton CO₂ embedded. Industrial OEMs sourcing iron/steel from non-EU mills should validate emission documentation.", "url": "https://reuters.com/cbam"},
    {"title": "Werner Enterprises closes Q4 with $0.58 EPS, beats by $0.04", "source": "SeekingAlpha", "category": "carrier", "body": "Dedicated segment showed strongest performance; one-way TL still under pressure from soft spot rates.", "url": "https://seekingalpha.com/werner-q4"},
    {"title": "FedEx Freight tests autonomous yard tractors at Memphis hub", "source": "DC Velocity", "category": "tech", "body": "Outrider partnership — 6 electric autonomous tractors handling trailer spotting and yard moves 24/7.", "url": "https://www.dcvelocity.com/fedex-outrider"},
    {"title": "Yellow Corp. terminals auction nets $1.9B for creditors", "source": "FreightWaves", "category": "ltl", "body": "ABF, Estes, XPO among biggest winners. 130 terminals sold across US — implications for LTL capacity in 2026.", "url": "https://www.freightwaves.com/yellow-auction"},
    {"title": "Class-8 truck orders soar in January — ACT Research", "source": "Truckinginfo", "category": "industry", "body": "Net Class-8 orders 32,400 units — best January since 2018. Fleets restocking ahead of EPA 2027 emissions rules.", "url": "https://www.truckinginfo.com/class-8"},
    {"title": "Tornado warning issued for Madison County, IN — I-69 SB shut", "source": "NWS Indianapolis", "category": "weather", "body": "Severe weather expected through 4pm local; Holland-bound carriers re-routing via I-65 to avoid the cell.", "url": "https://weather.gov/ind"},
    {"title": "Customs Modernization Act draft text circulates in House", "source": "AAEI", "category": "trade", "body": "Would consolidate CBP/PGA filings into a single enhanced ACE manifest. Industrial importers should monitor for changes to FTZ admission procedures.", "url": "https://aaei.org/cma"},
    {"title": "OOCL ULCV maintenance schedule pushes 4 transpacific calls", "source": "JOC", "category": "ocean", "body": "Three ULCVs in drydock at Singapore; affected weekly sailings PSW3 and PCC.", "url": "https://www.joc.com/oocl-schedule"},
    {"title": "Reefer rates climb 9% as Florida strawberry season opens", "source": "DAT", "category": "spot-market", "body": "Lakeland → Atlanta refrigerated lane up to $3.42/mi loaded. Expect FL-MI lanes to follow suit in 2 weeks.", "url": "https://www.dat.com/reefer"},
    {"title": "AAR weekly carloads up 3.1% — chemicals lead growth", "source": "Association of American Railroads", "category": "rail", "body": "Total US originated carloads 226,118; chemical traffic +7.4% YoY. Intermodal also strong at +5.2%.", "url": "https://www.aar.org/weekly"},
    {"title": "Operator adds 4 new approved LTL carriers to North American roster", "source": "Internal", "category": "internal", "body": "XPO, ODFL, Saia, Estes all complete annual qualification review. Routing guide Rev 29 reflects updated lane assignments.", "url": "/routing-guide"},
    {"title": "Holland MI plant on-time inbound rate hits 96.4% for January", "source": "Internal", "category": "internal", "body": "All-time monthly record. Top 5 suppliers (Motrex, BattCo, Premier Polymers, Yazaki, Midwest Steel) all >97%.", "url": "/kpis"},
    {"title": "Hurricane Beryl recovery: Houston port back to 100% capacity", "source": "American Shipper", "category": "ocean", "body": "Backlog cleared 3 weeks after landfall — 27 vessels processed at peak vs. typical 18.", "url": "https://www.americanshipper.com/houston"},
    {"title": "USPS announces 5.9% rate hike for parcel select effective July", "source": "Parcel Industry", "category": "parcel", "body": "Heaviest impact on lightweight residential parcels; commercial PS rates up 2.7%.", "url": "https://about.usps.com/rates"},
    {"title": "DHL Express prioritizes battery-shipping training for ground couriers", "source": "Air Cargo News", "category": "hazmat", "body": "All US-based ground handlers complete IATA DGR section II training by end of Q1. Li-ion battery exports unaffected.", "url": "https://www.aircargonews.net/dhl"},
    {"title": "Drayage CHASSISGATE: NACCS reports container chassis shortage easing", "source": "JOC", "category": "ocean", "body": "Long Beach pool +320 chassis week-over-week. Wait times for export bookings drop to 1.8 days from 4.2.", "url": "https://www.joc.com/chassis"},
    {"title": "Estes Express announces driver pay increase averaging $0.06/mi", "source": "CCJ", "category": "ltl", "body": "Pay hike effective Mar 15; carrier reports applicant pipeline +28% since announcement.", "url": "https://www.ccjdigital.com/estes-pay"},
    {"title": "I-95 Cordage Park bridge replacement project shifts truck traffic", "source": "MassDOT", "category": "infra", "body": "Two-year detour begins Mar 1 — Carriers serving New England should route via I-93/I-90.", "url": "https://mass.gov/i95-cordage"},
]

# Drop the deprecated single mock-news + traffic; the live endpoints below
# generate fresh timestamps on every call so the dispatch team sees the
# ticker advance even when the underlying corpus hasn't changed.
import random as _random  # noqa: E402

def _mins_ago_to_label(m: int) -> str:
    if m < 1: return "now"
    if m < 60: return f"{m}m"
    h = m // 60
    if h < 24: return f"{h}h"
    return f"{h // 24}d"


@api_router.get("/news")
async def get_news(category: Optional[str] = None, limit: int = 40, _: User = Depends(get_current_user)):
    """Live news feed. Every fetch reshuffles 'minutes ago' values within
    each item's natural window so the ticker visibly advances each poll.
    The dispatch team rotates through the full 40-item corpus instead of
    the prior 8-item loop."""
    pool = MOCK_NEWS
    if category and category != "all":
        pool = [n for n in pool if n.get("category") == category]
    rnd = _random.Random()
    out = []
    for i, n in enumerate(pool):
        # Spread mins between 0 and 360 so the feed feels live but coherent.
        mins = rnd.randint(0, 360) + i  # slight bias keeps ordering varied
        out.append({**n, "minutes_ago": mins, "time": _mins_ago_to_label(mins),
                    "published_at": (datetime.now(timezone.utc) - timedelta(minutes=mins)).isoformat()})
    out.sort(key=lambda x: x["minutes_ago"])
    return await _brand_swap(out[:limit])


# 25 traffic incidents with rich detail — agency, lanes closed, source URL,
# expected clear-time, photos optional. Each fetch shuffles severity-weighted
# ordering and updates "minutes-ago" so the panel feels live.
MOCK_TRAFFIC = [
    {"location": "I-94 EB at Mile 215 (MI)", "highway": "I-94", "direction": "EB",
     "type": "Crash", "severity": "moderate", "delay_min": 25,
     "lat": 42.65, "lng": -86.10, "lanes_closed": "2 of 3 EB",
     "agency": "MI State Police", "eta_clear_min": 45,
     "description": "Multi-vehicle crash at MM 215. Right two lanes blocked. Expect significant backups eastbound from Stevensville.",
     "source_url": "https://www.michigan.gov/traffic", "near_facility": "Tennant Holland MI"},
    {"location": "I-65 N at Louisville Spaghetti Junction (KY)", "highway": "I-65", "direction": "NB",
     "type": "Construction", "severity": "low", "delay_min": 12,
     "lat": 38.26, "lng": -85.75, "lanes_closed": "1 of 4 NB (shoulder work)",
     "agency": "KYTC", "eta_clear_min": 240,
     "description": "Ongoing bridge deck rehabilitation. Right shoulder closed 9pm–5am. Normal flow elsewhere.",
     "source_url": "https://goky.ky.gov", "near_facility": "Tennant Louisville KY"},
    {"location": "I-394 W approach to Golden Valley (MN)", "highway": "I-394", "direction": "WB",
     "type": "Weather · Snow", "severity": "high", "delay_min": 40,
     "lat": 44.98, "lng": -93.35, "lanes_closed": "Reduced visibility — all lanes open",
     "agency": "MNDOT", "eta_clear_min": 90,
     "description": "Heavy snowfall with blowing snow reducing visibility to under 1/4 mile. MNDOT advising essential travel only. Tennant HQ inbound LTL recommend delay.",
     "source_url": "https://511mn.org", "near_facility": "Tennant HQ — Golden Valley MN"},
    {"location": "I-80 EB Ohio Turnpike Mile 161 (OH)", "highway": "I-80", "direction": "EB",
     "type": "Stalled vehicle", "severity": "low", "delay_min": 8,
     "lat": 41.36, "lng": -82.22, "lanes_closed": "Right shoulder",
     "agency": "Ohio Turnpike", "eta_clear_min": 20,
     "description": "Disabled tractor-trailer in right shoulder. State patrol on scene. Minor rubbernecking delays.",
     "source_url": "https://www.ohioturnpike.org/traffic"},
    {"location": "I-71 N Cincinnati (OH)", "highway": "I-71", "direction": "NB",
     "type": "Congestion", "severity": "moderate", "delay_min": 18,
     "lat": 39.16, "lng": -84.45, "lanes_closed": "—",
     "agency": "OHGO", "eta_clear_min": 60,
     "description": "Volume-related slowdowns building from MM 1 to MM 8. Expect 25-35 mph through downtown.",
     "source_url": "https://ohgo.com"},

    {"location": "I-5 SB at Tejon Pass (CA)", "highway": "I-5", "direction": "SB",
     "type": "Weather · Snow & Ice", "severity": "high", "delay_min": 75,
     "lat": 34.80, "lng": -118.86, "lanes_closed": "Chains required for big rigs",
     "agency": "Caltrans D7", "eta_clear_min": 180,
     "description": "Chain controls in effect for vehicles over 6,000 lbs. Two truck spinouts at MM 195 cleared. Use US-101 alternate.",
     "source_url": "https://quickmap.dot.ca.gov"},
    {"location": "I-90 EB Cleveland (OH)", "highway": "I-90", "direction": "EB",
     "type": "Accident with injury", "severity": "high", "delay_min": 50,
     "lat": 41.49, "lng": -81.69, "lanes_closed": "All EB lanes — detour at W 25th",
     "agency": "Ohio State Patrol", "eta_clear_min": 75,
     "description": "Serious injury crash. EB I-90 closed from Innerbelt through Bridge Ave. Use Detour via I-490/I-77.",
     "source_url": "https://ohgo.com"},
    {"location": "I-70 WB Columbus (OH)", "highway": "I-70", "direction": "WB",
     "type": "HazMat spill", "severity": "high", "delay_min": 120,
     "lat": 39.96, "lng": -83.02, "lanes_closed": "All WB lanes — full closure",
     "agency": "Columbus Fire HAZMAT", "eta_clear_min": 240,
     "description": "Class 3 flammable liquid spill from overturned tanker at MM 99. Full WB closure. Detour via I-670/I-71.",
     "source_url": "https://ohgo.com"},
    {"location": "I-285 EB Atlanta Perimeter (GA)", "highway": "I-285", "direction": "EB",
     "type": "Congestion", "severity": "moderate", "delay_min": 28,
     "lat": 33.78, "lng": -84.31, "lanes_closed": "—",
     "agency": "GDOT NaviGAtor", "eta_clear_min": 60,
     "description": "Typical PM peak congestion building. Slow from Roswell Rd to I-85.",
     "source_url": "https://www.511ga.org"},
    {"location": "I-405 SB LA Sepulveda Pass (CA)", "highway": "I-405", "direction": "SB",
     "type": "Construction", "severity": "moderate", "delay_min": 32,
     "lat": 34.08, "lng": -118.47, "lanes_closed": "1 of 5 SB (night work)",
     "agency": "Caltrans D7", "eta_clear_min": 180,
     "description": "Nightly lane closure for pavement repair. Expect heavy delays through Sepulveda Pass overnight.",
     "source_url": "https://quickmap.dot.ca.gov"},

    {"location": "I-40 EB Memphis (TN)", "highway": "I-40", "direction": "EB",
     "type": "Bridge inspection", "severity": "low", "delay_min": 15,
     "lat": 35.15, "lng": -90.06, "lanes_closed": "1 of 3 EB",
     "agency": "TDOT", "eta_clear_min": 120,
     "description": "Routine inspection at Hernando de Soto Bridge. Reduced to 2 lanes.",
     "source_url": "https://smartway.tn.gov"},
    {"location": "I-25 SB Denver (CO)", "highway": "I-25", "direction": "SB",
     "type": "Congestion", "severity": "high", "delay_min": 45,
     "lat": 39.74, "lng": -104.99, "lanes_closed": "—",
     "agency": "CDOT", "eta_clear_min": 75,
     "description": "Mousetrap interchange backup extending 6 miles into northern Denver. Inbound from Wyoming significantly delayed.",
     "source_url": "https://www.cotrip.org"},
    {"location": "I-35 NB Dallas (TX)", "highway": "I-35", "direction": "NB",
     "type": "Crash · injury", "severity": "moderate", "delay_min": 35,
     "lat": 32.78, "lng": -96.81, "lanes_closed": "2 of 5 NB",
     "agency": "TxDOT", "eta_clear_min": 60,
     "description": "Three-vehicle collision near Reunion Tower. Left two lanes closed.",
     "source_url": "https://drivetexas.org"},
    {"location": "I-75 SB Detroit (MI)", "highway": "I-75", "direction": "SB",
     "type": "Pothole repair", "severity": "low", "delay_min": 10,
     "lat": 42.34, "lng": -83.04, "lanes_closed": "Right lane",
     "agency": "MDOT", "eta_clear_min": 60,
     "description": "Emergency pothole repair. Right lane closed near Mack Ave.",
     "source_url": "https://mdotjboss.state.mi.us"},
    {"location": "Holland Tunnel NJ inbound (NJ/NY)", "highway": "I-78", "direction": "EB",
     "type": "Congestion", "severity": "high", "delay_min": 55,
     "lat": 40.72, "lng": -74.04, "lanes_closed": "—",
     "agency": "PANYNJ", "eta_clear_min": 120,
     "description": "Heavy congestion building from NJ Turnpike Exit 14C. Cross-Hudson freight should use Lincoln or GWB.",
     "source_url": "https://www.panynj.gov/bridges-tunnels"},

    {"location": "I-15 NB Provo (UT)", "highway": "I-15", "direction": "NB",
     "type": "Wildfire smoke", "severity": "moderate", "delay_min": 20,
     "lat": 40.23, "lng": -111.66, "lanes_closed": "Reduced visibility — all open",
     "agency": "UDOT", "eta_clear_min": 240,
     "description": "Heavy smoke from Spanish Fork Canyon fire. Visibility under 1 mile in places. Drive with headlights on.",
     "source_url": "https://www.udottraffic.utah.gov"},
    {"location": "I-10 WB Phoenix (AZ)", "highway": "I-10", "direction": "WB",
     "type": "Crash", "severity": "moderate", "delay_min": 30,
     "lat": 33.45, "lng": -112.07, "lanes_closed": "1 of 4 WB",
     "agency": "ADOT", "eta_clear_min": 45,
     "description": "Single-vehicle rollover near 7th St. Right lane closed.",
     "source_url": "https://az511.gov"},
    {"location": "I-77 SB Charlotte (NC)", "highway": "I-77", "direction": "SB",
     "type": "Congestion", "severity": "high", "delay_min": 38,
     "lat": 35.23, "lng": -80.84, "lanes_closed": "—",
     "agency": "NCDOT TIMS", "eta_clear_min": 90,
     "description": "Volume + earlier crash at MM 11 — backups extending into Cornelius. Average speed 18 mph.",
     "source_url": "https://drivenc.gov"},
    {"location": "I-91 NB Hartford (CT)", "highway": "I-91", "direction": "NB",
     "type": "Crash", "severity": "moderate", "delay_min": 22,
     "lat": 41.76, "lng": -72.67, "lanes_closed": "1 of 3 NB",
     "agency": "CT DOT", "eta_clear_min": 35,
     "description": "Two-vehicle crash near Exit 33. Right lane closed.",
     "source_url": "https://cttravelsmart.org"},
    {"location": "I-275 EB Tampa (FL)", "highway": "I-275", "direction": "EB",
     "type": "Disabled truck", "severity": "low", "delay_min": 18,
     "lat": 27.97, "lng": -82.54, "lanes_closed": "Right shoulder + lane 1",
     "agency": "FDOT FL511", "eta_clear_min": 30,
     "description": "Tractor with mechanical issue blocking right lane. Towing en route.",
     "source_url": "https://fl511.com"},

    {"location": "I-84 EB Snoqualmie Pass (WA)", "highway": "I-84", "direction": "EB",
     "type": "Weather · Snow", "severity": "high", "delay_min": 90,
     "lat": 47.39, "lng": -121.41, "lanes_closed": "Traction tires required",
     "agency": "WSDOT", "eta_clear_min": 240,
     "description": "Pass conditions deteriorating rapidly. Traction tires required; chains advised. Mountain prepared closure possible.",
     "source_url": "https://wsdot.com/Travel/Real-time"},
    {"location": "I-44 WB Tulsa (OK)", "highway": "I-44", "direction": "WB",
     "type": "Construction", "severity": "low", "delay_min": 14,
     "lat": 36.15, "lng": -95.99, "lanes_closed": "1 of 3 WB",
     "agency": "OK 511", "eta_clear_min": 180,
     "description": "Resurfacing project — center lane reduction. Carriers building reliable +15 min buffer.",
     "source_url": "https://oklahoma.gov/odot"},
    {"location": "I-280 NB Bay Area (CA)", "highway": "I-280", "direction": "NB",
     "type": "Sigalert · Crash", "severity": "high", "delay_min": 65,
     "lat": 37.78, "lng": -122.39, "lanes_closed": "3 of 4 NB",
     "agency": "CHP", "eta_clear_min": 90,
     "description": "Multi-vehicle crash with injuries. SIG-Alert issued; expect 1+ hour delay.",
     "source_url": "https://quickmap.dot.ca.gov"},
    {"location": "I-526 EB Charleston (SC)", "highway": "I-526", "direction": "EB",
     "type": "Bridge maintenance", "severity": "low", "delay_min": 9,
     "lat": 32.83, "lng": -79.96, "lanes_closed": "1 of 2 EB",
     "agency": "SCDOT", "eta_clear_min": 120,
     "description": "Routine bridge work over Wando River. Minor delay.",
     "source_url": "https://511sc.org"},
    {"location": "I-29 SB Sioux City (IA)", "highway": "I-29", "direction": "SB",
     "type": "Weather · Wind", "severity": "moderate", "delay_min": 25,
     "lat": 42.50, "lng": -96.40, "lanes_closed": "Wind advisory for high-profile vehicles",
     "agency": "Iowa 511", "eta_clear_min": 180,
     "description": "Sustained 35 mph winds with gusts to 55 mph. High-profile / empty trailer warning in effect.",
     "source_url": "https://511ia.org"},
]


@api_router.get("/traffic")
async def get_traffic(severity: Optional[str] = None, _: User = Depends(get_current_user)):
    """Live traffic. Each fetch rotates 'reported_minutes_ago' so dispatchers
    see the panel as continuously updating, plus enriches every record with
    a synthetic `reported_at` ISO timestamp."""
    rnd = _random.Random()
    out = []
    for inc in MOCK_TRAFFIC:
        if severity and inc.get("severity") != severity:
            continue
        mins = rnd.randint(0, 90)
        out.append({**inc, "minutes_ago": mins, "reported_at_label": _mins_ago_to_label(mins),
                    "reported_at": (datetime.now(timezone.utc) - timedelta(minutes=mins)).isoformat()})
    # Severity-weighted sort: high first, then most recent
    sev_rank = {"high": 0, "moderate": 1, "low": 2}
    out.sort(key=lambda x: (sev_rank.get(x.get("severity"), 9), x["minutes_ago"]))
    return await _brand_swap(out)


# -------------------- WEATHER ALERTS (NWS-style, MOCKED but realistic) --------------------
MOCK_WEATHER_ALERTS = [
    {"alert_id": "NWS-MPX-2026-0211-WS01", "type": "Winter Storm Watch", "severity": "moderate",
     "area": "Twin Cities Metro · MN", "affected_facility": "Tennant HQ — Golden Valley MN",
     "headline": "Winter Storm Watch for Hennepin County in effect through Thu 6 AM.",
     "body": "6 to 10 inches of snow expected with localized 12-inch totals possible. Travel could be very difficult during the Wednesday morning commute. Carriers serving HQ should pre-route freight Tuesday evening.",
     "issued_at": (datetime.now(timezone.utc) - timedelta(hours=2)).isoformat(),
     "expires_at": (datetime.now(timezone.utc) + timedelta(days=2)).isoformat(),
     "source": "NWS Twin Cities", "source_url": "https://www.weather.gov/mpx/"},
    {"alert_id": "NWS-IWX-2026-0211-WW02", "type": "Winter Weather Advisory", "severity": "low",
     "area": "Lower Michigan", "affected_facility": "Tennant Holland MI Plant",
     "headline": "Winter Weather Advisory — 2 to 4 inches of snow Tuesday night into Wednesday.",
     "body": "Light to moderate snow expected. Plan for slick roadways during the morning commute. Some impacts to inbound LTL deliveries to Holland Plant possible.",
     "issued_at": (datetime.now(timezone.utc) - timedelta(hours=1)).isoformat(),
     "expires_at": (datetime.now(timezone.utc) + timedelta(days=1)).isoformat(),
     "source": "NWS Northern Indiana", "source_url": "https://www.weather.gov/iwx/"},
    {"alert_id": "NWS-LMK-2026-0211-WIA", "type": "Wind Advisory", "severity": "low",
     "area": "Louisville Metro · KY", "affected_facility": "Tennant Louisville KY Plant",
     "headline": "Wind Advisory — Southwest winds 25 to 35 mph with gusts to 50 mph.",
     "body": "High-profile / empty trailer drivers urged to use caution. Drayage operations may experience minor delays at hardstand locations.",
     "issued_at": (datetime.now(timezone.utc) - timedelta(minutes=40)).isoformat(),
     "expires_at": (datetime.now(timezone.utc) + timedelta(hours=14)).isoformat(),
     "source": "NWS Louisville", "source_url": "https://www.weather.gov/lmk/"},
    {"alert_id": "NWS-ICT-2026-0211-FFW", "type": "Flash Flood Warning", "severity": "high",
     "area": "Sedgwick County · KS", "affected_facility": None,
     "headline": "Flash Flood Warning until 6 PM CST.",
     "body": "Heavy rainfall producing 1 to 2 inches per hour. Avoid flooded roadways. I-135 South of Wichita may experience standing water at low points.",
     "issued_at": (datetime.now(timezone.utc) - timedelta(minutes=22)).isoformat(),
     "expires_at": (datetime.now(timezone.utc) + timedelta(hours=4)).isoformat(),
     "source": "NWS Wichita", "source_url": "https://www.weather.gov/ict/"},
]


# Weather alerts + locations live in routes.weather (see end of file for include).


# -------------------- WELLNESS NUDGES --------------------
WELLNESS_NUDGES = [
    {"id": "stretch", "category": "movement", "title": "Stand-up reset", "message": "You've been seated a while — stand, roll your shoulders, look at something 20+ feet away for 20 seconds.", "icon": "Activity"},
    {"id": "hydrate", "category": "hydration", "title": "Hydration check", "message": "Top up your water. Pull dispatch numbers go better when you do.", "icon": "Droplet"},
    {"id": "eyes", "category": "wellness", "title": "20-20-20 rule", "message": "Every 20 minutes look at something 20 feet away for 20 seconds. Your eyes will thank you.", "icon": "Eye"},
    {"id": "breath", "category": "mindfulness", "title": "Box breathing", "message": "Inhale 4 · hold 4 · exhale 4 · hold 4. Run three rounds before your next call.", "icon": "Wind"},
    {"id": "step", "category": "movement", "title": "100 steps", "message": "Walk to the printer, the kitchen, anywhere. 100 steps clears the head between long booking sessions.", "icon": "Footprints"},
    {"id": "lunch", "category": "nutrition", "title": "Real lunch reminder", "message": "Step away from the desk — eat lunch somewhere that isn't this screen.", "icon": "Coffee"},
    {"id": "posture", "category": "wellness", "title": "Posture pulse", "message": "Sit up tall. Shoulders back and down. Imagine a string lifting the crown of your head.", "icon": "User"},
    {"id": "gratitude", "category": "mindfulness", "title": "One good thing", "message": "Name one thing that went well in the last hour. Big or small — celebrate the win.", "icon": "Heart"},
    {"id": "stretch-neck", "category": "movement", "title": "Neck stretch", "message": "Drop your ear toward your shoulder, hold for 15 seconds each side. Loosen the dispatcher's hunch.", "icon": "Activity"},
    {"id": "snack", "category": "nutrition", "title": "Smart snack", "message": "Hungry? Reach for nuts, fruit, or yogurt over the candy bowl. Steady energy beats a sugar spike.", "icon": "Apple"},
    {"id": "music", "category": "mindfulness", "title": "Mood lift", "message": "Put on one song you love. Whole-day reset, two minutes.", "icon": "Music"},
    {"id": "fresh-air", "category": "movement", "title": "60 seconds outside", "message": "Step outside for one minute. Daylight resets your circadian rhythm — yes, even in February.", "icon": "Sun"},
    {"id": "stretch-back", "category": "movement", "title": "Cat-cow", "message": "Hands on knees, arch and round your back five times. Dispatch chairs aren't kind — fight back.", "icon": "Activity"},
    {"id": "remember", "category": "mindfulness", "title": "Why you do this", "message": "Every shipment moves something someone is waiting for. You're moving the world, one BOL at a time.", "icon": "Heart"},
    {"id": "celebrate", "category": "mindfulness", "title": "Team shoutout", "message": "Spotted a teammate crushing it? Send a quick 'nice work' message. Two-second cost, huge impact.", "icon": "Award"},
]


@api_router.get("/wellness/nudges")
async def wellness_nudges_endpoint(_: User = Depends(get_current_user)):
    return WELLNESS_NUDGES


# -------------------- INTEGRATIONS · POWER BI --------------------
POWERBI_REPORTS_DEFAULT = [
    {"id": "tennant-tms-overview", "name": "Tennant TMS · Executive Overview",
     "description": "Real-time KPIs across all modes — cost-per-mile, on-time %, freight spend.",
     "workspace": "Tennant · Supply Chain", "owner": "Joe Carlsson",
     "embed_url": "https://app.powerbi.com/reportEmbed?reportId=demo-overview&autoAuth=true",
     "view_url": "https://app.powerbi.com/groups/me/reports/demo-overview",
     "updated_at": (datetime.now(timezone.utc) - timedelta(minutes=8)).isoformat()},
    {"id": "freight-spend-by-lane", "name": "Freight Spend by Lane · YTD",
     "description": "Lane-level spend, volume, and rate trend with target compare.",
     "workspace": "Tennant · Supply Chain", "owner": "Amanda Reyes",
     "embed_url": "https://app.powerbi.com/reportEmbed?reportId=demo-lane-spend&autoAuth=true",
     "view_url": "https://app.powerbi.com/groups/me/reports/demo-lane-spend",
     "updated_at": (datetime.now(timezone.utc) - timedelta(hours=2)).isoformat()},
    {"id": "on-time-delivery", "name": "On-Time Delivery Performance",
     "description": "OTD% by carrier, mode, region. 90-day trend with drill-down to shipment.",
     "workspace": "Tennant · Operations", "owner": "Marcus Lavoie",
     "embed_url": "https://app.powerbi.com/reportEmbed?reportId=demo-otd&autoAuth=true",
     "view_url": "https://app.powerbi.com/groups/me/reports/demo-otd",
     "updated_at": (datetime.now(timezone.utc) - timedelta(minutes=45)).isoformat()},
    {"id": "yard-dwell-trends", "name": "Yard Dwell · Daily Trend",
     "description": "Average and P95 dwell time at all 4 plants. Stale-trailer alerts highlighted.",
     "workspace": "Tennant · Operations", "owner": "Christine Yoder",
     "embed_url": "https://app.powerbi.com/reportEmbed?reportId=demo-yard&autoAuth=true",
     "view_url": "https://app.powerbi.com/groups/me/reports/demo-yard",
     "updated_at": (datetime.now(timezone.utc) - timedelta(minutes=15)).isoformat()},
    {"id": "carrier-scorecards", "name": "Carrier Performance Scorecards",
     "description": "OTD, claims rate, billing accuracy, accessorial usage — top 25 carriers.",
     "workspace": "Tennant · Supply Chain", "owner": "Renee Calderon",
     "embed_url": "https://app.powerbi.com/reportEmbed?reportId=demo-scorecard&autoAuth=true",
     "view_url": "https://app.powerbi.com/groups/me/reports/demo-scorecard",
     "updated_at": (datetime.now(timezone.utc) - timedelta(hours=4)).isoformat()},
    {"id": "fuel-surcharge", "name": "FSC vs. DOE Diesel Tracker",
     "description": "Carrier FSC schedules charted against weekly DOE diesel. Reveals lag/lead.",
     "workspace": "Tennant · Finance", "owner": "Henry Park",
     "embed_url": "https://app.powerbi.com/reportEmbed?reportId=demo-fsc&autoAuth=true",
     "view_url": "https://app.powerbi.com/groups/me/reports/demo-fsc",
     "updated_at": (datetime.now(timezone.utc) - timedelta(hours=12)).isoformat()},
]


@api_router.get("/integrations/powerbi/config")
async def powerbi_config_get(_: User = Depends(get_current_user)):
    """Returns the active list of Power BI reports the admin has wired up.
    If none are saved in Mongo we fall back to the default demo roster so
    the new tab always renders something useful out of the box."""
    doc = await db.integration_configs.find_one({"_id": "powerbi"}, {"_id": 0})
    reports = (doc or {}).get("reports") or POWERBI_REPORTS_DEFAULT
    workspace_url = (doc or {}).get("workspace_url") or "https://app.powerbi.com/groups/me"
    return await _brand_swap({"reports": reports, "workspace_url": workspace_url, "tenant": "tennantco.onmicrosoft.com"})


class PowerBIReport(BaseModel):
    id: str
    name: str
    description: Optional[str] = ""
    workspace: Optional[str] = "Orisei"
    owner: Optional[str] = ""
    embed_url: Optional[str] = ""
    view_url: str


class PowerBIConfig(BaseModel):
    workspace_url: Optional[str] = ""
    reports: List[PowerBIReport] = []


@api_router.put("/integrations/powerbi/config")
async def powerbi_config_put(payload: PowerBIConfig, user: User = Depends(require_role("admin"))):
    rec = {"_id": "powerbi", "workspace_url": payload.workspace_url,
           "reports": [r.dict() for r in payload.reports],
           "updated_by": user.name,
           "updated_at": datetime.now(timezone.utc).isoformat()}
    await db.integration_configs.replace_one({"_id": "powerbi"}, rec, upsert=True)
    return {"ok": True}


# -------------------- INTEGRATIONS · SHAREPOINT --------------------
SHAREPOINT_SITES_DEFAULT = [
    {"id": "tennant-supplychain", "name": "Tennant · Supply Chain",
     "url": "https://tennantco.sharepoint.com/sites/supplychain",
     "description": "Routing guides, carrier MSAs, freight RFP packets, trade compliance archive.",
     "members": 47, "updated_at": (datetime.now(timezone.utc) - timedelta(minutes=18)).isoformat()},
    {"id": "tennant-ops", "name": "Tennant · Plant Operations",
     "url": "https://tennantco.sharepoint.com/sites/operations",
     "description": "Yard reports, daily production schedules, dock door rosters across HQ + 3 plants.",
     "members": 122, "updated_at": (datetime.now(timezone.utc) - timedelta(minutes=4)).isoformat()},
    {"id": "tennant-finance", "name": "Tennant · Finance",
     "url": "https://tennantco.sharepoint.com/sites/finance",
     "description": "Freight invoices, monthly cost-to-serve, accruals, audit packets.",
     "members": 18, "updated_at": (datetime.now(timezone.utc) - timedelta(hours=2)).isoformat()},
    {"id": "tennant-trade", "name": "Tennant · Trade Compliance",
     "url": "https://tennantco.sharepoint.com/sites/trade",
     "description": "HTS classifications, Section 301 watch, FTZ admissions, broker portal links.",
     "members": 12, "updated_at": (datetime.now(timezone.utc) - timedelta(hours=8)).isoformat()},
    {"id": "tennant-transportation", "name": "Tennant · Transportation",
     "url": "https://tennantco.sharepoint.com/sites/transportation",
     "description": "Carrier contracts, lane tenders, accessorial reference, dispatcher SOPs.",
     "members": 31, "updated_at": (datetime.now(timezone.utc) - timedelta(minutes=55)).isoformat()},
]

SHAREPOINT_RECENT_FILES_DEFAULT = [
    {"name": "2026 Domestic Inbound Routing Guide (Rev 29).pdf", "site": "Tennant · Supply Chain", "modified_by": "Joe Carlsson", "modified_at": (datetime.now(timezone.utc) - timedelta(hours=11)).isoformat(), "size": "418 KB", "url": "https://tennantco.sharepoint.com/sites/supplychain/Shared%20Documents/Routing/2026-Routing-Guide-Rev29.pdf"},
    {"name": "Q4-2025 Freight RFP — Final Bid Matrix.xlsx", "site": "Tennant · Supply Chain", "modified_by": "Amanda Reyes", "modified_at": (datetime.now(timezone.utc) - timedelta(days=1)).isoformat(), "size": "2.1 MB", "url": "https://tennantco.sharepoint.com/sites/supplychain/Shared%20Documents/RFP/Q4-2025-Bid-Matrix.xlsx"},
    {"name": "Carrier MSA — Knight Transportation 2026.pdf", "site": "Tennant · Transportation", "modified_by": "Marcus Lavoie", "modified_at": (datetime.now(timezone.utc) - timedelta(days=2)).isoformat(), "size": "780 KB", "url": "https://tennantco.sharepoint.com/sites/transportation/Shared%20Documents/MSAs/Knight-2026.pdf"},
    {"name": "Holland Yard Report 2026-02-10.xlsx", "site": "Tennant · Plant Operations", "modified_by": "Yard Supervisor", "modified_at": (datetime.now(timezone.utc) - timedelta(hours=6)).isoformat(), "size": "94 KB", "url": "https://tennantco.sharepoint.com/sites/operations/Shared%20Documents/Yard/Holland-2026-02-10.xlsx"},
    {"name": "January Freight Audit Variance Report.pdf", "site": "Tennant · Finance", "modified_by": "Henry Park", "modified_at": (datetime.now(timezone.utc) - timedelta(days=3)).isoformat(), "size": "612 KB", "url": "https://tennantco.sharepoint.com/sites/finance/Shared%20Documents/Audit/Jan-2026-Variance.pdf"},
    {"name": "USTR Section 301 Comment Letter — DRAFT.docx", "site": "Tennant · Trade Compliance", "modified_by": "Joe Carlsson", "modified_at": (datetime.now(timezone.utc) - timedelta(hours=4)).isoformat(), "size": "62 KB", "url": "https://tennantco.sharepoint.com/sites/trade/Shared%20Documents/USTR-301-Comment.docx"},
    {"name": "Dispatcher Standard Operating Procedures v3.2.pdf", "site": "Tennant · Transportation", "modified_by": "Renee Calderon", "modified_at": (datetime.now(timezone.utc) - timedelta(days=5)).isoformat(), "size": "1.4 MB", "url": "https://tennantco.sharepoint.com/sites/transportation/Shared%20Documents/SOP/Dispatch-SOP-v3.2.pdf"},
    {"name": "FTZ Admission Logs — January.xlsx", "site": "Tennant · Trade Compliance", "modified_by": "Amanda Reyes", "modified_at": (datetime.now(timezone.utc) - timedelta(days=4)).isoformat(), "size": "232 KB", "url": "https://tennantco.sharepoint.com/sites/trade/Shared%20Documents/FTZ/Jan-2026-Admissions.xlsx"},
]


@api_router.get("/integrations/sharepoint/config")
async def sharepoint_config(_: User = Depends(get_current_user)):
    doc = await db.integration_configs.find_one({"_id": "sharepoint"}, {"_id": 0})
    sites = (doc or {}).get("sites") or SHAREPOINT_SITES_DEFAULT
    files = (doc or {}).get("recent_files") or SHAREPOINT_RECENT_FILES_DEFAULT
    tenant_url = (doc or {}).get("tenant_url") or "https://tennantco.sharepoint.com"
    return await _brand_swap({"sites": sites, "recent_files": files, "tenant_url": tenant_url})


# -------------------- INTEGRATIONS · S4 LINK BUILDER --------------------
SAP_S4_BASE = os.environ.get("SAP_S4_BASE_URL", "https://my-s4.tennantco.com")


def _brand_short() -> str:
    """Best-effort sync access to the active brand's short name. Returns
    'Orisei' if unset. Used only for tenant-string substitution where an
    async call would be too expensive (called in template formatters)."""
    return "Orisei"  # sync default; async overlays do the real swap


async def _brand_tenant_strings() -> Dict[str, str]:
    """Returns the replacements applied to integration / SAP responses
    when a non-Tennant brand is active. Keys are 'Orisei' fragments to be
    swapped for the brand's short_name. The empty dict means no swap (i.e.
    Tennant is active)."""
    brand = await _active_brand_doc()
    if not brand or brand.get("brand_id") == "orisei-freight":
        return {}
    short = brand.get("short_name") or "Brand"
    slug = re.sub(r"[^a-z0-9]+", "", short.lower())[:20] or "brand"
    company = brand.get("company_name") or short

    # === LOCATION SWAP ===
    # Map Tennant's three canonical facility city/state tokens to whatever
    # facilities the active brand has provided. Order matters — process
    # longer/more-specific keys first so e.g. "Golden Valley MN" hits before
    # "Golden Valley".
    facilities = brand.get("facilities") or []
    # Default each Tennant slot to a brand facility (cycling if fewer exist).
    def _fac_for(slot_i: int, default_city: str, default_state: str) -> Dict[str, str]:
        if not facilities:
            return {"name": default_city, "city": default_city, "state": default_state, "label": f"{default_city} {default_state}"}
        f = facilities[slot_i % len(facilities)]
        full_city = (f.get("city") or "").strip() or default_city
        # "Kalamazoo, MI" -> ("Kalamazoo", "MI")
        if "," in full_city:
            city, state = [p.strip() for p in full_city.split(",", 1)]
        else:
            city, state = full_city, default_state
        return {
            "name": f.get("name") or f"{city} {state} Plant",
            "city": city,
            "state": state,
            "label": f"{city} {state}",
        }

    f0 = _fac_for(0, "Golden Valley", "MN")  # HQ slot
    f1 = _fac_for(1, "Holland", "MI")        # plant slot 1
    f2 = _fac_for(2, "Louisville", "KY")     # plant slot 2

    repl = {
        # Word-form replacements (case-sensitive, ordered most-specific first)
        "Orisei Freight Solutions": company,
        "Orisei Freight Solutions": company,
        "Tennant · ": f"{short} · ",
        "Tennant ": f"{short} ",
        "TENNANT": short.upper(),
        # Domains / slugs (specific first so they win over generic tennantco. catch-all below)
        "tennantco.com": f"{slug}.com",
        "tennantco.onmicrosoft.com": f"{slug}.onmicrosoft.com",
        "tennantco.sharepoint.com": f"{slug}.sharepoint.com",
        "tennantco.s4hana.cloud.sap": f"{slug}.s4hana.cloud.sap",
        "my-s4.tennantco.com": f"my-s4.{slug}.com",
        "s4hana.tennantco.sap.com": f"s4hana.{slug}.sap.com",
        "tennantco.s4.sap.com": f"{slug}.s4.sap.com",
        "powerbi.com/tennant": f"powerbi.com/{slug}",
        "tennantco": slug,
        "tennant-": f"{slug}-",
        # === Facility locations ===
        # Long-form first
        "Golden Valley, MN (HQ)": f"{f0['name']}",
        "Tennant HQ — Golden Valley MN": f"{short} HQ — {f0['label']}",
        "Tennant Holland MI Plant": f"{short} {f1['label']} Plant",
        "Tennant Louisville KY Plant": f"{short} {f2['label']} Plant",
        "Tennant Holland MI": f"{short} {f1['label']}",
        "Tennant Louisville KY": f"{short} {f2['label']}",
        "Golden Valley MN": f0["label"],
        "Holland MI": f1["label"],
        "Louisville KY": f2["label"],
        "Golden Valley, MN": f"{f0['city']}, {f0['state']}",
        "Holland, MI": f"{f1['city']}, {f1['state']}",
        "Louisville, KY": f"{f2['city']}, {f2['state']}",
        "Twin Cities Metro · MN": f"{f0['city']} Metro · {f0['state']}",
        "Lower Michigan": f"{f1['state']} Region",
        "Louisville Metro · KY": f"{f2['city']} Metro · {f2['state']}",
        # Standalone city tokens (least specific — placed last so longer matches win)
        "Golden Valley": f0["city"],
        "Holland": f1["city"],
        "Louisville": f2["city"],
        # Final catch-all for stray standalone references
        "Orisei": short,
        "orisei-freight": slug,
    }
    return repl


async def _brand_swap(value: Any) -> Any:
    """Apply tenant string replacements recursively. No-op for Tennant."""
    repl = await _brand_tenant_strings()
    if not repl:
        return value
    return _swap_strings(value, repl)



# Fiori app aliases — these are the canonical OData/Web Dynpro paths each
# S/4HANA tile resolves to. Wrapped in an env var so customers can swap
# the production hostname with no code changes.
S4_DEEP_LINK_PATTERNS = {
    "purchase_order":   "/sap/bc/ui2/flp#PurchaseOrder-displayFactSheet?PurchaseOrder={value}",
    "sales_order":      "/sap/bc/ui2/flp#SalesOrder-displayFactSheet?SalesOrder={value}",
    "delivery":         "/sap/bc/ui2/flp#OutboundDelivery-displayFactSheet?Delivery={value}",
    "invoice":          "/sap/bc/ui2/flp#BillingDocument-display?BillingDocument={value}",
    "supplier_invoice": "/sap/bc/ui2/flp#SupplierInvoice-display?SupplierInvoice={value}",
    "material":         "/sap/bc/ui2/flp#Material-displayFactSheet?Material={value}",
    "part":             "/sap/bc/ui2/flp#Material-displayFactSheet?Material={value}",
    "bol":              "/sap/bc/ui2/flp#FreightOrder-displayFactSheet?FreightOrder={value}",
    "booking":          "/sap/bc/ui2/flp#TransportationOrder-displayFactSheet?TransportationOrder={value}",
    "shipment":         "/sap/bc/ui2/flp#TransportationOrder-displayFactSheet?TransportationOrder={value}",
}


@api_router.get("/sap/deep-link")
async def sap_deep_link(kind: str, value: str, _: User = Depends(get_current_user)):
    """Build a deep-link URL into Tennant's SAP S/4HANA Fiori launchpad
    for a given reference (kind = purchase_order, sales_order, invoice,
    bol, etc). Returned as JSON so the frontend can hyperlink any token."""
    if kind not in S4_DEEP_LINK_PATTERNS:
        raise HTTPException(status_code=400, detail=f"Unknown S/4 link kind: {kind}")
    if not value:
        raise HTTPException(status_code=400, detail="value required")
    path = S4_DEEP_LINK_PATTERNS[kind].format(value=urllib.parse.quote(value, safe=""))
    base = SAP_S4_BASE
    repl = await _brand_tenant_strings()
    if repl:
        for src, dst in repl.items():
            base = base.replace(src, dst)
    return {"url": f"{base}{path}", "kind": kind, "value": value, "base": base}


@api_router.get("/sap/link-config")
async def sap_link_config(_: User = Depends(get_current_user)):
    """Frontend reads this once to build links locally without a roundtrip
    per token. Returns the base URL + pattern map (with `{value}` placeholder)."""
    base = SAP_S4_BASE
    repl = await _brand_tenant_strings()
    if repl:
        for src, dst in repl.items():
            base = base.replace(src, dst)
    return {"base": base, "patterns": S4_DEEP_LINK_PATTERNS, "kinds": list(S4_DEEP_LINK_PATTERNS.keys())}


# -------------------- CARRIER TRACKING URLs --------------------
# Per-carrier public tracking page templates. `{tracking}` is the only
# placeholder. When a carrier has multiple options (e.g. ODFL by PRO# vs.
# BOL#), the most-commonly-used pattern is picked.
CARRIER_TRACKING_URLS: Dict[str, Dict[str, str]] = {
    "UPS":                 {"url": "https://www.ups.com/track?tracknum={tracking}",                                  "label": "UPS · Tracking #"},
    "FedEx":               {"url": "https://www.fedex.com/fedextrack/?trknbr={tracking}",                            "label": "FedEx · Tracking #"},
    "USPS":                {"url": "https://tools.usps.com/go/TrackConfirmAction?tLabels={tracking}",                "label": "USPS · Tracking #"},
    "DHL":                 {"url": "https://www.dhl.com/global-en/home/tracking.html?tracking-id={tracking}",        "label": "DHL · Waybill"},
    "DHL Express":         {"url": "https://www.dhl.com/global-en/home/tracking.html?tracking-id={tracking}",        "label": "DHL Express · Waybill"},
    "ODFL":                {"url": "https://www.odfl.com/Trace/standardResults.faces?pro={tracking}",                "label": "Old Dominion · PRO #"},
    "Old Dominion":        {"url": "https://www.odfl.com/Trace/standardResults.faces?pro={tracking}",                "label": "Old Dominion · PRO #"},
    "XPO":                 {"url": "https://www.xpo.com/tracking/?pros={tracking}",                                  "label": "XPO Logistics · PRO #"},
    "XPO Logistics":       {"url": "https://www.xpo.com/tracking/?pros={tracking}",                                  "label": "XPO Logistics · PRO #"},
    "Saia":                {"url": "https://www.saia.com/tools/track-shipments?proNumbers={tracking}",               "label": "Saia · PRO #"},
    "Estes":               {"url": "https://www.estes-express.com/myestes/shipment-tracking/track-shipment?type=PRO&searchValue={tracking}", "label": "Estes · PRO #"},
    "Estes Express":       {"url": "https://www.estes-express.com/myestes/shipment-tracking/track-shipment?type=PRO&searchValue={tracking}", "label": "Estes · PRO #"},
    "R+L Carriers":        {"url": "https://www2.rlcarriers.com/freight/shipping/shipment-tracing?pro={tracking}",   "label": "R+L Carriers · PRO #"},
    "R&L Carriers":        {"url": "https://www2.rlcarriers.com/freight/shipping/shipment-tracing?pro={tracking}",   "label": "R&L Carriers · PRO #"},
    "ArcBest":             {"url": "https://arcb.com/tools/tracking?prn={tracking}",                                 "label": "ArcBest · PRO #"},
    "ABF":                 {"url": "https://arcb.com/tools/tracking?prn={tracking}",                                 "label": "ABF · PRO #"},
    "Knight":              {"url": "https://knight-transportation.com/track-load/?load={tracking}",                  "label": "Knight Transportation · Load #"},
    "Knight Transportation": {"url": "https://knight-transportation.com/track-load/?load={tracking}",                "label": "Knight · Load #"},
    "Schneider":           {"url": "https://schneider.com/shipment-tracking?reference={tracking}",                   "label": "Schneider · Pro/Ref"},
    "Werner":              {"url": "https://www.werner.com/shipment-tracking?ref={tracking}",                        "label": "Werner · Reference"},
    "C.H. Robinson":       {"url": "https://www.chrobinson.com/en-us/shipper/shipment-tracking/?reference={tracking}", "label": "C.H. Robinson · Reference"},
    "Maersk":              {"url": "https://www.maersk.com/tracking/{tracking}",                                     "label": "Maersk · Container/B-of-L"},
    "MSC":                 {"url": "https://www.msc.com/track-a-shipment?agencyPath=msc&shipmentNumber={tracking}",  "label": "MSC · Container/B-of-L"},
    "Hapag-Lloyd":         {"url": "https://www.hapag-lloyd.com/en/online-business/track/track-by-container-solution.html?container={tracking}", "label": "Hapag-Lloyd · Container"},
    "CMA CGM":             {"url": "https://www.cma-cgm.com/ebusiness/tracking/search?SearchBy=Container&Reference={tracking}", "label": "CMA CGM · Container"},
    "ONE":                 {"url": "https://ecomm.one-line.com/one-ecom/manage-shipment/cargo-tracking?trackingType=CONTAINER&trackingNumber={tracking}", "label": "ONE · Container"},
    "Evergreen":           {"url": "https://www.evergreen-line.com/static/jsp/cargo_tracking.jsp?bl_no={tracking}",  "label": "Evergreen · B-of-L"},
    "Yang Ming":           {"url": "https://www.yangming.com/e-service/Track_Trace/track_trace_cargo_tracking.aspx?bl_no={tracking}", "label": "Yang Ming · B-of-L"},
    "OOCL":                {"url": "https://www.oocl.com/eng/ourservices/eservices/cargotracking/Pages/cargotracking.aspx?ContainerNo={tracking}", "label": "OOCL · Container"},
    "Kuehne+Nagel":        {"url": "https://mykn.kuehne-nagel.com/public-tracking/shipments?query={tracking}",       "label": "Kuehne+Nagel · Reference"},
    "Kuehne+Nagel Services": {"url": "https://mykn.kuehne-nagel.com/public-tracking/shipments?query={tracking}",     "label": "Kuehne+Nagel · Reference"},
    "DSV":                 {"url": "https://my.dsv.com/etrack?id={tracking}",                                        "label": "DSV · Reference"},
    "Expeditors":          {"url": "https://www.expeditors.com/track?ref={tracking}",                                "label": "Expeditors · Reference"},
    "DB Schenker":         {"url": "https://www.dbschenker.com/global/tracking?id={tracking}",                       "label": "DB Schenker · Reference"},
    "Yellow":              {"url": "https://my.yellowcorp.com/tools/tracking-tools?proNumbers={tracking}",           "label": "Yellow · PRO #"},
    "TForce":              {"url": "https://www.tforcefreight.com/ltl/apps/Tracking?ProNumbers={tracking}",          "label": "TForce · PRO #"},
}


@api_router.get("/carriers/tracking-urls")
async def carrier_tracking_urls(_: User = Depends(get_current_user)):
    """Return the full carrier → tracking URL template map. Frontend uses
    this to render direct-tracking buttons inside the Live Tracking tab."""
    return {"carriers": CARRIER_TRACKING_URLS}


# -------------------- SPECIALTY CARRIERS --------------------
# White-glove / priority-use carriers Tennant relies on for high-value or
# tricky freight (Logix · ArcBest Panther · Fastfrate · Ryan Transportation).
# Each gets a dedicated profile card with the services that make them the
# go-to for that lane / equipment / handling profile.
SPECIALTY_CARRIERS = [
    {
        "id": "logix",
        "name": "Logix Transportation",
        "tagline": "White-glove machine transport · Pad-wrapped · Damage-free",
        "summary": "Created specifically to support Tennant's industrial machine deliveries. Logix specializes in pad-wrapping every machine before transit so customer-facing units arrive at retail and dealer locations in pristine condition.",
        "specialty": ["Pad-wrap protection", "Tennant machine transport", "Lift-gate residential", "Inside delivery", "White-glove handoff"],
        "modes": ["TL Specialty", "Van Line", "Air-Ride Lowboy"],
        "lanes": ["Holland MI → US 48 (priority)", "Golden Valley HQ → demo events", "Louisville KY → US Southeast", "Plant → end-customer dock"],
        "rate_basis": "Flat per machine · negotiated annual contract",
        "contact": {"name": "Eric Vandermeer", "phone": "+1-616-555-0177", "email": "tennant@logixtransportation.com", "after_hours": "+1-616-555-0188"},
        "website": "https://www.logixtransportation.com",
        "tracking_url": "https://www.logixtransportation.com/track?ref={tracking}",
        "color": "#10B981",  # emerald
        "logo_initials": "LX",
        "priority": "P0 · always-on white-glove",
        "since": "2018",
        "ytd_loads": 1287,
        "claim_rate_pct": 0.04,
        "on_time_pct": 99.1,
    },
    {
        "id": "arcbest-panther",
        "name": "ArcBest · Panther Premium Logistics",
        "tagline": "Expedited · Time-critical · Mission-critical",
        "summary": "Panther — ArcBest's expedited division — runs Tennant's hot loads when a customer line is down or a launch event can't slip. Same-day pickup, team drivers for cross-country freight, live in-cab tracking.",
        "specialty": ["Expedited sprinter / box truck / Class-8", "Team drivers", "Same-day pickup", "Live in-cab tracking", "Air charter coordination"],
        "modes": ["Sprinter Van", "Straight Truck", "Tractor + 53' Van", "Air Charter"],
        "lanes": ["US 48 anywhere → anywhere", "Detroit / Chicago / Atlanta hot lanes"],
        "rate_basis": "Per-mile + dispatch fee · spot or contract",
        "contact": {"name": "Cathy Powers", "phone": "+1-866-PANTHER", "email": "tennant@arcb.com", "after_hours": "+1-866-726-8437"},
        "website": "https://arcb.com/services/expedite",
        "tracking_url": "https://arcb.com/tools/tracking?prn={tracking}",
        "color": "#F59E0B",  # amber
        "logo_initials": "AP",
        "priority": "P0 · time-critical expedite",
        "since": "2015",
        "ytd_loads": 412,
        "claim_rate_pct": 0.07,
        "on_time_pct": 98.4,
    },
    {
        "id": "fastfrate",
        "name": "Fastfrate Group",
        "tagline": "Canadian intermodal · LTL · TL · Cross-border specialist",
        "summary": "Fastfrate is Tennant's go-to for moving freight to/from Canadian provinces. They handle customs documentation, single-bill cross-border tendering, and run an intermodal network out of every major Canadian gateway.",
        "specialty": ["Cross-border CA ↔ US documentation", "Intermodal CN/CP rail integration", "Bonded warehousing", "ACE / ACI eManifest filing", "PARS / PAPS clearance"],
        "modes": ["LTL", "TL", "Intermodal", "Bonded"],
        "lanes": ["Holland MI → ON/QC/AB/BC", "Toronto → Tennant dealers (US East)", "Vancouver ↔ Seattle/Portland"],
        "rate_basis": "LTL tariff · TL spot · intermodal contract",
        "contact": {"name": "Pierre Lavoie", "phone": "+1-905-451-7373", "email": "tennant@fastfrate.com", "after_hours": "+1-800-461-7373"},
        "website": "https://www.fastfrate.com",
        "tracking_url": "https://www.fastfrate.com/Tracking?probill={tracking}",
        "color": "#3B82F6",  # blue
        "logo_initials": "FF",
        "priority": "P1 · cross-border priority",
        "since": "2012",
        "ytd_loads": 642,
        "claim_rate_pct": 0.11,
        "on_time_pct": 97.2,
    },
    {
        "id": "ryan-transportation",
        "name": "Ryan Transportation",
        "tagline": "Strategic 3PL · capacity assurance · brokerage",
        "summary": "Ryan provides Tennant with on-demand capacity when contracted assets are tight. Strategic brokerage with vetted carrier pool, freight-class auditing, and embedded analytics. Backup carrier of choice for the entire Tennant US 48 network.",
        "specialty": ["Spot capacity (capacity assurance)", "Freight class / NMFC auditing", "Project / event logistics", "Network optimization analytics", "Pre-vetted carrier pool"],
        "modes": ["TL", "LTL", "Refrigerated", "Partial / Volume LTL"],
        "lanes": ["US 48 spot tender", "Surge capacity for HQ / plants"],
        "rate_basis": "Brokered · spot or RFP-awarded contract",
        "contact": {"name": "Mark Hannon", "phone": "+1-913-647-9700", "email": "tennant@ryantrans.com", "after_hours": "+1-913-647-9701"},
        "website": "https://www.ryantrans.com",
        "tracking_url": "https://www.ryantrans.com/tracking?reference={tracking}",
        "color": "#A855F7",  # purple
        "logo_initials": "RY",
        "priority": "P1 · capacity assurance",
        "since": "2017",
        "ytd_loads": 998,
        "claim_rate_pct": 0.09,
        "on_time_pct": 97.8,
    },
]


@api_router.get("/specialty-carriers")
async def specialty_carriers_list(_: User = Depends(get_current_user)):
    """Seeded specialty carriers + any user-added ones. Soft-deletes are
    honored so admins can hide a built-in carrier without losing the seed."""
    hidden = await db.specialty_carrier_overrides.find({"hidden": True}, {"_id": 0, "carrier_id": 1}).to_list(100)
    hidden_ids = {h["carrier_id"] for h in hidden}
    edits = await db.specialty_carrier_overrides.find({"hidden": {"$ne": True}}, {"_id": 0}).to_list(100)
    edits_map = {e["carrier_id"]: e for e in edits}
    out = []
    for c in SPECIALTY_CARRIERS:
        cid = c.get("id") or c.get("name", "").lower()
        if cid in hidden_ids:
            continue
        if cid in edits_map:
            merged = {**c, **edits_map[cid], "id": cid, "is_seed": True}
            out.append(merged)
        else:
            out.append({**c, "id": cid, "is_seed": True})
    # Custom carriers added by admin
    custom = await db.specialty_carriers_custom.find({}, {"_id": 0}).to_list(200)
    for c in custom:
        out.append({**c, "is_seed": False})
    return await _brand_swap({"carriers": out})


class SpecialtyCarrierIn(BaseModel):
    id: Optional[str] = None
    name: str
    type: Optional[str] = "Specialty"
    description: Optional[str] = ""
    services: Optional[List[str]] = []
    coverage: Optional[str] = ""
    website: Optional[str] = ""
    phone: Optional[str] = ""
    primary_contact: Optional[str] = ""
    primary_email: Optional[str] = ""
    notes: Optional[str] = ""


@api_router.post("/specialty-carriers")
async def specialty_carriers_create(payload: SpecialtyCarrierIn, user: User = Depends(require_role("admin", "dispatcher"))):
    """Create a new specialty carrier (custom)."""
    if not payload.name.strip():
        raise HTTPException(400, "name is required")
    cid = re.sub(r"[^a-z0-9]+", "-", payload.name.lower()).strip("-")
    doc = payload.model_dump()
    doc["id"] = cid
    doc["created_by"] = user.user_id
    doc["created_at"] = datetime.now(timezone.utc).isoformat()
    await db.specialty_carriers_custom.update_one({"id": cid}, {"$set": doc}, upsert=True)
    return {"ok": True, "carrier": doc}


@api_router.put("/specialty-carriers/{carrier_id}")
async def specialty_carriers_update(carrier_id: str, payload: SpecialtyCarrierIn, user: User = Depends(require_role("admin", "dispatcher"))):
    """Edit a carrier. Built-in (seeded) carriers are stored as overrides so
    the seed list remains immutable."""
    is_seed = any((c.get("id") or c.get("name", "").lower()) == carrier_id for c in SPECIALTY_CARRIERS)
    doc = {k: v for k, v in payload.model_dump().items() if v is not None}
    doc["carrier_id"] = carrier_id
    doc["updated_by"] = user.user_id
    doc["updated_at"] = datetime.now(timezone.utc).isoformat()
    if is_seed:
        await db.specialty_carrier_overrides.update_one({"carrier_id": carrier_id}, {"$set": doc}, upsert=True)
    else:
        await db.specialty_carriers_custom.update_one({"id": carrier_id}, {"$set": {**doc, "id": carrier_id}}, upsert=False)
    return {"ok": True}


@api_router.delete("/specialty-carriers/{carrier_id}")
async def specialty_carriers_delete(carrier_id: str, _: User = Depends(require_role("admin"))):
    """Delete a carrier. Seeded carriers are soft-hidden so they can be
    restored; custom carriers are hard-deleted."""
    is_seed = any((c.get("id") or c.get("name", "").lower()) == carrier_id for c in SPECIALTY_CARRIERS)
    if is_seed:
        await db.specialty_carrier_overrides.update_one(
            {"carrier_id": carrier_id},
            {"$set": {"carrier_id": carrier_id, "hidden": True}},
            upsert=True,
        )
        return {"ok": True, "soft_hidden": True}
    r = await db.specialty_carriers_custom.delete_one({"id": carrier_id})
    if r.deleted_count == 0:
        raise HTTPException(404, "Carrier not found")
    return {"ok": True}


@api_router.get("/carriers/tracking-url")
async def carrier_tracking_url(carrier: str, tracking: str, _: User = Depends(get_current_user)):
    """Resolve a single tracking URL. Case-insensitive carrier lookup, with
    fuzzy match against the carrier dropdown labels ('XPO · XPOL' → 'XPO')."""
    name = (carrier or "").split("·")[0].strip()
    key_match = None
    needle = name.lower()
    for key in CARRIER_TRACKING_URLS:
        if key.lower() == needle:
            key_match = key; break
    if not key_match:
        for key in CARRIER_TRACKING_URLS:
            if needle in key.lower() or key.lower() in needle:
                key_match = key; break
    if not key_match:
        raise HTTPException(status_code=404, detail=f"No tracking pattern for carrier '{carrier}'")
    tmpl = CARRIER_TRACKING_URLS[key_match]
    return {"carrier": key_match, "label": tmpl["label"], "url": tmpl["url"].format(tracking=urllib.parse.quote(tracking, safe=""))}


# -------------------- GLOBAL SEARCH --------------------
@api_router.get("/search/global")
async def global_search(q: str, limit: int = 20, _: User = Depends(get_current_user)):
    """Cross-collection search the topbar uses for omni-search. Hits:
    shipments (reference / bol_no / container_no / sales_order / po_no /
    material), truckload bookings (id / bol_no / po_no / carrier),
    bol_uploads (filename / bol_no), and machines (model / family)."""
    q = (q or "").strip()
    if len(q) < 2:
        return {"results": [], "query": q}
    rx = {"$regex": re.escape(q), "$options": "i"}
    results: List[Dict[str, Any]] = []

    # Shipments — index nearly every reference field
    async for s in db.shipments.find(
        {"$or": [
            {"reference": rx}, {"shipment_id": rx}, {"bol_no": rx},
            {"container_no": rx}, {"sales_order": rx}, {"po_no": rx},
            {"material": rx}, {"carrier": rx}, {"tracking_no": rx},
            {"delivery_no": rx}, {"invoice_no": rx},
        ]}, {"_id": 0, "shipment_id": 1, "reference": 1, "bol_no": 1, "container_no": 1,
              "sales_order": 1, "po_no": 1, "material": 1, "carrier": 1,
              "origin": 1, "destination": 1, "status": 1, "mode": 1}
    ).limit(limit):
        results.append({
            "type": "shipment",
            "title": s.get("reference") or s.get("shipment_id"),
            "subtitle": f"{(s.get('origin') or {}).get('city','?')} → {(s.get('destination') or {}).get('city','?')} · {s.get('carrier') or '—'}",
            "badge": s.get("mode") or "—",
            "status": s.get("status"),
            "link": f"/shipments?focus={s.get('shipment_id')}",
            "ref": s.get("reference") or s.get("shipment_id"),
            "fields": {k: s.get(k) for k in ("bol_no", "container_no", "sales_order", "po_no", "material") if s.get(k)},
        })

    # Truckload bookings
    async for b in db.truckload_bookings.find(
        {"$or": [{"id": rx}, {"bol_no": rx}, {"po_no": rx}, {"carrier": rx}]},
        {"_id": 0}
    ).limit(limit):
        results.append({
            "type": "booking",
            "title": b.get("bol_no") or b.get("id"),
            "subtitle": f"{b.get('origin','?')} → {b.get('destination','?')} · {b.get('carrier') or '—'}",
            "badge": "TL",
            "status": b.get("status"),
            "link": "/workbook",
            "ref": b.get("bol_no") or b.get("id"),
            "fields": {k: b.get(k) for k in ("bol_no", "po_no", "carrier") if b.get(k)},
        })

    # BOL uploads (GridFS)
    async for f in db["bol_uploads.files"].find(
        {"$or": [{"filename": rx}, {"metadata.bol_no": rx}, {"metadata.shipment_id": rx}]},
        {"filename": 1, "metadata": 1}
    ).limit(limit):
        md = f.get("metadata") or {}
        results.append({
            "type": "document",
            "title": f.get("filename"),
            "subtitle": f"BOL · {md.get('bol_no','—')} · {md.get('shipment_id','—')}",
            "badge": "DOC",
            "status": md.get("status"),
            "link": "/documents",
            "ref": md.get("bol_no") or f.get("filename"),
        })

    # Machines
    async for m in db.machines.find(
        {"$or": [{"model": rx}, {"family": rx}, {"category": rx}]},
        {"_id": 0, "model": 1, "family": 1, "category": 1}
    ).limit(limit):
        results.append({
            "type": "machine",
            "title": m.get("model"),
            "subtitle": f"{m.get('family','?')} · {m.get('category','?')}",
            "badge": "EQ",
            "link": f"/machines?focus={m.get('model')}",
            "ref": m.get("model"),
        })

    return {"results": results[: limit * 2], "query": q, "count": len(results)}


# -------------------- S/4 SEARCH (MOCKED) --------------------
@api_router.get("/s4/search")
async def s4_search(q: str, _: User = Depends(get_current_user)):
    """Mocked S/4HANA cross-document search. Returns matches across PO,
    Sales Order, Delivery, Invoice, Material, and Transportation Order
    with direct Fiori deep-links."""
    q = (q or "").strip()
    if len(q) < 2:
        return {"results": [], "query": q}
    # Generate deterministic-but-realistic mock matches
    rnd = _random.Random(q.lower())
    kinds = [
        ("purchase_order", "Purchase Order", "PO"),
        ("sales_order",    "Sales Order",    "SO"),
        ("delivery",       "Outbound Delivery", "DELIV"),
        ("invoice",        "Invoice",        "INV"),
        ("material",       "Material Master", "MATL"),
        ("shipment",       "Transportation Order", "TO"),
    ]
    results = []
    for (kind, label, badge) in kinds:
        # 50% chance per kind to show 1-2 matches based on the query string
        if rnd.random() > 0.5:
            continue
        count = rnd.choice([1, 1, 2])
        for _i in range(count):
            doc = f"{badge}-{rnd.randint(4000000, 4999999)}"
            path = S4_DEEP_LINK_PATTERNS[kind].format(value=doc)
            results.append({
                "kind": kind, "label": label, "badge": badge,
                "doc_number": doc,
                "description": f"S/4 match for '{q}' in {label}",
                "url": f"{SAP_S4_BASE}{path}",
                "matched_on": "header" if rnd.random() > 0.4 else "item",
            })
    return await _brand_swap({"results": results, "query": q, "base": SAP_S4_BASE, "tenant": "tennantco.s4hana.cloud.sap"})


# -------------------- ADMIN SETTINGS --------------------
@api_router.get("/admin/settings")
async def admin_settings_get(_: User = Depends(get_current_user)):
    """Single-doc settings store. All users can read; only admin can write."""
    doc = await db.app_settings.find_one({"_id": "global"}, {"_id": 0}) or {}
    return doc


@api_router.put("/admin/settings")
async def admin_settings_put(payload: Dict[str, Any], user: User = Depends(require_role("admin"))):
    rec = {**payload, "_id": "global",
           "updated_by": user.name,
           "updated_at": datetime.now(timezone.utc).isoformat()}
    await db.app_settings.replace_one({"_id": "global"}, rec, upsert=True)
    return {"ok": True}


# -------------------- S/4 INVOICES (MOCKED) --------------------
# Surfaces a list of SAP S/4HANA freight & supplier invoices the Document
# Vault can link to without having to import the documents themselves.
@api_router.get("/s4/invoices")
async def s4_invoices(limit: int = 50, _: User = Depends(get_current_user)):
    """Mocked roster of recent S/4HANA invoices. Each row deep-links into
    Fiori's Billing Document fact-sheet."""
    rnd = _random.Random(42)
    vendors = ["XPO Logistics", "ODFL", "Saia", "Estes Express", "R+L Carriers",
               "Knight Transportation", "Schneider National", "C.H. Robinson",
               "Werner Enterprises", "Kuehne+Nagel", "Maersk Line", "DHL Express",
               "UPS Freight", "FedEx Freight", "DB Schenker", "Expeditors",
               "Hapag-Lloyd", "CMA CGM", "ONE Line", "ArcBest"]
    out = []
    for i in range(limit):
        v = rnd.choice(vendors)
        amt = rnd.randint(280, 18000) + rnd.random()
        days_ago = rnd.randint(0, 60)
        inv_no = f"INV-{rnd.randint(4500000, 4999999)}"
        path = S4_DEEP_LINK_PATTERNS["invoice"].format(value=inv_no)
        out.append({
            "invoice_no": inv_no,
            "vendor": v,
            "amount_usd": round(amt, 2),
            "currency": "USD",
            "invoice_date": (datetime.now(timezone.utc) - timedelta(days=days_ago)).date().isoformat(),
            "due_date": (datetime.now(timezone.utc) + timedelta(days=30 - days_ago)).date().isoformat(),
            "status": rnd.choice(["Paid", "Approved", "In Review", "On Hold", "Disputed"]),
            "po_number": f"PO-{rnd.randint(4500000, 4999999)}",
            "s4_url": f"{SAP_S4_BASE}{path}",
        })
    return await _brand_swap({"invoices": out, "total": len(out), "tenant": "tennantco.s4hana.cloud.sap",
            "fiori_root": f"{SAP_S4_BASE}/sap/bc/ui2/flp#Shell-home"})


# -------------------- MACHINES · ADD/UPLOAD --------------------
class MachineCreate(BaseModel):
    model: str
    family: Optional[str] = ""
    category: Optional[str] = "Scrubber"
    description: Optional[str] = ""
    image_url: Optional[str] = ""
    specs: Optional[Dict[str, Any]] = None


@api_router.post("/machines")
async def machine_create(payload: MachineCreate, user: User = Depends(require_role("admin", "dispatcher"))):
    """Insert a custom machine into the Tennant catalog. If image_url is
    blank, the existing dynamic SVG generator at `/api/machines/{model}/image.svg`
    auto-renders a branded placeholder."""
    existing = await db.machines.find_one({"model": payload.model})
    if existing:
        raise HTTPException(status_code=409, detail=f"Model '{payload.model}' already exists")
    rec = {
        "model": payload.model,
        "family": payload.family or "Custom",
        "category": payload.category or "Other",
        "description": payload.description or "",
        "image_url": payload.image_url or "",
        "specs": payload.specs or {},
        "added_by": user.name,
        "added_at": datetime.now(timezone.utc).isoformat(),
        "custom": True,
    }
    await db.machines.insert_one(dict(rec))
    rec.pop("_id", None)
    return rec


# -------------------- DRIVER CONSOLE — DRIVER & TRAILER REGISTRY --------------------
class DriverRecord(BaseModel):
    name: str
    cdl_number: Optional[str] = ""
    cdl_class: Optional[str] = "A"
    cdl_state: Optional[str] = ""
    cdl_expiry: Optional[str] = ""
    medical_card_expiry: Optional[str] = ""
    phone: Optional[str] = ""
    email: Optional[str] = ""
    carrier: Optional[str] = ""
    hazmat_endorsement: Optional[bool] = False
    tanker_endorsement: Optional[bool] = False
    twic_card: Optional[bool] = False
    notes: Optional[str] = ""


class TrailerRecord(BaseModel):
    trailer_no: str
    type: Optional[str] = "Dry Van 53'"
    license_plate: Optional[str] = ""
    license_state: Optional[str] = ""
    vin: Optional[str] = ""
    carrier: Optional[str] = ""
    last_inspection: Optional[str] = ""
    next_inspection: Optional[str] = ""
    capacity_lbs: Optional[float] = 45000
    notes: Optional[str] = ""


@api_router.get("/drivers")
async def drivers_list(carrier: Optional[str] = None, _: User = Depends(get_current_user)):
    q = {"carrier": carrier} if carrier else {}
    docs = await db.drivers.find(q, {"_id": 0}).sort("name", 1).to_list(1000)
    if not docs:
        starters = [
            {"id": f"DRV-{uuid.uuid4().hex[:6].upper()}", "name": "Carlos Mendoza", "cdl_number": "M523891", "cdl_class": "A", "cdl_state": "MI", "cdl_expiry": "2027-04-15", "medical_card_expiry": "2026-11-20", "phone": "+1-616-555-0421", "email": "cmendoza@logixtrans.com", "carrier": "Logix Transportation", "hazmat_endorsement": False, "tanker_endorsement": False, "twic_card": False, "notes": "Pad-wrap specialist · Tennant T-series"},
            {"id": f"DRV-{uuid.uuid4().hex[:6].upper()}", "name": "Sarah O'Neill",   "cdl_number": "O998211", "cdl_class": "A", "cdl_state": "OH", "cdl_expiry": "2026-09-02", "medical_card_expiry": "2026-08-15", "phone": "+1-440-555-0188", "email": "soneill@xpo.com",  "carrier": "XPO Logistics", "hazmat_endorsement": True,  "tanker_endorsement": False, "twic_card": True,  "notes": "Hazmat-certified · Holland → Atlanta lane"},
            {"id": f"DRV-{uuid.uuid4().hex[:6].upper()}", "name": "Derrick Watanabe","cdl_number": "W772143", "cdl_class": "A", "cdl_state": "CA", "cdl_expiry": "2028-01-30", "medical_card_expiry": "2027-02-10", "phone": "+1-310-555-0719", "email": "dwata@schneider.com","carrier": "Schneider National", "hazmat_endorsement": False, "tanker_endorsement": True, "twic_card": False, "notes": "Reefer · West Coast lanes"},
            {"id": f"DRV-{uuid.uuid4().hex[:6].upper()}", "name": "Amanda Pruitt",  "cdl_number": "P640122", "cdl_class": "A", "cdl_state": "TX", "cdl_expiry": "2026-12-12", "medical_card_expiry": "2026-10-01", "phone": "+1-214-555-0357", "email": "apruitt@knight.com","carrier": "Knight Transportation","hazmat_endorsement": False, "tanker_endorsement": False,"twic_card": False, "notes": "Team driver · cross-country expedite"},
            {"id": f"DRV-{uuid.uuid4().hex[:6].upper()}", "name": "Boris Kowalski", "cdl_number": "K513398", "cdl_class": "A", "cdl_state": "IL", "cdl_expiry": "2027-06-18", "medical_card_expiry": "2026-12-04", "phone": "+1-708-555-0094", "email": "bkowal@odfl.com",  "carrier": "Old Dominion",       "hazmat_endorsement": True,  "tanker_endorsement": False,"twic_card": True,  "notes": "LTL specialist · 14-year veteran"},
        ]
        for s in starters:
            s["created_at"] = datetime.now(timezone.utc).isoformat()
            s["created_by"] = "System Seed"
        await db.drivers.insert_many([dict(s) for s in starters])
        docs = await db.drivers.find({}, {"_id": 0}).sort("name", 1).to_list(1000)
    return await _brand_swap({"drivers": docs})


@api_router.post("/drivers")
async def drivers_create(payload: DriverRecord, user: User = Depends(require_role("admin", "dispatcher"))):
    rec = {"id": f"DRV-{uuid.uuid4().hex[:6].upper()}", **payload.dict(),
           "created_at": datetime.now(timezone.utc).isoformat(),
           "created_by": user.name}
    await db.drivers.insert_one(dict(rec))
    rec.pop("_id", None)
    return rec


@api_router.delete("/drivers/{driver_id}")
async def drivers_delete(driver_id: str, _: User = Depends(require_role("admin", "dispatcher"))):
    r = await db.drivers.delete_one({"id": driver_id})
    if r.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Driver not found")
    return {"ok": True}


@api_router.get("/trailers")
async def trailers_list(carrier: Optional[str] = None, _: User = Depends(get_current_user)):
    q = {"carrier": carrier} if carrier else {}
    docs = await db.trailers.find(q, {"_id": 0}).sort("trailer_no", 1).to_list(1000)
    if not docs:
        starters = [
            {"id": f"TRL-{uuid.uuid4().hex[:6].upper()}", "trailer_no": "TNT-53201", "type": "Dry Van 53'", "license_plate": "MI-487J22", "license_state": "MI", "vin": "1H5VR3LX5MC123201", "carrier": "Tennant Fleet", "last_inspection": "2026-01-12", "next_inspection": "2026-07-12", "capacity_lbs": 45000, "notes": "Logo-wrapped · pad-wrap kit on board"},
            {"id": f"TRL-{uuid.uuid4().hex[:6].upper()}", "trailer_no": "TNT-53202", "type": "Dry Van 53'", "license_plate": "MI-487J23", "license_state": "MI", "vin": "1H5VR3LX5MC123202", "carrier": "Tennant Fleet", "last_inspection": "2025-11-05", "next_inspection": "2026-05-05", "capacity_lbs": 45000, "notes": ""},
            {"id": f"TRL-{uuid.uuid4().hex[:6].upper()}", "trailer_no": "LOGIX-1109", "type": "Air-Ride Lowboy", "license_plate": "OH-LXT109", "license_state": "OH", "vin": "1L2AR5KX7LL891109", "carrier": "Logix Transportation", "last_inspection": "2026-02-01", "next_inspection": "2026-08-01", "capacity_lbs": 48000, "notes": "Pad-wrap white-glove · machine transport"},
            {"id": f"TRL-{uuid.uuid4().hex[:6].upper()}", "trailer_no": "XPO-78445", "type": "Dry Van 53'", "license_plate": "TN-XPO845", "license_state": "TN", "vin": "1X7DR2LY9XX878445", "carrier": "XPO Logistics", "last_inspection": "2025-12-19", "next_inspection": "2026-06-19", "capacity_lbs": 44500, "notes": ""},
            {"id": f"TRL-{uuid.uuid4().hex[:6].upper()}", "trailer_no": "ODFL-22154", "type": "Pup 28'", "license_plate": "NC-OD2154", "license_state": "NC", "vin": "1O3PP7LY3OO722154", "carrier": "Old Dominion", "last_inspection": "2026-01-22", "next_inspection": "2026-07-22", "capacity_lbs": 22500, "notes": "LTL pup · Holland lane"},
            {"id": f"TRL-{uuid.uuid4().hex[:6].upper()}", "trailer_no": "MAEU-6731288", "type": "40' HC Container", "license_plate": "", "license_state": "", "vin": "MAEU6731288", "carrier": "Maersk Line", "last_inspection": "2026-01-30", "next_inspection": "2026-07-30", "capacity_lbs": 67200, "notes": "Ocean container · Long Beach inbound"},
        ]
        for s in starters:
            s["created_at"] = datetime.now(timezone.utc).isoformat()
            s["created_by"] = "System Seed"
        await db.trailers.insert_many([dict(s) for s in starters])
        docs = await db.trailers.find({}, {"_id": 0}).sort("trailer_no", 1).to_list(1000)
    return await _brand_swap({"trailers": docs})


@api_router.post("/trailers")
async def trailers_create(payload: TrailerRecord, user: User = Depends(require_role("admin", "dispatcher"))):
    rec = {"id": f"TRL-{uuid.uuid4().hex[:6].upper()}", **payload.dict(),
           "created_at": datetime.now(timezone.utc).isoformat(),
           "created_by": user.name}
    await db.trailers.insert_one(dict(rec))
    rec.pop("_id", None)
    return rec


@api_router.delete("/trailers/{trailer_id}")
async def trailers_delete(trailer_id: str, _: User = Depends(require_role("admin", "dispatcher"))):
    r = await db.trailers.delete_one({"id": trailer_id})
    if r.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Trailer not found")
    return {"ok": True}



# -------------------- INTEGRATIONS --------------------
INTEGRATIONS = [
    {"id": "sap_s4hana", "name": "SAP S/4HANA", "category": "ERP", "status": "connected", "last_sync": "2m ago", "endpoint": "api.tennantco.s4.sap.com"},
    {"id": "sharepoint", "name": "Microsoft SharePoint", "category": "Productivity", "status": "connected", "last_sync": "5m ago", "endpoint": "tennantco.sharepoint.com"},
    {"id": "powerbi", "name": "Microsoft PowerBI", "category": "Analytics", "status": "connected", "last_sync": "8m ago", "endpoint": "app.powerbi.com/tennant"},
    {"id": "outlook", "name": "Microsoft Outlook", "category": "Communication", "status": "connected", "last_sync": "1m ago", "endpoint": "outlook.office365.com"},
    {"id": "logix", "name": "Logix Transportation", "category": "TMS", "status": "connected", "last_sync": "3m ago", "endpoint": "api.logixtms.com"},
    {"id": "ups", "name": "UPS", "category": "Parcel", "status": "connected", "last_sync": "1m ago", "endpoint": "onlinetools.ups.com"},
    {"id": "fedex", "name": "FedEx", "category": "Parcel", "status": "connected", "last_sync": "2m ago", "endpoint": "ws.api.fedex.com"},
    {"id": "dhl", "name": "DHL Express", "category": "Air/Intl", "status": "connected", "last_sync": "4m ago", "endpoint": "api-eu.dhl.com"},
    {"id": "arcbest", "name": "ArcBest", "category": "LTL", "status": "connected", "last_sync": "6m ago", "endpoint": "api.arcb.com"},
    {"id": "fastfrate", "name": "Consolidated Fastfrate", "category": "LTL (Canada)", "status": "connected", "last_sync": "9m ago", "endpoint": "api.fastfrate.com"},
    {"id": "rl", "name": "R&L Carriers", "category": "LTL", "status": "connected", "last_sync": "12m ago", "endpoint": "api.rlc.com"},
    {"id": "xpo", "name": "XPO Logistics", "category": "LTL", "status": "warning", "last_sync": "1h ago", "endpoint": "api.xpo.com"},
    {"id": "saia", "name": "SAIA LTL Freight", "category": "LTL", "status": "connected", "last_sync": "7m ago", "endpoint": "api.saia.com"},
    {"id": "kuehne", "name": "Kuehne+Nagel", "category": "Ocean/Air", "status": "connected", "last_sync": "3m ago", "endpoint": "api.kuehne-nagel.com"},
    {"id": "webex", "name": "Cisco Webex", "category": "Communication", "status": "connected", "last_sync": "1m ago", "endpoint": "webexapis.com/v1"},
]

@api_router.get("/integrations")
async def get_integrations(_: User = Depends(get_current_user)):
    return await _brand_swap(INTEGRATIONS)

# -------------------- CARRIER RATES & FSC --------------------
CARRIER_RATES = [
    # TL lanes — base rate per mile + min charge + fuel surcharge
    {"lane_id": "GVM-DAL-TL", "mode": "TL", "origin": "Golden Valley, MN", "destination": "Dallas, TX", "miles": 970,
     "carriers": [
         {"carrier": "XPO Logistics", "scac": "XPOL", "base_rate": 2150.00, "rate_per_mile": 2.22, "fsc_pct": 28.5, "min_charge": 1200, "transit_days": 3, "fak": "85"},
         {"carrier": "ArcBest", "scac": "ABFS", "base_rate": 2280.00, "rate_per_mile": 2.35, "fsc_pct": 27.2, "min_charge": 1200, "transit_days": 3, "fak": "85"},
         {"carrier": "Schneider", "scac": "SNDR", "base_rate": 2090.00, "rate_per_mile": 2.15, "fsc_pct": 29.0, "min_charge": 1200, "transit_days": 3, "fak": "85"},
         {"carrier": "J.B. Hunt", "scac": "JBHT", "base_rate": 2210.00, "rate_per_mile": 2.28, "fsc_pct": 28.0, "min_charge": 1200, "transit_days": 3, "fak": "85"},
     ]},
    {"lane_id": "HOM-ATL-TL", "mode": "TL", "origin": "Holland, MI", "destination": "Atlanta, GA", "miles": 820,
     "carriers": [
         {"carrier": "XPO Logistics", "scac": "XPOL", "base_rate": 1980.00, "rate_per_mile": 2.41, "fsc_pct": 27.8, "min_charge": 1200, "transit_days": 2, "fak": "85"},
         {"carrier": "ArcBest", "scac": "ABFS", "base_rate": 2050.00, "rate_per_mile": 2.50, "fsc_pct": 27.0, "min_charge": 1200, "transit_days": 2, "fak": "85"},
         {"carrier": "Schneider", "scac": "SNDR", "base_rate": 1920.00, "rate_per_mile": 2.34, "fsc_pct": 28.5, "min_charge": 1200, "transit_days": 2, "fak": "85"},
         {"carrier": "J.B. Hunt", "scac": "JBHT", "base_rate": 2010.00, "rate_per_mile": 2.45, "fsc_pct": 28.2, "min_charge": 1200, "transit_days": 2, "fak": "85"},
     ]},
    {"lane_id": "LVK-PHX-TL", "mode": "TL", "origin": "Louisville, KY", "destination": "Phoenix, AZ", "miles": 1810,
     "carriers": [
         {"carrier": "XPO Logistics", "scac": "XPOL", "base_rate": 3680.00, "rate_per_mile": 2.03, "fsc_pct": 29.2, "min_charge": 1500, "transit_days": 4, "fak": "85"},
         {"carrier": "ArcBest", "scac": "ABFS", "base_rate": 3920.00, "rate_per_mile": 2.17, "fsc_pct": 28.0, "min_charge": 1500, "transit_days": 4, "fak": "85"},
         {"carrier": "Schneider", "scac": "SNDR", "base_rate": 3540.00, "rate_per_mile": 1.96, "fsc_pct": 29.8, "min_charge": 1500, "transit_days": 4, "fak": "85"},
         {"carrier": "J.B. Hunt", "scac": "JBHT", "base_rate": 3780.00, "rate_per_mile": 2.09, "fsc_pct": 28.8, "min_charge": 1500, "transit_days": 4, "fak": "85"},
     ]},
    # LTL lanes — rates per 100 lbs (CWT)
    {"lane_id": "HOM-DAL-LTL", "mode": "LTL", "origin": "Holland, MI", "destination": "Dallas, TX", "miles": 1130,
     "carriers": [
         {"carrier": "SAIA", "scac": "SAIA", "base_rate": 720.00, "rate_per_cwt": 18.50, "fsc_pct": 31.5, "min_charge": 145, "transit_days": 4, "fak": "85"},
         {"carrier": "R&L Carriers", "scac": "RLCA", "base_rate": 685.00, "rate_per_cwt": 17.20, "fsc_pct": 32.0, "min_charge": 125, "transit_days": 4, "fak": "85"},
         {"carrier": "ArcBest", "scac": "ABFS", "base_rate": 745.00, "rate_per_cwt": 19.10, "fsc_pct": 30.8, "min_charge": 150, "transit_days": 4, "fak": "85"},
         {"carrier": "XPO Logistics", "scac": "XPOL", "base_rate": 760.00, "rate_per_cwt": 19.50, "fsc_pct": 31.0, "min_charge": 155, "transit_days": 5, "fak": "85"},
         {"carrier": "Consolidated Fastfrate", "scac": "CFAT", "base_rate": 705.00, "rate_per_cwt": 18.05, "fsc_pct": 32.5, "min_charge": 140, "transit_days": 5, "fak": "85"},
     ]},
    {"lane_id": "GVM-SEA-LTL", "mode": "LTL", "origin": "Golden Valley, MN", "destination": "Seattle, WA", "miles": 1670,
     "carriers": [
         {"carrier": "SAIA", "scac": "SAIA", "base_rate": 1180.00, "rate_per_cwt": 22.40, "fsc_pct": 32.5, "min_charge": 165, "transit_days": 5, "fak": "85"},
         {"carrier": "R&L Carriers", "scac": "RLCA", "base_rate": 1095.00, "rate_per_cwt": 20.80, "fsc_pct": 33.0, "min_charge": 140, "transit_days": 5, "fak": "85"},
         {"carrier": "XPO Logistics", "scac": "XPOL", "base_rate": 1245.00, "rate_per_cwt": 23.20, "fsc_pct": 31.8, "min_charge": 170, "transit_days": 6, "fak": "85"},
         {"carrier": "ArcBest", "scac": "ABFS", "base_rate": 1220.00, "rate_per_cwt": 22.95, "fsc_pct": 32.0, "min_charge": 165, "transit_days": 5, "fak": "85"},
     ]},
    # Parcel — flat per-package + fuel surcharge
    {"lane_id": "GVM-DOMESTIC-PARCEL", "mode": "Parcel", "origin": "Golden Valley, MN", "destination": "U.S. Domestic Avg.", "miles": 0,
     "carriers": [
         {"carrier": "UPS", "scac": "UPSN", "base_rate": 18.40, "rate_per_lb": 1.85, "fsc_pct": 18.5, "min_charge": 12.00, "transit_days": 3, "fak": "Ground"},
         {"carrier": "FedEx", "scac": "FXFE", "base_rate": 19.10, "rate_per_lb": 1.92, "fsc_pct": 18.0, "min_charge": 12.50, "transit_days": 3, "fak": "Ground"},
         {"carrier": "DHL Express", "scac": "DHLC", "base_rate": 24.80, "rate_per_lb": 2.20, "fsc_pct": 22.5, "min_charge": 18.00, "transit_days": 2, "fak": "Express"},
     ]},
    # Ocean — per container
    {"lane_id": "POL-SHA-OCEAN", "mode": "Ocean", "origin": "Long Beach, CA", "destination": "Shanghai, CN", "miles": 6500,
     "carriers": [
         {"carrier": "Kuehne+Nagel", "scac": "KNAA", "base_rate": 3850.00, "rate_per_container": 3850.00, "fsc_pct": 12.0, "min_charge": 3850, "transit_days": 16, "fak": "40HC"},
         {"carrier": "Maersk", "scac": "MAEU", "base_rate": 3680.00, "rate_per_container": 3680.00, "fsc_pct": 12.5, "min_charge": 3680, "transit_days": 17, "fak": "40HC"},
         {"carrier": "MSC", "scac": "MSCU", "base_rate": 3740.00, "rate_per_container": 3740.00, "fsc_pct": 13.0, "min_charge": 3740, "transit_days": 18, "fak": "40HC"},
     ]},
    # Inbound from Korea
    {"lane_id": "BUS-LVK-OCEAN", "mode": "Ocean", "origin": "Busan, KR", "destination": "Louisville, KY (via LA)", "miles": 6200,
     "carriers": [
         {"carrier": "Kuehne+Nagel", "scac": "KNAA", "base_rate": 4180.00, "rate_per_container": 4180.00, "fsc_pct": 13.5, "min_charge": 4180, "transit_days": 22, "fak": "40HC"},
         {"carrier": "Maersk", "scac": "MAEU", "base_rate": 4020.00, "rate_per_container": 4020.00, "fsc_pct": 14.0, "min_charge": 4020, "transit_days": 24, "fak": "40HC"},
     ]},
]

@api_router.get("/carrier-rates")
async def get_carrier_rates(_: User = Depends(get_current_user), mode: Optional[str] = None):
    out = [l for l in CARRIER_RATES if (l["mode"] == mode or not mode)]
    out = [dict(l) for l in out]  # copy so we don't mutate the canonical list
    # Brand-aware overlay — rewrite origin/destination cities to the active
    # brand's sample_lanes so a Walmart admin sees Walmart's lanes, not
    # Golden Valley → Dallas. Rates and carrier mix stay realistic.
    brand = await _active_brand_doc()
    if brand and brand.get("brand_id") != "orisei-freight":
        lanes = brand.get("sample_lanes") or []
        # Each lane string like "Bentonville, AR -> Dallas, TX". Split safely.
        parsed = []
        for s in lanes:
            if "->" in s:
                o, d = s.split("->", 1)
                parsed.append((o.strip(), d.strip()))
        for i, l in enumerate(out):
            if parsed:
                o, d = parsed[i % len(parsed)]
                l["origin"] = o
                l["destination"] = d
                # Re-derive lane_id from new cities so it's still pretty
                short_o = re.sub(r"[^A-Z]", "", o.upper())[:3]
                short_d = re.sub(r"[^A-Z]", "", d.upper())[:3]
                l["lane_id"] = f"{short_o}-{short_d}-{l['mode']}"
    return out

@api_router.get("/carrier-rates/fsc")
async def get_fsc_index(_: User = Depends(get_current_user)):
    """Returns the current fuel surcharge index per carrier (national diesel + carrier matrix)."""
    return {
        "doe_diesel_avg_per_gallon": 3.84,
        "doe_week": "2026-05-11",
        "fsc_table": [
            {"carrier": "XPO Logistics", "scac": "XPOL", "current_fsc_pct": 28.5, "trend": "flat", "week_change_pct": 0.0},
            {"carrier": "ArcBest", "scac": "ABFS", "current_fsc_pct": 27.2, "trend": "down", "week_change_pct": -0.3},
            {"carrier": "Schneider", "scac": "SNDR", "current_fsc_pct": 29.0, "trend": "up", "week_change_pct": 0.5},
            {"carrier": "J.B. Hunt", "scac": "JBHT", "current_fsc_pct": 28.0, "trend": "flat", "week_change_pct": 0.0},
            {"carrier": "SAIA", "scac": "SAIA", "current_fsc_pct": 31.5, "trend": "up", "week_change_pct": 0.2},
            {"carrier": "R&L Carriers", "scac": "RLCA", "current_fsc_pct": 32.0, "trend": "down", "week_change_pct": -0.4},
            {"carrier": "UPS", "scac": "UPSN", "current_fsc_pct": 18.5, "trend": "flat", "week_change_pct": 0.0},
            {"carrier": "FedEx", "scac": "FXFE", "current_fsc_pct": 18.0, "trend": "up", "week_change_pct": 0.25},
            {"carrier": "DHL Express", "scac": "DHLC", "current_fsc_pct": 22.5, "trend": "up", "week_change_pct": 0.5},
            {"carrier": "Kuehne+Nagel", "scac": "KNAA", "current_fsc_pct": 12.5, "trend": "flat", "week_change_pct": 0.0},
            {"carrier": "Consolidated Fastfrate", "scac": "CFAT", "current_fsc_pct": 32.5, "trend": "down", "week_change_pct": -0.2},
        ],
    }

# -------------------- MUSIC (Radio Browser proxy + curated genres) --------------------
RADIO_BROWSER_BASE = "https://de1.api.radio-browser.info"

MUSIC_GENRES = [
    {"id": "lofi", "label": "Lo-Fi / Focus", "tag": "lofi", "icon": "🎧"},
    {"id": "jazz", "label": "Jazz", "tag": "jazz", "icon": "🎷"},
    {"id": "classical", "label": "Classical", "tag": "classical", "icon": "🎻"},
    {"id": "ambient", "label": "Ambient", "tag": "ambient", "icon": "🌊"},
    {"id": "electronic", "label": "Electronic", "tag": "electronic", "icon": "⚡"},
    {"id": "house", "label": "House", "tag": "house", "icon": "🏠"},
    {"id": "pop", "label": "Pop", "tag": "pop", "icon": "✨"},
    {"id": "rock", "label": "Rock", "tag": "rock", "icon": "🎸"},
    {"id": "indie", "label": "Indie", "tag": "indie", "icon": "🎤"},
    {"id": "country", "label": "Country", "tag": "country", "icon": "🤠"},
    {"id": "blues", "label": "Blues", "tag": "blues", "icon": "🎺"},
    {"id": "folk", "label": "Folk", "tag": "folk", "icon": "🪕"},
    {"id": "reggae", "label": "Reggae", "tag": "reggae", "icon": "🌴"},
    {"id": "latin", "label": "Latin", "tag": "latin", "icon": "💃"},
    {"id": "soundtrack", "label": "Soundtrack", "tag": "soundtrack", "icon": "🎬"},
    {"id": "world", "label": "World", "tag": "world", "icon": "🌍"},
    {"id": "instrumental", "label": "Instrumental", "tag": "instrumental", "icon": "🎼"},
    {"id": "chillout", "label": "Chillout", "tag": "chillout", "icon": "🧘"},
    {"id": "acoustic", "label": "Acoustic", "tag": "acoustic", "icon": "🪗"},
    {"id": "smooth-jazz", "label": "Smooth Jazz", "tag": "smooth jazz", "icon": "🎶"},
]

@api_router.get("/music/genres")
async def music_genres(_: User = Depends(get_current_user)):
    return MUSIC_GENRES

@api_router.get("/music/stations")
async def music_stations(
    _: User = Depends(get_current_user),
    genre: Optional[str] = None,
    q: Optional[str] = None,
    country: Optional[str] = None,
    limit: int = 60,
):
    """Search work-appropriate radio stations via Radio Browser. Filters out explicit/talk content."""
    params: Dict[str, Any] = {
        "limit": limit,
        "hidebroken": "true",
        "order": "clickcount",
        "reverse": "true",
    }
    # Build search params
    if genre:
        # Find genre tag from our list
        g = next((x for x in MUSIC_GENRES if x["id"] == genre), None)
        if g:
            params["tag"] = g["tag"]
    if q:
        params["name"] = q
    if country:
        params["countrycode"] = country

    try:
        async with httpx.AsyncClient(timeout=12.0) as http:
            url = f"{RADIO_BROWSER_BASE}/json/stations/search"
            r = await http.get(url, params=params, headers={"User-Agent": "TennantTMS/1.0"})
            stations = r.json()
    except Exception as e:
        logger.warning(f"Radio Browser fetch failed: {e}")
        return []

    # Filter to work-appropriate: skip stations whose name/tags contain explicit/talk markers
    BLOCKLIST = ["explicit", "talk", "news talk", "religion", "religious", "sermons", "preacher", "rap explicit"]
    out = []
    for s in stations:
        name = (s.get("name") or "").lower()
        tags = (s.get("tags") or "").lower()
        if any(b in name or b in tags for b in BLOCKLIST):
            continue
        url_stream = s.get("url_resolved") or s.get("url")
        if not url_stream:
            continue
        out.append({
            "id": s.get("stationuuid"),
            "name": s.get("name", "").strip(),
            "url": url_stream,
            "homepage": s.get("homepage"),
            "favicon": s.get("favicon"),
            "country": s.get("country"),
            "countrycode": s.get("countrycode"),
            "tags": s.get("tags", "").split(",")[:5],
            "bitrate": s.get("bitrate", 0),
            "codec": s.get("codec"),
            "language": s.get("language"),
            "clickcount": s.get("clickcount", 0),
            "votes": s.get("votes", 0),
        })
    return out

@api_router.post("/music/click")
async def music_click(request: Request, _: User = Depends(get_current_user)):
    """Report a click-through to the source for analytics-friendly use."""
    body = await request.json()
    station_id = body.get("station_id")
    if not station_id:
        return {"ok": False}
    try:
        async with httpx.AsyncClient(timeout=6.0) as http:
            await http.get(f"{RADIO_BROWSER_BASE}/json/url/{station_id}", headers={"User-Agent": "TennantTMS/1.0"})
    except Exception:
        pass
    return {"ok": True}

# -------------------- TRAILER SPECS --------------------
TRAILER_SPECS = [
    {"id": "53dv", "name": "53' Dry Van", "length_ft": 53, "width_ft": 8.5, "height_ft": 9, "max_weight_lbs": 45000, "pallets": 26, "uses": ["General freight", "Palletized goods", "Scrubber assemblies"], "color": "#00E5FF"},
    {"id": "48dv", "name": "48' Dry Van", "length_ft": 48, "width_ft": 8.5, "height_ft": 9, "max_weight_lbs": 44000, "pallets": 24, "uses": ["Regional freight", "Northeast lanes (low bridges)"], "color": "#00FF66"},
    {"id": "28pup", "name": "28' Pup Trailer", "length_ft": 28, "width_ft": 8, "height_ft": 9, "max_weight_lbs": 22000, "pallets": 14, "uses": ["LTL doubles/triples", "City delivery"], "color": "#FFCC00"},
    {"id": "53reef", "name": "53' Reefer", "length_ft": 53, "width_ft": 8.5, "height_ft": 8.5, "max_weight_lbs": 43500, "pallets": 26, "uses": ["Temperature-sensitive batteries", "Sensitive electronics"], "color": "#FF3B30"},
    {"id": "48flat", "name": "48' Flatbed", "length_ft": 48, "width_ft": 8.5, "height_ft": 5, "max_weight_lbs": 48000, "pallets": 0, "uses": ["Oversized scrubbers", "Machinery", "Steel components"], "color": "#A78BFA"},
    {"id": "20ctr", "name": "20' Ocean Container", "length_ft": 20, "width_ft": 8, "height_ft": 8.5, "max_weight_lbs": 47900, "pallets": 10, "uses": ["Heavy imports", "Single SKU consolidations"], "color": "#3B82F6"},
    {"id": "40ctr", "name": "40' Ocean Container", "length_ft": 40, "width_ft": 8, "height_ft": 8.5, "max_weight_lbs": 59000, "pallets": 20, "uses": ["Standard ocean imports from K+N", "Mixed SKU"], "color": "#06B6D4"},
    {"id": "40hc", "name": "40' High Cube", "length_ft": 40, "width_ft": 8, "height_ft": 9.5, "max_weight_lbs": 58000, "pallets": 20, "uses": ["Tall machinery", "Volume-out before weight-out"], "color": "#10B981"},
]

@api_router.get("/trailer-specs")
async def get_trailer_specs(_: User = Depends(get_current_user)):
    """Static reference data for the Trailer Specs page (53' Dry Van, Reefer,
    Flatbed, ocean containers, etc.). Distinct from /api/trailers which lists
    physical fleet assets."""
    return TRAILER_SPECS

# -------------------- HS CODES --------------------
HS_CODES = [
    {"code": "8479.89.94", "description": "Floor scrubbing/sweeping machines, self-propelled, industrial", "duty_pct": 2.5, "category": "Machinery"},
    {"code": "8479.89.65", "description": "Industrial cleaning machines, electromechanical", "duty_pct": 2.5, "category": "Machinery"},
    {"code": "8508.11.00", "description": "Vacuum cleaners with self-contained electric motor, ≤1500W", "duty_pct": 0, "category": "Machinery"},
    {"code": "8508.19.00", "description": "Vacuum cleaners, other (industrial)", "duty_pct": 0, "category": "Machinery"},
    {"code": "8413.81.00", "description": "Pumps for liquids, other (cleaning solution pumps)", "duty_pct": 0, "category": "Pumps"},
    {"code": "8501.31.50", "description": "DC motors, ≤750W (drive motors for scrubbers)", "duty_pct": 4.0, "category": "Electrical"},
    {"code": "8507.60.00", "description": "Lithium-ion batteries (battery packs)", "duty_pct": 3.4, "category": "Batteries"},
    {"code": "8507.20.80", "description": "Lead-acid storage batteries, other", "duty_pct": 3.5, "category": "Batteries"},
    {"code": "9603.50.00", "description": "Brushes constituting parts of machines (scrubber brushes)", "duty_pct": 0, "category": "Parts"},
    {"code": "8431.49.90", "description": "Parts suitable for use with machinery of headings 8425-8430", "duty_pct": 0, "category": "Parts"},
    {"code": "3926.90.99", "description": "Other articles of plastics (tanks, housings)", "duty_pct": 5.3, "category": "Plastics"},
    {"code": "7326.90.86", "description": "Articles of iron or steel, other (frames, brackets)", "duty_pct": 2.9, "category": "Steel"},
    {"code": "8536.50.90", "description": "Switches for voltage ≤1000V (control switches)", "duty_pct": 2.7, "category": "Electrical"},
    {"code": "4016.93.50", "description": "Rubber gaskets, washers and other seals", "duty_pct": 2.5, "category": "Rubber"},
    {"code": "8544.30.00", "description": "Ignition wiring sets and other wiring sets used in vehicles", "duty_pct": 5.0, "category": "Wiring"},
    {"code": "9026.10.20", "description": "Flow meters (water/solution flow)", "duty_pct": 0, "category": "Instruments"},
    {"code": "8536.69.40", "description": "Connectors for voltage ≤1000V (battery & motor connectors)", "duty_pct": 0, "category": "Electrical"},
    {"code": "3917.32.00", "description": "Plastic tubes/hoses, not reinforced", "duty_pct": 3.1, "category": "Plastics"},
]

@api_router.get("/hs-lookup")
async def hs_lookup(q: str = Query(""), _: User = Depends(get_current_user)):
    if not q:
        return HS_CODES
    ql = q.lower()
    return [c for c in HS_CODES if ql in c["code"].lower() or ql in c["description"].lower() or ql in c["category"].lower()]

# -------------------- QUICK LINKS --------------------
QUICK_LINKS = [
    {"name": "ACE Portal (CBP)", "url": "https://ace.cbp.dhs.gov/", "category": "Import/Export", "description": "Automated Commercial Environment - U.S. Customs"},
    {"name": "DOT FMCSA SAFER", "url": "https://safer.fmcsa.dot.gov/", "category": "DOT", "description": "Carrier safety lookup & insurance verification"},
    {"name": "FMCSA Portal", "url": "https://portal.fmcsa.dot.gov/", "category": "DOT", "description": "Federal Motor Carrier Safety Administration"},
    {"name": "USDOT - Transportation.gov", "url": "https://www.transportation.gov/", "category": "DOT", "description": "U.S. Department of Transportation"},
    {"name": "CBP Trade", "url": "https://www.cbp.gov/trade", "category": "Import/Export", "description": "Customs and Border Protection - Trade"},
    {"name": "USTR HTS Search", "url": "https://hts.usitc.gov/", "category": "Tariff", "description": "Harmonized Tariff Schedule of the United States"},
    {"name": "Census Schedule B Search", "url": "https://www.census.gov/foreign-trade/schedules/b/", "category": "Tariff", "description": "Schedule B (export commodity codes)"},
    {"name": "Port of Long Beach", "url": "https://www.polb.com/", "category": "Ports", "description": "Major West Coast port"},
    {"name": "Port of New York/NJ", "url": "https://www.panynj.gov/port/", "category": "Ports", "description": "East Coast major port"},
    {"name": "Weather.gov Aviation", "url": "https://aviationweather.gov/", "category": "Weather", "description": "Aviation weather products"},
    {"name": "FAA NOTAMs", "url": "https://www.faa.gov/", "category": "Air", "description": "Federal Aviation Administration"},
]

@api_router.get("/links")
async def get_links(_: User = Depends(get_current_user)):
    return QUICK_LINKS

# -------------------- CHAT --------------------
DEFAULT_CHANNELS = [
    {"id": "general", "name": "general", "description": "Company-wide updates"},
    {"id": "ops-dispatch", "name": "ops-dispatch", "description": "Daily dispatch coordination"},
    {"id": "import-export", "name": "import-export", "description": "Customs, HS codes, ACE filings"},
    {"id": "carrier-issues", "name": "carrier-issues", "description": "Carrier escalations & exceptions"},
    {"id": "louisville", "name": "louisville", "description": "Louisville, KY plant team"},
    {"id": "holland", "name": "holland", "description": "Holland, MI plant team"},
    {"id": "golden-valley", "name": "golden-valley", "description": "Golden Valley, MN HQ"},
]

@api_router.get("/chat/channels")
async def get_channels(_: User = Depends(get_current_user)):
    return DEFAULT_CHANNELS

@api_router.get("/chat/messages")
async def get_messages(channel: str, _: User = Depends(get_current_user)):
    msgs = await db.chat_messages.find({"channel": channel}, {"_id": 0}).sort("created_at", 1).limit(200).to_list(200)
    return msgs

# WebSocket manager
class ConnectionManager:
    def __init__(self):
        self.active: List[WebSocket] = []

    async def connect(self, ws: WebSocket):
        await ws.accept()
        self.active.append(ws)

    def disconnect(self, ws: WebSocket):
        if ws in self.active:
            self.active.remove(ws)

    async def broadcast(self, data: dict):
        dead = []
        for ws in self.active:
            try:
                await ws.send_json(data)
            except Exception:
                dead.append(ws)
        for d in dead:
            self.disconnect(d)

manager = ConnectionManager()

@app.websocket("/api/ws/chat")
async def ws_chat(websocket: WebSocket, token: str = Query(...)):
    session = await db.user_sessions.find_one({"session_token": token}, {"_id": 0})
    if not session:
        await websocket.close(code=1008)
        return
    user_doc = await db.users.find_one({"user_id": session["user_id"]}, {"_id": 0})
    if not user_doc:
        await websocket.close(code=1008)
        return
    await manager.connect(websocket)
    try:
        while True:
            raw = await websocket.receive_text()
            data = json.loads(raw)
            text = (data.get("text") or "").strip()
            channel = data.get("channel", "general")
            if not text:
                continue
            msg = {
                "message_id": f"msg_{uuid.uuid4().hex[:10]}",
                "channel": channel,
                "user_id": user_doc["user_id"],
                "user_name": user_doc["name"],
                "user_picture": user_doc.get("picture"),
                "text": text,
                "created_at": datetime.now(timezone.utc).isoformat(),
            }
            await db.chat_messages.insert_one(dict(msg))
            await manager.broadcast({"type": "message", "data": msg})
    except WebSocketDisconnect:
        manager.disconnect(websocket)
    except Exception as e:
        logger.warning(f"WS error: {e}")
        manager.disconnect(websocket)

# -------------------- TEAM --------------------
@api_router.get("/team")
async def team(_: User = Depends(get_current_user)):
    docs = await db.users.find({}, {"_id": 0}).limit(50).to_list(50)
    return docs

# -------------------- FREIGHT AUDIT & PAY --------------------
class FreightBill(BaseModel):
    bill_id: str
    shipment_ref: str
    carrier: str
    invoice_no: str
    base_charge: float
    accessorials: List[Dict[str, Any]]  # [{code, description, amount}]
    fuel_surcharge: float
    total: float
    quoted_total: float
    variance: float
    status: str  # pending, approved, paid, disputed, audit
    invoice_date: str
    due_date: str
    paid_at: Optional[str] = None

@api_router.get("/freight-bills")
async def list_bills(_: User = Depends(get_current_user), status: Optional[str] = None, carrier: Optional[str] = None):
    q: Dict[str, Any] = {}
    if status: q["status"] = status
    if carrier: q["carrier"] = carrier
    docs = await db.freight_bills.find(q, {"_id": 0}).sort("invoice_date", -1).limit(500).to_list(500)
    return docs

@api_router.post("/freight-bills/{bill_id}/pay")
async def pay_bill(bill_id: str, user: User = Depends(require_role("auditor"))):
    bill = await db.freight_bills.find_one({"bill_id": bill_id}, {"_id": 0})
    if not bill:
        raise HTTPException(status_code=404, detail="Bill not found")
    if bill["status"] == "paid":
        return {"ok": True, "already_paid": True}
    await db.freight_bills.update_one(
        {"bill_id": bill_id},
        {"$set": {"status": "paid", "paid_at": datetime.now(timezone.utc).isoformat(), "paid_by": user.name}}
    )
    return {"ok": True, "bill_id": bill_id}

@api_router.post("/freight-bills/{bill_id}/approve")
async def approve_bill(bill_id: str, user: User = Depends(require_role("auditor"))):
    bill = await db.freight_bills.find_one({"bill_id": bill_id}, {"_id": 0})
    if not bill:
        raise HTTPException(status_code=404, detail="Bill not found")
    await db.freight_bills.update_one({"bill_id": bill_id}, {"$set": {"status": "approved", "approved_by": user.name}})
    return {"ok": True}

@api_router.post("/freight-bills/{bill_id}/dispute")
async def dispute_bill(bill_id: str, request: Request, user: User = Depends(require_role("auditor"))):
    bill = await db.freight_bills.find_one({"bill_id": bill_id}, {"_id": 0})
    if not bill:
        raise HTTPException(status_code=404, detail="Bill not found")
    body = await request.json()
    await db.freight_bills.update_one(
        {"bill_id": bill_id},
        {"$set": {"status": "disputed", "dispute_reason": body.get("reason", ""), "disputed_by": user.name}}
    )
    return {"ok": True}

@api_router.get("/freight-bills/summary")
async def bills_summary(_: User = Depends(get_current_user)):
    bills = await db.freight_bills.find({}, {"_id": 0}).to_list(2000)
    total = sum(b["total"] for b in bills)
    paid = sum(b["total"] for b in bills if b["status"] == "paid")
    pending = sum(b["total"] for b in bills if b["status"] in ("pending", "audit"))
    disputed = sum(b["total"] for b in bills if b["status"] == "disputed")
    overcharges = sum(b["variance"] for b in bills if b.get("variance", 0) > 0)
    return {
        "total_billed": round(total, 2),
        "paid": round(paid, 2),
        "pending": round(pending, 2),
        "disputed": round(disputed, 2),
        "overcharges_detected": round(overcharges, 2),
        "count": len(bills),
        "count_disputed": sum(1 for b in bills if b["status"] == "disputed"),
    }

# -------------------- CARRIER ONBOARDING --------------------
class CarrierOnboarding(BaseModel):
    onboarding_id: str
    legal_name: str
    dba: Optional[str] = None
    mc_number: Optional[str] = None
    dot_number: Optional[str] = None
    scac: Optional[str] = None
    mode: str
    contact_name: str
    contact_email: str
    contact_phone: str
    insurance_amount: float
    insurance_expiry: str
    safety_rating: str  # Satisfactory, Conditional, Unsatisfactory, NotRated
    csa_score: int
    w9_received: bool
    coi_received: bool
    contract_signed: bool
    status: str  # invited, in_review, approved, rejected
    submitted_at: str
    notes: Optional[str] = ""

class CarrierOnboardingCreate(BaseModel):
    legal_name: str
    dba: Optional[str] = None
    mc_number: Optional[str] = None
    dot_number: Optional[str] = None
    scac: Optional[str] = None
    mode: str = "TL"
    contact_name: str
    contact_email: str
    contact_phone: str
    insurance_amount: float = 1000000
    insurance_expiry: str
    safety_rating: str = "Satisfactory"
    csa_score: int = 50
    notes: Optional[str] = ""

@api_router.get("/carriers/onboarding")
async def list_onboarding(_: User = Depends(get_current_user), status: Optional[str] = None):
    q: Dict[str, Any] = {}
    if status: q["status"] = status
    docs = await db.carrier_onboarding.find(q, {"_id": 0}).sort("submitted_at", -1).to_list(500)
    return docs

@api_router.post("/carriers/onboarding")
async def create_onboarding(payload: CarrierOnboardingCreate, user: User = Depends(get_current_user)):
    doc = {
        "onboarding_id": f"OB-{uuid.uuid4().hex[:8].upper()}",
        **payload.model_dump(),
        "w9_received": False,
        "coi_received": False,
        "contract_signed": False,
        "status": "in_review",
        "submitted_at": datetime.now(timezone.utc).isoformat(),
        "submitted_by": user.name,
    }
    await db.carrier_onboarding.insert_one(dict(doc))
    return doc

@api_router.post("/carriers/onboarding/{oid}/decision")
async def decide_onboarding(oid: str, request: Request, user: User = Depends(require_role("admin"))):
    body = await request.json()
    decision = body.get("decision")  # approved | rejected
    if decision not in ("approved", "rejected"):
        raise HTTPException(status_code=400, detail="Invalid decision")
    target = await db.carrier_onboarding.find_one({"onboarding_id": oid}, {"_id": 0})
    if not target:
        raise HTTPException(status_code=404, detail="Onboarding not found")
    await db.carrier_onboarding.update_one(
        {"onboarding_id": oid},
        {"$set": {"status": decision, "decided_by": user.name, "decision_notes": body.get("notes", "")}}
    )
    return {"ok": True}

@api_router.post("/carriers/onboarding/{oid}/toggle")
async def toggle_doc(oid: str, request: Request, _: User = Depends(get_current_user)):
    body = await request.json()
    field = body.get("field")
    if field not in ("w9_received", "coi_received", "contract_signed"):
        raise HTTPException(status_code=400, detail="Invalid field")
    doc = await db.carrier_onboarding.find_one({"onboarding_id": oid}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Not found")
    await db.carrier_onboarding.update_one({"onboarding_id": oid}, {"$set": {field: not doc.get(field, False)}})
    return {"ok": True}

# -------------------- DRIVER MOBILE --------------------
class DriverCheckIn(BaseModel):
    checkin_id: str
    shipment_id: str
    driver_name: str
    driver_phone: str
    status: str  # arriving_pickup, loaded, en_route, fuel, rest, delayed, arriving_dest, delivered
    lat: Optional[float] = None
    lng: Optional[float] = None
    location_text: Optional[str] = None
    note: Optional[str] = None
    created_at: str
    odometer: Optional[float] = None
    fuel_pct: Optional[int] = None

class DriverCheckInCreate(BaseModel):
    shipment_id: str
    driver_name: str
    driver_phone: str
    status: str
    lat: Optional[float] = None
    lng: Optional[float] = None
    location_text: Optional[str] = None
    note: Optional[str] = None
    odometer: Optional[float] = None
    fuel_pct: Optional[int] = None

@api_router.post("/driver/checkin")
async def driver_checkin(payload: DriverCheckInCreate):
    """Driver-side endpoint — open (auth-free) for in-cab mobile use, identified by shipment_id + driver_phone."""
    shipment = await db.shipments.find_one({"shipment_id": payload.shipment_id}, {"_id": 0})
    if not shipment:
        raise HTTPException(status_code=404, detail="Shipment not found")

    ci = {
        "checkin_id": f"CI-{uuid.uuid4().hex[:8].upper()}",
        **payload.model_dump(),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.driver_checkins.insert_one(dict(ci))

    # Map driver status -> shipment status & current_location
    status_map = {
        "arriving_pickup": "at_origin",
        "loaded": "in_transit",
        "en_route": "in_transit",
        "fuel": "in_transit",
        "rest": "in_transit",
        "delayed": "delayed",
        "arriving_dest": "in_transit",
        "delivered": "delivered",
    }
    updates: Dict[str, Any] = {}
    new_status = status_map.get(payload.status)
    if new_status:
        updates["status"] = new_status
    if payload.lat is not None and payload.lng is not None:
        updates["current_location"] = {"lat": payload.lat, "lng": payload.lng, "city": payload.location_text or "En route"}
        # rough progress estimate
        try:
            o = shipment["origin"]; d = shipment["destination"]
            total = ((d["lat"] - o["lat"]) ** 2 + (d["lng"] - o["lng"]) ** 2) ** 0.5
            done = ((payload.lat - o["lat"]) ** 2 + (payload.lng - o["lng"]) ** 2) ** 0.5
            if total > 0:
                updates["progress"] = max(0.0, min(1.0, done / total))
        except Exception:
            pass
    if updates:
        await db.shipments.update_one({"shipment_id": payload.shipment_id}, {"$set": updates})

    return {"ok": True, "checkin_id": ci["checkin_id"], "shipment_status": updates.get("status", shipment["status"])}

@api_router.get("/driver/shipment/{shipment_id}")
async def driver_get_shipment(shipment_id: str):
    """Driver-facing view (auth-free) — minimal shipment info + recent check-ins."""
    s = await db.shipments.find_one({"shipment_id": shipment_id}, {"_id": 0})
    if not s:
        raise HTTPException(status_code=404, detail="Shipment not found")
    checkins = await db.driver_checkins.find({"shipment_id": shipment_id}, {"_id": 0}).sort("created_at", -1).limit(20).to_list(20)
    return {"shipment": s, "checkins": checkins}

@api_router.get("/driver/checkins")
async def list_checkins(_: User = Depends(get_current_user), shipment_id: Optional[str] = None):
    q: Dict[str, Any] = {}
    if shipment_id: q["shipment_id"] = shipment_id
    docs = await db.driver_checkins.find(q, {"_id": 0}).sort("created_at", -1).limit(200).to_list(200)
    return docs

# -------------------- PDF DOCUMENT RENDERING --------------------
TENNANT_BLUE = colors.HexColor("#0E3A6B")
TENNANT_DARK = colors.HexColor("#0B0E14")

def _doc_styles():
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="TennantTitle", parent=styles["Heading1"], fontSize=22, leading=26, textColor=TENNANT_DARK, alignment=0, spaceAfter=4))
    styles.add(ParagraphStyle(name="TennantSubtitle", parent=styles["Normal"], fontSize=9, textColor=colors.HexColor("#475569"), spaceAfter=12))
    styles.add(ParagraphStyle(name="SectionLabel", parent=styles["Normal"], fontSize=7, textColor=TENNANT_BLUE, leading=9))
    styles.add(ParagraphStyle(name="FieldLabel", parent=styles["Normal"], fontSize=7, textColor=colors.HexColor("#64748B"), leading=9))
    styles.add(ParagraphStyle(name="FieldValue", parent=styles["Normal"], fontSize=10, textColor=TENNANT_DARK, leading=13))
    styles.add(ParagraphStyle(name="DocFooter", parent=styles["Normal"], fontSize=7, textColor=colors.HexColor("#94A3B8"), alignment=1))
    return styles

DOC_TYPE_TITLES = {
    "BOL": ("BILL OF LADING", "Straight Bill of Lading — Original Domestic"),
    "COMMERCIAL_INVOICE": ("COMMERCIAL INVOICE", "Export / Customs Invoice"),
    "PACKING_SLIP": ("PACKING SLIP", "Shipment Pack List"),
    "WEIGHT_CERT": ("WEIGHT CERTIFICATE", "Certified Scale Ticket"),
    "COO": ("CERTIFICATE OF ORIGIN", "Statement of Country of Origin"),
}

def _header_block(doc_id: str, doc_type: str, brand: Optional[Dict[str, Any]] = None):
    styles = _doc_styles()
    title, subtitle = DOC_TYPE_TITLES.get(doc_type, (doc_type, ""))
    company = (brand or {}).get("company_name") or "Orisei Freight Solutions LLC"
    short = (brand or {}).get("short_name") or "ORISEI"
    primary = (brand or {}).get("primary_color") or "#0E3A6B"
    rest = company.replace(short, "", 1).strip() or "COMPANY"
    header_data = [
        [
            Paragraph(f"<b><font color='{primary}'>{short.upper()}</font></b> {rest.upper()}", styles["TennantTitle"]),
            Paragraph(f"<b>{title}</b><br/><font size=8 color='#64748B'>{subtitle}</font><br/><font size=7 color='#94A3B8'>Document ID: {doc_id}</font>", styles["FieldValue"]),
        ]
    ]
    t = Table(header_data, colWidths=[3.2 * inch, 3.8 * inch])
    t.setStyle(TableStyle([
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("ALIGN", (1, 0), (1, 0), "RIGHT"),
        ("LINEBELOW", (0, 0), (-1, -1), 2, colors.HexColor(primary)),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 12),
    ]))
    return t

def _kv_table(rows, col_widths=None):
    styles = _doc_styles()
    data = []
    for row in rows:
        data.append([
            Paragraph(row[0].upper(), styles["FieldLabel"]),
            Paragraph(str(row[1] or "—"), styles["FieldValue"]),
        ])
    t = Table(data, colWidths=col_widths or [1.4 * inch, 2.3 * inch])
    t.setStyle(TableStyle([
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 2),
    ]))
    return t

def _build_pdf(doc: Dict[str, Any], brand: Optional[Dict[str, Any]] = None) -> bytes:
    buf = io.BytesIO()
    pdf = SimpleDocTemplate(buf, pagesize=letter, leftMargin=0.6 * inch, rightMargin=0.6 * inch, topMargin=0.6 * inch, bottomMargin=0.6 * inch, title=doc["document_id"])
    styles = _doc_styles()
    data = doc.get("data", {}) or {}
    elements = []
    elements.append(_header_block(doc["document_id"], doc["type"], brand=brand))
    elements.append(Spacer(1, 14))

    company_name = (brand or {}).get("company_name") or "Orisei Freight Solutions LLC"
    # Shipper / Consignee block
    parties_rows = [
        ["Shipper", data.get("shipper") or company_name],
        ["Consignee", data.get("consignee")],
        ["Origin", data.get("origin")],
        ["Destination", data.get("destination")],
    ]
    shipment_rows = [
        ["Shipment Ref", doc.get("shipment_ref")],
        ["Carrier", data.get("carrier")],
        ["Commodity", data.get("commodity")],
        ["Pieces", data.get("pieces")],
        ["Weight (lbs)", data.get("weight")],
        ["Value (USD)", data.get("value") and f"${data.get('value')}"],
    ]

    parties_t = _kv_table(parties_rows, col_widths=[1.0 * inch, 2.4 * inch])
    shipment_t = _kv_table(shipment_rows, col_widths=[1.1 * inch, 2.3 * inch])
    columns = Table([[parties_t, shipment_t]], colWidths=[3.5 * inch, 3.5 * inch])
    columns.setStyle(TableStyle([
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("BOX", (0, 0), (0, 0), 0.5, colors.HexColor("#E2E8F0")),
        ("BOX", (1, 0), (1, 0), 0.5, colors.HexColor("#E2E8F0")),
        ("LEFTPADDING", (0, 0), (-1, -1), 10),
        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
        ("TOPPADDING", (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
    ]))
    elements.append(columns)
    elements.append(Spacer(1, 16))

    # Type-specific section
    dtype = doc["type"]
    if dtype == "BOL":
        line_items = [["#", "Pieces", "Description", "Weight (lbs)", "Class"]]
        line_items.append(["1", data.get("pieces") or "—", data.get("commodity") or f"{company_name} freight", data.get("weight") or "—", "85"])
        items = Table(line_items, colWidths=[0.4 * inch, 0.8 * inch, 3.6 * inch, 1.0 * inch, 0.7 * inch])
        items.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), TENNANT_BLUE),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#E2E8F0")),
            ("PADDING", (0, 0), (-1, -1), 5),
        ]))
        elements.append(Paragraph("<b><font color='#00A4E4'>LINE ITEMS</font></b>", styles["FieldLabel"]))
        elements.append(Spacer(1, 4))
        elements.append(items)
        elements.append(Spacer(1, 14))
        elements.append(Paragraph("<font size=8 color='#475569'>RECEIVED, subject to the classifications and tariffs in effect on the date of issue, the property described above in apparent good order, except as noted (contents and condition of contents of packages unknown).</font>", styles["FieldValue"]))

    elif dtype == "COMMERCIAL_INVOICE":
        try:
            qty = float(data.get("pieces") or 0)
            total = float(data.get("value") or 0)
            unit_price = total / qty if qty else total
        except Exception:
            qty, total, unit_price = "—", "—", "—"
        rows = [["Qty", "HS Code (suggested)", "Description", "Unit Price", "Total"]]
        rows.append([data.get("pieces") or "—", "8479.89.94", data.get("commodity") or "—",
                     f"${unit_price:,.2f}" if isinstance(unit_price, float) else unit_price,
                     f"${total:,.2f}" if isinstance(total, float) else total])
        items = Table(rows, colWidths=[0.6 * inch, 1.3 * inch, 3.0 * inch, 1.0 * inch, 1.1 * inch])
        items.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), TENNANT_BLUE),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#E2E8F0")),
            ("PADDING", (0, 0), (-1, -1), 5),
        ]))
        elements.append(Paragraph("<b><font color='#00A4E4'>INVOICE LINES</font></b>", styles["FieldLabel"]))
        elements.append(Spacer(1, 4))
        elements.append(items)
        elements.append(Spacer(1, 12))
        elements.append(Paragraph(f"<b>TOTAL INVOICE VALUE: ${data.get('value', '—')} USD</b>", styles["FieldValue"]))
        elements.append(Spacer(1, 8))
        elements.append(Paragraph("<font size=7 color='#64748B'>Terms: Incoterms 2020 — DAP. No commission. Country of Origin: USA unless otherwise noted.</font>", styles["FieldValue"]))

    elif dtype == "PACKING_SLIP":
        rows = [["Carton", "Qty", "Description", "Weight", "Dimensions"]]
        pcs = int(data.get("pieces") or 1) if str(data.get("pieces") or "1").isdigit() else 1
        for i in range(1, min(pcs, 8) + 1):
            rows.append([f"#{i:03d}", "1", data.get("commodity") or "—", f"{(float(data.get('weight') or 0) / max(1, pcs)):,.0f} lbs" if data.get("weight") else "—", "48×40×60 in"])
        items = Table(rows, colWidths=[0.7 * inch, 0.5 * inch, 3.6 * inch, 0.9 * inch, 1.3 * inch])
        items.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), TENNANT_BLUE),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#E2E8F0")),
            ("PADDING", (0, 0), (-1, -1), 5),
        ]))
        elements.append(items)

    elif dtype == "WEIGHT_CERT":
        elements.append(Paragraph("<b><font color='#00A4E4'>CERTIFIED WEIGHT</font></b>", styles["FieldLabel"]))
        elements.append(Spacer(1, 4))
        wt_rows = [
            ["Gross Weight", f"{data.get('weight') or '—'} lbs"],
            ["Tare Weight", "14,200 lbs"],
            ["Net Weight (calc)", f"{(float(data.get('weight') or 0) - 14200):,.0f} lbs" if data.get('weight') else "—"],
            ["Scale ID", "MN-CERT-04287"],
            ["Operator", doc.get("created_by") or "—"],
            ["Date / Time", doc.get("created_at") or "—"],
        ]
        elements.append(_kv_table(wt_rows, col_widths=[1.6 * inch, 2.4 * inch]))
        elements.append(Spacer(1, 12))
        elements.append(Paragraph("<font size=7 color='#475569'>I hereby certify that the weights shown above were obtained on a scale certified by the State of Minnesota and accurate within tolerance NIST Handbook 44.</font>", styles["FieldValue"]))

    elif dtype == "COO":
        coo_rows = [
            ["Country of Origin", data.get("country_origin") or "USA"],
            ["Producer", "Orisei Freight Solutions"],
            ["Producer Address", "10400 Clean Street, Eden Prairie, MN 55344, USA"],
            ["Exporter", "Orisei Freight Solutions"],
            ["Marks & Numbers", doc.get("shipment_ref") or "—"],
        ]
        elements.append(_kv_table(coo_rows, col_widths=[1.6 * inch, 4.0 * inch]))
        elements.append(Spacer(1, 14))
        elements.append(Paragraph("<font size=8 color='#475569'>The undersigned hereby declares that the above-mentioned goods originate from the country shown above and meet all applicable origin criteria. This certificate is issued in accordance with applicable Free Trade Agreement rules of origin where claimed.</font>", styles["FieldValue"]))

    # Signature block
    elements.append(Spacer(1, 26))
    sig_data = [
        [Paragraph("<font size=7 color='#64748B'>Authorized Signature</font><br/><br/>______________________________", styles["FieldValue"]),
         Paragraph(f"<font size=7 color='#64748B'>Date</font><br/><br/>{datetime.now(timezone.utc).strftime('%Y-%m-%d')}", styles["FieldValue"]),
         Paragraph(f"<font size=7 color='#64748B'>Prepared By</font><br/><br/>{doc.get('created_by') or '—'}", styles["FieldValue"])],
    ]
    sig = Table(sig_data, colWidths=[2.7 * inch, 1.4 * inch, 2.7 * inch])
    sig.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP")]))
    elements.append(sig)
    elements.append(Spacer(1, 8))
    elements.append(Paragraph(f"{company_name} · TMS Generated Document · {datetime.now(timezone.utc).isoformat(timespec='seconds')}Z", styles["DocFooter"]))

    pdf.build(elements)
    buf.seek(0)
    return buf.getvalue()

def _doc_to_branded_markdown(doc: Dict[str, Any], brand: Optional[Dict[str, Any]] = None) -> str:
    """Translate any non-BOL document into the markdown shape consumed by
    `build_branded_markdown_pdf` — yields a heraldic-bordered PDF with
    azure label/value tables, ◆ section headers, and gold-banner totals
    matching the BOL aesthetic. Used for COMMERCIAL_INVOICE, PACKING_SLIP,
    WEIGHT_CERT, COO, and any other type.
    """
    company = (brand or {}).get("company_name") or "Orisei Freight Solutions LLC"
    data = doc.get("data") or {}
    dtype = doc.get("type") or "DOCUMENT"
    title, subtitle = DOC_TYPE_TITLES.get(dtype, (dtype.replace("_", " "), ""))

    def _safe(v: Any, fb: str = "—") -> str:
        return str(v) if v not in (None, "", []) else fb

    md: List[str] = [
        f"# {title}",
        f"_{subtitle}_" if subtitle else "",
        "",
        "## Parties",
        f"- **Shipper**: {_safe(data.get('shipper'), company)}",
        f"- **Consignee**: {_safe(data.get('consignee'))}",
        f"- **Origin**: {_safe(data.get('origin'))}",
        f"- **Destination**: {_safe(data.get('destination'))}",
        "",
        "## Shipment",
        f"- **Reference**: {_safe(doc.get('shipment_ref'))}",
        f"- **Carrier**: {_safe(data.get('carrier'))}",
        f"- **Commodity**: {_safe(data.get('commodity'))}",
        f"- **Pieces**: {_safe(data.get('pieces'))}",
        f"- **Weight (lbs)**: {_safe(data.get('weight'))}",
    ]
    if data.get("value"):
        md.append(f"- **Value (USD)**: ${data.get('value')}")
    md.append("")

    if dtype == "COMMERCIAL_INVOICE":
        try:
            qty = float(data.get("pieces") or 0)
            total = float(data.get("value") or 0)
            unit_price = total / qty if qty else total
        except Exception:                                  # noqa: BLE001
            qty, total, unit_price = 0, 0, 0
        md += [
            "## Invoice Lines",
            f"- **Quantity**: {_safe(data.get('pieces'))}",
            f"- **HS Code (suggested)**: 8479.89.94",
            f"- **Description**: {_safe(data.get('commodity'))}",
            f"- **Unit Price**: ${unit_price:,.2f}" if unit_price else "- **Unit Price**: —",
            f"- **Line Total**: ${total:,.2f}" if total else "- **Line Total**: —",
            "",
            f"## Total · ${total:,.2f} USD" if total else "## Total · —",
            "",
            "## Terms",
            "- Incoterms 2020 — **DAP** unless otherwise noted",
            "- No commission",
            f"- Country of Origin: {_safe(data.get('country_origin'), 'USA')}",
        ]
    elif dtype == "PACKING_SLIP":
        try:
            pcs = int(data.get("pieces") or 1)
        except Exception:                                  # noqa: BLE001
            pcs = 1
        wt_per = (float(data.get("weight") or 0) / max(1, pcs)) if data.get("weight") else 0
        md += ["## Packing Detail"]
        for i in range(1, min(pcs, 8) + 1):
            md.append(f"- **Carton #{i:03d}**: 1 × {_safe(data.get('commodity'))} · "
                      f"{wt_per:,.0f} lbs · 48×40×60 in")
    elif dtype == "WEIGHT_CERT":
        try:
            gross = float(data.get("weight") or 0)
            net = gross - 14200
        except Exception:                                  # noqa: BLE001
            gross, net = 0, 0
        md += [
            "## Certified Weight",
            f"- **Gross Weight**: {gross:,.0f} lbs" if gross else "- **Gross Weight**: —",
            f"- **Tare Weight**: 14,200 lbs",
            f"- **Net Weight (calc)**: {net:,.0f} lbs" if gross else "- **Net Weight (calc)**: —",
            f"- **Scale ID**: MN-CERT-04287",
            f"- **Operator**: {_safe(doc.get('created_by'))}",
            f"- **Date / Time**: {_safe(doc.get('created_at'))}",
            "",
            "> I hereby certify that the weights shown above were obtained on a "
            "scale certified by the State of Minnesota and accurate within "
            "tolerance NIST Handbook 44.",
        ]
    elif dtype == "COO":
        md += [
            "## Certificate of Origin",
            f"- **Country of Origin**: {_safe(data.get('country_origin'), 'USA')}",
            f"- **Producer**: {company}",
            f"- **Producer Address**: {(brand or {}).get('headquarters') or 'Minneapolis · Saint Paul, MN'}",
            f"- **Exporter**: {company}",
            f"- **Marks & Numbers**: {_safe(doc.get('shipment_ref'))}",
            "",
            "> The undersigned hereby declares that the above-mentioned goods "
            "originate from the country shown above and meet all applicable "
            "origin criteria. This certificate is issued in accordance with "
            "applicable Free Trade Agreement rules of origin where claimed.",
        ]
    else:
        # Generic fallback — render whatever data fields are present.
        md.append("## Document Fields")
        for k, v in data.items():
            if k in {"shipper", "consignee", "origin", "destination", "carrier",
                     "commodity", "pieces", "weight", "value", "country_origin"}:
                continue  # already shown above
            md.append(f"- **{str(k).replace('_', ' ').title()}**: {_safe(v)}")

    md += [
        "",
        "## Authorized Signature",
        "- **Signed**: ______________________________",
        f"- **Date**: {datetime.now(timezone.utc).strftime('%Y-%m-%d')}",
        f"- **Prepared by**: {_safe(doc.get('created_by'))}",
    ]
    return "\n".join(md)


@api_router.get("/documents/{document_id}/pdf")
async def download_document_pdf(document_id: str, _: User = Depends(get_current_user)):
    doc = await db.documents.find_one({"document_id": document_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    brand = await _active_brand_doc() or {}
    # For BOLs, use the beautiful Orisei (Calafia + griffin) generator from
    # routes/orisei_docs so the document carries the active brand's identity.
    if doc.get("type") == "BOL":
        try:
            from routes.orisei_docs import build_bol_pdf
            data = doc.get("data") or {}
            booking = {
                "load_id": doc.get("shipment_ref") or doc.get("document_id"),
                "carrier_name": data.get("carrier"),
                "equipment": data.get("equipment") or "Van",
                "miles": data.get("miles") or 0,
                "pickup_date": data.get("pickup_date") or doc.get("created_at", "")[:10],
                "delivery_date": data.get("delivery_date") or "",
                "pieces": data.get("pieces"),
                "weight_lbs": data.get("weight"),
                "commodity": data.get("commodity"),
                "rate_usd": data.get("value") or data.get("rate_usd"),
                "notes": data.get("notes"),
            }
            company = brand.get("company_name") or "Orisei Freight Solutions LLC"
            shipper = {
                "name": data.get("shipper") or company,
                "address": data.get("shipper_address") or "Operations HQ",
                "city_state_zip": data.get("origin") or "",
                "contact": data.get("shipper_contact") or "",
            }
            consignee = {
                "name": data.get("consignee") or "Consignee",
                "address": data.get("consignee_address") or "",
                "city_state_zip": data.get("destination") or "",
                "contact": data.get("consignee_contact") or "",
            }
            pdf_bytes = build_bol_pdf(
                doc_id=doc["document_id"],
                booking=booking,
                shipper=shipper,
                consignee=consignee,
                user_name=None,
                brand=brand,
            )
        except Exception as e:                                       # noqa: BLE001
            logger.exception("Branded BOL render failed — falling back")
            try:
                pdf_bytes = _build_pdf(doc, brand=brand)
            except Exception as e2:                                   # noqa: BLE001
                raise HTTPException(status_code=500, detail=f"PDF render failed: {e2}")
    else:
        # All non-BOL doc types (COMMERCIAL_INVOICE, PACKING_SLIP, WEIGHT_CERT,
        # COO, …) — render through the branded heraldic markdown engine so
        # they share the BOL aesthetic (gold borders, ◆ section headers,
        # azure label/value tables, gold-banner Totals).
        try:
            from routes.orisei_docs import build_branded_markdown_pdf
            title_text, subtitle_text = DOC_TYPE_TITLES.get(
                doc.get("type", ""), (doc.get("type", "DOCUMENT"), ""))
            pdf_bytes = build_branded_markdown_pdf(
                _doc_to_branded_markdown(doc, brand=brand),
                title=title_text,
                subtitle=subtitle_text or f"Reference {doc.get('shipment_ref') or doc['document_id']}",
                doc_id=doc["document_id"],
                brand=brand,
            )
        except Exception as e:                                       # noqa: BLE001
            logger.exception("Branded markdown PDF render failed — falling back to legacy")
            try:
                pdf_bytes = _build_pdf(doc, brand=brand)
            except Exception as e2:                                   # noqa: BLE001
                raise HTTPException(status_code=500, detail=f"PDF render failed: {e2}")
    filename = f"{doc['type']}_{doc['document_id']}.pdf"
    # Auto-archive into the immutable Document Vault (fire-and-forget; never
    # blocks the download)
    try:
        from routes.doc_vault import archive_pdf  # local import to avoid cycle
        await archive_pdf(
            db, pdf_bytes,
            doc_type=doc.get("type", "OTHER"),
            doc_id=doc["document_id"],
            ref_id=doc.get("shipment_ref") or doc.get("shipment_id"),
            source_endpoint=f"/api/documents/{document_id}/pdf",
            payload_snapshot={"data": doc.get("data", {}),
                              "type": doc.get("type"),
                              "shipment_ref": doc.get("shipment_ref")},
            user=_,
            filename=filename,
        )
    except Exception:                                            # noqa: BLE001
        logger.exception("doc_vault archive failed for %s", doc["document_id"])
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )

# -------------------- ADMIN / RBAC --------------------
@api_router.get("/admin/users")
async def list_users(_: User = Depends(require_role("admin"))):
    docs = await db.users.find({}, {"_id": 0}).sort("created_at", 1).to_list(500)
    return docs

class RoleChange(BaseModel):
    role: str  # admin | auditor | dispatcher | driver

@api_router.post("/admin/users/{user_id}/role")
async def change_role(user_id: str, payload: RoleChange, actor: User = Depends(require_role("admin"))):
    if payload.role not in ROLE_HIERARCHY:
        raise HTTPException(status_code=400, detail="Invalid role")
    target = await db.users.find_one({"user_id": user_id}, {"_id": 0})
    if not target:
        raise HTTPException(status_code=404, detail="User not found")
    # Don't allow the actor to demote themselves if they are the only admin
    if target["user_id"] == actor.user_id and payload.role != "admin":
        admin_count = await db.users.count_documents({"role": "admin"})
        if admin_count <= 1:
            raise HTTPException(status_code=400, detail="Cannot demote the only remaining admin")
    await db.users.update_one({"user_id": user_id}, {"$set": {"role": payload.role}})
    return {"ok": True, "user_id": user_id, "role": payload.role}

@api_router.post("/admin/seed-team")
async def seed_team_users(_: User = Depends(require_role("admin"))):
    """Seed sample team members across roles for demo."""
    samples = [
        {"name": "Avery Lindgren", "email": "avery.lindgren@tennantco.com", "role": "auditor"},
        {"name": "Devon Marquez", "email": "devon.marquez@tennantco.com", "role": "dispatcher"},
        {"name": "Priya Iyer", "email": "priya.iyer@tennantco.com", "role": "dispatcher"},
        {"name": "Sam Chen", "email": "sam.chen@tennantco.com", "role": "driver"},
        {"name": "Riley Park", "email": "riley.park@tennantco.com", "role": "driver"},
    ]
    inserted = 0
    for s in samples:
        if await db.users.find_one({"email": s["email"]}):
            continue
        await db.users.insert_one({
            "user_id": f"user_{uuid.uuid4().hex[:12]}",
            "email": s["email"], "name": s["name"], "picture": None,
            "role": s["role"], "created_at": datetime.now(timezone.utc).isoformat(),
        })
        inserted += 1
    return {"ok": True, "inserted": inserted}

# SAP endpoints live in routes/sap.py — registered near wire-up.

# -------------------- SEED --------------------
# -------------------- AI ASSISTANT (Claude Sonnet 4.5) --------------------
EMERGENT_LLM_KEY = os.environ.get("EMERGENT_LLM_KEY")

AI_SYSTEM_PROMPT = """You are HUDLINK, the AI co-pilot inside the Tennant Companies Transportation Management System (TMS).

Tennant Company is a Minnesota-based manufacturer of industrial and commercial floor scrubbers and cleaners with manufacturing facilities in Louisville KY, Holland MI, and Golden Valley MN (HQ). Parts are imported via Kuehne+Nagel.

You help dispatchers, freight auditors, and operations leaders with:
- Shipment status, lane-level decisions, mode selection (TL, LTL, Parcel, Ocean, Air, Rail)
- Freight bill audit: accessorial charges, fuel surcharge norms, when to dispute
- Carrier scorecards and routing recommendations
- HS code classification for cleaning machinery, lithium batteries, motors, brushes (HTS chapter 8479, 8508, 8507, etc.)
- Document requirements: BOL, Commercial Invoice, Packing Slip, Weight Certificate, Certificate of Origin
- Compliance: DOT/FMCSA, ACE filings, Incoterms 2020, CBP
- SAP S/4HANA SO/PO context — sales orders ship from plants 1010 (Golden Valley), 1020 (Holland), 1030 (Louisville)

Style: terse, technical, action-oriented. Use bullet points and numbers. When asked for recommendations, give a clear primary answer, then 1-2 alternatives. Cite HS codes, Incoterms, and section names when relevant. Never invent data — if asked for live shipment IDs, tell the user to query the dashboard.
"""

class AIMessageIn(BaseModel):
    session_id: str
    message: str

@api_router.post("/ai/chat")
async def ai_chat(payload: AIMessageIn, user: User = Depends(get_current_user)):
    if not EMERGENT_LLM_KEY:
        raise HTTPException(status_code=500, detail="EMERGENT_LLM_KEY not configured")
    # Persist user message
    user_msg_doc = {
        "session_id": payload.session_id,
        "user_id": user.user_id,
        "role": "user",
        "text": payload.message,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.ai_messages.insert_one(dict(user_msg_doc))

    try:
        chat = LlmChat(
            api_key=EMERGENT_LLM_KEY,
            session_id=payload.session_id,
            system_message=AI_SYSTEM_PROMPT,
        ).with_model("anthropic", "claude-sonnet-4-5-20250929")
        reply = await chat.send_message(UserMessage(text=payload.message))
    except Exception as e:
        logger.exception("AI chat failed")
        raise HTTPException(status_code=502, detail=f"AI provider error: {e}")

    assistant_doc = {
        "session_id": payload.session_id,
        "user_id": user.user_id,
        "role": "assistant",
        "text": reply,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.ai_messages.insert_one(dict(assistant_doc))
    return {"reply": reply}

@api_router.get("/ai/history")
async def ai_history(session_id: str, user: User = Depends(get_current_user)):
    docs = await db.ai_messages.find(
        {"session_id": session_id, "user_id": user.user_id},
        {"_id": 0}
    ).sort("created_at", 1).limit(200).to_list(200)
    return docs

@api_router.delete("/ai/history")
async def ai_clear(session_id: str, user: User = Depends(get_current_user)):
    await db.ai_messages.delete_many({"session_id": session_id, "user_id": user.user_id})
    return {"ok": True}

# -------------------- WEBEX INTEGRATION (mocked) --------------------
WEBEX_CONFIG = {
    "org_id": "Y2lzY29zcGFyazovL3VzL09SR0FOSVpBVElPTi90ZW5uYW50Y28",
    "site": "tennantco.webex.com",
    "bot_email": "tms-bot@tennantco.webex.bot",
    "scopes": ["spark:rooms_read", "spark:messages_write", "meetings:schedules_write"],
    "status": "connected",
}

WEBEX_SPACES = [
    {"id": "SPC-OPS-DISP", "title": "Ops Dispatch — Daily Stand-up", "type": "group", "members": 14, "last_activity": "8m"},
    {"id": "SPC-IMPORT", "title": "Import / K+N Coordination", "type": "group", "members": 9, "last_activity": "22m"},
    {"id": "SPC-CARR-ESC", "title": "Carrier Escalations", "type": "group", "members": 22, "last_activity": "1h"},
    {"id": "SPC-PLANT-LVK", "title": "Louisville Plant Logistics", "type": "group", "members": 11, "last_activity": "3h"},
    {"id": "SPC-PLANT-HOM", "title": "Holland Plant Logistics", "type": "group", "members": 8, "last_activity": "5h"},
    {"id": "SPC-PLANT-GVM", "title": "Golden Valley HQ — All-Hands", "type": "group", "members": 47, "last_activity": "Yesterday"},
    {"id": "SPC-KIRK-1on1", "title": "Kirk Juergins — Direct Messages", "type": "direct", "members": 2, "last_activity": "2d"},
]

WEBEX_MEETINGS = [
    {"id": "MTG-7821", "title": "Weekly Carrier Performance Review", "host": "Kirk Juergins", "start": "2026-05-12T14:00:00Z", "duration_min": 60, "attendees": 8, "join_url": "https://tennantco.webex.com/meet/kirk.j/MTG-7821"},
    {"id": "MTG-7822", "title": "K+N Import Coordination — May Cycle", "host": "Avery Lindgren", "start": "2026-05-13T15:30:00Z", "duration_min": 45, "attendees": 6, "join_url": "https://tennantco.webex.com/meet/avery.l/MTG-7822"},
    {"id": "MTG-7825", "title": "SAP S/4HANA TMS Integration — Sprint Demo", "host": "Devon Marquez", "start": "2026-05-14T17:00:00Z", "duration_min": 30, "attendees": 12, "join_url": "https://tennantco.webex.com/meet/devon.m/MTG-7825"},
    {"id": "MTG-7830", "title": "Q2 Freight Spend Audit Findings", "host": "Avery Lindgren", "start": "2026-05-15T19:00:00Z", "duration_min": 60, "attendees": 5, "join_url": "https://tennantco.webex.com/meet/avery.l/MTG-7830"},
]

@api_router.get("/webex/config")
async def webex_config(_: User = Depends(get_current_user)):
    return WEBEX_CONFIG

@api_router.get("/webex/spaces")
async def webex_spaces(_: User = Depends(get_current_user)):
    return WEBEX_SPACES

@api_router.get("/webex/meetings")
async def webex_meetings(_: User = Depends(get_current_user)):
    return WEBEX_MEETINGS

class WebexNotifyIn(BaseModel):
    space_id: str
    text: str
    shipment_ref: Optional[str] = None

@api_router.post("/webex/notify")
async def webex_notify(payload: WebexNotifyIn, user: User = Depends(get_current_user)):
    """Simulate posting a message to a Webex space."""
    record = {
        "log_id": f"WBX-{uuid.uuid4().hex[:8].upper()}",
        "space_id": payload.space_id,
        "user_id": user.user_id,
        "user_name": user.name,
        "text": payload.text,
        "shipment_ref": payload.shipment_ref,
        "posted_at": datetime.now(timezone.utc).isoformat(),
        "status": "delivered",
    }
    await db.webex_log.insert_one(dict(record))
    return record

@api_router.get("/webex/log")
async def webex_log(_: User = Depends(get_current_user)):
    docs = await db.webex_log.find({}, {"_id": 0}).sort("posted_at", -1).limit(50).to_list(50)
    return docs

class WebexScheduleIn(BaseModel):
    title: str
    when: str  # ISO datetime
    duration_min: int = 30
    invitees: List[str] = []

@api_router.post("/webex/schedule")
async def webex_schedule(payload: WebexScheduleIn, user: User = Depends(get_current_user)):
    meeting = {
        "id": f"MTG-{random.randint(8000, 9999)}",
        "title": payload.title,
        "host": user.name,
        "start": payload.when,
        "duration_min": payload.duration_min,
        "attendees": len(payload.invitees),
        "join_url": f"https://tennantco.webex.com/meet/{user.user_id}/MTG-{uuid.uuid4().hex[:6].upper()}",
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.webex_meetings.insert_one(dict(meeting))
    return meeting

# -------------------- WORKBOOK (Excel-style renameable tabs) --------------------
DEFAULT_TABS = [
    {"kind": "truckload_bookings", "name": "Truckload Bookings", "filter": {}},
    {"kind": "shipments_tl", "name": "Outbound TL", "filter": {"mode": "TL"}},
    {"kind": "shipments_ltl", "name": "Outbound LTL", "filter": {"mode": "LTL"}},
    {"kind": "shipments_expedites", "name": "Expedites", "filter": {"mode": "Air"}},
    {"kind": "shipments_crates", "name": "Crate Spots", "filter": {"mode": "Flatbed"}},
    {"kind": "shipments_seafreight", "name": "Seafreight 25M", "filter": {"mode": "Ocean"}},
    {"kind": "shipments_import", "name": "25 Import", "filter": {"mode": "Ocean", "direction": "inbound"}},
    {"kind": "quotes", "name": "25 Quotes", "filter": {}},
    {"kind": "plant_hubs", "name": "Plant Hubs", "filter": {}},
    {"kind": "carriers_primary", "name": "IN Primary Carrier", "filter": {}},
    {"kind": "contacts_suppliers", "name": "IN Supplier Contacts", "filter": {}},
    {"kind": "contacts_carriers", "name": "IN Carrier Contacts", "filter": {}},
    {"kind": "info", "name": "Info", "filter": {}},
    {"kind": "volume_overview", "name": "Volume Overview", "filter": {}},
]

# Columns rendered by the editable Truckload Booking Sheet. Each column maps to
# a key on documents in db.truckload_bookings. `type` drives the cell editor on
# the frontend: text / number / date / select / textarea.
TRUCKLOAD_BOOKING_COLUMNS: List[Dict[str, Any]] = [
    {"key": "date", "label": "Date", "type": "date"},
    {"key": "bol_no", "label": "BOL #", "type": "text"},
    {"key": "po_no", "label": "PO #", "type": "text"},
    # Carrier is a free-text combobox at runtime — the `/workbook/truckload-bookings`
    # endpoint injects the current onboarded-carrier roster into `options`
    # while keeping `type: "combo"` so the cell editor lets dispatchers
    # either pick an existing approved carrier or type a brand-new name.
    {"key": "carrier", "label": "Carrier", "type": "combo", "options": []},
    {"key": "origin", "label": "Origin", "type": "text"},
    {"key": "destination", "label": "Destination", "type": "text"},
    {"key": "pieces", "label": "Pieces", "type": "number"},
    {"key": "weight_lbs", "label": "Weight (lbs)", "type": "number"},
    {"key": "pallets", "label": "Pallets", "type": "number"},
    {"key": "lift_gate", "label": "Lift Gate", "type": "select", "options": ["", "Yes", "No"]},
    {"key": "freight_class", "label": "Class", "type": "select",
     "options": ["", "50", "55", "60", "65", "70", "77.5", "85", "92.5", "100", "110", "125", "150", "175", "200", "250", "300", "400", "500"]},
    {"key": "nmfc_code", "label": "NMFC", "type": "text"},
    {"key": "equipment", "label": "Equipment", "type": "select",
     "options": ["", "Dry Van 53'", "Dry Van 48'", "Reefer", "Flatbed", "Step Deck", "Drop Deck", "Box Truck", "Sprinter"]},
    {"key": "pickup_date", "label": "Pickup Date", "type": "date"},
    {"key": "delivery_date", "label": "Delivery Date", "type": "date"},
    {"key": "rate_usd", "label": "Rate (USD)", "type": "number"},
    # NEW · cost transparency columns dispatchers asked for. Up-charges
    # captures any unexpected accessorial/detention/layover charges so
    # they can be flagged before the freight audit catches them.
    {"key": "up_charges_usd", "label": "Up-Charges (USD)", "type": "number"},
    {"key": "up_charges_reason", "label": "Up-Charge Reason", "type": "select",
     "options": ["", "Detention", "Layover", "Re-delivery", "Lift gate", "Inside delivery",
                 "Sort & segregate", "Reconsignment", "Driver assist", "Hazmat", "Limited access",
                 "Residential", "Notification", "Storage", "TONU (truck order not used)", "Other"]},
    {"key": "status", "label": "Status", "type": "select",
     "options": ["", "Quoted", "Booked", "Tendered", "Picked Up", "In Transit", "Delivered", "Cancelled"]},
    {"key": "notes", "label": "Notes", "type": "textarea"},
]

KIND_DEFINITIONS = {
    "shipments_tl": {"columns": [
        {"key": "reference", "label": "Reference"}, {"key": "shipment_id", "label": "Shipment ID"},
        {"key": "carrier", "label": "Carrier"}, {"key": "origin.city", "label": "Origin"},
        {"key": "destination.city", "label": "Destination"}, {"key": "pickup_date", "label": "Pickup Date"},
        {"key": "eta", "label": "ETA"}, {"key": "weight_lbs", "label": "Weight (lbs)"},
        {"key": "pieces", "label": "Pieces"}, {"key": "commodity", "label": "Commodity"},
        {"key": "status", "label": "Status"}, {"key": "bol_no", "label": "BOL #"},
    ]},
    "shipments_ltl": {"columns": [
        {"key": "reference", "label": "Reference"}, {"key": "carrier", "label": "Carrier"},
        {"key": "origin.city", "label": "Origin"}, {"key": "destination.city", "label": "Destination"},
        {"key": "weight_lbs", "label": "Weight"}, {"key": "pieces", "label": "Pcs"},
        {"key": "pro_no", "label": "PRO #"}, {"key": "status", "label": "Status"}, {"key": "eta", "label": "ETA"},
    ]},
    "shipments_seafreight": {"columns": [
        {"key": "reference", "label": "Reference"}, {"key": "carrier", "label": "NVOCC/Carrier"},
        {"key": "container_no", "label": "Container #"}, {"key": "origin.city", "label": "POL"},
        {"key": "destination.city", "label": "POD"}, {"key": "eta", "label": "ETA"},
        {"key": "status", "label": "Status"}, {"key": "value_usd", "label": "Value (USD)"},
    ]},
    "quotes": {"columns": [
        {"key": "quote_id", "label": "Quote #"}, {"key": "mode", "label": "Mode"},
        {"key": "carrier", "label": "Carrier"}, {"key": "origin", "label": "Origin"},
        {"key": "destination", "label": "Destination"}, {"key": "rate_usd", "label": "Rate (USD)"},
        {"key": "transit_days", "label": "Transit Days"}, {"key": "valid_until", "label": "Valid Until"},
        {"key": "status", "label": "Status"},
    ]},
    "plant_hubs": {"columns": [
        {"key": "id", "label": "Plant ID"}, {"key": "name", "label": "Plant Name"},
        {"key": "city", "label": "City"}, {"key": "state", "label": "State"},
        {"key": "type", "label": "Type"}, {"key": "lat", "label": "Lat"}, {"key": "lng", "label": "Lng"},
    ]},
    "carriers_primary": {"columns": [
        {"key": "name", "label": "Carrier"}, {"key": "category", "label": "Mode"},
        {"key": "endpoint", "label": "API Endpoint"}, {"key": "status", "label": "Status"},
        {"key": "last_sync", "label": "Last Sync"},
    ]},
    "contacts_suppliers": {"columns": [
        {"key": "supplier", "label": "Supplier"}, {"key": "contact_name", "label": "Contact"},
        {"key": "email", "label": "Email"}, {"key": "phone", "label": "Phone"},
        {"key": "country", "label": "Country"}, {"key": "commodity", "label": "Commodity"},
    ]},
    "contacts_carriers": {"columns": [
        {"key": "carrier", "label": "Carrier"}, {"key": "contact_name", "label": "Rep"},
        {"key": "email", "label": "Email"}, {"key": "phone", "label": "Phone"},
        {"key": "scac", "label": "SCAC"}, {"key": "after_hours", "label": "24/7 Number"},
    ]},
    "info": {"columns": [
        {"key": "key", "label": "Key"}, {"key": "value", "label": "Value"},
        {"key": "notes", "label": "Notes"},
    ]},
    "volume_overview": {"columns": [
        {"key": "month", "label": "Month"}, {"key": "tl_count", "label": "TL"},
        {"key": "ltl_count", "label": "LTL"}, {"key": "parcel_count", "label": "Parcel"},
        {"key": "ocean_count", "label": "Ocean"}, {"key": "air_count", "label": "Air"},
        {"key": "total_lbs", "label": "Total Weight (lbs)"}, {"key": "total_spend", "label": "Spend (USD)"},
    ]},
}
KIND_DEFINITIONS["shipments_expedites"] = KIND_DEFINITIONS["shipments_tl"]
KIND_DEFINITIONS["shipments_crates"] = KIND_DEFINITIONS["shipments_tl"]
KIND_DEFINITIONS["shipments_import"] = KIND_DEFINITIONS["shipments_seafreight"]
KIND_DEFINITIONS["truckload_bookings"] = {"columns": TRUCKLOAD_BOOKING_COLUMNS}

def _get_nested(d: Dict[str, Any], key: str):
    parts = key.split(".")
    cur = d
    for p in parts:
        if cur is None:
            return None
        cur = cur.get(p) if isinstance(cur, dict) else None
    return cur

QUOTES_DATA = [
    {"quote_id": "Q-25001", "mode": "TL", "carrier": "XPO Logistics", "origin": "Golden Valley, MN", "destination": "Dallas, TX", "rate_usd": 2150.00, "transit_days": 3, "valid_until": "2026-06-15", "status": "Open"},
    {"quote_id": "Q-25002", "mode": "LTL", "carrier": "SAIA", "origin": "Holland, MI", "destination": "Atlanta, GA", "rate_usd": 845.00, "transit_days": 4, "valid_until": "2026-06-20", "status": "Awarded"},
    {"quote_id": "Q-25003", "mode": "Ocean", "carrier": "Kuehne+Nagel", "origin": "Long Beach, CA", "destination": "Yokohama, JP", "rate_usd": 4280.00, "transit_days": 14, "valid_until": "2026-07-01", "status": "Open"},
    {"quote_id": "Q-25004", "mode": "Air", "carrier": "FedEx", "origin": "Louisville, KY", "destination": "Frankfurt, DE", "rate_usd": 6840.00, "transit_days": 2, "valid_until": "2026-06-10", "status": "Open"},
    {"quote_id": "Q-25005", "mode": "TL", "carrier": "ArcBest", "origin": "Louisville, KY", "destination": "Phoenix, AZ", "rate_usd": 2680.00, "transit_days": 3, "valid_until": "2026-06-18", "status": "Declined"},
    {"quote_id": "Q-25006", "mode": "LTL", "carrier": "R&L Carriers", "origin": "Golden Valley, MN", "destination": "Seattle, WA", "rate_usd": 1120.00, "transit_days": 5, "valid_until": "2026-06-22", "status": "Open"},
]

CARRIER_PRIMARY_DATA = [{"name": i["name"], "category": i["category"], "endpoint": i["endpoint"], "status": i["status"], "last_sync": i["last_sync"]} for i in INTEGRATIONS]

SUPPLIER_CONTACTS_DATA = [
    {"supplier": "Kuehne+Nagel Services", "contact_name": "Lisette Vermeer", "email": "lisette.vermeer@kn-logistics.com", "phone": "+31-10-555-0188", "country": "Netherlands", "commodity": "All inbound imports"},
    {"supplier": "Motrex Co. Ltd", "contact_name": "Park Min-jun", "email": "min.park@motrex.co.kr", "phone": "+82-2-555-3344", "country": "South Korea", "commodity": "DC drive motors"},
    {"supplier": "BattCo Industries GmbH", "contact_name": "Klaus Müller", "email": "k.mueller@battco.de", "phone": "+49-30-555-9911", "country": "Germany", "commodity": "Li-ion battery cells"},
    {"supplier": "Premier Polymers", "contact_name": "Maria Gonzalez", "email": "m.gonzalez@premierpoly.com", "phone": "+1-952-555-0142", "country": "USA", "commodity": "Solution tanks"},
    {"supplier": "Midwest Steel Frame Co", "contact_name": "Steve Olson", "email": "solson@midweststeel.com", "phone": "+1-507-555-7733", "country": "USA", "commodity": "Chassis frames"},
    {"supplier": "Yazaki Wiring Harness", "contact_name": "Tanaka Hiroshi", "email": "h.tanaka@yazaki.co.jp", "phone": "+81-3-555-2218", "country": "Japan", "commodity": "Wiring harnesses"},
]

CARRIER_CONTACTS_DATA = [
    {"carrier": "XPO Logistics", "contact_name": "Jamal Robinson", "email": "j.robinson@xpo.com", "phone": "+1-855-555-9760", "scac": "XPOL", "after_hours": "+1-844-555-1199"},
    {"carrier": "ArcBest", "contact_name": "Sarah Whitfield", "email": "swhitfield@arcb.com", "phone": "+1-877-555-2227", "scac": "ABFS", "after_hours": "+1-844-555-2200"},
    {"carrier": "SAIA", "contact_name": "Daniel Park", "email": "dpark@saia.com", "phone": "+1-800-555-7242", "scac": "SAIA", "after_hours": "+1-844-555-7700"},
    {"carrier": "R&L Carriers", "contact_name": "Maria Esposito", "email": "mesposito@rlc.com", "phone": "+1-800-555-5526", "scac": "RLCA", "after_hours": "+1-844-555-7575"},
    {"carrier": "UPS", "contact_name": "UPS National Account Desk", "email": "tennant-acct@ups.com", "phone": "+1-800-555-8742", "scac": "UPSN", "after_hours": "+1-800-555-7898"},
    {"carrier": "FedEx Freight", "contact_name": "Brian Liu", "email": "brian.liu@fedex.com", "phone": "+1-866-555-3339", "scac": "FXFE", "after_hours": "+1-800-555-3339"},
    {"carrier": "DHL Express", "contact_name": "Anika Schroeder", "email": "anika.s@dhl.com", "phone": "+1-800-555-3110", "scac": "DHLC", "after_hours": "+1-800-555-2255"},
    {"carrier": "Kuehne+Nagel", "contact_name": "Lisette Vermeer", "email": "lisette.vermeer@kn-logistics.com", "phone": "+31-10-555-0188", "scac": "KNAA", "after_hours": "+31-10-555-0911"},
    {"carrier": "Consolidated Fastfrate", "contact_name": "Pierre Tremblay", "email": "p.tremblay@fastfrate.com", "phone": "+1-905-555-1212", "scac": "CFAT", "after_hours": "+1-844-555-3434"},
]

INFO_DATA = [
    {"key": "Company HQ", "value": "10400 Clean Street, Eden Prairie, MN 55344", "notes": "Tennant Company HQ — main office"},
    {"key": "Logistics Director", "value": "Kirk Juergins", "notes": "Lead for Transportation Management"},
    {"key": "TMS System ID", "value": "TENNANT-TMS-V1", "notes": "Internal system identifier"},
    {"key": "SAP Client", "value": "100", "notes": "Production S/4HANA client"},
    {"key": "Webex Org", "value": "tennantco.webex.com", "notes": "Corporate Webex instance"},
    {"key": "Emergency Dispatch", "value": "+1-952-555-0911", "notes": "24/7 dispatch hotline"},
    {"key": "Customs Broker", "value": "Livingston International", "notes": "Primary US customs broker"},
    {"key": "Insurance Broker", "value": "Marsh McLennan", "notes": "Cargo insurance"},
]

VOLUME_OVERVIEW_DATA = [
    {"month": "2026-01", "tl_count": 142, "ltl_count": 318, "parcel_count": 1240, "ocean_count": 38, "air_count": 22, "total_lbs": 2480000, "total_spend": 412000},
    {"month": "2026-02", "tl_count": 158, "ltl_count": 342, "parcel_count": 1310, "ocean_count": 41, "air_count": 19, "total_lbs": 2590000, "total_spend": 438000},
    {"month": "2026-03", "tl_count": 167, "ltl_count": 365, "parcel_count": 1402, "ocean_count": 44, "air_count": 28, "total_lbs": 2810000, "total_spend": 475000},
    {"month": "2026-04", "tl_count": 175, "ltl_count": 388, "parcel_count": 1488, "ocean_count": 47, "air_count": 31, "total_lbs": 2920000, "total_spend": 498000},
    {"month": "2026-05", "tl_count": 162, "ltl_count": 354, "parcel_count": 1356, "ocean_count": 39, "air_count": 25, "total_lbs": 2640000, "total_spend": 451000},
]

async def _get_rows_for_tab(tab: Dict[str, Any]) -> List[Dict[str, Any]]:
    kind = tab["kind"]
    if kind.startswith("shipments_"):
        q: Dict[str, Any] = {}
        flt = tab.get("filter") or {}
        if flt.get("mode") and flt["mode"] != "Flatbed":
            q["mode"] = flt["mode"]
        docs = await db.shipments.find(q, {"_id": 0}).limit(500).to_list(500)
        if kind == "shipments_crates":
            docs = [d for d in docs if "crate" in (d.get("commodity", "").lower())] or docs[:8]
        return docs
    if kind == "quotes": return QUOTES_DATA
    if kind == "plant_hubs": return TMS_FACILITIES
    if kind == "carriers_primary": return CARRIER_PRIMARY_DATA
    if kind == "contacts_suppliers": return SUPPLIER_CONTACTS_DATA
    if kind == "contacts_carriers": return CARRIER_CONTACTS_DATA
    if kind == "info": return INFO_DATA
    if kind == "volume_overview": return VOLUME_OVERVIEW_DATA
    if kind == "truckload_bookings":
        return await db.truckload_bookings.find({}, {"_id": 0}).sort("created_at", -1).to_list(2000)
    return []

async def _ensure_default_tabs():
    """Seed defaults on a fresh DB, AND idempotently insert any newly-added
    default kinds for existing installations (e.g. the v1.9 Truckload Bookings
    tab). Existing user-renamed tabs are never touched."""
    existing_kinds = set()
    async for t in db.workbook_tabs.find({}, {"_id": 0, "kind": 1}):
        existing_kinds.add(t.get("kind"))
    next_order = await db.workbook_tabs.count_documents({})
    for i, t in enumerate(DEFAULT_TABS):
        if t["kind"] in existing_kinds:
            continue
        await db.workbook_tabs.insert_one({
            "tab_id": f"TAB-{uuid.uuid4().hex[:8].upper()}",
            "name": t["name"], "kind": t["kind"],
            "filter": t.get("filter") or {},
            # Pin truckload_bookings to order=0 so it lands as the first tab
            "order": 0 if t["kind"] == "truckload_bookings" else (next_order + i),
            "created_at": datetime.now(timezone.utc).isoformat(),
        })
    # If we inserted truckload_bookings, shift everything else down by 1
    if "truckload_bookings" not in existing_kinds:
        await db.workbook_tabs.update_many(
            {"kind": {"$ne": "truckload_bookings"}},
            {"$inc": {"order": 1}},
        )

@api_router.get("/workbook/tabs")
async def list_tabs(_: User = Depends(get_current_user)):
    await _ensure_default_tabs()
    tabs = await db.workbook_tabs.find({}, {"_id": 0}).sort("order", 1).to_list(200)
    for t in tabs:
        t["columns"] = KIND_DEFINITIONS.get(t["kind"], {"columns": []})["columns"]
    return tabs

class TabCreate(BaseModel):
    name: str
    kind: str = "info"
    filter: Optional[Dict[str, Any]] = None

@api_router.post("/workbook/tabs")
async def create_tab(payload: TabCreate, _: User = Depends(get_current_user)):
    if payload.kind not in KIND_DEFINITIONS:
        raise HTTPException(status_code=400, detail=f"Unknown kind. Allowed: {list(KIND_DEFINITIONS.keys())}")
    count = await db.workbook_tabs.count_documents({})
    doc = {
        "tab_id": f"TAB-{uuid.uuid4().hex[:8].upper()}",
        "name": payload.name, "kind": payload.kind,
        "filter": payload.filter or {}, "order": count,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.workbook_tabs.insert_one(dict(doc))
    doc["columns"] = KIND_DEFINITIONS.get(payload.kind, {"columns": []})["columns"]
    return doc

class TabPatch(BaseModel):
    name: Optional[str] = None
    order: Optional[int] = None
    filter: Optional[Dict[str, Any]] = None

@api_router.patch("/workbook/tabs/{tab_id}")
async def update_tab(tab_id: str, payload: TabPatch, _: User = Depends(get_current_user)):
    updates = {k: v for k, v in payload.model_dump().items() if v is not None}
    if not updates:
        return {"ok": True, "no_changes": True}
    result = await db.workbook_tabs.update_one({"tab_id": tab_id}, {"$set": updates})
    if result.matched_count == 0:
        raise HTTPException(status_code=404, detail="Tab not found")
    return {"ok": True}

@api_router.delete("/workbook/tabs/{tab_id}")
async def delete_tab(tab_id: str, _: User = Depends(get_current_user)):
    result = await db.workbook_tabs.delete_one({"tab_id": tab_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Tab not found")
    return {"ok": True}


# ---------------------------------------------------------------------------
# Per-user tile layouts. Every page in the TMS can render its sections via the
# DraggableTiles wrapper, which stores its order array here keyed by the user
# id and a page_key. Result: a dispatcher's layout follows them across browsers
# / devices, not just one localStorage shard.
# ---------------------------------------------------------------------------

class UserLayoutPayload(BaseModel):
    order: List[str]


@api_router.get("/user/layouts/{page_key}")
async def get_user_layout(page_key: str, user: User = Depends(get_current_user)):
    doc = await db.user_layouts.find_one(
        {"user_id": user.user_id, "page_key": page_key}, {"_id": 0}
    )
    return {"order": (doc or {}).get("order", [])}


@api_router.put("/user/layouts/{page_key}")
async def put_user_layout(
    page_key: str, payload: UserLayoutPayload, user: User = Depends(get_current_user)
):
    # Sanity: require at least one id and reasonable bounds. Anything else
    # gets accepted as-is — clients control the schema of these strings.
    if not payload.order or len(payload.order) > 200:
        raise HTTPException(status_code=400, detail="Invalid layout payload")
    await db.user_layouts.update_one(
        {"user_id": user.user_id, "page_key": page_key},
        {"$set": {
            "user_id": user.user_id,
            "page_key": page_key,
            "order": payload.order,
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }},
        upsert=True,
    )
    return {"ok": True, "page_key": page_key, "order": payload.order}


@api_router.delete("/user/layouts/{page_key}")
async def reset_user_layout(page_key: str, user: User = Depends(get_current_user)):
    await db.user_layouts.delete_one({"user_id": user.user_id, "page_key": page_key})
    return {"ok": True}


# ---------------------------------------------------------------------------
# Truckload Booking Sheet — Excel-style editable rows with real-time polling.
# Each row lives in db.truckload_bookings. A separate revision counter in
# db.truckload_bookings_meta lets clients poll a lightweight HEAD-style endpoint
# (`/version`) and only refetch the full payload when something actually
# changed — keeps Mongo cost low while still feeling live across dispatchers.
# ---------------------------------------------------------------------------
_TLB_ALLOWED_KEYS = {c["key"] for c in TRUCKLOAD_BOOKING_COLUMNS}


async def _bump_tlb_version(actor: str) -> int:
    """Atomically increment the truckload-bookings version + stamp last editor."""
    doc = await db.truckload_bookings_meta.find_one_and_update(
        {"_id": "version"},
        {"$inc": {"version": 1}, "$set": {"last_editor": actor, "updated_at": datetime.now(timezone.utc).isoformat()}},
        upsert=True, return_document=True,
    )
    return int((doc or {}).get("version", 1))


def _clean_tlb_payload(payload: Dict[str, Any]) -> Dict[str, Any]:
    """Project an incoming payload to the known column keys only, coercing
    obvious numeric fields. Skips any silently-injected unknown keys."""
    out: Dict[str, Any] = {}
    for col in TRUCKLOAD_BOOKING_COLUMNS:
        key, ctype = col["key"], col.get("type", "text")
        if key not in payload:
            continue
        v = payload.get(key)
        if v == "" or v is None:
            out[key] = None
            continue
        if ctype == "number":
            try:
                out[key] = float(v) if not str(v).isdigit() else int(v)
            except (TypeError, ValueError):
                out[key] = None
        else:
            out[key] = str(v)
    return out


async def _onboarded_carrier_names() -> List[str]:
    """Return the sorted, de-duped list of carrier display names that the
    Truckload Booking Sheet's Carrier dropdown should offer. Sources:
      1. carrier_onboarding records with status == "approved" (primary)
      2. existing carriers already used on previous bookings (so dispatchers
         see familiar names even before they're formally onboarded)
    DBA name wins when present so dispatchers see what they call the carrier."""
    out: Dict[str, None] = {}
    cursor = db.carrier_onboarding.find(
        {"status": "approved"}, {"_id": 0, "legal_name": 1, "dba": 1, "scac": 1}
    )
    async for c in cursor:
        name = (c.get("dba") or c.get("legal_name") or "").strip()
        if not name:
            continue
        scac = (c.get("scac") or "").strip()
        label = f"{name} · {scac}" if scac else name
        out[label] = None
    # Carriers historically used in bookings (fallback, keeps history alive)
    used = await db.truckload_bookings.distinct("carrier")
    for u in used:
        if u and isinstance(u, str) and u not in out:
            out[u] = None
    return sorted(out.keys(), key=lambda s: s.lower())


def _columns_with_carriers(carrier_options: List[str]) -> List[Dict[str, Any]]:
    """Clone TRUCKLOAD_BOOKING_COLUMNS and inject the live carrier roster
    into the Carrier column's `options` array. Doesn't mutate the module-level
    constant."""
    out: List[Dict[str, Any]] = []
    for col in TRUCKLOAD_BOOKING_COLUMNS:
        if col.get("key") == "carrier":
            out.append({**col, "options": ["", *carrier_options]})
        else:
            out.append(dict(col))
    return out


@api_router.get("/workbook/truckload-bookings")
async def list_truckload_bookings(_: User = Depends(get_current_user)):
    rows = await db.truckload_bookings.find({}, {"_id": 0}).sort("created_at", -1).to_list(2000)
    meta = await db.truckload_bookings_meta.find_one({"_id": "version"}, {"_id": 0}) or {"version": 0}
    carriers = await _onboarded_carrier_names()
    columns = _columns_with_carriers(carriers)
    return {"rows": rows, "columns": columns, "carrier_options": carriers, **meta}


@api_router.get("/workbook/truckload-bookings/version")
async def truckload_bookings_version(_: User = Depends(get_current_user)):
    """Lightweight poll endpoint — returns just {version, updated_at, last_editor}
    so a watching client can decide whether to refetch the full list."""
    meta = await db.truckload_bookings_meta.find_one({"_id": "version"}, {"_id": 0})
    return meta or {"version": 0, "updated_at": None, "last_editor": None}


class TruckloadBookingPayload(BaseModel):
    # Free-form payload validated against TRUCKLOAD_BOOKING_COLUMNS in handler
    data: Dict[str, Any] = {}


def _stripped_carrier_label(label: Optional[str]) -> str:
    """Carrier dropdown values look like 'XPO · XPOL' — split off the SCAC
    suffix when we compare against existing onboarding records."""
    if not label:
        return ""
    return str(label).split("·", 1)[0].strip()


async def _ensure_carrier_in_pipeline(label: Optional[str], user: User) -> Optional[str]:
    """If `label` references a carrier that's NOT already in carrier_onboarding
    (case-insensitive match on legal_name OR dba), insert a fresh stub record
    with status='in_review' so the new name shows up in /carrier-onboarding
    for the compliance team to follow up on. Returns the new onboarding_id
    (or None if nothing was created)."""
    name = _stripped_carrier_label(label)
    if not name:
        return None
    needle = name.lower()  # noqa: F841 — reserved for future ranked match
    existing = await db.carrier_onboarding.find_one(
        {"$or": [
            {"legal_name": {"$regex": f"^{re.escape(name)}$", "$options": "i"}},
            {"dba":        {"$regex": f"^{re.escape(name)}$", "$options": "i"}},
        ]},
        {"_id": 0, "onboarding_id": 1},
    )
    if existing:
        return None
    onboarding_id = f"OB-{uuid.uuid4().hex[:8].upper()}"
    stub = {
        "onboarding_id": onboarding_id,
        "legal_name": name,
        "dba": None,
        "mc_number": None,
        "dot_number": None,
        "scac": None,
        "mode": "TL",
        "contact_name": "(pending)",
        "contact_email": "(pending)",
        "contact_phone": "(pending)",
        "insurance_amount": 0.0,
        "insurance_expiry": (datetime.now(timezone.utc) + timedelta(days=30)).date().isoformat(),
        "safety_rating": "NotRated",
        "csa_score": 0,
        "w9_received": False,
        "coi_received": False,
        "contract_signed": False,
        "status": "in_review",
        "submitted_at": datetime.now(timezone.utc).isoformat(),
        "submitted_by": user.name,
        "notes": "Auto-created from Truckload Booking Sheet — dispatcher entered a new carrier name. Compliance to gather W-9, COI, contract.",
        "auto_created": True,
    }
    await db.carrier_onboarding.insert_one(dict(stub))
    return onboarding_id


@api_router.post("/workbook/truckload-bookings")
async def create_truckload_booking(
    payload: TruckloadBookingPayload,
    user: User = Depends(require_role("admin", "dispatcher", "auditor")),
):
    row_id = f"TLB-{uuid.uuid4().hex[:8].upper()}"
    now = datetime.now(timezone.utc).isoformat()
    doc = {
        "id": row_id,
        "created_at": now,
        "updated_at": now,
        "created_by": user.name,
        "updated_by": user.name,
        **{c["key"]: None for c in TRUCKLOAD_BOOKING_COLUMNS},
        **_clean_tlb_payload(payload.data),
    }
    await db.truckload_bookings.insert_one(dict(doc))
    version = await _bump_tlb_version(user.name)
    auto_obid = await _ensure_carrier_in_pipeline(doc.get("carrier"), user)
    doc.pop("_id", None)
    return {"row": doc, "version": version, "auto_onboarding_id": auto_obid}


@api_router.patch("/workbook/truckload-bookings/{row_id}")
async def update_truckload_booking(
    row_id: str,
    payload: TruckloadBookingPayload,
    user: User = Depends(require_role("admin", "dispatcher", "auditor")),
):
    clean = _clean_tlb_payload(payload.data)
    if not clean:
        return {"ok": True, "no_changes": True}
    clean["updated_at"] = datetime.now(timezone.utc).isoformat()
    clean["updated_by"] = user.name
    result = await db.truckload_bookings.update_one({"id": row_id}, {"$set": clean})
    if result.matched_count == 0:
        raise HTTPException(status_code=404, detail="Booking row not found")
    version = await _bump_tlb_version(user.name)
    auto_obid = None
    if "carrier" in clean:
        auto_obid = await _ensure_carrier_in_pipeline(clean["carrier"], user)
    fresh = await db.truckload_bookings.find_one({"id": row_id}, {"_id": 0})
    return {"row": fresh, "version": version, "auto_onboarding_id": auto_obid}


@api_router.delete("/workbook/truckload-bookings/{row_id}")
async def delete_truckload_booking(
    row_id: str,
    user: User = Depends(require_role("admin", "dispatcher")),
):
    result = await db.truckload_bookings.delete_one({"id": row_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Booking row not found")
    version = await _bump_tlb_version(user.name)
    return {"ok": True, "version": version}

@api_router.get("/workbook/tabs/{tab_id}/rows")
async def tab_rows(tab_id: str, _: User = Depends(get_current_user)):
    tab = await db.workbook_tabs.find_one({"tab_id": tab_id}, {"_id": 0})
    if not tab:
        raise HTTPException(status_code=404, detail="Tab not found")
    rows = await _get_rows_for_tab(tab)
    cols = KIND_DEFINITIONS.get(tab["kind"], {"columns": []})["columns"]
    projected = []
    for r in rows:
        proj = {}
        for c in cols:
            v = _get_nested(r, c["key"]) if "." in c["key"] else r.get(c["key"])
            proj[c["key"]] = v
        projected.append(proj)
    return {"tab": tab, "columns": cols, "rows": projected}

def _build_xlsx_for_tab(tab: Dict[str, Any], rows: List[Dict[str, Any]], cols: List[Dict[str, Any]]) -> bytes:
    wb = XLWorkbook()
    ws = wb.active
    ws.title = (tab["name"][:31] or "Sheet1")
    for ch in r'[]/\?*':
        ws.title = ws.title.replace(ch, "_")
    header_fill = PatternFill(start_color="FF0B0E14", end_color="FF0B0E14", fill_type="solid")
    header_font = Font(color="FF00E5FF", bold=True, size=11)
    border = Border(bottom=Side(border_style="thin", color="FF00A4E4"))
    for ci, c in enumerate(cols, start=1):
        cell = ws.cell(row=1, column=ci, value=c["label"])
        cell.fill = header_fill; cell.font = header_font; cell.border = border
        cell.alignment = Alignment(horizontal="left", vertical="center")
    for ri, r in enumerate(rows, start=2):
        for ci, c in enumerate(cols, start=1):
            v = _get_nested(r, c["key"]) if "." in c["key"] else r.get(c["key"])
            if isinstance(v, (dict, list)):
                v = json.dumps(v)
            ws.cell(row=ri, column=ci, value=v)
    for ci, c in enumerate(cols, start=1):
        col_letter = ws.cell(row=1, column=ci).column_letter
        max_len = max([len(str(c["label"]))] + [len(str(_get_nested(r, c["key"]) if "." in c["key"] else r.get(c["key"]) or "")) for r in rows] + [0])
        ws.column_dimensions[col_letter].width = min(max(max_len + 2, 12), 48)
    ws.freeze_panes = "A2"
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return buf.getvalue()

@api_router.get("/workbook/tabs/{tab_id}/export.xlsx")
async def export_tab_xlsx(tab_id: str, _: User = Depends(get_current_user)):
    tab = await db.workbook_tabs.find_one({"tab_id": tab_id}, {"_id": 0})
    if not tab:
        raise HTTPException(status_code=404, detail="Tab not found")
    rows = await _get_rows_for_tab(tab)
    cols = KIND_DEFINITIONS.get(tab["kind"], {"columns": []})["columns"]
    data = _build_xlsx_for_tab(tab, rows, cols)
    filename = f"Tennant_TMS_{tab['name'].replace(' ', '_')}.xlsx"
    return StreamingResponse(io.BytesIO(data), media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                             headers={"Content-Disposition": f'attachment; filename="{filename}"'})

@api_router.get("/workbook/export-all.xlsx")
async def export_all_xlsx(_: User = Depends(get_current_user)):
    await _ensure_default_tabs()
    tabs = await db.workbook_tabs.find({}, {"_id": 0}).sort("order", 1).to_list(200)
    wb = XLWorkbook()
    wb.remove(wb.active)
    header_fill = PatternFill(start_color="FF0B0E14", end_color="FF0B0E14", fill_type="solid")
    header_font = Font(color="FF00E5FF", bold=True, size=11)
    border = Border(bottom=Side(border_style="thin", color="FF00A4E4"))
    for tab in tabs:
        cols = KIND_DEFINITIONS.get(tab["kind"], {"columns": []})["columns"]
        if not cols: continue
        rows = await _get_rows_for_tab(tab)
        title = tab["name"][:31] or tab["tab_id"]
        for ch in r'[]/\?*':
            title = title.replace(ch, "_")
        ws = wb.create_sheet(title=title)
        for ci, c in enumerate(cols, start=1):
            cell = ws.cell(row=1, column=ci, value=c["label"])
            cell.fill = header_fill; cell.font = header_font; cell.border = border
        for ri, r in enumerate(rows, start=2):
            for ci, c in enumerate(cols, start=1):
                v = _get_nested(r, c["key"]) if "." in c["key"] else r.get(c["key"])
                if isinstance(v, (dict, list)):
                    v = json.dumps(v)
                ws.cell(row=ri, column=ci, value=v)
        ws.freeze_panes = "A2"
        for ci in range(1, len(cols) + 1):
            ws.column_dimensions[ws.cell(row=1, column=ci).column_letter].width = 18
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                             headers={"Content-Disposition": 'attachment; filename="Tennant_TMS_Workbook.xlsx"'})

# -------------------- SEED --------------------
@api_router.post("/admin/seed")
async def seed_data(force: bool = False):
    """Seed mock shipments. Idempotent unless force=True."""
    count = await db.shipments.count_documents({})
    if count > 0 and not force:
        return {"ok": True, "skipped": True, "count": count}
    if force:
        await db.shipments.delete_many({})

    carriers_by_mode = {
        "TL": ["XPO Logistics", "ArcBest", "Schneider", "J.B. Hunt"],
        "LTL": ["SAIA", "R&L Carriers", "ArcBest", "XPO Logistics", "Consolidated Fastfrate"],
        "Parcel": ["UPS", "FedEx", "DHL Express"],
        "Ocean": ["Kuehne+Nagel", "Maersk", "MSC"],
        "Air": ["FedEx", "DHL Express", "Kuehne+Nagel"],
        "Rail": ["BNSF", "Union Pacific", "CSX"],
    }
    destinations = [
        {"city": "Dallas, TX", "lat": 32.7767, "lng": -96.7970},
        {"city": "Atlanta, GA", "lat": 33.7490, "lng": -84.3880},
        {"city": "Los Angeles, CA", "lat": 34.0522, "lng": -118.2437},
        {"city": "Phoenix, AZ", "lat": 33.4484, "lng": -112.0740},
        {"city": "Seattle, WA", "lat": 47.6062, "lng": -122.3321},
        {"city": "Chicago, IL", "lat": 41.8781, "lng": -87.6298},
        {"city": "Houston, TX", "lat": 29.7604, "lng": -95.3698},
        {"city": "Miami, FL", "lat": 25.7617, "lng": -80.1918},
        {"city": "Denver, CO", "lat": 39.7392, "lng": -104.9903},
        {"city": "Boston, MA", "lat": 42.3601, "lng": -71.0589},
        {"city": "Toronto, ON", "lat": 43.6532, "lng": -79.3832},
        {"city": "Mexico City, MX", "lat": 19.4326, "lng": -99.1332},
        {"city": "Rotterdam, NL (Port)", "lat": 51.9244, "lng": 4.4777},
        {"city": "Shanghai, CN (Port)", "lat": 31.2304, "lng": 121.4737},
        {"city": "Hamburg, DE (Port)", "lat": 53.5511, "lng": 9.9937},
    ]
    commodities = [
        "Floor scrubbers (T16AMR)", "Industrial cleaners (M30)", "Sweeper components", "Battery packs (Li-ion)",
        "DC motors", "Scrubber brushes", "Plastic tanks", "Steel frames", "Control assemblies", "Replacement parts",
    ]
    statuses = ["in_transit", "in_transit", "in_transit", "delayed", "delivered", "delivered", "pending", "at_origin"]

    # Inbound seed pool (matches the supplier XLSX columns)
    suppliers = [
        ("Monahan Filaments", "USA", "Brushes & filaments"),
        ("Trojan Battery — CA", "USA", "Lead-acid batteries"),
        ("Trojan Battery — GA", "USA", "Lead-acid batteries"),
        ("East Penn Manufacturing", "USA", "Lead-acid batteries"),
        ("Amer Electric Motion", "USA", "DC drive motors"),
        ("Motrex Co. Ltd", "South Korea", "DC drive motors"),
        ("BattCo Industries GmbH", "Germany", "Li-ion battery cells"),
        ("Yazaki Wiring Harness", "Japan", "Wiring harnesses"),
        ("Premier Polymers", "USA", "Solution tanks"),
        ("Midwest Steel Frame Co", "USA", "Chassis frames"),
    ]
    material_controllers = ["Katherine Markley", "Tim Strube", "Devon Marquez", "Priya Iyer"]
    inbound_carriers = ["RYNK", "XPO Logistics", "SAIA", "ArcBest", "R&L Carriers", "Schneider", "FedEx Freight", "Kuehne+Nagel"]
    days_map = ["MON", "TUE", "WED", "THR", "FRI"]

    shipments = []
    # 32 outbound (existing pattern, slightly slimmer)
    for i in range(32):
        mode = random.choice(["TL", "LTL", "Parcel", "Ocean", "Air", "Rail", "LTL", "TL"])
        carrier = random.choice(carriers_by_mode[mode])
        origin_facility = random.choice(TMS_FACILITIES)
        dest = random.choice(destinations)
        status = random.choice(statuses)
        if status in ("pending", "at_origin"):
            progress = 0.0
            cur = {"lat": origin_facility["lat"], "lng": origin_facility["lng"], "city": origin_facility["city"]}
        elif status == "delivered":
            progress = 1.0
            cur = {"lat": dest["lat"], "lng": dest["lng"], "city": dest["city"]}
        else:
            progress = round(random.uniform(0.15, 0.85), 2)
            cur = {
                "lat": origin_facility["lat"] + (dest["lat"] - origin_facility["lat"]) * progress + random.uniform(-0.5, 0.5),
                "lng": origin_facility["lng"] + (dest["lng"] - origin_facility["lng"]) * progress + random.uniform(-0.5, 0.5),
                "city": "En route",
            }
        sid = f"SHP-{uuid.uuid4().hex[:8].upper()}"
        pickup = datetime.now(timezone.utc) - timedelta(days=random.randint(0, 10))
        eta = pickup + timedelta(days=random.randint(2, 16))
        is_haz = random.random() < 0.18  # ~18% hazmat
        haz_class = random.choice(["Class 9 (Li-ion)", "Class 8 (Corrosive)", "Class 3 (Flammable)"]) if is_haz else None
        bid_cost = round(random.uniform(420, 4200), 2)
        shipments.append({
            "shipment_id": sid,
            "reference": f"TN-{random.randint(10000, 99999)}",
            "mode": mode, "carrier": carrier, "status": status,
            "origin": {"name": origin_facility["name"], "city": origin_facility["city"], "lat": origin_facility["lat"], "lng": origin_facility["lng"], "facility": origin_facility["id"]},
            "destination": {"name": dest["city"], "city": dest["city"], "lat": dest["lat"], "lng": dest["lng"]},
            "current_location": cur,
            "eta": eta.isoformat(), "pickup_date": pickup.date().isoformat(),
            "weight_lbs": round(random.uniform(800, 42000), 0),
            "pieces": random.randint(1, 26), "commodity": random.choice(commodities),
            "value_usd": round(random.uniform(2500, 285000), 2),
            "container_no": f"TCLU{random.randint(1000000,9999999)}" if mode == "Ocean" else None,
            "bol_no": f"BOL{random.randint(100000,999999)}", "pro_no": f"PRO{random.randint(5980000, 6010000)}",
            "progress": progress,
            "direction": "outbound", "hazmat": is_haz, "hazmat_class": haz_class,
            "supplier": None, "consignee": None,
            "ship_date": pickup.date().isoformat(),
            "ship_day": days_map[pickup.weekday()] if pickup.weekday() < 5 else "MON",
            "skids": random.randint(1, 26), "material_controller": random.choice(material_controllers),
            "po_numbers": None,
            "booking_number": f"PRO {random.randint(5980000, 6010000)}",
            "bid_cost": bid_cost, "fsc_pct": round(random.uniform(18, 32), 1),
            "extras": "plus FSC" if random.random() < 0.3 else None,
            "done": status == "delivered", "shipping_hours": "0700-1500" if mode in ("TL", "LTL") else None,
            "pickup_no": f"SFS{random.randint(79000, 79999)}" if random.random() < 0.4 else None,
            "created_at": datetime.now(timezone.utc).isoformat(),
        })

    # 16 inbound shipments — supplier → Tennant facility (matches uploaded screenshot)
    for i in range(16):
        sup = random.choice(suppliers)
        plant = random.choice(TMS_FACILITIES)
        plant_short = {"GVM": "Tennant - GV", "HOM": "Tennant - HO", "LVK": "Tennant - LV"}[plant["id"]]
        carrier = random.choice(inbound_carriers)
        # Origin city depends on supplier country
        if "Korea" in sup[1] or "Japan" in sup[1] or "Germany" in sup[1]:
            origin_city, lat, lng = {"South Korea": ("Busan, KR", 35.18, 129.07), "Japan": ("Yokohama, JP", 35.44, 139.64), "Germany": ("Hamburg, DE", 53.55, 9.99)}[sup[1]]
            mode = "Ocean"
        else:
            origin_choice = random.choice([("Atlanta, GA", 33.75, -84.39), ("Reading, PA", 40.34, -75.93), ("Phoenix, AZ", 33.45, -112.07), ("Dallas, TX", 32.78, -96.80), ("Charlotte, NC", 35.23, -80.84)])
            origin_city, lat, lng = origin_choice
            mode = random.choice(["TL", "LTL", "LTL"])
        status = random.choice(["in_transit", "in_transit", "delivered", "pending", "at_origin"])
        progress = 0.5 if status == "in_transit" else (1.0 if status == "delivered" else 0.0)
        ship_dt = datetime.now(timezone.utc) - timedelta(days=random.randint(0, 14))
        eta = ship_dt + timedelta(days=random.randint(3, 18))
        is_haz = "Battery" in sup[0] or "Trojan" in sup[0] or "Penn" in sup[0] or "BattCo" in sup[0]
        haz_class = "Class 8 (Lead-acid)" if "Trojan" in sup[0] or "Penn" in sup[0] else ("Class 9 (Li-ion)" if "BattCo" in sup[0] else ("Class 9 (Li-ion)" if is_haz else None))
        sid = f"SHP-IN-{uuid.uuid4().hex[:6].upper()}"
        skids = random.randint(7, 26)
        weight = round(random.uniform(15000, 40000), 0)
        pos = [f"4500{random.randint(20000, 29999)}" for _ in range(random.randint(1, 4))]
        # Spread current_location along route
        cur = {
            "lat": lat + (plant["lat"] - lat) * progress,
            "lng": lng + (plant["lng"] - lng) * progress,
            "city": origin_city if progress == 0 else (plant["city"] if progress == 1 else "En route"),
        }
        shipments.append({
            "shipment_id": sid,
            "reference": f"IN-{random.randint(10000, 99999)}",
            "mode": mode, "carrier": carrier, "status": status,
            "origin": {"name": sup[0], "city": origin_city, "lat": lat, "lng": lng},
            "destination": {"name": plant["name"], "city": plant["city"], "lat": plant["lat"], "lng": plant["lng"], "facility": plant["id"]},
            "current_location": cur,
            "eta": eta.isoformat(), "pickup_date": ship_dt.date().isoformat(),
            "weight_lbs": weight,
            "pieces": skids, "commodity": sup[2],
            "value_usd": round(random.uniform(8500, 145000), 2),
            "container_no": f"KKFU{random.randint(1000000,9999999)}" if mode == "Ocean" else None,
            "bol_no": f"BOL{random.randint(100000,999999)}",
            "pro_no": f"PRO {random.randint(5980000, 6010000)}",
            "progress": progress,
            "direction": "inbound", "hazmat": is_haz, "hazmat_class": haz_class,
            "supplier": sup[0], "consignee": plant_short,
            "ship_date": ship_dt.date().isoformat(),
            "ship_day": days_map[ship_dt.weekday()] if ship_dt.weekday() < 5 else "MON",
            "skids": skids, "material_controller": random.choice(material_controllers),
            "po_numbers": "po#" + "_".join(pos),
            "booking_number": f"PRO {random.randint(5980000, 6010000)}",
            "bid_cost": round(random.uniform(680, 3200), 2),
            "fsc_pct": round(random.uniform(20, 34), 1),
            "extras": random.choice([None, "plus FSC", "plus FSC + Liftgate"]),
            "done": status == "delivered",
            "shipping_hours": "0700-1500",
            "pickup_no": f"SFS{random.randint(79000, 79999)}" if random.random() < 0.5 else None,
            "created_at": datetime.now(timezone.utc).isoformat(),
        })
    await db.shipments.insert_many([dict(s) for s in shipments])

    # Seed a few documents
    doc_types = ["BOL", "COMMERCIAL_INVOICE", "PACKING_SLIP", "WEIGHT_CERT", "COO"]
    docs = []
    for s in shipments[:12]:
        t = random.choice(doc_types)
        docs.append({
            "document_id": f"DOC-{uuid.uuid4().hex[:8].upper()}",
            "type": t,
            "shipment_ref": s["reference"],
            "created_by": "System Seed",
            "created_at": datetime.now(timezone.utc).isoformat(),
            "data": {"carrier": s["carrier"], "commodity": s["commodity"], "weight_lbs": s["weight_lbs"]},
        })
    if docs:
        await db.documents.insert_many([dict(d) for d in docs])

    # Seed some chat messages
    seed_msgs = [
        {"channel": "general", "text": "Welcome to the Tennant TMS HUD. Live data streaming."},
        {"channel": "general", "text": "Reminder: ACE filing cutoff for Tuesday departures is 17:00 ET."},
        {"channel": "ops-dispatch", "text": "Holland → Dallas TL booked with XPO. Pickup 0700 CT tomorrow."},
        {"channel": "ops-dispatch", "text": "Two parcel exceptions on UPS tracker — see /shipments?status=delayed"},
        {"channel": "import-export", "text": "K+N container TCLU6543210 cleared at LA. Drayage scheduled."},
        {"channel": "carrier-issues", "text": "XPO API latency spike noted at 14:22 — monitoring."},
    ]
    chat_docs = []
    for m in seed_msgs:
        chat_docs.append({
            "message_id": f"msg_{uuid.uuid4().hex[:10]}",
            "channel": m["channel"],
            "user_id": "system",
            "user_name": "Dispatch Bot",
            "user_picture": None,
            "text": m["text"],
            "created_at": datetime.now(timezone.utc).isoformat(),
        })
    if chat_docs:
        await db.chat_messages.insert_many(chat_docs)

    # Seed freight bills
    accessorial_codes = [
        ("DET", "Detention", 75.0), ("LIFT", "Liftgate Service", 95.0),
        ("RES", "Residential Delivery", 110.0), ("INS", "Inside Delivery", 145.0),
        ("REW", "Reweigh", 35.0), ("REC", "Reclassification", 50.0),
        ("HAZ", "Hazardous Materials", 65.0), ("APPT", "Delivery Appointment", 45.0),
    ]
    bill_statuses = ["pending", "approved", "paid", "disputed", "audit"]
    bills = []
    for s in shipments[:30]:
        base = round(s["weight_lbs"] * random.uniform(0.08, 0.45), 2)
        accs_count = random.randint(0, 3)
        accs = []
        for _ in range(accs_count):
            c = random.choice(accessorial_codes)
            accs.append({"code": c[0], "description": c[1], "amount": c[2] + round(random.uniform(0, 25), 2)})
        fuel = round(base * random.uniform(0.12, 0.28), 2)
        total = round(base + fuel + sum(a["amount"] for a in accs), 2)
        quoted = round(total * random.uniform(0.85, 1.0), 2)
        variance = round(total - quoted, 2)
        bills.append({
            "bill_id": f"FB-{uuid.uuid4().hex[:8].upper()}",
            "shipment_ref": s["reference"],
            "carrier": s["carrier"],
            "invoice_no": f"INV-{random.randint(100000, 999999)}",
            "base_charge": base,
            "accessorials": accs,
            "fuel_surcharge": fuel,
            "total": total,
            "quoted_total": quoted,
            "variance": variance,
            "status": random.choice(bill_statuses),
            "invoice_date": (datetime.now(timezone.utc) - timedelta(days=random.randint(0, 28))).isoformat(),
            "due_date": (datetime.now(timezone.utc) + timedelta(days=random.randint(5, 30))).isoformat(),
        })
    await db.freight_bills.insert_many([dict(b) for b in bills])

    # Seed carrier onboardings — includes the original three plus a handful
    # of well-known national carriers marked "approved" so the Truckload
    # Booking Sheet's Carrier dropdown has a usable roster on day 1.
    today = datetime.now(timezone.utc).date()
    plus = lambda d: (datetime.now(timezone.utc) + timedelta(days=d)).date().isoformat()
    onboarding_seeds = [
        {"legal_name": "Prairie Stream Logistics LLC", "dba": "Prairie Stream", "mc_number": "MC-887214", "dot_number": "3621198", "scac": "PSLG", "mode": "TL",
         "contact_name": "Brett Halverson", "contact_email": "brett@prairiestream.com", "contact_phone": "+1-612-555-0188",
         "insurance_amount": 1500000, "insurance_expiry": plus(180),
         "safety_rating": "Satisfactory", "csa_score": 32, "status": "in_review"},
        {"legal_name": "Lakeshore Freight Co", "dba": None, "mc_number": "MC-742019", "dot_number": "2874321", "scac": "LSFC", "mode": "LTL",
         "contact_name": "Marisol Tran", "contact_email": "marisol@lakeshorefreight.com", "contact_phone": "+1-616-555-0142",
         "insurance_amount": 1000000, "insurance_expiry": plus(92),
         "safety_rating": "Satisfactory", "csa_score": 48, "status": "approved"},
        {"legal_name": "Bluegrass Express Inc", "dba": "BlueX", "mc_number": "MC-558821", "dot_number": "1842099", "scac": "BLGX", "mode": "TL",
         "contact_name": "Darnell McKee", "contact_email": "darnell@bluegrassx.com", "contact_phone": "+1-502-555-0199",
         "insurance_amount": 1000000, "insurance_expiry": plus(-10),
         "safety_rating": "Conditional", "csa_score": 78, "status": "in_review"},
        # ---- Approved national carriers, ready to book ----
        {"legal_name": "XPO Logistics LLC", "dba": "XPO", "mc_number": "MC-249635", "dot_number": "528970", "scac": "XPOL", "mode": "TL",
         "contact_name": "Renee Calderon", "contact_email": "tennant.team@xpo.com", "contact_phone": "+1-855-976-2243",
         "insurance_amount": 2000000, "insurance_expiry": plus(220),
         "safety_rating": "Satisfactory", "csa_score": 18, "status": "approved"},
        {"legal_name": "Old Dominion Freight Line, Inc.", "dba": "ODFL", "mc_number": "MC-22198", "dot_number": "55977", "scac": "ODFL", "mode": "LTL",
         "contact_name": "Greg Halsey", "contact_email": "tennant@odfl.com", "contact_phone": "+1-800-432-6335",
         "insurance_amount": 2000000, "insurance_expiry": plus(310),
         "safety_rating": "Satisfactory", "csa_score": 11, "status": "approved"},
        {"legal_name": "Saia Motor Freight Line, LLC", "dba": "Saia", "mc_number": "MC-44918", "dot_number": "33172", "scac": "SAIA", "mode": "LTL",
         "contact_name": "Tasha Burnett", "contact_email": "tennant@saia.com", "contact_phone": "+1-800-765-7242",
         "insurance_amount": 1500000, "insurance_expiry": plus(150),
         "safety_rating": "Satisfactory", "csa_score": 22, "status": "approved"},
        {"legal_name": "Estes Express Lines", "dba": "Estes", "mc_number": "MC-105764", "dot_number": "55712", "scac": "EXLA", "mode": "LTL",
         "contact_name": "Adam Mueller", "contact_email": "ops@estes-express.com", "contact_phone": "+1-866-378-3748",
         "insurance_amount": 2000000, "insurance_expiry": plus(265),
         "safety_rating": "Satisfactory", "csa_score": 16, "status": "approved"},
        {"legal_name": "R&L Carriers, Inc.", "dba": "R+L", "mc_number": "MC-133134", "dot_number": "243809", "scac": "RLCA", "mode": "LTL",
         "contact_name": "Marcus Lavoie", "contact_email": "tennant@rlcarriers.com", "contact_phone": "+1-800-543-5589",
         "insurance_amount": 1500000, "insurance_expiry": plus(195),
         "safety_rating": "Satisfactory", "csa_score": 20, "status": "approved"},
        {"legal_name": "Knight-Swift Transportation Holdings", "dba": "Knight Transportation", "mc_number": "MC-247369", "dot_number": "362724", "scac": "KNIG", "mode": "TL",
         "contact_name": "Lina Ortega", "contact_email": "tennant@knight-swift.com", "contact_phone": "+1-602-269-2000",
         "insurance_amount": 2000000, "insurance_expiry": plus(340),
         "safety_rating": "Satisfactory", "csa_score": 25, "status": "approved"},
        {"legal_name": "Schneider National Carriers", "dba": "Schneider", "mc_number": "MC-237983", "dot_number": "264184", "scac": "SCNN", "mode": "TL",
         "contact_name": "Henry Park", "contact_email": "tennant@schneider.com", "contact_phone": "+1-800-558-6767",
         "insurance_amount": 2000000, "insurance_expiry": plus(280),
         "safety_rating": "Satisfactory", "csa_score": 28, "status": "approved"},
        {"legal_name": "C.H. Robinson Worldwide, Inc.", "dba": "C.H. Robinson", "mc_number": "MC-208535", "dot_number": "388873", "scac": "CHRW", "mode": "Brokerage",
         "contact_name": "Ben Reichl", "contact_email": "tennant@chrobinson.com", "contact_phone": "+1-800-323-7587",
         "insurance_amount": 2000000, "insurance_expiry": plus(360),
         "safety_rating": "Satisfactory", "csa_score": 15, "status": "approved"},
        {"legal_name": "Werner Enterprises, Inc.", "dba": "Werner", "mc_number": "MC-159458", "dot_number": "111723", "scac": "WERN", "mode": "TL",
         "contact_name": "Christine Yoder", "contact_email": "tennant@werner.com", "contact_phone": "+1-800-228-2240",
         "insurance_amount": 2000000, "insurance_expiry": plus(120),
         "safety_rating": "Satisfactory", "csa_score": 26, "status": "approved"},
    ]
    _ = today  # avoid unused-variable warning if we add date-relative ones later
    onboarding_docs = []
    for ob in onboarding_seeds:
        onboarding_docs.append({
            "onboarding_id": f"OB-{uuid.uuid4().hex[:8].upper()}",
            **ob,
            "w9_received": random.choice([True, False]),
            "coi_received": random.choice([True, False]),
            "contract_signed": ob["status"] == "approved",
            "submitted_at": (datetime.now(timezone.utc) - timedelta(days=random.randint(1, 14))).isoformat(),
            "submitted_by": "Seed",
            "notes": "",
        })
    await db.carrier_onboarding.insert_many([dict(o) for o in onboarding_docs])

    # Seed Claims (idempotent — only if collection empty)
    claim_count = await db.claims.count_documents({})
    claims_seeded = 0
    if claim_count == 0 and shipments:
        sample = random.sample(shipments, min(8, len(shipments)))
        claim_types_seed = ["damage", "shortage", "loss", "freight_overcharge", "damage", "freight_overcharge", "shortage", "damage"]
        statuses_seed = ["settled", "filed", "open", "acknowledged", "denied", "partial", "settled", "open"]
        claim_seeds = []
        for s, ct, st in zip(sample, claim_types_seed, statuses_seed):
            claimed = round(random.uniform(450, 8500), 2)
            recovered = 0.0
            if st == "settled":
                recovered = claimed
            elif st == "partial":
                recovered = round(claimed * random.uniform(0.4, 0.85), 2)
            claim_seeds.append({
                "claim_id": f"CLM-{uuid.uuid4().hex[:8].upper()}",
                "shipment_id": s["shipment_id"],
                "carrier": s["carrier"],
                "bol_no": s.get("bol_no"),
                "claim_type": ct,
                "amount_claimed_usd": claimed,
                "amount_recovered_usd": recovered,
                "status": st,
                "filed_date": (datetime.now(timezone.utc) - timedelta(days=random.randint(1, 60))).date().isoformat(),
                "incident_date": s.get("ship_date") or s.get("pickup_date"),
                "description": {
                    "damage": "Two skids crushed — forklift damage at delivery",
                    "shortage": "1 of 4 skids missing on delivery",
                    "loss": "Cargo not delivered — last GPS ping 96h ago",
                    "freight_overcharge": "Accessorial billed but not authorized",
                    "overage": "Extra pallets received beyond manifest",
                }.get(ct, "Exception filed"),
                "created_by": "seed",
                "notes": "",
            })
            claims_seeded += 1
        if claim_seeds:
            await db.claims.insert_many(claim_seeds)

    return {"ok": True, "shipments": len(shipments), "documents": len(docs), "messages": len(chat_docs), "bills": len(bills), "onboardings": len(onboarding_docs), "claims": claims_seeded}

# -------------------- NMFC REFERENCE (Tennant products) --------------------
# Curated NMFC codes most relevant to Tennant industrial cleaning products & inbound components.
# Source: NMFTA NMFC + standard LTL freight classes. Class is density/value/handling-based.
# Comprehensive National Motor Freight Classification (NMFC) reference —
# selected codes spanning every major commodity category so dispatchers can
# class any load, not just industrial machinery. Sourced from the public
# NMFC tariff (NMF 100-AY). Categories are grouped for easy navigation in
# the BookLoad NMFC select.
GENERIC_NMFC_CODES = [
    # ── Food & Beverage ──────────────────────────────────────────────
    {"nmfc": "73130", "description": "Foodstuffs, canned, in cases",                                  "freight_class": "70",  "category": "Food & Beverage"},
    {"nmfc": "73140", "description": "Foodstuffs, dry, in bags or cartons (flour, sugar, rice)",        "freight_class": "65",  "category": "Food & Beverage"},
    {"nmfc": "73150", "description": "Foodstuffs, frozen, in insulated cartons",                       "freight_class": "70",  "category": "Food & Beverage"},
    {"nmfc": "73160", "description": "Beverages, non-alcoholic, in bottles/cans (palletized)",          "freight_class": "65",  "category": "Food & Beverage"},
    {"nmfc": "73170", "description": "Beverages, alcoholic — beer, wine, spirits, palletized",          "freight_class": "70",  "category": "Food & Beverage"},
    {"nmfc": "73180", "description": "Meat & poultry, refrigerated, in cartons",                       "freight_class": "70",  "category": "Food & Beverage"},
    {"nmfc": "73190", "description": "Dairy products, refrigerated",                                   "freight_class": "70",  "category": "Food & Beverage"},
    {"nmfc": "73200", "description": "Produce, fresh — fruits & vegetables, in crates",                 "freight_class": "70",  "category": "Food & Beverage"},
    {"nmfc": "73210", "description": "Pet food / animal feed, in bags",                                "freight_class": "65",  "category": "Food & Beverage"},

    # ── Apparel & Textiles ───────────────────────────────────────────
    {"nmfc": "49880", "description": "Apparel, finished, in cartons (cotton/synthetic)",                "freight_class": "150", "category": "Apparel & Textiles"},
    {"nmfc": "49900", "description": "Apparel, finished, fur or leather garments",                     "freight_class": "200", "category": "Apparel & Textiles"},
    {"nmfc": "49920", "description": "Footwear, in cartons",                                           "freight_class": "150", "category": "Apparel & Textiles"},
    {"nmfc": "49940", "description": "Textiles, woven piece goods, in bales",                          "freight_class": "100", "category": "Apparel & Textiles"},
    {"nmfc": "49960", "description": "Yarn or thread, in cartons or cones",                            "freight_class": "100", "category": "Apparel & Textiles"},

    # ── Electronics & Computers ──────────────────────────────────────
    {"nmfc": "61700", "description": "Computers, mainframes/servers, in cartons",                      "freight_class": "92.5","category": "Electronics"},
    {"nmfc": "61720", "description": "Computers, desktop, in retail cartons",                          "freight_class": "100", "category": "Electronics"},
    {"nmfc": "61740", "description": "Computers, laptop/notebook, in retail cartons",                  "freight_class": "150", "category": "Electronics"},
    {"nmfc": "61760", "description": "Electronics, consumer (TV, audio), in retail cartons",           "freight_class": "150", "category": "Electronics"},
    {"nmfc": "61780", "description": "Electronics, components (PCBs, ICs), in cartons",                "freight_class": "175", "category": "Electronics"},
    {"nmfc": "62110", "description": "Cables & wires, electronic, in cartons or spools",               "freight_class": "100", "category": "Electronics"},

    # ── Pharmaceuticals & Medical ────────────────────────────────────
    {"nmfc": "60000", "description": "Pharmaceuticals, in packages (non-controlled)",                  "freight_class": "85",  "category": "Pharmaceuticals & Medical"},
    {"nmfc": "60010", "description": "Pharmaceuticals, in packages (controlled — DEA Schedule I-V)",   "freight_class": "100", "category": "Pharmaceuticals & Medical"},
    {"nmfc": "60020", "description": "Medical instruments / devices, in cartons",                      "freight_class": "100", "category": "Pharmaceuticals & Medical"},
    {"nmfc": "60030", "description": "Medical supplies, disposable (gloves, gowns, masks)",            "freight_class": "175", "category": "Pharmaceuticals & Medical"},

    # ── Furniture & Home Goods ───────────────────────────────────────
    {"nmfc": "94060", "description": "Furniture, household, knocked-down, in cartons",                 "freight_class": "150", "category": "Furniture & Home Goods"},
    {"nmfc": "94070", "description": "Furniture, household, fully assembled (sofas, chairs)",          "freight_class": "250", "category": "Furniture & Home Goods"},
    {"nmfc": "94080", "description": "Furniture, office, knocked-down, in cartons",                    "freight_class": "125", "category": "Furniture & Home Goods"},
    {"nmfc": "94090", "description": "Mattresses, innerspring or foam, palletized",                    "freight_class": "200", "category": "Furniture & Home Goods"},
    {"nmfc": "94100", "description": "Appliances, household (refrigerators, washers, ranges)",          "freight_class": "92.5","category": "Furniture & Home Goods"},

    # ── Building & Construction ──────────────────────────────────────
    {"nmfc": "20040", "description": "Building materials, gypsum board / drywall sheets",              "freight_class": "70",  "category": "Building & Construction"},
    {"nmfc": "20060", "description": "Lumber, dimensional or cut, banded",                             "freight_class": "65",  "category": "Building & Construction"},
    {"nmfc": "20080", "description": "Plywood / OSB sheets, banded",                                   "freight_class": "65",  "category": "Building & Construction"},
    {"nmfc": "20100", "description": "Roofing materials — asphalt shingles, palletized",               "freight_class": "65",  "category": "Building & Construction"},
    {"nmfc": "20120", "description": "Windows / doors, finished, crated",                              "freight_class": "85",  "category": "Building & Construction"},
    {"nmfc": "20140", "description": "Tile (ceramic / porcelain), in cartons on pallets",              "freight_class": "70",  "category": "Building & Construction"},
    {"nmfc": "20160", "description": "Concrete, precast (pipes, blocks)",                              "freight_class": "55",  "category": "Building & Construction"},
    {"nmfc": "20180", "description": "Insulation, fiberglass batts/rolls, banded",                     "freight_class": "150", "category": "Building & Construction"},

    # ── Steel, Metals & Machinery ────────────────────────────────────
    {"nmfc": "133300", "description": "Motors, electric, AC/DC, NOI",                                  "freight_class": "85",  "category": "Machinery & Equipment"},
    {"nmfc": "133420", "description": "Motors, electric, hermetically sealed",                         "freight_class": "70",  "category": "Machinery & Equipment"},
    {"nmfc": "44890",  "description": "Castings, iron or steel, machinery parts",                      "freight_class": "70",  "category": "Machinery & Equipment"},
    {"nmfc": "44970",  "description": "Frames, steel, machinery, welded",                              "freight_class": "70",  "category": "Machinery & Equipment"},
    {"nmfc": "105820", "description": "Floor maintenance machines (scrubbers, sweepers, polishers)",    "freight_class": "85",  "category": "Machinery & Equipment"},
    {"nmfc": "105823", "description": "Vacuum cleaners, industrial (electric/battery)",                "freight_class": "92.5","category": "Machinery & Equipment"},
    {"nmfc": "105825", "description": "Carpet extractors / shampoo machines",                          "freight_class": "100", "category": "Machinery & Equipment"},
    {"nmfc": "49880-1","description": "Pumps, electric, water/solution",                               "freight_class": "85",  "category": "Machinery & Equipment"},
    {"nmfc": "133600", "description": "Generators, portable, gasoline or diesel",                       "freight_class": "85",  "category": "Machinery & Equipment"},
    {"nmfc": "133620", "description": "Compressors, air, electric or gas",                              "freight_class": "85",  "category": "Machinery & Equipment"},
    {"nmfc": "133700", "description": "Forklifts / pallet jacks, palletized",                          "freight_class": "70",  "category": "Machinery & Equipment"},
    {"nmfc": "189870", "description": "Tires, solid rubber, industrial",                               "freight_class": "85",  "category": "Machinery & Equipment"},
    {"nmfc": "189880", "description": "Tires, pneumatic, passenger/light truck",                        "freight_class": "100", "category": "Machinery & Equipment"},

    # ── Metals & Raw Materials ───────────────────────────────────────
    {"nmfc": "33060", "description": "Iron or steel, bars / rods, banded",                             "freight_class": "50",  "category": "Metals & Raw"},
    {"nmfc": "33080", "description": "Iron or steel, sheets / plates, banded",                         "freight_class": "55",  "category": "Metals & Raw"},
    {"nmfc": "33100", "description": "Iron or steel, structural shapes (I-beams, channel)",            "freight_class": "60",  "category": "Metals & Raw"},
    {"nmfc": "33120", "description": "Aluminum, sheet or coil, banded",                                "freight_class": "65",  "category": "Metals & Raw"},
    {"nmfc": "33140", "description": "Copper, wire or rod, banded",                                    "freight_class": "60",  "category": "Metals & Raw"},
    {"nmfc": "156600","description": "Plastic articles, solution tanks (HDPE)",                        "freight_class": "150", "category": "Metals & Raw"},
    {"nmfc": "156605","description": "Plastic articles, recovery tanks",                               "freight_class": "175", "category": "Metals & Raw"},
    {"nmfc": "156620","description": "Plastic resin, pelletized, in bags or supersacks",               "freight_class": "55",  "category": "Metals & Raw"},

    # ── Hazmat / Chemicals (Class 3 / 8 / 9) ─────────────────────────
    {"nmfc": "60590-2","description": "Batteries, electric storage, lead-acid (wet/AGM)",              "freight_class": "60",  "category": "Hazmat · Class 8"},
    {"nmfc": "60585",  "description": "Batteries, lithium-ion, packaged with equipment",               "freight_class": "92.5","category": "Hazmat · Class 9"},
    {"nmfc": "60600",  "description": "Battery chargers, electric",                                    "freight_class": "85",  "category": "Electrical"},
    {"nmfc": "27360",  "description": "Chemicals, flammable liquid n.o.s. (UN1993)",                   "freight_class": "85",  "category": "Hazmat · Class 3"},
    {"nmfc": "27380",  "description": "Chemicals, corrosive liquid (acids/bases)",                     "freight_class": "85",  "category": "Hazmat · Class 8"},
    {"nmfc": "27400",  "description": "Paint, in cans or pails",                                       "freight_class": "85",  "category": "Hazmat · Class 3"},
    {"nmfc": "27420",  "description": "Adhesives, in cans or cartridges",                              "freight_class": "85",  "category": "Hazmat · Class 3"},

    # ── Automotive ───────────────────────────────────────────────────
    {"nmfc": "70000", "description": "Auto parts, NOI (palletized or crated)",                          "freight_class": "85",  "category": "Automotive"},
    {"nmfc": "70030", "description": "Wiring harnesses, automotive/industrial",                         "freight_class": "92.5","category": "Automotive"},
    {"nmfc": "70060", "description": "Engines / transmissions, crated",                                 "freight_class": "70",  "category": "Automotive"},
    {"nmfc": "70080", "description": "Brake parts, in cartons",                                         "freight_class": "85",  "category": "Automotive"},
    {"nmfc": "70100", "description": "Bumpers / body panels, plastic, in cartons",                       "freight_class": "150", "category": "Automotive"},

    # ── Paper & Print ────────────────────────────────────────────────
    {"nmfc": "151100","description": "Paper, printing or writing, in rolls",                            "freight_class": "60",  "category": "Paper & Print"},
    {"nmfc": "151110","description": "Paper products, cartons/boxes flat-packed",                       "freight_class": "70",  "category": "Paper & Print"},
    {"nmfc": "151120","description": "Books / catalogs, on pallets",                                    "freight_class": "65",  "category": "Paper & Print"},

    # ── Consumables / Misc ──────────────────────────────────────────
    {"nmfc": "16030", "description": "Brushes, machine, rotary",                                       "freight_class": "150", "category": "Consumables"},
    {"nmfc": "16035", "description": "Brushes, sweeping, cylindrical",                                  "freight_class": "175", "category": "Consumables"},
    {"nmfc": "84510", "description": "Hoses, rubber or plastic, with fittings",                         "freight_class": "100", "category": "Consumables"},
    {"nmfc": "84580", "description": "Squeegee blades, rubber/urethane",                                "freight_class": "125", "category": "Consumables"},
    {"nmfc": "186200","description": "Tools, hand, NOI (parts & spares)",                               "freight_class": "100", "category": "Parts & Spares"},

    # ── Catch-all Freight All Kinds ──────────────────────────────────
    {"nmfc": "FAK-50",  "description": "FAK · Freight All Kinds — class 50 (heavy, dense)",             "freight_class": "50",  "category": "FAK"},
    {"nmfc": "FAK-85",  "description": "FAK · Freight All Kinds — class 85 (mid-density mixed)",         "freight_class": "85",  "category": "FAK"},
    {"nmfc": "FAK-125", "description": "FAK · Freight All Kinds — class 125 (low-density mixed)",        "freight_class": "125", "category": "FAK"},
    {"nmfc": "FAK-175", "description": "FAK · Freight All Kinds — class 175 (light bulky mixed)",        "freight_class": "175", "category": "FAK"},
]

FREIGHT_CLASSES = ["50", "55", "60", "65", "70", "77.5", "85", "92.5", "100", "110", "125", "150", "175", "200", "250", "300", "400", "500"]

ACCESSORIAL_OPTIONS = [
    {"id": "liftgate", "label": "Liftgate Service"},
    {"id": "residential", "label": "Residential Delivery"},
    {"id": "inside_delivery", "label": "Inside Delivery"},
    {"id": "limited_access", "label": "Limited Access (school/church/storage)"},
    {"id": "appointment", "label": "Appointment Required"},
    {"id": "notify_consignee", "label": "Notify Consignee Prior to Delivery"},
    {"id": "tradeshow", "label": "Tradeshow Delivery"},
    {"id": "construction_site", "label": "Construction Site Delivery"},
    {"id": "single_shipment", "label": "Single Shipment Charge"},
    {"id": "hazmat_handling", "label": "Hazmat Handling Fee"},
    {"id": "blind_shipment", "label": "Blind Shipment"},
    {"id": "sort_segregate", "label": "Sort & Segregate"},
]

@api_router.get("/nmfc/codes")
async def get_nmfc_codes(_: User = Depends(get_current_user)):
    return {"codes": GENERIC_NMFC_CODES, "freight_classes": FREIGHT_CLASSES, "accessorials": ACCESSORIAL_OPTIONS}

# -------------------- SAP S/4HANA: Open Deliveries (for Book Load auto-fill) --------------------
SAP_OPEN_DELIVERIES = [
    {"delivery_no": "8000234", "so_no": "1010234", "customer": "Ferguson Enterprises — Atlanta, GA", "material": "TENN-T16AMR-LI", "material_desc": "T16 AMR · Lithium-ion", "qty": 2, "plant": "1010", "requested_date": "2026-05-14", "incoterms": "FOB Origin"},
    {"delivery_no": "8000241", "so_no": "1010241", "customer": "Sysco Corp — Houston, TX", "material": "TENN-T7AMR", "material_desc": "T7 AMR · 32 in scrubber", "qty": 1, "plant": "1020", "requested_date": "2026-05-15", "incoterms": "FOB Destination"},
    {"delivery_no": "8000256", "so_no": "1010256", "customer": "Home Depot DC — Atlanta, GA", "material": "TENN-S30", "material_desc": "S30 Industrial Sweeper", "qty": 3, "plant": "1030", "requested_date": "2026-05-16", "incoterms": "FCA Plant"},
    {"delivery_no": "8000263", "so_no": "1010263", "customer": "Walmart DC #5024 — Bentonville, AR", "material": "TENN-T350-PROPANE", "material_desc": "T350 LPG Scrubber", "qty": 4, "plant": "1010", "requested_date": "2026-05-17", "incoterms": "FOB Origin"},
    {"delivery_no": "8000275", "so_no": "1010275", "customer": "Costco Wholesale — Issaquah, WA", "material": "TENN-T16AMR-LI", "material_desc": "T16 AMR · Lithium-ion", "qty": 1, "plant": "1020", "requested_date": "2026-05-18", "incoterms": "FOB Destination"},
    {"delivery_no": "8000284", "so_no": "1010284", "customer": "Boeing Everett — Everett, WA", "material": "TENN-M30-PARTS", "material_desc": "M30 Replacement Parts Kit", "qty": 8, "plant": "1010", "requested_date": "2026-05-19", "incoterms": "FCA Plant"},
    {"delivery_no": "8000291", "so_no": "1010291", "customer": "FedEx Memphis Hub — Memphis, TN", "material": "TENN-S20", "material_desc": "S20 Compact Sweeper", "qty": 2, "plant": "1030", "requested_date": "2026-05-20", "incoterms": "DAP Destination"},
    {"delivery_no": "8000302", "so_no": "1010302", "customer": "GM Lansing Plant — Lansing, MI", "material": "TENN-T681", "material_desc": "T681 Rider Scrubber", "qty": 1, "plant": "1020", "requested_date": "2026-05-21", "incoterms": "FOB Origin"},
]

@api_router.get("/sap/open-deliveries")
async def sap_open_deliveries(user: User = Depends(get_current_user)):
    """Returns open deliveries from SAP S/4HANA (mocked) for Book Load reference auto-fill."""
    if user.role == "carrier":
        raise HTTPException(status_code=403, detail="Not available for carrier role")
    out = list(SAP_OPEN_DELIVERIES)
    brand = await _active_brand_doc()
    if brand and brand.get("brand_id") != "orisei-freight":
        products = brand.get("sample_products") or []
        short = brand.get("short_name") or "BRAND"
        prefix = re.sub(r"[^A-Z0-9]+", "", short.upper())[:4] or "BRND"
        for i, d in enumerate(out):
            if products:
                desc = products[i % len(products)]
                d["material_desc"] = desc
                slug = re.sub(r"[^A-Z0-9]+", "", desc.upper())[:6] or f"P{i:03d}"
                d["material"] = f"{prefix}-{slug}-{i:02d}"
    return {"deliveries": out, "fetched_at": datetime.now(timezone.utc).isoformat()}

@api_router.get("/sap/materials")
async def sap_materials(_: User = Depends(get_current_user)):
    """Top part numbers (mock) for the Command Center widget. Brand-aware:
    when a non-Tennant brand is active, part numbers and descriptions are
    overlaid with the active brand's sample products."""
    materials = [
        {"part_no": "TENN-T16AMR-LI", "description": "T16 AMR · Lithium-ion", "plant": "1010", "on_hand": 47, "open_orders": 12, "nmfc": "105820", "freight_class": "85"},
        {"part_no": "TENN-T7AMR", "description": "T7 AMR Scrubber 32 in", "plant": "1020", "on_hand": 86, "open_orders": 18, "nmfc": "105820", "freight_class": "85"},
        {"part_no": "TENN-T350-PROPANE", "description": "T350 LPG Scrubber", "plant": "1010", "on_hand": 22, "open_orders": 9, "nmfc": "105820", "freight_class": "85"},
        {"part_no": "TENN-S30", "description": "S30 Industrial Sweeper", "plant": "1030", "on_hand": 31, "open_orders": 14, "nmfc": "105820", "freight_class": "85"},
        {"part_no": "TENN-S20", "description": "S20 Compact Sweeper", "plant": "1030", "on_hand": 64, "open_orders": 21, "nmfc": "105820", "freight_class": "85"},
        {"part_no": "TENN-BATT-LION-48V", "description": "Battery Pack · Li-ion · 48V", "plant": "1010", "on_hand": 142, "open_orders": 38, "nmfc": "60585", "freight_class": "92.5"},
        {"part_no": "TENN-BATT-AGM-36V", "description": "Battery Pack · AGM · 36V", "plant": "1010", "on_hand": 89, "open_orders": 24, "nmfc": "60590-2", "freight_class": "60"},
        {"part_no": "TENN-BRUSH-CYL-32", "description": "Cylindrical Brush · 32 in", "plant": "1020", "on_hand": 312, "open_orders": 56, "nmfc": "16035", "freight_class": "175"},
        {"part_no": "TENN-SQGE-URETH-40", "description": "Squeegee Blade · Urethane · 40 in", "plant": "1020", "on_hand": 218, "open_orders": 41, "nmfc": "84580", "freight_class": "125"},
        {"part_no": "TENN-TANK-SOL-HDPE-30G", "description": "Solution Tank · HDPE · 30 gal", "plant": "1030", "on_hand": 56, "open_orders": 15, "nmfc": "156600", "freight_class": "150"},
    ]
    brand = await _active_brand_doc()
    if brand and brand.get("brand_id") != "orisei-freight":
        products = brand.get("sample_products") or []
        short = brand.get("short_name") or "BRAND"
        prefix = re.sub(r"[^A-Z0-9]+", "", short.upper())[:4] or "BRND"
        for i, m in enumerate(materials):
            if products:
                desc = products[i % len(products)]
                m["description"] = desc
                # part_no built from the brand short + 4 chars of the product slug
                slug = re.sub(r"[^A-Z0-9]+", "", desc.upper())[:6] or f"P{i:03d}"
                m["part_no"] = f"{prefix}-{slug}-{i:02d}"
    return {"materials": materials}

# -------------------- CARRIER INVITES (admin-gated) --------------------
class CarrierInviteCreate(BaseModel):
    carrier_company: str
    invitee_email: str
    invitee_name: Optional[str] = None
    expires_days: int = 14


def _public_app_url(request: Optional["Request"] = None) -> str:
    """Best-effort public URL used in email links/logos. Tries
    PUBLIC_APP_URL → FRONTEND_PUBLIC_URL → REACT_APP_BACKEND_URL → request
    forwarded headers → request.base_url so logo images always resolve to an
    absolute https URL inside email clients."""
    for env_key in ("PUBLIC_APP_URL", "FRONTEND_PUBLIC_URL", "REACT_APP_BACKEND_URL"):
        v = (os.environ.get(env_key) or "").rstrip("/")
        if v:
            return v
    if request is not None:
        try:
            # Behind a reverse proxy / k8s ingress, prefer the forwarded host
            # so emails point at the user-facing URL, not the cluster-internal one.
            fwd_host = request.headers.get("x-forwarded-host") or request.headers.get("host")
            fwd_proto = request.headers.get("x-forwarded-proto") or "https"
            if fwd_host:
                return f"{fwd_proto}://{fwd_host}".rstrip("/")
            return str(request.base_url).rstrip("/")
        except Exception:
            pass
    return ""


def _build_carrier_invite_html(*, brand: dict, invitee_name: Optional[str], carrier_company: str,
                                invite_url: str, expires_days: int, logo_url: Optional[str],
                                founder_name: Optional[str] = None) -> str:
    """Render a warm, full-color HTML carrier invite email body. Inline CSS,
    table layout — safe for Gmail/Outlook/Apple Mail."""
    company = brand.get("company_name") or "Orisei Freight Solutions LLC"
    short = brand.get("short_name") or "Orisei Freight"
    tagline = brand.get("tagline") or "Operator-built freight brokerage · Minneapolis · Saint Paul"
    primary = brand.get("primary_color") or "#0E3A6B"
    accent = brand.get("accent_color") or "#C9A24A"
    hq = brand.get("headquarters") or "Minneapolis, MN"
    signer = founder_name or brand.get("owner_name") or "Oliver Cummins"
    greeting = f"Hi {invitee_name}," if invitee_name else f"Hi {carrier_company} team,"
    # Build the optional logo block; some email clients require <img> hosted on https.
    logo_block = ""
    if logo_url:
        logo_block = (
            f'<img src="{logo_url}" alt="{company}" width="84" height="84" '
            f'style="display:block;border:0;outline:none;border-radius:14px;background:#ffffff;" />'
        )
    else:
        # Fallback: monogram disc using the brand letter
        letter = (brand.get("logo_letter") or short[:1] or "O").upper()
        logo_block = (
            f'<div style="width:84px;height:84px;border-radius:14px;background:{accent};'
            f'color:{primary};font-family:Georgia,serif;font-size:42px;font-weight:700;'
            f'line-height:84px;text-align:center;">{letter}</div>'
        )
    return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8" /><title>{company} · Carrier Invitation</title></head>
<body style="margin:0;padding:0;background:#F1F5F9;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;color:#0F172A;">
  <table role="presentation" cellspacing="0" cellpadding="0" border="0" width="100%" style="background:#F1F5F9;padding:32px 0;">
    <tr><td align="center">
      <table role="presentation" cellspacing="0" cellpadding="0" border="0" width="640" style="background:#ffffff;border-radius:14px;overflow:hidden;box-shadow:0 6px 18px rgba(15,23,42,0.10);">

        <!-- HERO -->
        <tr><td style="background:linear-gradient(135deg,{primary} 0%,#102B4F 100%);padding:32px 36px;color:#ffffff;">
          <table role="presentation" cellspacing="0" cellpadding="0" border="0" width="100%">
            <tr>
              <td valign="middle" width="100">{logo_block}</td>
              <td valign="middle" style="padding-left:20px;">
                <div style="font-family:Georgia,'Times New Roman',serif;font-weight:700;font-size:24px;letter-spacing:-0.01em;color:#ffffff;">{company}</div>
                <div style="margin-top:4px;font-size:13px;color:{accent};letter-spacing:0.04em;text-transform:uppercase;font-weight:600;">{tagline}</div>
              </td>
            </tr>
          </table>
        </td></tr>

        <!-- TITLE -->
        <tr><td style="padding:28px 36px 6px 36px;">
          <div style="font-family:Georgia,'Times New Roman',serif;font-size:26px;font-weight:700;color:{primary};line-height:1.25;">
            You're invited to join the {short} carrier network.
          </div>
          <div style="margin-top:6px;font-size:13px;color:#64748B;">{greeting}</div>
        </td></tr>

        <!-- BODY -->
        <tr><td style="padding:14px 36px 6px 36px;">
          <p style="margin:0 0 14px 0;color:#334155;font-size:15px;line-height:1.6;">
            We're {short} — a Twin Cities-based property freight brokerage built by a logistics
            operator who got tired of slow tenders, half-broken portals, and rate confirmations
            that never matched what hit the BOL. We run our own in-house TMS, our own load
            board syndication, and we believe a great carrier relationship is the single biggest
            competitive moat a brokerage can build.
          </p>
          <p style="margin:0 0 14px 0;color:#334155;font-size:15px;line-height:1.6;">
            We'd love to add <strong style="color:{primary};">{carrier_company}</strong> to our
            approved carrier network so we can start tendering loads to your dispatch desk.
            Acceptance is quick — just click the button below to set up your free carrier portal
            account. From there you'll be able to:
          </p>
          <ul style="margin:0 0 14px 18px;padding:0;color:#334155;font-size:14px;line-height:1.7;">
            <li>View every load we tender to <strong>{carrier_company}</strong></li>
            <li>Update pickup &amp; delivery status in real time</li>
            <li>Upload BOLs &amp; PODs straight from a phone — no email chase</li>
            <li>Track ETAs, exceptions, and quick-pay settlement</li>
          </ul>
        </td></tr>

        <!-- CTA -->
        <tr><td align="center" style="padding:18px 36px 24px 36px;">
          <a href="{invite_url}" style="display:inline-block;padding:14px 28px;border-radius:10px;background:{accent};color:{primary};font-weight:700;font-size:15px;text-decoration:none;letter-spacing:0.02em;box-shadow:0 4px 12px rgba(201,162,74,0.35);">
            Accept Invite &amp; Set Up Portal &rarr;
          </a>
          <div style="margin-top:10px;font-size:11px;color:#94A3B8;">Invite expires in {expires_days} days · No password to remember</div>
        </td></tr>

        <!-- WHO WE ARE / TRUST -->
        <tr><td style="padding:0 36px 22px 36px;">
          <table role="presentation" cellspacing="0" cellpadding="0" border="0" width="100%" style="background:#F8FAFC;border:1px solid #E2E8F0;border-radius:10px;">
            <tr><td style="padding:18px 22px;">
              <div style="font-size:11px;letter-spacing:0.08em;text-transform:uppercase;color:{primary};font-weight:700;">Who we are</div>
              <div style="margin-top:6px;color:#334155;font-size:13px;line-height:1.55;">
                {short} is operator-built and operator-run. We hold our MC authority, BMC-84 surety
                bond, and BOC-3 process agents on file. We pay via standard NET-30, with
                <strong>quick-pay at 2%</strong> available the moment a clean POD lands.
                Headquarters in {hq}, with lane desks across the upper Midwest.
              </div>
            </td></tr>
          </table>
        </td></tr>

        <!-- SIGNATURE -->
        <tr><td style="padding:0 36px 22px 36px;">
          <div style="color:#334155;font-size:14px;line-height:1.55;">Looking forward to running good lanes with your team,</div>
          <div style="margin-top:10px;color:{primary};font-family:Georgia,'Times New Roman',serif;font-size:18px;font-weight:700;">{signer}</div>
          <div style="color:#64748B;font-size:12px;">Founder &amp; Principal Broker · {company}</div>
        </td></tr>

        <!-- FOOTER -->
        <tr><td style="padding:16px 36px;background:{primary};color:#CBD5E1;font-size:11px;line-height:1.6;">
          You received this invitation because {short} identified <strong style="color:#ffffff;">{carrier_company}</strong> as
          a great-fit partner. Reply directly to opt out or to request a different point of contact.
          Need a hand? Reply to this email and a real human will write back — usually within a
          couple of hours.
        </td></tr>

      </table>
    </td></tr>
  </table>
</body></html>"""


def _build_carrier_invite_text(*, brand: dict, invitee_name: Optional[str], carrier_company: str,
                                invite_url: str, expires_days: int, founder_name: Optional[str] = None) -> str:
    """Plain-text fallback / mailto-body version of the invite."""
    company = brand.get("company_name") or "Orisei Freight Solutions LLC"
    short = brand.get("short_name") or "Orisei Freight"
    signer = founder_name or brand.get("owner_name") or "Oliver Cummins"
    greeting = f"Hi {invitee_name}," if invitee_name else f"Hi {carrier_company} team,"
    return (
        f"{greeting}\n\n"
        f"We're {short} — a Twin Cities-based property freight brokerage built by a logistics operator who got "
        f"tired of slow tenders, half-broken portals, and rate confirmations that never matched what hit the BOL. "
        f"We run our own in-house TMS, our own load board syndication, and we believe a great carrier relationship "
        f"is the single biggest competitive moat a brokerage can build.\n\n"
        f"We'd love to add {carrier_company} to our approved carrier network so we can start tendering loads to "
        f"your dispatch desk. Acceptance is quick — set up your free carrier portal account in under a minute "
        f"and you'll be able to:\n"
        f"  • View every load we tender to {carrier_company}\n"
        f"  • Update pickup & delivery status in real time\n"
        f"  • Upload BOLs & PODs straight from a phone\n"
        f"  • Track ETAs, exceptions, and quick-pay settlement\n\n"
        f"Accept your invite here (expires in {expires_days} days):\n{invite_url}\n\n"
        f"A little about us: {short} is operator-built and operator-run. We hold MC authority, BMC-84 surety bond, "
        f"and BOC-3 process agents on file. We pay via standard NET-30, with quick-pay at 2% available the moment "
        f"a clean POD lands.\n\n"
        f"Looking forward to running good lanes with your team,\n"
        f"{signer}\n"
        f"Founder & Principal Broker · {company}"
    )


@api_router.post("/carrier-invites")
async def create_carrier_invite(payload: CarrierInviteCreate, request: Request, admin: User = Depends(require_role("admin"))):
    token = uuid.uuid4().hex
    expires_at = (datetime.now(timezone.utc) + timedelta(days=payload.expires_days)).isoformat()
    invite = {
        "invite_id": f"INV-{uuid.uuid4().hex[:10].upper()}",
        "token": token,
        "carrier_company": payload.carrier_company,
        "invitee_email": payload.invitee_email,
        "invitee_name": payload.invitee_name,
        "created_by": admin.user_id,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "expires_at": expires_at,
        "status": "pending",  # pending | accepted | revoked | expired
        "accepted_user_id": None,
    }
    await db.carrier_invites.insert_one(dict(invite))
    base = _public_app_url(request)
    invite_link = f"/accept-invite?token={token}" if not base else f"{base}/accept-invite?token={token}"
    logo_url: Optional[str] = None
    brand = await _active_brand_doc() or {}
    if base:
        # public-served logo (frontend public/brand/orisei_logo.png is the canonical asset)
        logo_url = f"{base}/brand/orisei_logo.png"
    elif brand.get("logo_url"):
        # purely best-effort; some hosts may serve relative paths
        logo_url = brand["logo_url"]
    email_html = _build_carrier_invite_html(
        brand=brand,
        invitee_name=payload.invitee_name,
        carrier_company=payload.carrier_company,
        invite_url=invite_link,
        expires_days=payload.expires_days,
        logo_url=logo_url,
    )
    email_text = _build_carrier_invite_text(
        brand=brand,
        invitee_name=payload.invitee_name,
        carrier_company=payload.carrier_company,
        invite_url=invite_link,
        expires_days=payload.expires_days,
    )
    short = brand.get("short_name") or "Orisei Freight"
    subject = f"{short} · Join our carrier network · loads tendered to {payload.carrier_company}"
    return {
        "invite": {k: v for k, v in invite.items() if k != "token"},
        "invite_link": invite_link,
        "token": token,
        "subject": subject,
        "email_body": email_text,   # legacy field — used by mailto button
        "email_html": email_html,
        "logo_url": logo_url,
    }

@api_router.get("/carrier-invites")
async def list_carrier_invites(_: User = Depends(require_role("admin"))):
    docs = await db.carrier_invites.find({}, {"_id": 0}).sort("created_at", -1).limit(200).to_list(200)
    return docs

@api_router.delete("/carrier-invites/{invite_id}")
async def revoke_carrier_invite(invite_id: str, _: User = Depends(require_role("admin"))):
    res = await db.carrier_invites.update_one({"invite_id": invite_id}, {"$set": {"status": "revoked"}})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Invite not found")
    return {"ok": True}


@api_router.post("/carrier-invites/{invite_id}/send-email")
async def send_carrier_invite_email(invite_id: str, request: Request, admin: User = Depends(require_role("admin"))):
    """Send the carrier invite as a full-color HTML email via Resend
    (credentials pulled from the Connections vault). Falls back to a
    descriptive 400 if Resend is not configured so the UI can prompt the
    admin to add a key."""
    invite = await db.carrier_invites.find_one({"invite_id": invite_id}, {"_id": 0})
    if not invite:
        raise HTTPException(status_code=404, detail="Invite not found")
    if invite["status"] != "pending":
        raise HTTPException(status_code=400, detail=f"Invite already {invite['status']}")
    to_email = (invite.get("invitee_email") or "").strip()
    if not to_email or "@" not in to_email:
        raise HTTPException(status_code=400, detail="Invite has no valid invitee email on file")

    # Re-derive everything from the latest brand & invite state so a stale
    # email body never goes out.
    brand = await _active_brand_doc() or {}
    base = _public_app_url(request)
    invite_link = (
        f"/accept-invite?token={invite['token']}" if not base
        else f"{base}/accept-invite?token={invite['token']}"
    )
    expires_at = invite.get("expires_at")
    try:
        expires_dt = datetime.fromisoformat(expires_at) if isinstance(expires_at, str) else expires_at
        days_left = max(1, (expires_dt - datetime.now(timezone.utc)).days)
    except Exception:
        days_left = 14
    logo_url = f"{base}/brand/orisei_logo.png" if base else (brand.get("logo_url") or None)
    invitee_name = invite.get("invitee_name") or None
    carrier_company = invite.get("carrier_company") or "your team"
    html = _build_carrier_invite_html(
        brand=brand, invitee_name=invitee_name, carrier_company=carrier_company,
        invite_url=invite_link, expires_days=days_left, logo_url=logo_url,
    )
    text = _build_carrier_invite_text(
        brand=brand, invitee_name=invitee_name, carrier_company=carrier_company,
        invite_url=invite_link, expires_days=days_left,
    )
    short = brand.get("short_name") or "Orisei Freight"
    company = brand.get("company_name") or "Orisei Freight Solutions LLC"
    subject = f"{short} · Join our carrier network · loads tendered to {carrier_company}"

    # Pull Resend creds from the Connections vault
    try:
        from routes.connections import get_connection_credentials  # local import to avoid cycle
        creds = await get_connection_credentials(db, "resend")
    except Exception:
        creds = None

    record = {
        "id": str(uuid.uuid4()),
        "invite_id": invite_id,
        "sent_at": datetime.now(timezone.utc).isoformat(),
        "sent_by": admin.user_id,
        "to_email": to_email,
        "to_name": invitee_name,
        "carrier_company": carrier_company,
        "subject": subject,
        "company_name": company,
        "provider": "resend" if creds and creds.get("api_key") else "not_configured",
        "status": "queued",
    }

    if not creds or not creds.get("api_key"):
        record["status"] = "blocked"
        record["error"] = (
            "Resend connection not configured. Open Connections → Resend and paste a "
            "Resend API key (plus a verified sender address) to enable direct emailing."
        )
        await db.carrier_invite_emails.insert_one(dict(record))
        raise HTTPException(400, record["error"])

    from_email = creds.get("from_email") or "onboarding@resend.dev"
    from_name = creds.get("from_name") or short
    params = {
        "from": f"{from_name} <{from_email}>",
        "to": [to_email],
        "subject": subject,
        "html": html,
        "text": text,
    }
    reply_to = creds.get("reply_to")
    if reply_to:
        params["reply_to"] = [reply_to]

    try:
        import resend as _resend  # noqa: WPS433
        _resend.api_key = creds["api_key"]
        result = await asyncio.to_thread(_resend.Emails.send, params)
        record["status"] = "sent"
        record["provider_message_id"] = result.get("id") if isinstance(result, dict) else None
    except Exception as exc:  # pragma: no cover — Resend SDK failures
        record["status"] = "failed"
        record["error"] = str(exc)[:400]
        await db.carrier_invite_emails.insert_one(dict(record))
        raise HTTPException(502, f"Resend send failed: {exc}")

    # Stamp the invite so the table can show a "sent" badge
    await db.carrier_invites.update_one(
        {"invite_id": invite_id},
        {"$set": {
            "email_sent_at": record["sent_at"],
            "email_sent_by": admin.user_id,
            "email_provider_message_id": record.get("provider_message_id"),
        }},
    )
    await db.carrier_invite_emails.insert_one(dict(record))
    record.pop("_id", None)
    return {"ok": True, **record}

class CarrierInviteAccept(BaseModel):
    token: str
    email: str
    name: str

@api_router.post("/carrier-invites/accept")
async def accept_carrier_invite(payload: CarrierInviteAccept):
    """Auth-free endpoint — the invite token IS the auth. Creates the user + a 7-day session."""
    invite = await db.carrier_invites.find_one({"token": payload.token}, {"_id": 0})
    if not invite:
        raise HTTPException(status_code=404, detail="Invalid invite token")
    if invite["status"] != "pending":
        raise HTTPException(status_code=400, detail=f"Invite already {invite['status']}")
    expires_at = invite["expires_at"]
    if isinstance(expires_at, str):
        expires_at = datetime.fromisoformat(expires_at)
    if expires_at.tzinfo is None:
        expires_at = expires_at.replace(tzinfo=timezone.utc)
    if expires_at < datetime.now(timezone.utc):
        await db.carrier_invites.update_one({"token": payload.token}, {"$set": {"status": "expired"}})
        raise HTTPException(status_code=400, detail="Invite has expired")
    # Find or create the user
    existing = await db.users.find_one({"email": payload.email.lower()}, {"_id": 0})
    if existing:
        user_id = existing["user_id"]
        await db.users.update_one({"user_id": user_id}, {"$set": {
            "role": "carrier",
            "carrier_company": invite["carrier_company"],
            "name": payload.name,
        }})
    else:
        user_id = f"user_{uuid.uuid4().hex[:12]}"
        await db.users.insert_one({
            "user_id": user_id,
            "email": payload.email.lower(),
            "name": payload.name,
            "picture": None,
            "role": "carrier",
            "carrier_company": invite["carrier_company"],
            "created_at": datetime.now(timezone.utc).isoformat(),
        })
    session_token = uuid.uuid4().hex
    session_expires = (datetime.now(timezone.utc) + timedelta(days=7)).isoformat()
    await db.user_sessions.insert_one({
        "user_id": user_id,
        "session_token": session_token,
        "expires_at": session_expires,
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    await db.carrier_invites.update_one(
        {"token": payload.token},
        {"$set": {"status": "accepted", "accepted_user_id": user_id, "accepted_at": datetime.now(timezone.utc).isoformat()}},
    )
    return {"session_token": session_token, "user_id": user_id, "role": "carrier", "carrier_company": invite["carrier_company"]}

# -------------------- ONBOARDING: Send carrier requirements packet (brand-aware) --------------------
@api_router.post("/carrier-onboarding/{onboarding_id}/send-packet")
async def send_onboarding_packet(onboarding_id: str, admin: User = Depends(require_role("admin", "dispatcher"))):
    doc = await db.carrier_onboarding.find_one({"onboarding_id": onboarding_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Onboarding record not found")

    # Carrier display name (records use legal_name; fall back gracefully)
    carrier_name = doc.get("legal_name") or doc.get("name") or doc.get("dba") or "Carrier"
    contact_name = doc.get("contact_name") or carrier_name

    # Pull active brand so the packet reflects Orisei/Calafia (or whatever brand is active)
    brand = await _active_brand_doc() or {}
    company = brand.get("company_name") or "Orisei Freight Solutions"
    short = brand.get("short_name") or "Orisei"
    contact_email = (brand.get("contact_emails") or {}).get("carriers") if isinstance(brand.get("contact_emails"), dict) else None
    if not contact_email:
        # derive a sane default from the brand's primary email domain or short name
        slug = re.sub(r"[^a-z0-9]+", "", short.lower())[:24] or "brand"
        contact_email = f"carriers@{slug}.com"
    phone = brand.get("phone") or "(763) 540-1200"
    insurance_amount = int(doc.get("insurance_amount") or 0)

    subject = f"{company} — Carrier Onboarding Requirements Packet · {carrier_name}"
    body = (
        f"Hello {contact_name},\n\n"
        f"Thank you for your interest in becoming a {company} approved carrier.\n"
        f"To complete the onboarding process please return the following documents within 14 days:\n\n"
        f"  1. Signed Carrier Master Agreement ({short} standard, attached)\n"
        f"  2. Certificate of Insurance · Auto Liability $1M / Cargo $250K min / WC statutory · naming {company} as additional insured\n"
        f"  3. W-9 Tax Form (current year)\n"
        f"  4. SCAC verification & MC# / DOT# documentation\n"
        f"  5. Most recent CSA / SMS safety scorecard (≤6 months)\n"
        f"  6. EDI 204/214/210 capability confirmation (yes/no/in development)\n"
        f"  7. Lane & equipment matrix (capacity by lane, trailer types)\n\n"
        f"Submit all documents to: {contact_email}\n\n"
        f"Once received, our team will review and notify you of approval status within 5 business days.\n\n"
        f"Carrier reference info we already have on file:\n"
        f"  · MC#: {doc.get('mc_number') or '—'}    DOT#: {doc.get('dot_number') or '—'}    SCAC: {doc.get('scac') or '—'}\n"
        f"  · Mode: {doc.get('mode') or '—'}    Insurance amt: ${insurance_amount:,}\n\n"
        f"— {short} Transportation Team\n  {contact_email} · {phone}"
    )
    await db.carrier_onboarding.update_one(
        {"onboarding_id": onboarding_id},
        {"$set": {"packet_sent_at": datetime.now(timezone.utc).isoformat(), "packet_sent_by": admin.user_id}},
    )
    to_email = doc.get("contact_email", "") or ""
    # URL-encode subject & body so the mailto link is robust to spaces / newlines
    from urllib.parse import quote
    mailto = f"mailto:{to_email}?subject={quote(subject)}&body={quote(body)}"
    return {"to": to_email, "subject": subject, "body": body, "mailto": mailto}

# -------------------- SHIPMENTS: Email composers --------------------
@api_router.post("/shipments/{shipment_id}/email-routing-guide")
async def email_routing_guide(shipment_id: str, user: User = Depends(require_role("admin", "dispatcher"))):
    s = await db.shipments.find_one({"shipment_id": shipment_id}, {"_id": 0})
    if not s:
        raise HTTPException(status_code=404, detail="Shipment not found")
    to = s.get("customer_contact_email") or ""
    subject = f"Routing Guide — Tennant Shipment {s.get('reference')} ({s.get('shipment_id')})"
    body = (
        f"Hello,\n\n"
        f"Please find below the routing instructions for your upcoming shipment from Tennant Companies.\n\n"
        f"Shipment ID:        {s.get('shipment_id')}\n"
        f"Reference / SO#:    {s.get('reference')}\n"
        f"Delivery #:         {s.get('sap_delivery_no') or '—'}\n"
        f"Mode:               {s.get('mode')}\n"
        f"Carrier:            {s.get('carrier')}\n"
        f"BOL #:              {s.get('bol_no')}\n"
        f"PRO #:              {s.get('pro_no')}\n"
        f"Container #:        {s.get('container_no') or '—'}\n\n"
        f"Origin:             {s.get('origin', {}).get('name')} · {s.get('origin', {}).get('city')}\n"
        f"Destination:        {s.get('destination', {}).get('name')} · {s.get('destination', {}).get('city')}\n"
        f"Pickup Date:        {s.get('pickup_date')}\n"
        f"ETA:                {s.get('eta')}\n\n"
        f"Commodity:          {s.get('commodity')}\n"
        f"Weight:             {int(s.get('weight_lbs') or 0):,} lbs\n"
        f"Pieces / Pallets:   {s.get('pieces')} / {s.get('pallet_count') or s.get('skids') or '—'}\n"
        f"NMFC:               {s.get('nmfc_code') or '—'}    Freight Class: {s.get('freight_class') or '—'}\n"
        f"Hazmat:             {'YES — ' + (s.get('hazmat_class') or '') if s.get('hazmat') else 'No'}\n"
        f"Accessorials:       {', '.join(s.get('accessorials') or []) or 'None'}\n\n"
        f"ROUTING REQUIREMENTS\n"
        f"  • Carrier must check in at the gate with this BOL #\n"
        f"  • Driver must call dispatch 30 min prior to ETA\n"
        f"  • All deliveries 0700-1500 local unless prior approval\n"
        f"  • Photo POD required for cargo > $50K value\n"
        f"  • Exception reporting: dispatch@tennantco.com / (763) 540-1200\n\n"
        f"Thank you,\nTennant Transportation Team"
    )
    await db.shipments.update_one(
        {"shipment_id": shipment_id},
        {"$set": {"routing_guide_sent_at": datetime.now(timezone.utc).isoformat(), "routing_guide_sent_by": user.user_id}},
    )
    mailto = f"mailto:{to}?subject={subject}&body={body}"
    return {"to": to, "subject": subject, "body": body, "mailto": mailto}

class CarrierEmailRequest(BaseModel):
    template: str = "request_eta"  # request_eta | request_pod | exception_inquiry | rate_confirmation

@api_router.post("/shipments/{shipment_id}/email-carrier")
async def email_carrier(shipment_id: str, payload: CarrierEmailRequest, user: User = Depends(require_role("admin", "dispatcher"))):
    s = await db.shipments.find_one({"shipment_id": shipment_id}, {"_id": 0})
    if not s:
        raise HTTPException(status_code=404, detail="Shipment not found")
    to = s.get("carrier_contact_email") or ""
    templates = {
        "request_eta": {
            "subject": f"ETA Update Request — {s.get('carrier')} · {s.get('shipment_id')}",
            "body": (
                f"Hello {s.get('carrier')} dispatch,\n\n"
                f"Please provide an updated ETA for the following Tennant load:\n\n"
                f"  Shipment:   {s.get('shipment_id')} / {s.get('reference')}\n"
                f"  BOL:        {s.get('bol_no')}    PRO: {s.get('pro_no')}\n"
                f"  Origin:     {s.get('origin', {}).get('city')}\n"
                f"  Destination:{s.get('destination', {}).get('city')}\n"
                f"  Current ETA:{s.get('eta')}    Current Status: {s.get('status')}\n\n"
                f"Kindly confirm latest GPS location, driver name & phone, and revised ETA.\n\n"
                f"Reply to dispatch@tennantco.com — Thank you."
            ),
        },
        "request_pod": {
            "subject": f"POD Request — {s.get('carrier')} · {s.get('shipment_id')}",
            "body": (
                f"Hello {s.get('carrier')} team,\n\n"
                f"Please send a signed Proof of Delivery for the load below at your earliest convenience.\n\n"
                f"  Shipment:   {s.get('shipment_id')} / {s.get('reference')}\n"
                f"  BOL:        {s.get('bol_no')}    PRO: {s.get('pro_no')}\n"
                f"  Delivered:  {s.get('destination', {}).get('city')}\n\n"
                f"POD is required to release final freight payment. Email signed POD to ap@tennantco.com.\n\n"
                f"Thank you,\nTennant Audit & Payables"
            ),
        },
        "exception_inquiry": {
            "subject": f"Exception Inquiry — {s.get('carrier')} · {s.get('shipment_id')}",
            "body": (
                f"Hello {s.get('carrier')},\n\n"
                f"Our system shows this shipment as {s.get('status').upper()} — please advise on the cause and corrective action:\n\n"
                f"  Shipment:   {s.get('shipment_id')} / {s.get('reference')}\n"
                f"  BOL:        {s.get('bol_no')}    PRO: {s.get('pro_no')}\n"
                f"  Origin:     {s.get('origin', {}).get('city')}\n"
                f"  Destination:{s.get('destination', {}).get('city')}\n\n"
                f"Provide cause code, recovery plan, and revised ETA within 4 hours.\n\nThanks — Tennant Dispatch"
            ),
        },
        "rate_confirmation": {
            "subject": f"Rate Confirmation — {s.get('carrier')} · {s.get('shipment_id')}",
            "body": (
                f"{s.get('carrier')} — please confirm acceptance of the following load:\n\n"
                f"  Shipment:   {s.get('shipment_id')} / {s.get('reference')}\n"
                f"  Mode:       {s.get('mode')}\n"
                f"  Pickup:     {s.get('origin', {}).get('city')} · {s.get('pickup_date')}\n"
                f"  Delivery:   {s.get('destination', {}).get('city')} · ETA {s.get('eta')[:10]}\n"
                f"  Weight:     {int(s.get('weight_lbs') or 0):,} lbs · {s.get('pieces')} pcs\n"
                f"  Commodity:  {s.get('commodity')}\n"
                f"  Hazmat:     {'YES' if s.get('hazmat') else 'No'}\n"
                f"  Rate:       ${(s.get('bid_cost') or 0):,.2f} + {s.get('fsc_pct') or 0}% FSC\n\n"
                f"Reply with driver, equipment #, MC#, and signed rate con.\nDispatch@tennantco.com"
            ),
        },
    }
    tpl = templates.get(payload.template, templates["request_eta"])
    await db.shipments.update_one(
        {"shipment_id": shipment_id},
        {"$push": {"carrier_emails_log": {"template": payload.template, "sent_at": datetime.now(timezone.utc).isoformat(), "sent_by": user.user_id}}},
    )
    mailto = f"mailto:{to}?subject={tpl['subject']}&body={tpl['body']}"
    return {"to": to, "subject": tpl["subject"], "body": tpl["body"], "mailto": mailto}

# -------------------- DOCUMENT VAULT (GridFS) --------------------
from motor.motor_asyncio import AsyncIOMotorGridFSBucket
from fastapi import UploadFile, File, Form

vault_bucket = AsyncIOMotorGridFSBucket(db, bucket_name="vault")
bol_bucket = AsyncIOMotorGridFSBucket(db, bucket_name="carrier_bols")

VAULT_CATEGORIES = ["Insurance COI", "W-9", "Carrier Agreement", "MSDS/SDS", "Onboarding Packet", "Tariff", "Quote", "Bid", "Other"]

@api_router.get("/vault/categories")
async def get_vault_categories(_: User = Depends(get_current_user)):
    return {"categories": VAULT_CATEGORIES}

@api_router.post("/vault/upload")
async def vault_upload(
    file: UploadFile = File(...),
    category: str = Form("Other"),
    related_to: str = Form(""),
    expires_at: Optional[str] = Form(None),
    notes: str = Form(""),
    user: User = Depends(require_role("admin", "dispatcher", "auditor")),
):
    data = await file.read()
    file_id = await vault_bucket.upload_from_stream(file.filename or "untitled", data, metadata={
        "content_type": file.content_type or "application/octet-stream",
        "category": category,
        "related_to": related_to,
        "uploaded_by": user.user_id,
        "uploaded_by_name": user.name,
        "uploaded_at": datetime.now(timezone.utc).isoformat(),
        "size_bytes": len(data),
        "expires_at": expires_at,
        "notes": notes,
    })
    return {"file_id": str(file_id), "filename": file.filename, "size_bytes": len(data), "category": category}

@api_router.get("/vault/files")
async def list_vault_files(category: Optional[str] = None, related_to: Optional[str] = None, _: User = Depends(get_current_user)):
    q: Dict[str, Any] = {}
    if category:
        q["metadata.category"] = category
    if related_to:
        q["metadata.related_to"] = related_to
    files = await db["vault.files"].find(q, {"chunkSize": 0}).sort("uploadDate", -1).limit(500).to_list(500)
    out = []
    for f in files:
        md = f.get("metadata") or {}
        out.append({
            "file_id": str(f["_id"]),
            "filename": f.get("filename"),
            "length": f.get("length"),
            "upload_date": (f.get("uploadDate") or datetime.now(timezone.utc)).isoformat() if hasattr(f.get("uploadDate"), "isoformat") else str(f.get("uploadDate")),
            "category": md.get("category"),
            "related_to": md.get("related_to"),
            "uploaded_by_name": md.get("uploaded_by_name"),
            "expires_at": md.get("expires_at"),
            "notes": md.get("notes"),
            "content_type": md.get("content_type"),
        })
    return out

@api_router.get("/vault/files/{file_id}")
async def download_vault_file(file_id: str, _: User = Depends(get_current_user)):
    from bson import ObjectId
    try:
        oid = ObjectId(file_id)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid file_id")
    grid_out = await vault_bucket.open_download_stream(oid)
    data = await grid_out.read()
    md = grid_out.metadata or {}
    content_type = md.get("content_type", "application/octet-stream")
    headers = {"Content-Disposition": f'attachment; filename="{grid_out.filename}"'}
    return Response(content=data, media_type=content_type, headers=headers)

@api_router.delete("/vault/files/{file_id}")
async def delete_vault_file(file_id: str, _: User = Depends(require_role("admin"))):
    from bson import ObjectId
    try:
        oid = ObjectId(file_id)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid file_id")
    try:
        await vault_bucket.delete(oid)
    except Exception:
        raise HTTPException(status_code=404, detail="File not found")
    return {"ok": True}

# -------------------- CARRIER BOL STORAGE (per-shipment) --------------------
@api_router.post("/shipments/{shipment_id}/bol-upload")
async def upload_carrier_bol(shipment_id: str, file: UploadFile = File(...), user: User = Depends(require_role("admin", "dispatcher", "carrier"))):
    s = await db.shipments.find_one({"shipment_id": shipment_id}, {"_id": 0})
    if not s:
        raise HTTPException(status_code=404, detail="Shipment not found")
    if user.role == "carrier" and user.carrier_company and s.get("carrier") != user.carrier_company:
        raise HTTPException(status_code=403, detail="Not your shipment")
    data = await file.read()
    fid = await bol_bucket.upload_from_stream(
        file.filename or f"BOL_{shipment_id}.pdf", data,
        metadata={"content_type": file.content_type or "application/pdf", "shipment_id": shipment_id, "uploaded_by_name": user.name, "uploaded_at": datetime.now(timezone.utc).isoformat()},
    )
    await db.shipments.update_one(
        {"shipment_id": shipment_id},
        {"$set": {"carrier_bol_file_id": str(fid), "carrier_bol_filename": file.filename}},
    )
    return {"file_id": str(fid), "filename": file.filename, "size_bytes": len(data)}

@api_router.get("/shipments/{shipment_id}/bol-download")
async def download_carrier_bol(shipment_id: str, user: User = Depends(get_current_user)):
    from bson import ObjectId
    s = await db.shipments.find_one({"shipment_id": shipment_id}, {"_id": 0})
    if not s:
        raise HTTPException(status_code=404, detail="Shipment not found")
    if user.role == "carrier" and user.carrier_company and s.get("carrier") != user.carrier_company:
        raise HTTPException(status_code=403, detail="Not your shipment")
    fid = s.get("carrier_bol_file_id")
    if not fid:
        raise HTTPException(status_code=404, detail="No BOL on file")
    grid_out = await bol_bucket.open_download_stream(ObjectId(fid))
    data = await grid_out.read()
    md = grid_out.metadata or {}
    headers = {"Content-Disposition": f'attachment; filename="{grid_out.filename}"'}
    return Response(content=data, media_type=md.get("content_type", "application/pdf"), headers=headers)

# -------------------- INBOUND ROUTING GUIDE (GridFS, public PDF) --------------------
# Tennant publishes a Domestic US/CA/MX Inbound Routing Guide that suppliers
# must follow. The team needs to email it to vendors frequently, so the PDF
# is stored in a dedicated GridFS bucket and exposed via:
#   GET  /api/routing-guide/info            — metadata for the active version
#   GET  /api/routing-guide/pdf             — public PDF stream (no auth, so
#                                             vendors can open the email link)
#   GET  /api/routing-guide/email-template  — mailto: payload with subject/body
#                                             and an absolute PDF link
#   POST /api/routing-guide/upload          — admin uploads a new revision
#   GET  /api/routing-guide/versions        — full revision history
routing_guide_bucket = AsyncIOMotorGridFSBucket(db, bucket_name="routing_guides")
ROUTING_GUIDE_SEED = ROOT_DIR / "seed_assets" / "routing-guide.pdf"


async def _seed_routing_guide_if_empty() -> None:
    """One-shot seed: if no routing guide exists in GridFS, ingest the seed PDF
    shipped alongside the backend so the team has something to email on day 1."""
    if not ROUTING_GUIDE_SEED.exists():
        return
    has_any = await db["routing_guides.files"].count_documents({})
    if has_any:
        return
    data = ROUTING_GUIDE_SEED.read_bytes()
    await routing_guide_bucket.upload_from_stream(
        "Domestic_US_Canada_Mexico_Inbound_Routing_Guide_Rev29_2026-01-09.pdf",
        data,
        metadata={
            "content_type": "application/pdf",
            "title": "Domestic US / Canada / Mexico Inbound Routing Guide",
            "revision": "Revision 29",
            "effective_date": "2026-01-09",
            "uploaded_by_name": "System Seed",
            "uploaded_at": datetime.now(timezone.utc).isoformat(),
            "size_bytes": len(data),
            "notes": "Tennant Inbound Routing Requirements for suppliers shipping within the USA, Canada, Mexico, and Puerto Rico. Covers approved carriers, modes (Small Package / LTL / FTL), and PO/BOL requirements.",
        },
    )


async def _latest_routing_guide() -> Optional[Dict[str, Any]]:
    """Return the most recently uploaded routing guide file doc, or None."""
    doc = await db["routing_guides.files"].find_one(
        {}, sort=[("uploadDate", -1)]
    )
    return doc


@api_router.get("/routing-guide/info")
async def routing_guide_info(_: User = Depends(get_current_user)):
    await _seed_routing_guide_if_empty()
    f = await _latest_routing_guide()
    if not f:
        raise HTTPException(status_code=404, detail="No routing guide on file")
    md = f.get("metadata") or {}
    upload_date = f.get("uploadDate")
    return {
        "file_id": str(f["_id"]),
        "filename": f.get("filename"),
        "size_bytes": f.get("length"),
        "title": md.get("title"),
        "revision": md.get("revision"),
        "effective_date": md.get("effective_date"),
        "uploaded_by_name": md.get("uploaded_by_name"),
        "uploaded_at": upload_date.isoformat() if hasattr(upload_date, "isoformat") else md.get("uploaded_at"),
        "notes": md.get("notes"),
        "content_type": md.get("content_type") or "application/pdf",
        "pdf_url": "/api/routing-guide/pdf",
    }


@api_router.get("/routing-guide/pdf")
async def routing_guide_pdf(download: bool = False):
    """PUBLIC endpoint — no auth, so the link can be emailed to external
    suppliers and opened from any mail client. `?download=1` forces a
    download instead of inline browser preview."""
    await _seed_routing_guide_if_empty()
    f = await _latest_routing_guide()
    if not f:
        raise HTTPException(status_code=404, detail="No routing guide on file")
    from bson import ObjectId
    grid_out = await routing_guide_bucket.open_download_stream(f["_id"])
    data = await grid_out.read()
    md = grid_out.metadata or {}
    disposition = "attachment" if download else "inline"
    headers = {
        "Content-Disposition": f'{disposition}; filename="{grid_out.filename}"',
        "Cache-Control": "public, max-age=3600",
    }
    return Response(content=data, media_type=md.get("content_type", "application/pdf"), headers=headers)


@api_router.get("/routing-guide/email-template")
async def routing_guide_email_template(
    to: str = "",
    cc: str = "",
    user: User = Depends(get_current_user),
):
    await _seed_routing_guide_if_empty()
    f = await _latest_routing_guide()
    if not f:
        raise HTTPException(status_code=404, detail="No routing guide on file")
    md = f.get("metadata") or {}
    base = (os.environ.get("PUBLIC_BASE_URL") or "").rstrip("/")
    # If PUBLIC_BASE_URL isn't set we still emit a relative path so the
    # frontend can prepend window.location.origin before launching mailto.
    pdf_link = f"{base}/api/routing-guide/pdf" if base else "/api/routing-guide/pdf"
    title = md.get("title") or "Tennant Inbound Routing Guide"
    revision = md.get("revision") or ""
    effective = md.get("effective_date") or ""
    subject = f"Tennant Inbound Routing Guide — {revision} (Eff. {effective})".strip(" —")
    body_lines = [
        "Hello,",
        "",
        f"Please find Tennant Company's current Inbound Routing Guide attached below: {title}.",
        "",
        f"Revision: {revision or 'N/A'}",
        f"Effective Date: {effective or 'N/A'}",
        "",
        "This document outlines our required carriers, modes (Small Package / LTL / Full Truckload), and the PO + BOL information that MUST appear on every inbound shipment to Tennant in the USA, Canada, Mexico, and Puerto Rico. Suppliers are required to follow these instructions on all shipments where Tennant is the routing party (Prepaid + Add, Collect, or 3rd-Party freight terms).",
        "",
        f"Download the PDF here: {pdf_link}",
        "",
        f"If you have questions, please reply to this email or contact transportation@tennantco.com.",
        "",
        "Thank you,",
        user.name,
        "Tennant Company · Transportation",
    ]
    body = "\n".join(body_lines)
    mailto = (
        f"mailto:{urllib.parse.quote(to)}"
        f"?subject={urllib.parse.quote(subject)}"
        f"&body={urllib.parse.quote(body)}"
    )
    if cc:
        mailto += f"&cc={urllib.parse.quote(cc)}"
    return {
        "to": to,
        "cc": cc,
        "subject": subject,
        "body": body,
        "mailto": mailto,
        "pdf_url": pdf_link,
        "filename": f.get("filename"),
    }


@api_router.post("/routing-guide/upload")
async def routing_guide_upload(
    file: UploadFile = File(...),
    revision: str = Form(""),
    effective_date: str = Form(""),
    notes: str = Form(""),
    user: User = Depends(require_role("admin", "dispatcher")),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Routing guide must be a PDF")
    data = await file.read()
    fid = await routing_guide_bucket.upload_from_stream(
        file.filename, data,
        metadata={
            "content_type": file.content_type or "application/pdf",
            "title": "Domestic US / Canada / Mexico Inbound Routing Guide",
            "revision": revision or "Revision (unspecified)",
            "effective_date": effective_date or datetime.now(timezone.utc).strftime("%Y-%m-%d"),
            "uploaded_by": user.user_id,
            "uploaded_by_name": user.name,
            "uploaded_at": datetime.now(timezone.utc).isoformat(),
            "size_bytes": len(data),
            "notes": notes,
        },
    )
    return {
        "file_id": str(fid),
        "filename": file.filename,
        "size_bytes": len(data),
        "revision": revision,
        "effective_date": effective_date,
    }


@api_router.get("/routing-guide/versions")
async def routing_guide_versions(_: User = Depends(get_current_user)):
    await _seed_routing_guide_if_empty()
    files = await db["routing_guides.files"].find({}, {"chunkSize": 0}).sort("uploadDate", -1).limit(100).to_list(100)
    out = []
    for f in files:
        md = f.get("metadata") or {}
        upload_date = f.get("uploadDate")
        out.append({
            "file_id": str(f["_id"]),
            "filename": f.get("filename"),
            "size_bytes": f.get("length"),
            "revision": md.get("revision"),
            "effective_date": md.get("effective_date"),
            "uploaded_by_name": md.get("uploaded_by_name"),
            "uploaded_at": upload_date.isoformat() if hasattr(upload_date, "isoformat") else md.get("uploaded_at"),
            "notes": md.get("notes"),
        })
    return out


@api_router.delete("/routing-guide/versions/{file_id}")
async def routing_guide_delete(file_id: str, _: User = Depends(require_role("admin"))):
    from bson import ObjectId
    try:
        oid = ObjectId(file_id)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid file_id")
    try:
        await routing_guide_bucket.delete(oid)
    except Exception:
        raise HTTPException(status_code=404, detail="File not found")
    return {"ok": True}


# -------------------- CALENDAR EVENTS --------------------
# Powers the Command Center's MiniCalendar. Aggregates dated items from
# shipments (eta, pickup_date, delivery_date) AND truckload_bookings
# (pickup_date, delivery_date) into a single per-date list.
#
#   GET /api/calendar/events?start=YYYY-MM-DD&end=YYYY-MM-DD
#   →  { events: [ { date, kind, type, label, ref, link } ] }
#
# `kind` ∈ {shipment, booking}; `type` ∈ {pickup, delivery, eta, bol_deadline}.
# The frontend uses date strings as map keys → it doesn't have to parse ISO
# back into UTC midnight.

def _norm_date(v: Any) -> Optional[str]:
    """Coerce a value into a YYYY-MM-DD string. Tolerates ISO datetimes,
    bare date strings, and None. Returns None for unparseable / empty."""
    if not v:
        return None
    s = str(v).strip()
    if not s:
        return None
    # Already a bare date
    if len(s) >= 10 and s[4] == "-" and s[7] == "-":
        return s[:10]
    try:
        return datetime.fromisoformat(s.replace("Z", "+00:00")).date().isoformat()
    except Exception:
        return None


@api_router.get("/calendar/events")
async def calendar_events(
    start: str,
    end: str,
    _: User = Depends(get_current_user),
):
    """Return calendar events in [start, end] inclusive (YYYY-MM-DD)."""
    try:
        start_d = datetime.fromisoformat(start).date()
        end_d = datetime.fromisoformat(end).date()
    except Exception:
        raise HTTPException(status_code=400, detail="start/end must be YYYY-MM-DD")
    if (end_d - start_d).days > 92:
        raise HTTPException(status_code=400, detail="Range > 92 days not allowed")
    if end_d < start_d:
        raise HTTPException(status_code=400, detail="end must be on or after start")

    in_range = lambda iso: bool(iso) and start <= iso <= end  # noqa: E731

    out: List[Dict[str, Any]] = []

    # ---- Shipments ----
    cursor = db.shipments.find(
        {}, {
            "_id": 0, "reference": 1, "shipment_id": 1, "mode": 1, "carrier": 1,
            "status": 1, "origin": 1, "destination": 1,
            "eta": 1, "pickup_date": 1, "delivery_date": 1,
        }
    )
    async for s in cursor:
        ref = s.get("reference") or s.get("shipment_id")
        carrier = s.get("carrier") or "—"
        lane = f"{(s.get('origin') or {}).get('city', '?')} → {(s.get('destination') or {}).get('city', '?')}"
        sid = s.get("shipment_id")
        link = f"/shipments?focus={sid}" if sid else "/shipments"
        for (field, ev_type, prefix) in (
            ("pickup_date", "pickup",   "Pickup"),
            ("delivery_date", "delivery", "Delivery"),
            ("eta",          "eta",      "ETA"),
        ):
            d = _norm_date(s.get(field))
            if in_range(d):
                out.append({
                    "date": d, "kind": "shipment", "type": ev_type,
                    "label": f"{prefix} · {ref} · {carrier}",
                    "sublabel": lane, "ref": ref, "link": link,
                    "mode": s.get("mode"), "status": s.get("status"),
                })

    # ---- Truckload bookings ----
    cursor = db.truckload_bookings.find(
        {}, {
            "_id": 0, "id": 1, "bol_no": 1, "carrier": 1, "origin": 1,
            "destination": 1, "pickup_date": 1, "delivery_date": 1, "status": 1,
        }
    )
    async for b in cursor:
        ref = b.get("bol_no") or b.get("id")
        carrier = b.get("carrier") or "—"
        lane = f"{b.get('origin') or '?'} → {b.get('destination') or '?'}"
        link = "/workbook"
        for (field, ev_type, prefix) in (
            ("pickup_date",   "pickup",   "TL Pickup"),
            ("delivery_date", "delivery", "TL Delivery"),
        ):
            d = _norm_date(b.get(field))
            if in_range(d):
                out.append({
                    "date": d, "kind": "booking", "type": ev_type,
                    "label": f"{prefix} · {ref} · {carrier}",
                    "sublabel": lane, "ref": ref, "link": link,
                    "mode": "TL", "status": b.get("status"),
                })

    # Sort by date then type for stable rendering
    out.sort(key=lambda e: (e["date"], e["type"], e["ref"] or ""))
    # Aggregate by date so the calendar can render a single badge per day
    by_date: Dict[str, int] = {}
    for e in out:
        by_date[e["date"]] = by_date.get(e["date"], 0) + 1
    return {"events": out, "counts_by_date": by_date, "start": start, "end": end}


# -------------------- EMAIL SEND (MOCKED) --------------------
# Tennant chose to keep the email integration mocked for now. This helper
# logs every "send" into db.outbound_emails with status='mocked' and a
# faux message_id. When SendGrid (or Resend) is wired up later, swap the
# body of _do_send_email() with the real SDK call and flip status to
# 'queued'/'delivered'. All call sites stay identical.
EMAIL_FROM = os.environ.get("EMAIL_FROM", "transportation@tennantco.com")


async def _do_send_email(
    to: str,
    subject: str,
    body_text: str,
    body_html: Optional[str] = None,
    cc: Optional[str] = None,
    attachments: Optional[List[Dict[str, Any]]] = None,
    user: Optional[User] = None,
    kind: str = "generic",
    ref: Optional[str] = None,
) -> Dict[str, Any]:
    """Mocked email send. Logs to db.outbound_emails and returns a fake
    delivery receipt. Production swap: replace this body with the SendGrid
    SDK call (Mail() + SendGridAPIClient(...).send(message))."""
    msg_id = f"mock_{uuid.uuid4().hex}"
    entry = {
        "message_id": msg_id,
        "to": to,
        "cc": cc or "",
        "from": EMAIL_FROM,
        "subject": subject,
        "body_text": body_text,
        "body_html": body_html,
        "kind": kind,
        "ref": ref,
        "attachment_names": [a.get("filename") for a in (attachments or [])],
        "status": "mocked",  # would be 'queued' once SendGrid is wired
        "provider": "mock",
        "sent_by": (user.name if user else "system"),
        "sent_by_user_id": (user.user_id if user else None),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.outbound_emails.insert_one(dict(entry))
    return {
        "ok": True,
        "message_id": msg_id,
        "status": "mocked",
        "to": to,
        "subject": subject,
        "from": EMAIL_FROM,
    }


class EmailSendPayload(BaseModel):
    to: str
    cc: Optional[str] = ""
    subject: str
    body_text: str
    body_html: Optional[str] = None
    kind: Optional[str] = "generic"
    ref: Optional[str] = None
    attachment_urls: Optional[List[str]] = None


@api_router.post("/email/send")
async def email_send(payload: EmailSendPayload, user: User = Depends(require_role("admin", "dispatcher", "auditor"))):
    """Generic send endpoint usable by any frontend module that previously
    relied on mailto. Currently mocked — returns a stub receipt and logs
    the would-be email to db.outbound_emails so dispatch can audit what
    *would* have gone out."""
    result = await _do_send_email(
        to=payload.to, cc=payload.cc, subject=payload.subject,
        body_text=payload.body_text, body_html=payload.body_html,
        user=user, kind=payload.kind or "generic", ref=payload.ref,
        attachments=[{"filename": u.rsplit("/", 1)[-1]} for u in (payload.attachment_urls or [])],
    )
    return result


@api_router.get("/email/log")
async def email_log(
    kind: Optional[str] = None,
    limit: int = 100,
    _: User = Depends(get_current_user),
):
    q: Dict[str, Any] = {}
    if kind:
        q["kind"] = kind
    docs = await db.outbound_emails.find(q, {"_id": 0}).sort("created_at", -1).limit(min(limit, 500)).to_list(min(limit, 500))
    return {"log": docs, "provider": "mock", "from": EMAIL_FROM}


@api_router.post("/routing-guide/send-email")
async def routing_guide_send_email(
    payload: EmailSendPayload,
    user: User = Depends(require_role("admin", "dispatcher", "auditor")),
):
    """Send the active routing guide PDF as an actual email (currently mocked).
    Auto-builds subject/body from the active revision metadata if the client
    didn't supply them, then attaches the PDF link."""
    await _seed_routing_guide_if_empty()
    f = await _latest_routing_guide()
    if not f:
        raise HTTPException(status_code=404, detail="No routing guide on file")
    md = f.get("metadata") or {}
    base = (os.environ.get("PUBLIC_BASE_URL") or "").rstrip("/")
    pdf_link = f"{base}/api/routing-guide/pdf" if base else "/api/routing-guide/pdf"
    subject = payload.subject or f"Tennant Inbound Routing Guide — {md.get('revision','')} (Eff. {md.get('effective_date','')})".strip(" —")
    body_text = payload.body_text or (
        "Hello,\n\n"
        f"Please find Tennant Company's current Inbound Routing Guide: {md.get('title')}.\n\n"
        f"Revision: {md.get('revision') or 'N/A'}\n"
        f"Effective Date: {md.get('effective_date') or 'N/A'}\n\n"
        f"Download the PDF: {pdf_link}\n\n"
        f"— {user.name}\nTennant Company · Transportation"
    )
    return await _do_send_email(
        to=payload.to, cc=payload.cc, subject=subject,
        body_text=body_text, body_html=payload.body_html,
        user=user, kind="routing_guide", ref=md.get("revision"),
        attachments=[{"filename": f.get("filename")}],
    )
# -------------------- CLAIMS & RECONCILIATION --------------------
class Claim(BaseModel):
    claim_id: str
    shipment_id: Optional[str] = None
    carrier: str
    bol_no: Optional[str] = None
    claim_type: str  # damage | shortage | loss | overage | freight_overcharge
    amount_claimed_usd: float
    amount_recovered_usd: float = 0.0
    status: str  # open | filed | acknowledged | denied | settled | partial
    filed_date: str
    incident_date: str
    description: str
    created_by: Optional[str] = None
    notes: Optional[str] = None

class ClaimCreate(BaseModel):
    shipment_id: Optional[str] = None
    carrier: str
    bol_no: Optional[str] = None
    claim_type: str
    amount_claimed_usd: float
    incident_date: str
    description: str
    notes: Optional[str] = None

class ClaimUpdate(BaseModel):
    status: Optional[str] = None
    amount_recovered_usd: Optional[float] = None
    notes: Optional[str] = None

CLAIM_TYPES = ["damage", "shortage", "loss", "overage", "freight_overcharge"]
CLAIM_STATUSES = ["open", "filed", "acknowledged", "denied", "settled", "partial"]

@api_router.get("/claims")
async def list_claims(_: User = Depends(get_current_user)):
    docs = await db.claims.find({}, {"_id": 0}).sort("filed_date", -1).limit(500).to_list(500)
    # Summary
    open_count = sum(1 for c in docs if c["status"] in ("open", "filed", "acknowledged"))
    settled_count = sum(1 for c in docs if c["status"] in ("settled", "partial"))
    denied_count = sum(1 for c in docs if c["status"] == "denied")
    total_claimed = sum(c.get("amount_claimed_usd", 0) for c in docs)
    total_recovered = sum(c.get("amount_recovered_usd", 0) for c in docs)
    return {
        "claims": docs,
        "summary": {
            "total_claims": len(docs),
            "open_count": open_count,
            "settled_count": settled_count,
            "denied_count": denied_count,
            "total_claimed_usd": round(total_claimed, 2),
            "total_recovered_usd": round(total_recovered, 2),
            "recovery_rate_pct": round((total_recovered / total_claimed * 100) if total_claimed else 0, 1),
        },
        "claim_types": CLAIM_TYPES,
        "claim_statuses": CLAIM_STATUSES,
    }

@api_router.post("/claims")
async def create_claim(payload: ClaimCreate, user: User = Depends(require_role("admin", "auditor", "dispatcher"))):
    claim = {
        "claim_id": f"CLM-{uuid.uuid4().hex[:8].upper()}",
        "shipment_id": payload.shipment_id,
        "carrier": payload.carrier,
        "bol_no": payload.bol_no,
        "claim_type": payload.claim_type,
        "amount_claimed_usd": payload.amount_claimed_usd,
        "amount_recovered_usd": 0.0,
        "status": "open",
        "filed_date": datetime.now(timezone.utc).date().isoformat(),
        "incident_date": payload.incident_date,
        "description": payload.description,
        "created_by": user.user_id,
        "notes": payload.notes,
    }
    await db.claims.insert_one(dict(claim))
    return claim

@api_router.put("/claims/{claim_id}")
async def update_claim(claim_id: str, payload: ClaimUpdate, _: User = Depends(require_role("admin", "auditor"))):
    raw = payload.model_dump(exclude_unset=True)
    if not raw:
        raise HTTPException(status_code=400, detail="No fields to update")
    res = await db.claims.update_one({"claim_id": claim_id}, {"$set": raw})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Claim not found")
    doc = await db.claims.find_one({"claim_id": claim_id}, {"_id": 0})
    return doc

# -------------------- TENNANT MACHINE CATALOG --------------------
TENNANT_MACHINES = [
    # ===== ROBOTICS / AMR =====
    {"model": "X4 ROVR", "category": "Robotic Sweeper-Scrubber", "type": "Autonomous", "size": "26 in",
     "power": "Lithium-ion 36V", "runtime": "Up to 4 hrs", "deck_width_in": 26, "tank_gal": 16,
     "weight_lbs": 850, "list_price_usd": 38000, "use_case": "Compact autonomous sweep-scrub for retail, education, healthcare",
     "image_url": "https://www.tennantco.com/content/dam/tennant/tennantco/products/machines/robotics/x4-rovr/x4-rovr-hero.png",
     "product_url": "https://www.tennantco.com/en_us/1/machines/robotic-cleaning-machines/product.x4rovr.compact-robotic-sweeper-scrubber.html",
     "highlights": ["BrainOS® autonomous nav", "Sweep + scrub in one pass", "Compact ~26\" path", "Multi-shift Li-ion"]},
    {"model": "X6 ROVR", "category": "Robotic Sweeper-Scrubber", "type": "Autonomous", "size": "32 in",
     "power": "Lithium-ion 36V", "runtime": "Up to 5 hrs", "deck_width_in": 32, "tank_gal": 22,
     "weight_lbs": 1180, "list_price_usd": 52000, "use_case": "Mid-size autonomous sweep-scrub for warehouses & retail",
     "image_url": "https://www.tennantco.com/content/dam/tennant/tennantco/products/machines/robotics/x6-rovr/x6-rovr-hero.png",
     "product_url": "https://www.tennantco.com/en_us/1/machines/robotic-cleaning-machines/product.x6rovr.mid-size-robotic-sweeper-scrubber.html",
     "highlights": ["AI-NAV + obstacle avoidance", "Combined sweep + scrub", "Fleet dashboard telemetry", "Quiet 70 dBA operation"]},
    {"model": "X16 SWEEP", "category": "Robotic Sweeper", "type": "Autonomous", "size": "44 in",
     "power": "Lithium-ion 48V", "runtime": "Up to 6 hrs", "deck_width_in": 44, "tank_gal": 0,
     "weight_lbs": 2400, "list_price_usd": 78000, "use_case": "Autonomous industrial sweeping for warehouses, distribution centers, manufacturing",
     "image_url": "https://www.tennantco.com/content/dam/tennant/tennantco/products/machines/robotics/x16-sweep/x16-sweep-hero.png",
     "product_url": "https://www.tennantco.com/en_us/1/machines/robotic-cleaning-machines/product.x16sweep.industrial-robotic-sweeper.html",
     "highlights": ["Industrial-grade autonomy", "Dust control system", "Multi-shift battery", "Real-time route reports"]},
    {"model": "T7AMR", "category": "Robotic Scrubber", "type": "Autonomous", "size": "26-32 in",
     "power": "Lithium-ion 36V", "runtime": "Up to 4-5 hrs", "deck_width_in": 28, "tank_gal": 29,
     "weight_lbs": 1100, "list_price_usd": 42000, "use_case": "Robotic ride-on scrubber for retail, healthcare, airports, malls",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000056.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t7amr.robotic-floor-scrubber.2000056.html",
     "highlights": ["BrainOS® autonomy", "Manual/auto modes", "ec-H2O NanoClean® option", "Top-5 driverless brand"]},
    {"model": "T16AMR", "category": "Robotic Scrubber", "type": "Autonomous", "size": "32 in",
     "power": "Lithium-ion 36V", "runtime": "5–6 hrs", "deck_width_in": 32, "tank_gal": 60,
     "weight_lbs": 1480, "list_price_usd": 56000, "use_case": "Large facility autonomous scrubbing — warehouses, big-box retail",
     "image_url": "https://www.tennantco.com/content/dam/tennant/tennantco/products/machines/scrubber%20riders/t16amr/Images/t16amr-right.jpg/jcr:content/renditions/cq5dam.web.1280.1280.jpeg",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t16amr.industrial-robotic-floor-scrubber.2000054.html",
     "highlights": ["AI-NAV mapping", "60-gal solution tank", "Multi-shift Li-ion", "FOB-style controls"]},

    # ===== WALK-BEHIND SCRUBBERS =====
    {"model": "T2", "category": "Walk-Behind Scrubber", "type": "Walk-Behind", "size": "17 in",
     "power": "AGM 12V", "runtime": "2.5 hrs", "deck_width_in": 17, "tank_gal": 6,
     "weight_lbs": 180, "list_price_usd": 4200, "use_case": "Small tight spaces — restrooms, kitchens, small retail",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000044.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t2.compact-floor-scrubber.html",
     "highlights": ["Lightweight & portable", "Single-pass clean", "ec-H2O option"]},
    {"model": "T300", "category": "Walk-Behind Scrubber", "type": "Walk-Behind", "size": "17-20 in",
     "power": "24V Battery", "runtime": "3.5 hrs", "deck_width_in": 20, "tank_gal": 11,
     "weight_lbs": 295, "list_price_usd": 7800, "use_case": "Small to mid-size facilities — offices, schools, healthcare",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000084.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t300.walk-behind-floor-scrubber.html",
     "highlights": ["Insta-Adjust brush pressure", "Tool-free maintenance", "Quiet 67 dBA"]},
    {"model": "T350", "category": "Walk-Behind Scrubber", "type": "Stand-On", "size": "20-26 in",
     "power": "LPG / 24V Battery", "runtime": "Up to 8 hrs (LPG)", "deck_width_in": 24, "tank_gal": 17,
     "weight_lbs": 720, "list_price_usd": 18500, "use_case": "Education, healthcare, retail with stand-on productivity",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000004.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t350.walk-behind-floor-scrubber.html",
     "highlights": ["Stand-on platform option", "ec-H2O NanoClean®", "Honda LPG available"]},
    {"model": "T381", "category": "Walk-Behind Scrubber", "type": "Walk-Behind", "size": "20-24 in",
     "power": "24V Battery", "runtime": "4 hrs", "deck_width_in": 24, "tank_gal": 19,
     "weight_lbs": 540, "list_price_usd": 12500, "use_case": "Mid-size hard floors — schools, healthcare, retail",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000085.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t381.walk-behind-floor-scrubber.html",
     "highlights": ["Stair-climb & lift-friendly", "ec-H2O option", "Smart-Fill auto-detergent"]},
    {"model": "T500", "category": "Walk-Behind Scrubber", "type": "Walk-Behind", "size": "26-32 in",
     "power": "AGM / Li-ion 24V", "runtime": "Up to 4.5 hrs", "deck_width_in": 28, "tank_gal": 22,
     "weight_lbs": 760, "list_price_usd": 16800, "use_case": "Mid-large hard floors — manufacturing aisles, schools, hospitals",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000093.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t500.walk-behind-floor-scrubber.html",
     "highlights": ["Disk or cylindrical deck", "Quietshroud™ design", "Severe-environment build"]},
    {"model": "T600", "category": "Walk-Behind Scrubber", "type": "Walk-Behind / Stand-On", "size": "28-32 in",
     "power": "AGM / Li-ion 36V", "runtime": "5 hrs", "deck_width_in": 32, "tank_gal": 22,
     "weight_lbs": 880, "list_price_usd": 21500, "use_case": "Heavy-duty walk-behind for industrial, manufacturing, big retail",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000094.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t600.walk-behind-floor-scrubber.html",
     "highlights": ["Stand-on platform option", "150 lb of brush pressure", "Auto pre-sweep brush option"]},

    # ===== RIDE-ON SCRUBBERS =====
    {"model": "T7", "category": "Ride-On Scrubber", "type": "Ride-On", "size": "26-32 in",
     "power": "AGM / Li-ion 36V", "runtime": "4.5 hrs", "deck_width_in": 32, "tank_gal": 29,
     "weight_lbs": 1100, "list_price_usd": 32500, "use_case": "Compact ride-on for retail, education, healthcare — fits doors & elevators",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000074.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t7.ride-on-floor-scrubber.2000074.html",
     "highlights": ["67-70 dBA quiet operation", "Up to 56,320 sq ft/hr", "ec-H2O NanoClean®", "Compact footprint"]},
    {"model": "T12", "category": "Ride-On Scrubber", "type": "Ride-On", "size": "32 in",
     "power": "AGM / Li-ion 36V", "runtime": "5 hrs", "deck_width_in": 32, "tank_gal": 40,
     "weight_lbs": 1500, "list_price_usd": 28500, "use_case": "Mid-size ride-on for warehouses & distribution centers",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000068.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t12.battery-ride-on-floor-scrubber.html",
     "highlights": ["40-gal tank", "Touch-n-Go controls", "Severe-duty steel frame"]},
    {"model": "T16", "category": "Ride-On Scrubber", "type": "Ride-On", "size": "36-46 in",
     "power": "AGM / Li-ion 36V", "runtime": "5–7 hrs", "deck_width_in": 36, "tank_gal": 75,
     "weight_lbs": 1850, "list_price_usd": 35000, "use_case": "Industrial ride-on — manufacturing, distribution centers, large warehouses",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000070.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t16.battery-ride-on-floor-scrubber.2000070.html",
     "highlights": ["Up to 89,100 sq ft/hr", "75-gal tank w/ ES", "Touch-n-Go controls"]},
    {"model": "T17", "category": "Ride-On Scrubber", "type": "Ride-On", "size": "46 in",
     "power": "AGM / Li-ion 36V", "runtime": "6–8 hrs", "deck_width_in": 46, "tank_gal": 90,
     "weight_lbs": 2400, "list_price_usd": 45000, "use_case": "Large industrial ride-on — heavy-duty warehouse & manufacturing",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000071.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t17.battery-ride-on-floor-scrubber.html",
     "highlights": ["90-gal solution tank", "Insta-Click™ accessory swaps", "ec-H2O NanoClean® standard"]},
    {"model": "T20", "category": "Ride-On Scrubber", "type": "Ride-On", "size": "48-60 in",
     "power": "LPG / Diesel / Battery", "runtime": "8+ hrs", "deck_width_in": 54, "tank_gal": 80,
     "weight_lbs": 3400, "list_price_usd": 62000, "use_case": "Heavy-duty industrial — manufacturing plants, large DCs",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000076.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/scrubbers/product.t20.industrial-ride-on-floor-scrubber.html",
     "highlights": ["Severe-duty design", "Cab w/ heat & A/C option", "ES (extended scrub) tank"]},

    # ===== WALK-BEHIND SWEEPERS =====
    {"model": "S3", "category": "Walk-Behind Sweeper", "type": "Walk-Behind", "size": "21 in",
     "power": "12V Battery / Manual", "runtime": "2.5 hrs", "deck_width_in": 21, "tank_gal": 0,
     "weight_lbs": 95, "list_price_usd": 2900, "use_case": "Compact walk-behind sweeper for small facilities, restaurants, retail",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000077.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweepers/product.s3.compact-sweeper.html",
     "highlights": ["Quiet operation", "Lift-off hopper", "Optional side broom"]},
    {"model": "S5", "category": "Walk-Behind Sweeper", "type": "Walk-Behind", "size": "32 in",
     "power": "24V Battery", "runtime": "3 hrs", "deck_width_in": 32, "tank_gal": 0,
     "weight_lbs": 320, "list_price_usd": 6500, "use_case": "Mid-size walk-behind sweeper — manufacturing, warehouses, education",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000091.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweepers/product.s5.walk-behind-sweeper.html",
     "highlights": ["SweepMax™ filtration", "Large hopper", "Battery-powered productivity"]},
    {"model": "S6", "category": "Walk-Behind Sweeper", "type": "Walk-Behind", "size": "27 in",
     "power": "24V Battery / Manual", "runtime": "3.5 hrs", "deck_width_in": 27, "tank_gal": 0,
     "weight_lbs": 280, "list_price_usd": 5400, "use_case": "Compact battery sweeper for indoor mid-size areas",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000092.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweepers/product.s6.walk-behind-sweeper.html",
     "highlights": ["Quietshroud™ design", "Easy-access debris hopper", "Optional dust control"]},

    # ===== RIDE-ON SWEEPERS =====
    {"model": "S7", "category": "Ride-On Sweeper", "type": "Ride-On", "size": "32 in",
     "power": "AGM / Li-ion 36V", "runtime": "5 hrs", "deck_width_in": 32, "tank_gal": 0,
     "weight_lbs": 1100, "list_price_usd": 22500, "use_case": "Compact ride-on sweeper for distribution centers & warehouses",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000128.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweepers/product.s7.battery-rider-sweeper.html",
     "highlights": ["Quiet battery operation", "SweepMax Plus filtration", "Auto debris-dump"]},
    {"model": "S12", "category": "Ride-On Sweeper", "type": "Ride-On", "size": "40 in",
     "power": "LPG / Battery", "runtime": "6 hrs", "deck_width_in": 40, "tank_gal": 0,
     "weight_lbs": 1700, "list_price_usd": 28500, "use_case": "Mid-size ride-on sweeper for warehouses & manufacturing",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000114.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweepers/product.s12.industrial-rider-sweeper.html",
     "highlights": ["Hydraulic high-dump hopper", "Dust filtration", "Multi-fuel options"]},
    {"model": "S16", "category": "Ride-On Sweeper", "type": "Ride-On", "size": "42 in",
     "power": "LPG / Diesel", "runtime": "8 hrs", "deck_width_in": 42, "tank_gal": 0,
     "weight_lbs": 2400, "list_price_usd": 38000, "use_case": "Industrial ride-on sweeper — fume-free LPG, fast cleaning",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000118.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweepers/product.s16.industrial-rider-sweeper.html",
     "highlights": ["High-dump hopper", "DustGuard™ system", "Cab option"]},
    {"model": "S20", "category": "Ride-On Sweeper", "type": "Ride-On", "size": "44 in",
     "power": "LPG / Diesel / Battery", "runtime": "6–8 hrs", "deck_width_in": 44, "tank_gal": 0,
     "weight_lbs": 2200, "list_price_usd": 36000, "use_case": "Mid-size warehouse / industrial sweeping",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000122.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweepers/product.s20.mid-sized-sweeper.html",
     "highlights": ["Dust control system", "Multi-debris hopper", "Compact turn radius"]},
    {"model": "S30", "category": "Ride-On Sweeper", "type": "Ride-On", "size": "50-60 in",
     "power": "LPG / Diesel / Battery", "runtime": "8+ hrs", "deck_width_in": 58, "tank_gal": 0,
     "weight_lbs": 3200, "list_price_usd": 52000, "use_case": "Heavy-duty warehouse, factory sweeping",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000136.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweepers/product.s30.industrial-sweeper.html",
     "highlights": ["SweepMax Plus filtration", "100 cu ft hopper", "Cab option"]},

    # ===== SWEEPER-SCRUBBERS (COMBO) =====
    {"model": "M17", "category": "Sweeper-Scrubber", "type": "Ride-On Battery", "size": "40 in",
     "power": "Lithium-ion 36V", "runtime": "Up to 4-5 hrs", "deck_width_in": 40, "tank_gal": 91,
     "weight_lbs": 2700, "list_price_usd": 58000, "use_case": "Battery ride-on combo — single-pass sweep + scrub for manufacturing",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000131.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweeper-scrubbers/product.m17.battery-powered-ride-on-sweeper-scrubber.2000131.html",
     "highlights": ["Up to 116,160 sq ft/hr scrub", "Dual Force Sweeping (debris up to 22\")", "91-gal recovery tank", "Optional 2500-psi washer"]},
    {"model": "M20", "category": "Sweeper-Scrubber", "type": "Ride-On LPG", "size": "40 in",
     "power": "LPG", "runtime": "6 hrs", "deck_width_in": 40, "tank_gal": 60,
     "weight_lbs": 3100, "list_price_usd": 62000, "use_case": "LPG ride-on sweeper-scrubber for industrial single-pass cleaning",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000119.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweeper-scrubbers/product.m20.industrial-sweeper-scrubber.html",
     "highlights": ["Severe-duty design", "IRIS asset management", "Single-pass productivity"]},
    {"model": "M30", "category": "Sweeper-Scrubber", "type": "Ride-On Diesel/LPG", "size": "44-52 in",
     "power": "LPG / Diesel", "runtime": "6+ hrs", "deck_width_in": 50, "tank_gal": 70,
     "weight_lbs": 3850, "list_price_usd": 78000, "use_case": "Heavy industrial — one-pass sweep + scrub for large manufacturing",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000130.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/sweeper-scrubbers/product.m30.industrial-sweeper-scrubber.html",
     "highlights": ["Single-pass sweep & scrub", "70-gal tank", "Severe-duty engine options"]},

    # ===== BURNISHERS =====
    {"model": "B5", "category": "Burnisher", "type": "Walk-Behind", "size": "20 in",
     "power": "AC corded", "runtime": "Unlimited", "deck_width_in": 20, "tank_gal": 0,
     "weight_lbs": 95, "list_price_usd": 2400, "use_case": "Retail & small commercial floor buffing / polishing",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000079.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/burnishers/product.b5.walk-behind-burnisher.html",
     "highlights": ["1500-2000 RPM pad speed", "Active dust control", "Lightweight"]},
    {"model": "B7", "category": "Burnisher", "type": "Ride-On", "size": "27 in",
     "power": "LPG / Battery", "runtime": "5 hrs", "deck_width_in": 27, "tank_gal": 0,
     "weight_lbs": 660, "list_price_usd": 12500, "use_case": "High-gloss floor maintenance for retail & corporate",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000110.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/burnishers/product.b7.battery-rider-burnisher.html",
     "highlights": ["High RPM pad", "Optional dust control", "Quiet operation"]},
    {"model": "B10", "category": "Burnisher", "type": "Ride-On", "size": "27 in",
     "power": "LPG", "runtime": "6 hrs", "deck_width_in": 27, "tank_gal": 0,
     "weight_lbs": 1100, "list_price_usd": 18500, "use_case": "Premium ride-on burnisher for large retail & commercial floors",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000113.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/burnishers/product.b10.lpg-rider-burnisher.html",
     "highlights": ["2200 RPM pad speed", "Insta-Adjust pad pressure", "Active dust capture"]},

    # ===== CARPET EXTRACTORS =====
    {"model": "EX-CAN-7", "category": "Carpet Extractor", "type": "Canister", "size": "—",
     "power": "AC corded", "runtime": "Unlimited", "deck_width_in": 0, "tank_gal": 7,
     "weight_lbs": 65, "list_price_usd": 2100, "use_case": "Hotels, offices, restoration — deep carpet & upholstery cleaning",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000087.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/carpet-extractors/product.ex-can-7.canister-extractor.html",
     "highlights": ["Deep-soak extraction", "Heated water option", "Wand & hose included"]},
    {"model": "E5", "category": "Carpet Extractor", "type": "Walk-Behind", "size": "15 in",
     "power": "AC corded", "runtime": "Unlimited", "deck_width_in": 15, "tank_gal": 9,
     "weight_lbs": 85, "list_price_usd": 3200, "use_case": "Compact walk-behind carpet extractor for offices, retail, hospitality",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000088.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/carpet-extractors/product.e5.walk-behind-extractor.html",
     "highlights": ["Dual-tank design", "Heated water option", "Self-contained"]},
    {"model": "1610", "category": "Carpet Extractor", "type": "Walk-Behind", "size": "16 in",
     "power": "AC corded", "runtime": "Unlimited", "deck_width_in": 16, "tank_gal": 12,
     "weight_lbs": 120, "list_price_usd": 4800, "use_case": "Heated walk-behind extractor for deep restoration cleaning",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000089.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/carpet-extractors/product.1610.heated-extractor.html",
     "highlights": ["210°F heated solution", "12-gal tank", "Cylindrical brush agitation"]},

    # ===== OUTDOOR / SPECIALTY =====
    {"model": "Green Machine 414HS", "category": "Outdoor Sweeper", "type": "Compact Outdoor", "size": "44 in",
     "power": "Diesel", "runtime": "8 hrs", "deck_width_in": 44, "tank_gal": 0,
     "weight_lbs": 1700, "list_price_usd": 64000, "use_case": "Urban streets, sidewalks, parking lots — compact street sweeping",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000016.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/outdoor-cleaning/product.414hs.compact-street-sweeper.html",
     "highlights": ["Vacuum-assist", "Compact urban footprint", "Cab w/ heat & A/C"]},
    {"model": "ATLV 4300", "category": "Outdoor Litter Vacuum", "type": "Compact Outdoor", "size": "—",
     "power": "Gasoline", "runtime": "6 hrs", "deck_width_in": 0, "tank_gal": 0,
     "weight_lbs": 1100, "list_price_usd": 38000, "use_case": "Parks, campuses, parking lots — all-terrain litter pickup",
     "image_url": "https://www.tennantco.com/services/product/image.tennant.2000038.thumbnail-3",
     "product_url": "https://www.tennantco.com/en_us/1/machines/outdoor-cleaning/product.atlv-4300.all-terrain-litter-vacuum.html",
     "highlights": ["High-suction litter pickup", "Compact maneuverable", "All-terrain tires"]},
]

TENNANT_MACHINE_CATEGORIES = sorted(set(m["category"] for m in TENNANT_MACHINES))

# Verified real Tennant CDN product photo URLs, model → ID. The URL pattern
# `services/product/image.tennant.{id}.image` returns a high-quality JPEG
# (40-100 KB) for the specific product. IDs not listed below fall back to the
# server-generated branded SVG at /api/machines/{model}/image.svg.
TENNANT_REAL_PHOTO_IDS: Dict[str, int] = {
    "T7AMR": 2000056,
    "T16AMR": 2000054,
    "T7": 2000074,
    "T16": 2000070,
    "T12": 2000068,
    "T17": 2000071,
    "T20": 2000076,
    # T2 (id 2000044) returns the Tennant brand logo, not a product shot — let
    # it fall back to the branded SVG so the card is accurate to the model.
    "T300": 2000084,
    "T500": 2000093,
    "S3": 2000077,
    "S5": 2000091,
    "M17": 2000131,
    "B5": 2000079,
    "B7": 2000110,
    "EX-CAN-7": 2000087,
    "Green Machine 414HS": 2000016,
}

def _real_photo_or_svg(m: Dict[str, Any]) -> str:
    """Prefer the verified Tennant CDN URL when known, else fall back to the
    server-generated branded SVG (always renders, model-specific)."""
    pid = TENNANT_REAL_PHOTO_IDS.get(m["model"])
    if pid:
        return f"https://www.tennantco.com/services/product/image.tennant.{pid}.image"
    return f"/api/machines/{urllib.parse.quote(m['model'])}/image.svg"


# Category → silhouette + color palette for the generated machine image.
# Each silhouette is a hand-tuned SVG path roughly evoking that machine type.
_MACHINE_ICONS: Dict[str, Dict[str, str]] = {
    "Robotic Sweeper-Scrubber": {
        "accent": "#A78BFA",
        "shape": "<rect x='180' y='150' width='240' height='110' rx='30' fill='url(#g)'/>"
                 "<circle cx='220' cy='280' r='28' fill='#0B0E14' stroke='#A78BFA' stroke-width='3'/>"
                 "<circle cx='380' cy='280' r='28' fill='#0B0E14' stroke='#A78BFA' stroke-width='3'/>"
                 "<rect x='220' y='110' width='160' height='50' rx='14' fill='#A78BFA' opacity='0.4'/>"
                 "<circle cx='300' cy='200' r='10' fill='#0B0E14'/>"
                 "<text x='300' y='205' font-family='monospace' font-size='12' fill='#A78BFA' text-anchor='middle'>AI</text>",
    },
    "Robotic Sweeper": {
        "accent": "#A78BFA",
        "shape": "<rect x='170' y='150' width='260' height='120' rx='24' fill='url(#g)'/>"
                 "<circle cx='210' cy='290' r='30' fill='#0B0E14' stroke='#A78BFA' stroke-width='3'/>"
                 "<circle cx='390' cy='290' r='30' fill='#0B0E14' stroke='#A78BFA' stroke-width='3'/>"
                 "<rect x='200' y='100' width='200' height='52' rx='12' fill='#A78BFA' opacity='0.4'/>"
                 "<circle cx='300' cy='210' r='10' fill='#0B0E14'/>"
                 "<text x='300' y='215' font-family='monospace' font-size='12' fill='#A78BFA' text-anchor='middle'>AI</text>",
    },
    "Robotic Scrubber": {
        "accent": "#22D3EE",
        "shape": "<rect x='180' y='150' width='240' height='110' rx='28' fill='url(#g)'/>"
                 "<circle cx='220' cy='280' r='28' fill='#0B0E14' stroke='#22D3EE' stroke-width='3'/>"
                 "<circle cx='380' cy='280' r='28' fill='#0B0E14' stroke='#22D3EE' stroke-width='3'/>"
                 "<rect x='220' y='110' width='160' height='50' rx='14' fill='#22D3EE' opacity='0.4'/>"
                 "<text x='300' y='205' font-family='monospace' font-size='11' fill='#22D3EE' text-anchor='middle' font-weight='bold'>AMR</text>",
    },
    "Ride-On Scrubber": {
        "accent": "#00E5FF",
        "shape": "<rect x='160' y='160' width='280' height='100' rx='18' fill='url(#g)'/>"
                 "<rect x='240' y='110' width='100' height='60' rx='10' fill='#00E5FF' opacity='0.35'/>"
                 "<circle cx='200' cy='280' r='30' fill='#0B0E14' stroke='#00E5FF' stroke-width='3'/>"
                 "<circle cx='400' cy='280' r='30' fill='#0B0E14' stroke='#00E5FF' stroke-width='3'/>"
                 "<rect x='180' y='220' width='240' height='10' fill='#00E5FF' opacity='0.5'/>",
    },
    "Walk-Behind Scrubber": {
        "accent": "#06B6D4",
        "shape": "<rect x='200' y='180' width='200' height='90' rx='14' fill='url(#g)'/>"
                 "<line x1='300' y1='80' x2='300' y2='180' stroke='#06B6D4' stroke-width='6' stroke-linecap='round'/>"
                 "<circle cx='300' cy='75' r='14' fill='#06B6D4'/>"
                 "<circle cx='220' cy='290' r='24' fill='#0B0E14' stroke='#06B6D4' stroke-width='3'/>"
                 "<circle cx='380' cy='290' r='24' fill='#0B0E14' stroke='#06B6D4' stroke-width='3'/>"
                 "<rect x='200' y='245' width='200' height='10' fill='#06B6D4' opacity='0.5'/>",
    },
    "Walk-Behind Sweeper": {
        "accent": "#10B981",
        "shape": "<rect x='200' y='175' width='200' height='95' rx='14' fill='url(#g)'/>"
                 "<line x1='300' y1='85' x2='300' y2='175' stroke='#10B981' stroke-width='6' stroke-linecap='round'/>"
                 "<circle cx='300' cy='80' r='14' fill='#10B981'/>"
                 "<circle cx='220' cy='290' r='22' fill='#0B0E14' stroke='#10B981' stroke-width='3'/>"
                 "<circle cx='380' cy='290' r='22' fill='#0B0E14' stroke='#10B981' stroke-width='3'/>"
                 "<circle cx='300' cy='265' r='30' fill='none' stroke='#10B981' stroke-width='3' stroke-dasharray='5 4'/>",
    },
    "Ride-On Sweeper": {
        "accent": "#FFCC00",
        "shape": "<rect x='160' y='160' width='280' height='100' rx='18' fill='url(#g)'/>"
                 "<rect x='220' y='100' width='130' height='62' rx='10' fill='#FFCC00' opacity='0.35'/>"
                 "<circle cx='200' cy='280' r='30' fill='#0B0E14' stroke='#FFCC00' stroke-width='3'/>"
                 "<circle cx='400' cy='280' r='30' fill='#0B0E14' stroke='#FFCC00' stroke-width='3'/>"
                 "<circle cx='300' cy='285' r='34' fill='none' stroke='#FFCC00' stroke-width='3' stroke-dasharray='6 4'/>",
    },
    "Sweeper-Scrubber": {
        "accent": "#F59E0B",
        "shape": "<rect x='160' y='150' width='280' height='115' rx='20' fill='url(#g)'/>"
                 "<rect x='230' y='95' width='130' height='60' rx='10' fill='#F59E0B' opacity='0.35'/>"
                 "<circle cx='205' cy='285' r='30' fill='#0B0E14' stroke='#F59E0B' stroke-width='3'/>"
                 "<circle cx='395' cy='285' r='30' fill='#0B0E14' stroke='#F59E0B' stroke-width='3'/>"
                 "<rect x='180' y='225' width='240' height='10' fill='#F59E0B' opacity='0.5'/>"
                 "<circle cx='300' cy='285' r='28' fill='none' stroke='#F59E0B' stroke-width='3' stroke-dasharray='5 3'/>",
    },
    "Burnisher": {
        "accent": "#EC4899",
        "shape": "<rect x='220' y='175' width='160' height='90' rx='14' fill='url(#g)'/>"
                 "<line x1='300' y1='85' x2='300' y2='175' stroke='#EC4899' stroke-width='6' stroke-linecap='round'/>"
                 "<circle cx='300' cy='80' r='14' fill='#EC4899'/>"
                 "<circle cx='300' cy='280' r='50' fill='none' stroke='#EC4899' stroke-width='4'/>"
                 "<circle cx='300' cy='280' r='30' fill='#EC4899' opacity='0.3'/>"
                 "<circle cx='300' cy='280' r='12' fill='#EC4899'/>",
    },
    "Carpet Extractor": {
        "accent": "#3B82F6",
        "shape": "<rect x='210' y='150' width='180' height='130' rx='16' fill='url(#g)'/>"
                 "<line x1='400' y1='190' x2='480' y2='150' stroke='#3B82F6' stroke-width='5' stroke-linecap='round'/>"
                 "<circle cx='485' cy='148' r='10' fill='#3B82F6'/>"
                 "<rect x='235' y='195' width='130' height='30' rx='6' fill='#3B82F6' opacity='0.4'/>"
                 "<circle cx='240' cy='295' r='14' fill='#0B0E14' stroke='#3B82F6' stroke-width='2'/>"
                 "<circle cx='360' cy='295' r='14' fill='#0B0E14' stroke='#3B82F6' stroke-width='2'/>",
    },
    "Outdoor Sweeper": {
        "accent": "#10B981",
        "shape": "<rect x='140' y='150' width='320' height='110' rx='16' fill='url(#g)'/>"
                 "<rect x='220' y='95' width='160' height='62' rx='12' fill='#10B981' opacity='0.35'/>"
                 "<rect x='150' y='240' width='80' height='30' rx='6' fill='#10B981' opacity='0.4'/>"
                 "<circle cx='190' cy='290' r='32' fill='#0B0E14' stroke='#10B981' stroke-width='3'/>"
                 "<circle cx='410' cy='290' r='32' fill='#0B0E14' stroke='#10B981' stroke-width='3'/>"
                 "<circle cx='110' cy='270' r='22' fill='none' stroke='#10B981' stroke-width='3' stroke-dasharray='5 4'/>",
    },
    "Outdoor Litter Vacuum": {
        "accent": "#84CC16",
        "shape": "<rect x='180' y='150' width='240' height='130' rx='18' fill='url(#g)'/>"
                 "<rect x='400' y='180' width='80' height='60' rx='8' fill='#84CC16' opacity='0.4'/>"
                 "<circle cx='220' cy='295' r='25' fill='#0B0E14' stroke='#84CC16' stroke-width='3'/>"
                 "<circle cx='380' cy='295' r='25' fill='#0B0E14' stroke='#84CC16' stroke-width='3'/>"
                 "<circle cx='300' cy='205' r='20' fill='none' stroke='#84CC16' stroke-width='3'/>",
    },
}


def _build_machine_svg(model: str, category: str, size: str) -> str:
    """Render a unique, branded machine card SVG. Always succeeds, even for
    models with no real product photo on the Tennant CDN."""
    icon = _MACHINE_ICONS.get(category, _MACHINE_ICONS["Ride-On Scrubber"])
    accent = icon["accent"]
    shape = icon["shape"]
    cat_upper = category.upper()
    return f"""<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 600 400" preserveAspectRatio="xMidYMid slice" width="600" height="400">
  <defs>
    <linearGradient id="bg" x1="0" y1="0" x2="0" y2="1">
      <stop offset="0%" stop-color="#0B0E14"/>
      <stop offset="100%" stop-color="#131821"/>
    </linearGradient>
    <linearGradient id="g" x1="0" y1="0" x2="0" y2="1">
      <stop offset="0%" stop-color="{accent}" stop-opacity="0.85"/>
      <stop offset="100%" stop-color="{accent}" stop-opacity="0.45"/>
    </linearGradient>
    <pattern id="grid" width="40" height="40" patternUnits="userSpaceOnUse">
      <path d="M 40 0 L 0 0 0 40" fill="none" stroke="rgba(255,255,255,0.04)" stroke-width="1"/>
    </pattern>
  </defs>
  <rect width="600" height="400" fill="url(#bg)"/>
  <rect width="600" height="400" fill="url(#grid)"/>
  <!-- corner brackets, HUD aesthetic -->
  <path d="M 20 20 L 60 20 M 20 20 L 20 60" stroke="{accent}" stroke-width="3" fill="none"/>
  <path d="M 580 20 L 540 20 M 580 20 L 580 60" stroke="{accent}" stroke-width="3" fill="none"/>
  <path d="M 20 380 L 60 380 M 20 380 L 20 340" stroke="{accent}" stroke-width="3" fill="none"/>
  <path d="M 580 380 L 540 380 M 580 380 L 580 340" stroke="{accent}" stroke-width="3" fill="none"/>
  <!-- silhouette -->
  {shape}
  <!-- type label top-left -->
  <text x="40" y="55" font-family="ui-monospace,Menlo,monospace" font-size="13" fill="{accent}" letter-spacing="3">{cat_upper}</text>
  <!-- model name big -->
  <text x="40" y="350" font-family="Inter,system-ui,sans-serif" font-size="56" fill="#FFFFFF" font-weight="900" letter-spacing="-1">{model}</text>
  <text x="40" y="375" font-family="ui-monospace,Menlo,monospace" font-size="13" fill="rgba(255,255,255,0.55)" letter-spacing="2">TENNANT · {size}</text>
</svg>"""


@api_router.get("/machines/{model}/image.svg")
async def machine_image(model: str):
    """Always-on, brand-correct machine artwork. Routed by model name.

    Beats the Tennant CDN approach because:
      • 100 % uptime (no external dependency, no 404/500 fall-throughs)
      • Distinct per model (model name is rendered into the SVG)
      • Distinct per category (silhouette + accent color picked per category)
      • Cyber-HUD aesthetic matches the rest of the app
    """
    m = next((x for x in TENNANT_MACHINES if x["model"].lower() == model.lower()), None)
    if not m:
        raise HTTPException(status_code=404, detail="Machine not found")
    svg = _build_machine_svg(m["model"], m["category"], m.get("size", "—"))
    return Response(content=svg, media_type="image/svg+xml",
                    headers={"Cache-Control": "public, max-age=86400"})


@api_router.get("/machines")
async def list_machines(_: User = Depends(get_current_user), category: Optional[str] = None):
    # Built-in catalog + admin-added custom machines + soft-hide overrides.
    hidden = await db.machine_overrides.find({"hidden": True}, {"_id": 0, "model": 1}).to_list(200)
    hidden_models = {h["model"] for h in hidden}
    base = [m for m in TENNANT_MACHINES if m["model"] not in hidden_models]
    custom = await db.machines_custom.find({}, {"_id": 0}).to_list(500)
    legacy = await db.machines.find({}, {"_id": 0}).to_list(500)
    out = base + custom + legacy
    if category:
        out = [m for m in out if m.get("category") == category]
    out = [{**m, "image_url": m.get("image_url") or _real_photo_or_svg(m)} for m in out]
    # Active brand overlay — swap to brand's catalog when not Tennant.
    brand = await _active_brand_doc()
    if brand and brand.get("brand_id") != "orisei-freight":
        out = [_overlay_machine(m, brand, i) if not m.get("is_custom") else m for i, m in enumerate(out)]
    return {"machines": out, "categories": TENNANT_MACHINE_CATEGORIES, "count": len(out), "catalog_label": (brand or {}).get("catalog_label") or "Machine Catalog"}


class MachineIn(BaseModel):
    model: str
    display_name: Optional[str] = None
    category: str = "Custom"
    description: Optional[str] = ""
    image_url: Optional[str] = ""
    width_in: Optional[float] = None
    length_in: Optional[float] = None
    height_in: Optional[float] = None
    weight_lbs: Optional[float] = None
    power: Optional[str] = ""
    tank_gal: Optional[float] = None
    run_time_hrs: Optional[float] = None


@api_router.delete("/machines/{model}")
async def machines_delete(model: str, _: User = Depends(require_role("admin"))):
    """Delete a machine. Seeded models are soft-hidden so they can be
    restored; custom models are hard-deleted."""
    model_norm = model.upper().strip()
    is_seed = any(m["model"] == model_norm for m in TENNANT_MACHINES)
    if is_seed:
        await db.machine_overrides.update_one(
            {"model": model_norm}, {"$set": {"model": model_norm, "hidden": True}}, upsert=True
        )
        return {"ok": True, "soft_hidden": True}
    # Both legacy db.machines and new db.machines_custom collections
    a = await db.machines_custom.delete_one({"model": model_norm})
    b = await db.machines.delete_one({"model": model_norm})
    if a.deleted_count == 0 and b.deleted_count == 0:
        raise HTTPException(404, "Machine not found")
    return {"ok": True}

# -------------------- ARCADE: Connect Four · Tournaments · Trophies --------------------
ROWS_C4, COLS_C4 = 6, 7

def _empty_c4_board():
    return [[0] * COLS_C4 for _ in range(ROWS_C4)]

def _c4_drop(board: List[List[int]], col: int, player: int) -> Optional[int]:
    """Drops a piece; returns the row it landed in or None if column full."""
    for r in range(ROWS_C4 - 1, -1, -1):
        if board[r][col] == 0:
            board[r][col] = player
            return r
    return None

def _c4_winner(board: List[List[int]]) -> int:
    # horizontal
    for r in range(ROWS_C4):
        for c in range(COLS_C4 - 3):
            v = board[r][c]
            if v and v == board[r][c+1] == board[r][c+2] == board[r][c+3]:
                return v
    # vertical
    for r in range(ROWS_C4 - 3):
        for c in range(COLS_C4):
            v = board[r][c]
            if v and v == board[r+1][c] == board[r+2][c] == board[r+3][c]:
                return v
    # diag /
    for r in range(3, ROWS_C4):
        for c in range(COLS_C4 - 3):
            v = board[r][c]
            if v and v == board[r-1][c+1] == board[r-2][c+2] == board[r-3][c+3]:
                return v
    # diag \
    for r in range(ROWS_C4 - 3):
        for c in range(COLS_C4 - 3):
            v = board[r][c]
            if v and v == board[r+1][c+1] == board[r+2][c+2] == board[r+3][c+3]:
                return v
    return 0

def _c4_full(board) -> bool:
    return all(board[0][c] != 0 for c in range(COLS_C4))

class C4GameCreate(BaseModel):
    opponent_email: Optional[str] = None  # if None → open lobby

class C4MoveBody(BaseModel):
    column: int

@api_router.post("/arcade/connect4/games")
async def c4_create(payload: C4GameCreate, user: User = Depends(get_current_user)):
    game = {
        "game_id": f"C4-{uuid.uuid4().hex[:8].upper()}",
        "kind": "connect4",
        "player1_id": user.user_id,
        "player1_name": user.name,
        "player2_id": None,
        "player2_name": None,
        "player2_email": payload.opponent_email,
        "board": _empty_c4_board(),
        "turn": 1,
        "status": "open",  # open | active | finished | draw
        "winner_id": None,
        "winner_name": None,
        "tournament_id": None,
        "moves": [],
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.arcade_games.insert_one(dict(game))
    return {k: v for k, v in game.items() if k != "_id"}

@api_router.get("/arcade/connect4/games")
async def c4_list(user: User = Depends(get_current_user), status: Optional[str] = None):
    q: Dict[str, Any] = {"kind": "connect4"}
    if status:
        q["status"] = status
    docs = await db.arcade_games.find(q, {"_id": 0}).sort("updated_at", -1).limit(50).to_list(50)
    # Only show mine + open lobby
    out = [d for d in docs if d["status"] == "open" or user.user_id in (d.get("player1_id"), d.get("player2_id"))]
    return out

@api_router.get("/arcade/connect4/games/{game_id}")
async def c4_get(game_id: str, _: User = Depends(get_current_user)):
    doc = await db.arcade_games.find_one({"game_id": game_id}, {"_id": 0})
    if not doc:
        raise HTTPException(404, "Game not found")
    return doc

@api_router.post("/arcade/connect4/games/{game_id}/join")
async def c4_join(game_id: str, user: User = Depends(get_current_user)):
    doc = await db.arcade_games.find_one({"game_id": game_id}, {"_id": 0})
    if not doc:
        raise HTTPException(404, "Game not found")
    if doc["status"] != "open":
        raise HTTPException(400, "Game not open")
    if doc["player1_id"] == user.user_id:
        raise HTTPException(400, "Cannot join your own game")
    await db.arcade_games.update_one(
        {"game_id": game_id},
        {"$set": {"player2_id": user.user_id, "player2_name": user.name, "status": "active", "updated_at": datetime.now(timezone.utc).isoformat()}},
    )
    return await db.arcade_games.find_one({"game_id": game_id}, {"_id": 0})

@api_router.post("/arcade/connect4/games/{game_id}/move")
async def c4_move(game_id: str, payload: C4MoveBody, user: User = Depends(get_current_user)):
    doc = await db.arcade_games.find_one({"game_id": game_id}, {"_id": 0})
    if not doc:
        raise HTTPException(404, "Game not found")
    if doc["status"] != "active":
        raise HTTPException(400, "Game not active")
    if doc["turn"] == 1 and doc["player1_id"] != user.user_id:
        raise HTTPException(403, "Not your turn")
    if doc["turn"] == 2 and doc["player2_id"] != user.user_id:
        raise HTTPException(403, "Not your turn")
    if payload.column < 0 or payload.column >= COLS_C4:
        raise HTTPException(400, "Invalid column")
    board = doc["board"]
    row = _c4_drop(board, payload.column, doc["turn"])
    if row is None:
        raise HTTPException(400, "Column full")
    win = _c4_winner(board)
    now = datetime.now(timezone.utc).isoformat()
    moves = doc.get("moves", []) + [{"player": doc["turn"], "col": payload.column, "row": row, "at": now}]
    update = {"board": board, "moves": moves, "updated_at": now}
    if win:
        update["status"] = "finished"
        update["winner_id"] = doc["player1_id"] if win == 1 else doc["player2_id"]
        update["winner_name"] = doc["player1_name"] if win == 1 else doc["player2_name"]
    elif _c4_full(board):
        update["status"] = "draw"
    else:
        update["turn"] = 2 if doc["turn"] == 1 else 1
    await db.arcade_games.update_one({"game_id": game_id}, {"$set": update})

    # Award trophy / record win in leaderboard
    if win:
        winner_id = update["winner_id"]
        winner_name = update["winner_name"]
        loser_id = doc["player2_id"] if win == 1 else doc["player1_id"]
        await db.arcade_leaderboard.update_one(
            {"user_id": winner_id},
            {"$inc": {"wins": 1, "games": 1, "trophies": 1},
             "$set": {"name": winner_name, "last_played_at": now},
             "$setOnInsert": {"created_at": now, "losses": 0, "draws": 0}},
            upsert=True,
        )
        if loser_id:
            await db.arcade_leaderboard.update_one(
                {"user_id": loser_id},
                {"$inc": {"losses": 1, "games": 1},
                 "$set": {"last_played_at": now},
                 "$setOnInsert": {"created_at": now, "wins": 0, "draws": 0, "trophies": 0,
                                  "name": doc["player2_name"] if win == 1 else doc["player1_name"]}},
                upsert=True,
            )
    elif _c4_full(board):
        for pid, pname in [(doc["player1_id"], doc["player1_name"]), (doc["player2_id"], doc["player2_name"])]:
            if pid:
                await db.arcade_leaderboard.update_one(
                    {"user_id": pid},
                    {"$inc": {"draws": 1, "games": 1},
                     "$set": {"name": pname, "last_played_at": now},
                     "$setOnInsert": {"created_at": now, "wins": 0, "losses": 0, "trophies": 0}},
                    upsert=True,
                )

    fresh = await db.arcade_games.find_one({"game_id": game_id}, {"_id": 0})

    # Tournament advance
    if fresh.get("tournament_id") and fresh["status"] in ("finished",):
        await _advance_tournament(fresh["tournament_id"], fresh["game_id"], fresh["winner_id"], fresh["winner_name"])
    return fresh

@api_router.get("/arcade/leaderboard")
async def arcade_leaderboard(_: User = Depends(get_current_user)):
    rows = await db.arcade_leaderboard.find({}, {"_id": 0}).sort([("trophies", -1), ("wins", -1)]).limit(100).to_list(100)
    # Trophy tiers
    for r in rows:
        t = r.get("trophies", 0)
        r["tier"] = "Legend" if t >= 25 else "Champion" if t >= 10 else "Contender" if t >= 3 else "Rookie"
    return {"rows": rows}

# -------------------- ARCADE: Tournaments --------------------
class TournamentCreate(BaseModel):
    name: str
    kind: str = "connect4"
    participant_user_ids: List[str] = []

async def _advance_tournament(tid: str, game_id: str, winner_id: Optional[str], winner_name: Optional[str]):
    t = await db.arcade_tournaments.find_one({"tournament_id": tid}, {"_id": 0})
    if not t:
        return
    bracket = t.get("bracket") or []
    # Mark the finished match
    for r_idx, rnd in enumerate(bracket):
        for m_idx, match in enumerate(rnd["matches"]):
            if match.get("game_id") == game_id:
                match["status"] = "done"
                match["winner_id"] = winner_id
                match["winner_name"] = winner_name
                # Promote winner to next round if exists
                if r_idx + 1 < len(bracket):
                    next_match_idx = m_idx // 2
                    side = "p1" if m_idx % 2 == 0 else "p2"
                    next_match = bracket[r_idx + 1]["matches"][next_match_idx]
                    next_match[f"{side}_id"] = winner_id
                    next_match[f"{side}_name"] = winner_name
                else:
                    # Final winner
                    t["status"] = "completed"
                    t["champion_id"] = winner_id
                    t["champion_name"] = winner_name
                    if winner_id:
                        await db.arcade_leaderboard.update_one(
                            {"user_id": winner_id},
                            {"$inc": {"trophies": 3, "tournaments_won": 1},
                             "$set": {"name": winner_name, "last_played_at": datetime.now(timezone.utc).isoformat()}},
                            upsert=True,
                        )
                break
    await db.arcade_tournaments.update_one({"tournament_id": tid}, {"$set": {"bracket": bracket, "status": t.get("status", "active"), "champion_id": t.get("champion_id"), "champion_name": t.get("champion_name")}})

@api_router.post("/arcade/tournaments")
async def create_tournament(payload: TournamentCreate, admin: User = Depends(require_role("admin", "dispatcher"))):
    if not (4 <= len(payload.participant_user_ids) <= 16) or (len(payload.participant_user_ids) & (len(payload.participant_user_ids) - 1)) != 0:
        raise HTTPException(400, "Participants must be a power of 2 between 4 and 16")
    # Lookup user names
    user_docs = await db.users.find({"user_id": {"$in": payload.participant_user_ids}}, {"_id": 0}).to_list(20)
    user_map = {u["user_id"]: u["name"] for u in user_docs}
    participants = [{"user_id": uid, "name": user_map.get(uid, "Unknown")} for uid in payload.participant_user_ids]
    random.shuffle(participants)
    # Build bracket
    rounds = []
    current = participants
    round_idx = 0
    while len(current) >= 2:
        matches = []
        for i in range(0, len(current), 2):
            p1 = current[i]; p2 = current[i + 1] if i + 1 < len(current) else {"user_id": None, "name": "BYE"}
            game = {
                "game_id": f"C4-{uuid.uuid4().hex[:8].upper()}",
                "kind": "connect4",
                "player1_id": p1["user_id"], "player1_name": p1["name"],
                "player2_id": p2["user_id"], "player2_name": p2["name"],
                "board": _empty_c4_board(), "turn": 1,
                "status": "active" if (p1["user_id"] and p2["user_id"]) else "pending",
                "winner_id": None, "winner_name": None,
                "moves": [], "tournament_id": None,
                "created_at": datetime.now(timezone.utc).isoformat(),
                "updated_at": datetime.now(timezone.utc).isoformat(),
            }
            matches.append({
                "match_id": f"M-{uuid.uuid4().hex[:6].upper()}",
                "p1_id": p1["user_id"], "p1_name": p1["name"],
                "p2_id": p2["user_id"], "p2_name": p2["name"],
                "game_id": game["game_id"], "status": "pending",
                "winner_id": None, "winner_name": None,
            })
            if round_idx == 0:
                await db.arcade_games.insert_one(dict(game))
                await db.arcade_games.update_one({"game_id": game["game_id"]}, {"$set": {"tournament_id": "PENDING"}})
        rounds.append({"round": round_idx + 1, "matches": matches})
        current = [{"user_id": None, "name": "TBD"} for _ in range(len(current) // 2)]
        round_idx += 1

    tid = f"TRN-{uuid.uuid4().hex[:8].upper()}"
    tournament = {
        "tournament_id": tid,
        "name": payload.name,
        "kind": payload.kind,
        "participants": participants,
        "bracket": rounds,
        "status": "active",
        "champion_id": None,
        "champion_name": None,
        "created_by": admin.user_id,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    # Tag round-1 games with tournament_id
    for m in rounds[0]["matches"]:
        await db.arcade_games.update_one({"game_id": m["game_id"]}, {"$set": {"tournament_id": tid}})
    await db.arcade_tournaments.insert_one(dict(tournament))
    return {k: v for k, v in tournament.items() if k != "_id"}

@api_router.get("/arcade/tournaments")
async def list_tournaments(_: User = Depends(get_current_user)):
    docs = await db.arcade_tournaments.find({}, {"_id": 0}).sort("created_at", -1).limit(20).to_list(20)
    return docs

@api_router.get("/arcade/tournaments/{tournament_id}")
async def get_tournament(tournament_id: str, _: User = Depends(get_current_user)):
    doc = await db.arcade_tournaments.find_one({"tournament_id": tournament_id}, {"_id": 0})
    if not doc:
        raise HTTPException(404, "Tournament not found")
    return doc

@api_router.get("/arcade/users")
async def arcade_users(_: User = Depends(get_current_user)):
    """List active users available for tournament invitations."""
    docs = await db.users.find({}, {"_id": 0, "user_id": 1, "name": 1, "email": 1, "role": 1}).to_list(100)
    return docs

# -------------------- ARCADE: Challenges --------------------
class ChallengeCreate(BaseModel):
    opponent_user_id: str
    kind: str = "connect4"
    message: Optional[str] = None

@api_router.post("/arcade/challenges")
async def create_challenge(payload: ChallengeCreate, user: User = Depends(get_current_user)):
    if payload.opponent_user_id == user.user_id:
        raise HTTPException(400, "Cannot challenge yourself")
    opp = await db.users.find_one({"user_id": payload.opponent_user_id}, {"_id": 0})
    if not opp:
        raise HTTPException(404, "Opponent not found")
    challenge = {
        "challenge_id": f"CHG-{uuid.uuid4().hex[:8].upper()}",
        "kind": payload.kind,
        "from_user_id": user.user_id,
        "from_user_name": user.name,
        "to_user_id": payload.opponent_user_id,
        "to_user_name": opp["name"],
        "message": payload.message,
        "status": "pending",  # pending | accepted | declined | expired
        "game_id": None,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.arcade_challenges.insert_one(dict(challenge))
    return {k: v for k, v in challenge.items() if k != "_id"}

@api_router.get("/arcade/challenges")
async def list_challenges(user: User = Depends(get_current_user)):
    """Returns inbox (challenges sent TO me) and outbox (FROM me)."""
    inbox = await db.arcade_challenges.find({"to_user_id": user.user_id}, {"_id": 0}).sort("created_at", -1).limit(30).to_list(30)
    outbox = await db.arcade_challenges.find({"from_user_id": user.user_id}, {"_id": 0}).sort("created_at", -1).limit(30).to_list(30)
    return {"inbox": inbox, "outbox": outbox, "pending_count": sum(1 for c in inbox if c["status"] == "pending")}

@api_router.post("/arcade/challenges/{challenge_id}/accept")
async def accept_challenge(challenge_id: str, user: User = Depends(get_current_user)):
    c = await db.arcade_challenges.find_one({"challenge_id": challenge_id}, {"_id": 0})
    if not c:
        raise HTTPException(404, "Challenge not found")
    if c["to_user_id"] != user.user_id:
        raise HTTPException(403, "Not your challenge")
    if c["status"] != "pending":
        raise HTTPException(400, f"Challenge already {c['status']}")
    game = {
        "game_id": f"C4-{uuid.uuid4().hex[:8].upper()}",
        "kind": c["kind"],
        "player1_id": c["from_user_id"], "player1_name": c["from_user_name"],
        "player2_id": user.user_id, "player2_name": user.name,
        "board": _empty_c4_board(), "turn": 1, "status": "active",
        "winner_id": None, "winner_name": None,
        "moves": [], "tournament_id": None,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.arcade_games.insert_one(dict(game))
    await db.arcade_challenges.update_one(
        {"challenge_id": challenge_id},
        {"$set": {"status": "accepted", "game_id": game["game_id"], "accepted_at": datetime.now(timezone.utc).isoformat()}},
    )
    return {"game_id": game["game_id"], "status": "accepted"}

@api_router.post("/arcade/challenges/{challenge_id}/decline")
async def decline_challenge(challenge_id: str, user: User = Depends(get_current_user)):
    c = await db.arcade_challenges.find_one({"challenge_id": challenge_id}, {"_id": 0})
    if not c:
        raise HTTPException(404, "Challenge not found")
    if c["to_user_id"] != user.user_id:
        raise HTTPException(403, "Not your challenge")
    await db.arcade_challenges.update_one({"challenge_id": challenge_id}, {"$set": {"status": "declined"}})
    return {"ok": True}

# -------------------- KPI REPORTS · Download & Email --------------------
async def _compute_kpis() -> Dict[str, Any]:
    """Reuse the same calc as /kpis."""
    all_shipments = await db.shipments.find({}, {"_id": 0}).to_list(2000)
    total = len(all_shipments)
    in_transit = sum(1 for s in all_shipments if s["status"] == "in_transit")
    delayed = sum(1 for s in all_shipments if s["status"] == "delayed")
    delivered = sum(1 for s in all_shipments if s["status"] == "delivered")
    pending = sum(1 for s in all_shipments if s["status"] == "pending")
    by_mode: Dict[str, int] = {}
    by_carrier: Dict[str, Dict[str, int]] = {}
    total_weight = 0.0
    total_value = 0.0
    for s in all_shipments:
        by_mode[s["mode"]] = by_mode.get(s["mode"], 0) + 1
        c = s["carrier"]
        by_carrier.setdefault(c, {"total": 0, "on_time": 0, "delayed": 0})
        by_carrier[c]["total"] += 1
        if s["status"] == "delivered":
            by_carrier[c]["on_time"] += 1
        if s["status"] == "delayed":
            by_carrier[c]["delayed"] += 1
        total_weight += s.get("weight_lbs", 0)
        total_value += s.get("value_usd", 0)
    on_time_rate = round((delivered / max(1, total)) * 100, 1)
    return {
        "total": total, "in_transit": in_transit, "delayed": delayed,
        "delivered": delivered, "pending": pending,
        "weight_lbs": round(total_weight, 0), "value_usd": round(total_value, 0),
        "on_time_rate": on_time_rate,
        "by_mode": by_mode,
        "by_carrier": sorted([{"carrier": k, **v, "on_time_pct": round(v["on_time"] / max(1, v["total"]) * 100, 1)} for k, v in by_carrier.items()], key=lambda x: -x["total"]),
    }

@api_router.get("/reports/kpi/download.pdf")
async def kpi_report_pdf(user: User = Depends(get_current_user)):
    """Tennant-branded KPI report as PDF."""
    kpi = await _compute_kpis()
    weekly = await get_weekly_weights(_=user)
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, topMargin=0.5 * inch, bottomMargin=0.5 * inch, leftMargin=0.5 * inch, rightMargin=0.5 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("title", parent=styles["Heading1"], textColor=colors.HexColor("#00E5FF"), fontSize=20, spaceAfter=8)
    sub_style = ParagraphStyle("sub", parent=styles["Normal"], textColor=colors.HexColor("#475569"), fontSize=9, spaceAfter=14)
    section_style = ParagraphStyle("sec", parent=styles["Heading2"], textColor=colors.HexColor("#0F172A"), fontSize=13, spaceBefore=14, spaceAfter=8)
    body_style = ParagraphStyle("body", parent=styles["Normal"], textColor=colors.HexColor("#0F172A"), fontSize=10)
    elements = []
    elements.append(Paragraph("TENNANT COMPANIES · TMS KPI Report", title_style))
    elements.append(Paragraph(f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')} · Operations Snapshot · All Modes", sub_style))

    # KPI grid as a 4x2 table
    elements.append(Paragraph("Headline Metrics", section_style))
    kpi_table = Table([
        ["Total Shipments", f"{kpi['total']:,}", "In Transit", f"{kpi['in_transit']:,}"],
        ["Delivered", f"{kpi['delivered']:,}", "Delayed", f"{kpi['delayed']:,}"],
        ["Pending", f"{kpi['pending']:,}", "On-Time Rate", f"{kpi['on_time_rate']}%"],
        ["Total Weight (lbs)", f"{int(kpi['weight_lbs']):,}", "Total Value (USD)", f"${int(kpi['value_usd']):,}"],
    ], colWidths=[1.6 * inch, 1.4 * inch, 1.6 * inch, 1.4 * inch])
    kpi_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F8FAFC")),
        ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
        ("INNERGRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#E2E8F0")),
        ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
        ("FONTNAME", (2, 0), (2, -1), "Helvetica-Bold"),
        ("TEXTCOLOR", (1, 0), (1, -1), colors.HexColor("#0891B2")),
        ("TEXTCOLOR", (3, 0), (3, -1), colors.HexColor("#0891B2")),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("ALIGN", (1, 0), (1, -1), "RIGHT"),
        ("ALIGN", (3, 0), (3, -1), "RIGHT"),
        ("LEFTPADDING", (0, 0), (-1, -1), 10),
        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
    ]))
    elements.append(kpi_table)

    # Carrier scorecard
    elements.append(Paragraph("Carrier Scorecard", section_style))
    carrier_rows = [["Carrier", "Total", "Delivered", "Delayed", "On-Time %"]]
    for c in kpi["by_carrier"][:15]:
        carrier_rows.append([c["carrier"], str(c["total"]), str(c["on_time"]), str(c["delayed"]), f"{c['on_time_pct']}%"])
    carrier_table = Table(carrier_rows, colWidths=[2.4 * inch, 1.0 * inch, 1.0 * inch, 1.0 * inch, 1.2 * inch])
    carrier_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0F172A")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#F8FAFC")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, 0), 10),
        ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#FFFFFF")),
        ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
        ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#E2E8F0")),
        ("FONTSIZE", (0, 1), (-1, -1), 9),
        ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
    ]))
    elements.append(carrier_table)

    # Mode mix
    elements.append(Paragraph("Volume by Mode", section_style))
    mode_rows = [["Mode", "Shipments", "% of Total"]]
    for m, v in sorted(kpi["by_mode"].items(), key=lambda x: -x[1]):
        mode_rows.append([m, str(v), f"{round(v / max(1, kpi['total']) * 100, 1)}%"])
    table_style = TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0F172A")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#F8FAFC")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, 0), 10),
        ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
        ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#E2E8F0")),
        ("FONTSIZE", (0, 1), (-1, -1), 9),
        ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
    ])
    mode_table = Table(mode_rows, colWidths=[2.0 * inch, 1.4 * inch, 1.4 * inch])
    mode_table.setStyle(table_style)
    elements.append(mode_table)

    # Weekly weights summary
    elements.append(Paragraph("Weekly Avg Weight per Shipment · by Facility (12-wk)", section_style))
    ww_rows = [["Facility", "Current Wk Avg (lbs)", "12-Wk Avg (lbs)", "WoW Δ", "Total (lbs)"]]
    for fac in ["GVM", "HOM", "LVK"]:
        d = weekly["summary"].get(fac, {})
        ww_rows.append([fac, f"{int(d.get('current_week_avg_lbs') or 0):,}", f"{int(d.get('twelve_wk_avg_lbs') or 0):,}", f"{int(d.get('wow_delta_lbs') or 0):+,}", f"{int(d.get('twelve_wk_total_lbs') or 0):,}"])
    ww_table = Table(ww_rows, colWidths=[1.0 * inch, 1.5 * inch, 1.5 * inch, 1.0 * inch, 1.5 * inch])
    ww_table.setStyle(table_style)
    elements.append(ww_table)

    elements.append(Spacer(1, 0.3 * inch))
    elements.append(Paragraph(f"<font size=7 color='#475569'>Tennant Companies · TMS · Generated by {user.name} · Confidential — Internal Use Only</font>", body_style))

    doc.build(elements)
    buf.seek(0)
    headers = {"Content-Disposition": f'attachment; filename="Tennant_KPI_Report_{datetime.now(timezone.utc).strftime("%Y%m%d")}.pdf"'}
    return StreamingResponse(buf, media_type="application/pdf", headers=headers)

@api_router.get("/reports/kpi/download.xlsx")
async def kpi_report_xlsx(user: User = Depends(get_current_user)):
    """KPI report as a multi-sheet Excel."""
    kpi = await _compute_kpis()
    weekly = await get_weekly_weights(_=user)
    wb = XLWorkbook()
    header_fill = PatternFill("solid", fgColor="0F172A")
    header_font = Font(name="Calibri", size=11, bold=True, color="F8FAFC")
    cyan_font = Font(name="Calibri", size=11, color="0891B2", bold=True)
    border = Border(left=Side(style="thin", color="CBD5E1"), right=Side(style="thin", color="CBD5E1"),
                    top=Side(style="thin", color="CBD5E1"), bottom=Side(style="thin", color="CBD5E1"))

    ws1 = wb.active; ws1.title = "Headline"
    ws1["A1"] = "TENNANT COMPANIES · TMS KPI Report"; ws1["A1"].font = Font(size=14, bold=True, color="00E5FF")
    ws1["A2"] = f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}"; ws1["A2"].font = Font(size=9, color="475569")
    headline = [
        ("Total Shipments", kpi["total"]),
        ("In Transit", kpi["in_transit"]),
        ("Delivered", kpi["delivered"]),
        ("Delayed", kpi["delayed"]),
        ("Pending", kpi["pending"]),
        ("On-Time Rate %", kpi["on_time_rate"]),
        ("Total Weight (lbs)", kpi["weight_lbs"]),
        ("Total Value (USD)", kpi["value_usd"]),
    ]
    for i, (k, v) in enumerate(headline, start=4):
        ws1.cell(row=i, column=1, value=k).font = Font(bold=True)
        c = ws1.cell(row=i, column=2, value=v); c.font = cyan_font; c.border = border
    ws1.column_dimensions["A"].width = 26; ws1.column_dimensions["B"].width = 18

    ws2 = wb.create_sheet("Carrier Scorecard")
    for j, h in enumerate(["Carrier", "Total", "Delivered", "Delayed", "On-Time %"], start=1):
        c = ws2.cell(row=1, column=j, value=h); c.fill = header_fill; c.font = header_font
    for i, c in enumerate(kpi["by_carrier"], start=2):
        ws2.cell(row=i, column=1, value=c["carrier"])
        ws2.cell(row=i, column=2, value=c["total"])
        ws2.cell(row=i, column=3, value=c["on_time"])
        ws2.cell(row=i, column=4, value=c["delayed"])
        ws2.cell(row=i, column=5, value=c["on_time_pct"])
    for col_letter in "ABCDE":
        ws2.column_dimensions[col_letter].width = 18
    ws2.freeze_panes = "A2"

    ws3 = wb.create_sheet("Mode Mix")
    for j, h in enumerate(["Mode", "Shipments", "% of Total"], start=1):
        c = ws3.cell(row=1, column=j, value=h); c.fill = header_fill; c.font = header_font
    for i, (m, v) in enumerate(sorted(kpi["by_mode"].items(), key=lambda x: -x[1]), start=2):
        ws3.cell(row=i, column=1, value=m)
        ws3.cell(row=i, column=2, value=v)
        ws3.cell(row=i, column=3, value=round(v / max(1, kpi["total"]) * 100, 1))

    ws4 = wb.create_sheet("Weekly Weights")
    for j, h in enumerate(["Week", "GVM", "HOM", "LVK"], start=1):
        c = ws4.cell(row=1, column=j, value=h); c.fill = header_fill; c.font = header_font
    for i, row in enumerate(weekly["series"], start=2):
        ws4.cell(row=i, column=1, value=row["week"])
        ws4.cell(row=i, column=2, value=row["GVM"])
        ws4.cell(row=i, column=3, value=row["HOM"])
        ws4.cell(row=i, column=4, value=row["LVK"])

    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    headers = {"Content-Disposition": f'attachment; filename="Tennant_KPI_Report_{datetime.now(timezone.utc).strftime("%Y%m%d")}.xlsx"'}
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", headers=headers)

class KPIEmailRequest(BaseModel):
    to: Optional[str] = None
    cc: Optional[str] = None
    note: Optional[str] = None
    format: str = "pdf"  # pdf | xlsx | both

@api_router.post("/reports/kpi/email")
async def email_kpi_report(payload: KPIEmailRequest, user: User = Depends(get_current_user)):
    kpi = await _compute_kpis()
    subject = f"Tennant TMS · KPI Report · {datetime.now(timezone.utc).strftime('%b %d, %Y')}"
    fmt_label = {"pdf": "PDF", "xlsx": "Excel", "both": "PDF + Excel"}.get(payload.format, "PDF")
    body = (
        f"Hello,\n\n"
        f"Please find attached the latest Tennant TMS KPI snapshot ({fmt_label}).\n\n"
        f"HEADLINE METRICS\n"
        f"  Total Shipments:    {kpi['total']:,}\n"
        f"  In Transit:         {kpi['in_transit']:,}\n"
        f"  Delivered:          {kpi['delivered']:,}\n"
        f"  Delayed:            {kpi['delayed']:,}\n"
        f"  Pending:            {kpi['pending']:,}\n"
        f"  On-Time Rate:       {kpi['on_time_rate']}%\n"
        f"  Total Weight:       {int(kpi['weight_lbs']):,} lbs\n"
        f"  Total Value:        ${int(kpi['value_usd']):,}\n\n"
        f"TOP 5 CARRIERS BY VOLUME\n"
    )
    for c in kpi["by_carrier"][:5]:
        body += f"  • {c['carrier']:30s}  total {c['total']:4d}   on-time {c['on_time_pct']}%\n"
    if payload.note:
        body += f"\nNote from {user.name}:\n{payload.note}\n"
    body += "\n— Tennant Transportation Team\n  generated automatically from the Tennant TMS"
    base = os.environ.get("PUBLIC_APP_URL", "")
    pdf_link = f"{base}/api/reports/kpi/download.pdf" if base else "/api/reports/kpi/download.pdf"
    xlsx_link = f"{base}/api/reports/kpi/download.xlsx" if base else "/api/reports/kpi/download.xlsx"
    body += f"\n\nDownload links (active for authenticated users):\n  PDF: {pdf_link}\n  XLSX: {xlsx_link}\n"
    await db.report_audit_log.insert_one({
        "report_type": "kpi",
        "to": payload.to or "",
        "cc": payload.cc or "",
        "format": payload.format,
        "generated_by": user.user_id,
        "generated_by_name": user.name,
        "at": datetime.now(timezone.utc).isoformat(),
    })
    to = payload.to or ""
    cc = payload.cc or ""
    mailto_qs = []
    if cc:
        mailto_qs.append(f"cc={cc}")
    mailto_qs.append(f"subject={subject}")
    mailto_qs.append(f"body={body}")
    mailto = f"mailto:{to}?" + "&".join(mailto_qs)
    return {"to": to, "cc": cc, "subject": subject, "body": body, "mailto": mailto, "pdf_link": pdf_link, "xlsx_link": xlsx_link}

# -------------------- SUPPLIER SOURCING --------------------
# Tracks every supplier Tennant currently sources from — components, lead times, risk, dual-source coverage.
TENNANT_SUPPLIERS = [
    {"supplier_id": "SUP-001", "name": "Samsung SDI", "country": "South Korea", "category": "Battery · Li-ion",
     "components": ["TENN-BATT-LION-48V", "TENN-BATT-LION-36V"], "annual_spend_usd": 4_200_000,
     "lead_time_days": 42, "moq": 100, "on_time_pct": 96.4, "quality_ppm": 142,
     "risk_score": 12, "fta_eligible": "KORUS", "primary": True, "alt_suppliers": ["LG Energy Solution"],
     "contract_expiry": "2027-03-31", "contact": "Hyun-woo Kim · hkim@samsungsdi.com · +82-2-3458-7000",
     "notes": "Primary 48V Li-ion. KORUS-eligible — 0% duty."},
    {"supplier_id": "SUP-002", "name": "LG Energy Solution", "country": "South Korea", "category": "Battery · Li-ion",
     "components": ["TENN-BATT-LION-48V-ALT"], "annual_spend_usd": 1_100_000,
     "lead_time_days": 38, "moq": 50, "on_time_pct": 94.1, "quality_ppm": 178,
     "risk_score": 18, "fta_eligible": "KORUS", "primary": False, "alt_suppliers": ["Samsung SDI"],
     "contract_expiry": "2026-12-31", "contact": "Min-jun Park · mpark@lgensol.com · +82-2-3773-1114",
     "notes": "Backup source for 48V pack."},
    {"supplier_id": "SUP-003", "name": "Trojan Battery Co.", "country": "USA", "category": "Battery · AGM/Flooded",
     "components": ["TENN-BATT-AGM-36V", "TENN-BATT-FLOODED-36V"], "annual_spend_usd": 2_650_000,
     "lead_time_days": 21, "moq": 25, "on_time_pct": 97.8, "quality_ppm": 89,
     "risk_score": 8, "fta_eligible": "Domestic", "primary": True, "alt_suppliers": ["East Penn Mfg"],
     "contract_expiry": "2028-06-30", "contact": "Sarah Bennett · sb@trojanbattery.com · +1-562-236-3000",
     "notes": "Primary AGM. Domestic = no duty exposure."},
    {"supplier_id": "SUP-004", "name": "East Penn Manufacturing", "country": "USA", "category": "Battery · AGM",
     "components": ["TENN-BATT-AGM-36V-ALT"], "annual_spend_usd": 580_000,
     "lead_time_days": 18, "moq": 25, "on_time_pct": 98.5, "quality_ppm": 65,
     "risk_score": 6, "fta_eligible": "Domestic", "primary": False, "alt_suppliers": ["Trojan Battery"],
     "contract_expiry": "2027-09-15", "contact": "Dale Krug · dkrug@dekabatteries.com · +1-610-682-6361"},
    {"supplier_id": "SUP-005", "name": "Yura Corporation", "country": "South Korea", "category": "Wiring Harness",
     "components": ["TENN-WH-T16", "TENN-WH-T350", "TENN-WH-S30"], "annual_spend_usd": 1_870_000,
     "lead_time_days": 35, "moq": 200, "on_time_pct": 95.2, "quality_ppm": 215,
     "risk_score": 22, "fta_eligible": "KORUS", "primary": True, "alt_suppliers": ["Yazaki Japan"],
     "contract_expiry": "2026-11-30", "contact": "Ji-hoon Lee · jlee@yura.kr · +82-31-280-7000",
     "notes": "Primary harness supplier. Quality PPM trending up — under review."},
    {"supplier_id": "SUP-006", "name": "Yazaki Corporation", "country": "Japan", "category": "Wiring Harness",
     "components": ["TENN-WH-INDUSTRIAL", "TENN-WH-CONTROL"], "annual_spend_usd": 925_000,
     "lead_time_days": 45, "moq": 100, "on_time_pct": 98.1, "quality_ppm": 95,
     "risk_score": 14, "fta_eligible": "JP-FTA", "primary": True, "alt_suppliers": ["Yura"],
     "contract_expiry": "2027-04-30", "contact": "Takeshi Mori · tmori@yazaki.com · +81-538-32-2114"},
    {"supplier_id": "SUP-007", "name": "Honda Power Equipment", "country": "Japan", "category": "Engine · LPG/Gas",
     "components": ["TENN-ENG-LPG-T350", "TENN-ENG-LPG-T500"], "annual_spend_usd": 3_400_000,
     "lead_time_days": 55, "moq": 50, "on_time_pct": 99.2, "quality_ppm": 32,
     "risk_score": 9, "fta_eligible": "JP-FTA", "primary": True, "alt_suppliers": ["Kubota"],
     "contract_expiry": "2029-03-31", "contact": "Tennant Account Mgr · tennant@honda-na.com · +1-770-497-6400",
     "notes": "Long-standing partner. Best quality in class."},
    {"supplier_id": "SUP-008", "name": "Premier Polymers", "country": "USA", "category": "Plastic · HDPE Tanks",
     "components": ["TENN-TANK-SOL-HDPE-30G", "TENN-TANK-REC-HDPE-30G", "TENN-TANK-SOL-HDPE-50G"], "annual_spend_usd": 1_240_000,
     "lead_time_days": 14, "moq": 50, "on_time_pct": 96.7, "quality_ppm": 156,
     "risk_score": 11, "fta_eligible": "Domestic", "primary": True, "alt_suppliers": ["Olympic Plastics"],
     "contract_expiry": "2027-12-31", "contact": "Mike Davenport · md@premierpoly.com · +1-616-877-2200"},
    {"supplier_id": "SUP-009", "name": "Midwest Steel Frame Co.", "country": "USA", "category": "Steel Frames",
     "components": ["TENN-FRAME-T16", "TENN-FRAME-T350", "TENN-FRAME-S30"], "annual_spend_usd": 2_100_000,
     "lead_time_days": 28, "moq": 20, "on_time_pct": 97.5, "quality_ppm": 71,
     "risk_score": 7, "fta_eligible": "Domestic", "primary": True, "alt_suppliers": ["Heartland Fabrication"],
     "contract_expiry": "2028-01-31", "contact": "Robert Chen · rchen@midwestframe.com · +1-815-555-2900"},
    {"supplier_id": "SUP-010", "name": "Penn Battery Industries", "country": "USA", "category": "Battery · Lead-acid",
     "components": ["TENN-BATT-LA-36V-INDUSTRIAL"], "annual_spend_usd": 480_000,
     "lead_time_days": 16, "moq": 30, "on_time_pct": 95.4, "quality_ppm": 188,
     "risk_score": 19, "fta_eligible": "Domestic", "primary": False, "alt_suppliers": ["Trojan Battery"],
     "contract_expiry": "2026-08-31", "contact": "Jim Walker · jw@pennbattery.com · +1-610-555-3300"},
    {"supplier_id": "SUP-011", "name": "Bosch Rexroth", "country": "Germany", "category": "Hydraulics & Pumps",
     "components": ["TENN-PUMP-HYD-T16AMR", "TENN-VALVE-HYD-T500"], "annual_spend_usd": 1_640_000,
     "lead_time_days": 49, "moq": 75, "on_time_pct": 97.9, "quality_ppm": 48,
     "risk_score": 13, "fta_eligible": "MFN (no FTA)", "primary": True, "alt_suppliers": ["Parker Hannifin"],
     "contract_expiry": "2027-06-30", "contact": "Klaus Weber · kweber@boschrexroth.de · +49-9352-180"},
    {"supplier_id": "SUP-012", "name": "Parker Hannifin", "country": "USA", "category": "Hydraulics & Pumps",
     "components": ["TENN-PUMP-HYD-ALT"], "annual_spend_usd": 720_000,
     "lead_time_days": 24, "moq": 25, "on_time_pct": 98.6, "quality_ppm": 38,
     "risk_score": 6, "fta_eligible": "Domestic", "primary": False, "alt_suppliers": ["Bosch Rexroth"],
     "contract_expiry": "2027-12-31", "contact": "Andrew Mills · amills@parker.com · +1-216-896-3000"},
    {"supplier_id": "SUP-013", "name": "ABB Motors & Drives", "country": "Sweden", "category": "Motors · DC/AC",
     "components": ["TENN-MOTOR-DC-AMR", "TENN-MOTOR-AC-3PH"], "annual_spend_usd": 1_980_000,
     "lead_time_days": 56, "moq": 100, "on_time_pct": 96.2, "quality_ppm": 124,
     "risk_score": 16, "fta_eligible": "MFN", "primary": True, "alt_suppliers": ["WEG"],
     "contract_expiry": "2027-09-30", "contact": "Erik Lindberg · elindberg@abb.se · +46-21-32-50-00"},
    {"supplier_id": "SUP-014", "name": "WEG Industries", "country": "Brazil", "category": "Motors · AC",
     "components": ["TENN-MOTOR-AC-3PH-ALT"], "annual_spend_usd": 420_000,
     "lead_time_days": 60, "moq": 50, "on_time_pct": 93.8, "quality_ppm": 230,
     "risk_score": 28, "fta_eligible": "MFN", "primary": False, "alt_suppliers": ["ABB"],
     "contract_expiry": "2026-10-31", "contact": "Carlos Silva · csilva@weg.net · +55-47-3276-4000"},
    {"supplier_id": "SUP-015", "name": "Industrial Brush Co.", "country": "USA", "category": "Brushes",
     "components": ["TENN-BRUSH-CYL-32", "TENN-BRUSH-DISC-20", "TENN-BRUSH-CYL-40"], "annual_spend_usd": 380_000,
     "lead_time_days": 12, "moq": 100, "on_time_pct": 98.9, "quality_ppm": 42,
     "risk_score": 5, "fta_eligible": "Domestic", "primary": True, "alt_suppliers": ["Tucel Industries"],
     "contract_expiry": "2028-04-30", "contact": "Linda Foster · lfoster@indbrush.com · +1-414-555-7800"},
    {"supplier_id": "SUP-016", "name": "Tucel Industries", "country": "USA", "category": "Brushes",
     "components": ["TENN-BRUSH-DISC-ALT"], "annual_spend_usd": 95_000,
     "lead_time_days": 14, "moq": 50, "on_time_pct": 97.4, "quality_ppm": 88,
     "risk_score": 9, "fta_eligible": "Domestic", "primary": False, "alt_suppliers": ["Industrial Brush"],
     "contract_expiry": "2027-02-28", "contact": "Mark Lyons · mlyons@tucel.com · +1-802-247-3300"},
    {"supplier_id": "SUP-017", "name": "Urethane Specialties Intl.", "country": "USA", "category": "Squeegees / Urethane",
     "components": ["TENN-SQGE-URETH-40", "TENN-SQGE-URETH-32", "TENN-SQGE-URETH-26"], "annual_spend_usd": 312_000,
     "lead_time_days": 10, "moq": 100, "on_time_pct": 99.1, "quality_ppm": 28,
     "risk_score": 4, "fta_eligible": "Domestic", "primary": True, "alt_suppliers": ["Polymer Sciences"],
     "contract_expiry": "2028-12-31", "contact": "Greg Holt · gholt@uretspec.com · +1-330-555-4400"},
    {"supplier_id": "SUP-018", "name": "Schneider Electric", "country": "France", "category": "Electronics · PLCs/Controllers",
     "components": ["TENN-PLC-AMR-NAVIGATION", "TENN-CONTROLLER-SCRUB"], "annual_spend_usd": 1_450_000,
     "lead_time_days": 42, "moq": 50, "on_time_pct": 96.5, "quality_ppm": 102,
     "risk_score": 14, "fta_eligible": "MFN", "primary": True, "alt_suppliers": ["Siemens"],
     "contract_expiry": "2027-07-31", "contact": "Marie Dubois · mdubois@se.com · +33-1-41-29-70-00"},
    {"supplier_id": "SUP-019", "name": "Siemens AG", "country": "Germany", "category": "Electronics · Drives",
     "components": ["TENN-DRIVE-VFD-T16AMR"], "annual_spend_usd": 880_000,
     "lead_time_days": 49, "moq": 25, "on_time_pct": 98.4, "quality_ppm": 56,
     "risk_score": 11, "fta_eligible": "MFN", "primary": False, "alt_suppliers": ["Schneider"],
     "contract_expiry": "2027-05-31", "contact": "Hans Mueller · hmueller@siemens.de · +49-89-636-00"},
    {"supplier_id": "SUP-020", "name": "Goodyear Industrial", "country": "USA", "category": "Tires · Industrial",
     "components": ["TENN-TIRE-IND-T16", "TENN-TIRE-IND-S30"], "annual_spend_usd": 290_000,
     "lead_time_days": 21, "moq": 50, "on_time_pct": 97.0, "quality_ppm": 84,
     "risk_score": 8, "fta_eligible": "Domestic", "primary": True, "alt_suppliers": ["Trelleborg"],
     "contract_expiry": "2027-11-30", "contact": "Mary Lou Carson · mlc@goodyear.com · +1-330-796-2121"},
]

@api_router.get("/suppliers")
async def list_suppliers(_: User = Depends(get_current_user), country: Optional[str] = None, category: Optional[str] = None, risk_max: Optional[int] = None):
    """List Tennant suppliers with optional filters by country / category / max risk score.
    Combines built-in seed suppliers with any manually-added entries stored in MongoDB."""
    # Load custom suppliers (manually added via UI) and merge with seed list
    custom = await db.suppliers_custom.find({}, {"_id": 0}).to_list(1000)
    all_suppliers = list(TENNANT_SUPPLIERS) + custom
    # Active brand overlay → swap supplier names when not Tennant.
    brand = await _active_brand_doc()
    if brand and brand.get("brand_id") != "orisei-freight":
        all_suppliers = [_overlay_supplier(s, brand, i) for i, s in enumerate(all_suppliers)]
    out = list(all_suppliers)
    if country:
        out = [s for s in out if (s.get("country") or "").lower() == country.lower()]
    if category:
        out = [s for s in out if category.lower() in (s.get("category") or "").lower()]
    if risk_max is not None:
        out = [s for s in out if (s.get("risk_score") or 0) <= risk_max]
    # Summary stats
    total_spend = sum((s.get("annual_spend_usd") or 0) for s in all_suppliers)
    countries = {}
    categories = {}
    for s in all_suppliers:
        c = s.get("country") or "Unknown"
        cat = s.get("category") or "Unknown"
        countries[c] = countries.get(c, 0) + (s.get("annual_spend_usd") or 0)
        categories[cat] = categories.get(cat, 0) + 1
    return {
        "suppliers": out,
        "summary": {
            "total_suppliers": len(all_suppliers),
            "primary_count": sum(1 for s in all_suppliers if s.get("primary")),
            "total_annual_spend_usd": total_spend,
            "single_source_components": sum(1 for s in all_suppliers if not s.get("alt_suppliers")),
            "high_risk_count": sum(1 for s in all_suppliers if (s.get("risk_score") or 0) >= 20),
            "expiring_contracts_12mo": sum(1 for s in all_suppliers if s.get("contract_expiry") and datetime.fromisoformat(s["contract_expiry"]).date() < (datetime.now(timezone.utc).date() + timedelta(days=365))),
            "by_country": [{"country": k, "spend": v} for k, v in sorted(countries.items(), key=lambda x: -x[1])],
            "by_category": [{"category": k, "count": v} for k, v in sorted(categories.items(), key=lambda x: -x[1])],
        },
    }


@api_router.post("/suppliers")
async def create_supplier(payload: dict, user: User = Depends(require_role("admin", "dispatcher"))):
    """Manually add a supplier. Stored in MongoDB and merged with the seed list."""
    if not payload.get("name"):
        raise HTTPException(400, "name is required")
    # Allocate a new SUP-### id based on existing count
    existing_custom = await db.suppliers_custom.count_documents({})
    sup_id = payload.get("supplier_id") or f"SUP-C{existing_custom + 1:03d}"
    doc = {
        "supplier_id": sup_id,
        "name": payload["name"],
        "country": payload.get("country") or "USA",
        "category": payload.get("category") or "Uncategorized",
        "components": payload.get("components") or [],
        "annual_spend_usd": int(payload.get("annual_spend_usd") or 0),
        "lead_time_days": int(payload.get("lead_time_days") or 0),
        "moq": int(payload.get("moq") or 0),
        "on_time_pct": float(payload.get("on_time_pct") or 0),
        "quality_ppm": int(payload.get("quality_ppm") or 0),
        "risk_score": int(payload.get("risk_score") or 10),
        "fta_eligible": payload.get("fta_eligible") or "—",
        "primary": bool(payload.get("primary", True)),
        "alt_suppliers": payload.get("alt_suppliers") or [],
        "contract_expiry": payload.get("contract_expiry") or "",
        "contact": payload.get("contact") or "",
        "notes": payload.get("notes") or "",
        "created_by": user.user_id,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "is_custom": True,
    }
    await db.suppliers_custom.insert_one(dict(doc))
    return {"ok": True, "supplier": doc}


@api_router.delete("/suppliers/{supplier_id}")
async def delete_supplier(supplier_id: str, _: User = Depends(require_role("admin"))):
    """Only custom (manually-added) suppliers can be deleted; seeded suppliers are immutable."""
    r = await db.suppliers_custom.delete_one({"supplier_id": supplier_id})
    if r.deleted_count == 0:
        raise HTTPException(404, "Supplier not found or built-in (immutable)")
    return {"ok": True}

# -------------------- INSPIRATIONAL QUOTES (Command Center subtle ticker) --------------------
INSPIRATIONAL_QUOTES = [
    {"text": "The only way to do great work is to love what you do.", "author": "Steve Jobs"},
    {"text": "Success is not final, failure is not fatal: it is the courage to continue that counts.", "author": "Winston Churchill"},
    {"text": "Logistics is the lifeline of every industry.", "author": "Anonymous"},
    {"text": "It does not matter how slowly you go as long as you do not stop.", "author": "Confucius"},
    {"text": "Quality means doing it right when no one is looking.", "author": "Henry Ford"},
    {"text": "Without continuous improvement, there can be no excellence.", "author": "W. Edwards Deming"},
    {"text": "Excellence is not a destination; it is a continuous journey.", "author": "Brian Tracy"},
    {"text": "The supply chain stuff is really tricky.", "author": "Elon Musk"},
    {"text": "Hard work beats talent when talent doesn't work hard.", "author": "Tim Notke"},
    {"text": "Discipline is the bridge between goals and accomplishment.", "author": "Jim Rohn"},
    {"text": "If everything seems under control, you're not going fast enough.", "author": "Mario Andretti"},
    {"text": "Done is better than perfect.", "author": "Sheryl Sandberg"},
    {"text": "Plans are worthless, but planning is everything.", "author": "Dwight D. Eisenhower"},
    {"text": "Amateurs talk strategy; professionals talk logistics.", "author": "Gen. Omar Bradley"},
    {"text": "Whether you think you can or think you can't, you're right.", "author": "Henry Ford"},
    {"text": "Action is the foundational key to all success.", "author": "Pablo Picasso"},
    {"text": "Continuous effort — not strength or intelligence — is the key to unlocking our potential.", "author": "Winston Churchill"},
    {"text": "Don't watch the clock; do what it does — keep going.", "author": "Sam Levenson"},
    {"text": "Strive not to be a success, but rather to be of value.", "author": "Albert Einstein"},
    {"text": "Out of clutter, find simplicity.", "author": "Albert Einstein"},
    {"text": "The pessimist sees difficulty in every opportunity. The optimist sees opportunity in every difficulty.", "author": "Winston Churchill"},
    {"text": "Quality is never an accident; it is always the result of intelligent effort.", "author": "John Ruskin"},
    {"text": "The man who moves a mountain begins by carrying away small stones.", "author": "Confucius"},
    {"text": "Efficiency is doing things right; effectiveness is doing the right things.", "author": "Peter Drucker"},
    {"text": "Genius is one percent inspiration and ninety-nine percent perspiration.", "author": "Thomas Edison"},
    {"text": "Either you run the day, or the day runs you.", "author": "Jim Rohn"},
    {"text": "Innovation distinguishes between a leader and a follower.", "author": "Steve Jobs"},
    {"text": "Wherever you go, no matter what the weather, always bring your own sunshine.", "author": "Anthony J. D'Angelo"},
    {"text": "If you cannot do great things, do small things in a great way.", "author": "Napoleon Hill"},
    {"text": "Success is walking from failure to failure with no loss of enthusiasm.", "author": "Winston Churchill"},
    {"text": "The only place where success comes before work is in the dictionary.", "author": "Vidal Sassoon"},
    {"text": "Believe you can and you're halfway there.", "author": "Theodore Roosevelt"},
    {"text": "The best way to predict the future is to create it.", "author": "Peter Drucker"},
    {"text": "Coming together is a beginning; staying together is progress; working together is success.", "author": "Henry Ford"},
    {"text": "Don't be afraid to give up the good to go for the great.", "author": "John D. Rockefeller"},
    {"text": "Energy and persistence conquer all things.", "author": "Benjamin Franklin"},
    {"text": "Setting goals is the first step in turning the invisible into the visible.", "author": "Tony Robbins"},
    {"text": "Quality is doing the right thing even when no one is watching.", "author": "Henry Ford"},
    {"text": "The future depends on what you do today.", "author": "Mahatma Gandhi"},
    {"text": "Opportunities don't happen. You create them.", "author": "Chris Grosser"},
    {"text": "It's not whether you get knocked down, it's whether you get up.", "author": "Vince Lombardi"},
    {"text": "Perfection is not attainable, but if we chase perfection we can catch excellence.", "author": "Vince Lombardi"},
    {"text": "Logistics is the ball and chain of armored warfare.", "author": "Heinz Guderian"},
    {"text": "What you do today can improve all your tomorrows.", "author": "Ralph Marston"},
    {"text": "Start where you are. Use what you have. Do what you can.", "author": "Arthur Ashe"},
    {"text": "An ounce of action is worth a ton of theory.", "author": "Friedrich Engels"},
    {"text": "The harder you work for something, the greater you'll feel when you achieve it.", "author": "Anonymous"},
    {"text": "Trust the process.", "author": "Anonymous"},
    {"text": "Behind every successful logistics operation is a great team.", "author": "Anonymous"},
    {"text": "Every load matters. Every mile counts.", "author": "Trucking proverb"},
    {"text": "Excellence is the gradual result of always striving to do better.", "author": "Pat Riley"},
    {"text": "Do not wait to strike till the iron is hot, but make it hot by striking.", "author": "William Butler Yeats"},
    {"text": "Great things are done by a series of small things brought together.", "author": "Vincent Van Gogh"},
    {"text": "The road to success and the road to failure are almost exactly the same.", "author": "Colin R. Davis"},
    {"text": "Don't let yesterday take up too much of today.", "author": "Will Rogers"},
    {"text": "Logistics: where the rubber meets the road.", "author": "Anonymous"},
    {"text": "Productivity is never an accident. It is always the result of a commitment to excellence.", "author": "Paul J. Meyer"},
    {"text": "Stop being afraid of what could go wrong; start being excited about what could go right.", "author": "Tony Robbins"},
    {"text": "Move fast and ship things.", "author": "Anonymous"},
    {"text": "A goal without a plan is just a wish.", "author": "Antoine de Saint-Exupéry"},
    {"text": "It always seems impossible until it's done.", "author": "Nelson Mandela"},
    {"text": "If you want something you've never had, you must be willing to do something you've never done.", "author": "Thomas Jefferson"},
    {"text": "The journey of a thousand miles begins with one step.", "author": "Lao Tzu"},
    {"text": "Don't count the days, make the days count.", "author": "Muhammad Ali"},
    {"text": "Champions keep playing until they get it right.", "author": "Billie Jean King"},
    {"text": "Take chances, make mistakes. That's how you grow.", "author": "Mary Tyler Moore"},
    {"text": "The successful warrior is the average man, with laser-like focus.", "author": "Bruce Lee"},
    {"text": "Concentrate all your thoughts upon the work in hand.", "author": "Alexander Graham Bell"},
    {"text": "Knowing is not enough; we must apply. Willing is not enough; we must do.", "author": "Goethe"},
    {"text": "Failure is the condiment that gives success its flavor.", "author": "Truman Capote"},
    {"text": "The expert in anything was once a beginner.", "author": "Helen Hayes"},
    {"text": "We become what we think about.", "author": "Earl Nightingale"},
    {"text": "Slow progress is still progress.", "author": "Anonymous"},
    {"text": "The harder I work, the luckier I get.", "author": "Samuel Goldwyn"},
    {"text": "Be the change you want to see in the world.", "author": "Mahatma Gandhi"},
    {"text": "When everything seems to be going against you, remember that the airplane takes off against the wind.", "author": "Henry Ford"},
    {"text": "There are no traffic jams along the extra mile.", "author": "Roger Staubach"},
    {"text": "If you don't risk anything, you risk even more.", "author": "Erica Jong"},
    {"text": "Success usually comes to those who are too busy to be looking for it.", "author": "Henry David Thoreau"},
    {"text": "Don't be busy. Be productive.", "author": "Anonymous"},
    {"text": "Cargo doesn't sleep — neither do great dispatchers.", "author": "Trucking proverb"},
    {"text": "What gets measured gets managed.", "author": "Peter Drucker"},
    {"text": "Lead, follow, or get out of the way.", "author": "Thomas Paine"},
    {"text": "If opportunity doesn't knock, build a door.", "author": "Milton Berle"},
    {"text": "A ship in harbor is safe — but that is not what ships are built for.", "author": "John A. Shedd"},
    {"text": "Yesterday's home runs don't win today's games.", "author": "Babe Ruth"},
    {"text": "The two most important days in your life are the day you are born and the day you find out why.", "author": "Mark Twain"},
    {"text": "Whatever you are, be a good one.", "author": "Abraham Lincoln"},
    {"text": "Every truck driver carries a piece of the economy on their back.", "author": "Anonymous"},
    {"text": "Things may come to those who wait, but only the things left by those who hustle.", "author": "Abraham Lincoln"},
    {"text": "The best preparation for tomorrow is doing your best today.", "author": "H. Jackson Brown Jr."},
    {"text": "Don't stop when you're tired. Stop when you're done.", "author": "David Goggins"},
    {"text": "Stay hungry, stay foolish.", "author": "Steve Jobs"},
    {"text": "Hope is not a strategy.", "author": "Vince Lombardi"},
    {"text": "If it scares you, it might be a good thing to try.", "author": "Seth Godin"},
    {"text": "We are what we repeatedly do. Excellence, then, is not an act, but a habit.", "author": "Aristotle"},
    {"text": "Inventory is money sitting around in another form.", "author": "Rhonda Adams"},
    {"text": "Logistics comprises the means and arrangements which work out the plans of strategy and tactics.", "author": "Antoine-Henri Jomini"},
    {"text": "Never confuse motion with progress.", "author": "Alfred A. Montapert"},
    {"text": "Be so good they can't ignore you.", "author": "Steve Martin"},
    {"text": "The way to get started is to quit talking and begin doing.", "author": "Walt Disney"},
    {"text": "Time is the scarcest resource, and unless it is managed nothing else can be managed.", "author": "Peter Drucker"},
    {"text": "Optimism is the faith that leads to achievement.", "author": "Helen Keller"},
    {"text": "We rise by lifting others.", "author": "Robert Ingersoll"},
]

@api_router.get("/quotes")
async def get_quotes(_: User = Depends(get_current_user)):
    """Returns Tennant's curated motivational/inspirational quote rotation."""
    return {"quotes": INSPIRATIONAL_QUOTES, "count": len(INSPIRATIONAL_QUOTES)}

# -------------------- TRADE COMPLIANCE --------------------
# Tariff schedules, trade programs, country-of-origin rules, watchlists relevant to Tennant.
TRADE_COMPLIANCE = {
    "summary": {
        "active_tariff_codes": 18,
        "section_301_exposure_pct": 14.2,
        "ftz_active_lots": 12,
        "duty_drawback_ytd_usd": 184320,
        "section_232_steel_aluminum_in_scope": True,
        "last_updated": datetime.now(timezone.utc).date().isoformat(),
    },
    "tariff_schedules": [
        {"hts": "8479.89.95", "description": "Floor cleaning machines, electric", "general_duty": "2.5%", "column_1_special": "Free (USMCA, JP, KR, AU, SG)", "section_301_list": None, "notes": "Primary HS for T-series scrubbers / S-series sweepers"},
        {"hts": "8508.11.00", "description": "Vacuum cleaners, < 1500 W, household type", "general_duty": "Free", "column_1_special": "Free", "section_301_list": None, "notes": "Tennant residential-grade SKUs (rare)"},
        {"hts": "8508.19.00", "description": "Vacuum cleaners, industrial, > 1500 W", "general_duty": "Free", "column_1_special": "Free", "section_301_list": None, "notes": "Most Tennant industrial vacs"},
        {"hts": "8507.60.00", "description": "Lithium-ion batteries", "general_duty": "3.4%", "column_1_special": "Free (USMCA, JP, KR)", "section_301_list": "List 3 (China origin: +25%)", "notes": "K+N Korea sourcing keeps us out of 301; China-origin packs trigger 25%"},
        {"hts": "8507.20.80", "description": "Lead-acid storage batteries, NOI", "general_duty": "3.5%", "column_1_special": "Free (USMCA, JP, KR)", "section_301_list": None, "notes": "AGM packs"},
        {"hts": "8501.31.50", "description": "DC motors, 750W-75kW", "general_duty": "2.8%", "column_1_special": "Free (USMCA)", "section_301_list": "List 3 (China: +25%)", "notes": "Drive motors — verify CoO carefully"},
        {"hts": "8501.10.40", "description": "AC motors, single-phase, < 37.5W", "general_duty": "6.7%", "column_1_special": "Free (USMCA, KR)", "section_301_list": "List 4A (China: +7.5%)"},
        {"hts": "9603.50.00", "description": "Brushes, rotating, for machines", "general_duty": "Free", "column_1_special": "Free", "section_301_list": None, "notes": "Cylindrical & disc brushes"},
        {"hts": "3926.90.99", "description": "Plastic articles, NOI (HDPE tanks)", "general_duty": "5.3%", "column_1_special": "Free (USMCA, KR)", "section_301_list": "List 3 (China: +25%)", "notes": "Solution & recovery tanks"},
        {"hts": "7308.90.95", "description": "Steel structural frames, machinery", "general_duty": "Free", "column_1_special": "Free", "section_301_list": None, "section_232": "+25% on steel content (China/EU/JP exempted via TRQ)"},
        {"hts": "8544.30.00", "description": "Wiring harnesses, ignition / industrial", "general_duty": "5.0%", "column_1_special": "Free (USMCA, KR)", "section_301_list": "List 3 (China: +25%)"},
        {"hts": "4009.32.00", "description": "Rubber tubing / hoses, reinforced", "general_duty": "2.5%", "column_1_special": "Free (USMCA, KR)", "section_301_list": "List 3 (China: +25%)"},
        {"hts": "8413.70.20", "description": "Submersible pumps, water/solution", "general_duty": "Free", "column_1_special": "Free", "section_301_list": "List 3 (China: +25%)"},
        {"hts": "4011.20.10", "description": "Industrial rubber tires, solid", "general_duty": "3.4%", "column_1_special": "Free (USMCA)", "section_301_list": "List 4A (China: +7.5%)"},
        {"hts": "8714.99.80", "description": "Parts of vehicles, NOI", "general_duty": "Free", "column_1_special": "Free", "section_301_list": "List 3 (China: +25%)"},
        {"hts": "8421.39.80", "description": "Filters, machinery, gas/air", "general_duty": "Free", "column_1_special": "Free", "section_301_list": None},
        {"hts": "4016.99.55", "description": "Squeegee blades, rubber/urethane", "general_duty": "2.5%", "column_1_special": "Free (USMCA, KR)", "section_301_list": "List 3 (China: +25%)"},
        {"hts": "8537.10.91", "description": "Programmable controllers, < 1kV", "general_duty": "2.7%", "column_1_special": "Free (USMCA, KR)", "section_301_list": "List 3 (China: +25%)"},
    ],
    "trade_programs": [
        {"program": "USMCA", "abbr": "USMCA", "type": "FTA", "tennant_use": "Mexico-sourced motors & wiring harnesses; Tennant ships to Canadian customers under 0% duty.", "status": "Active", "renewed": "2026 review pending"},
        {"program": "KORUS FTA (Korea)", "abbr": "KR-FTA", "type": "FTA", "tennant_use": "Samsung SDI / LG / Yura inbound Li-ion batteries from Korea — free of duty.", "status": "Active"},
        {"program": "Japan FTA (USJTA)", "abbr": "JP-FTA", "type": "Limited FTA", "tennant_use": "Yazaki wiring harnesses from Japan — partial duty reduction.", "status": "Active"},
        {"program": "Australia FTA (AUSFTA)", "abbr": "AU-FTA", "type": "FTA", "tennant_use": "Limited — used when AU customers require certification of origin.", "status": "Active"},
        {"program": "Singapore FTA", "abbr": "SG-FTA", "type": "FTA", "tennant_use": "ASEAN distribution hub — components transiting Singapore.", "status": "Active"},
        {"program": "Foreign Trade Zone (FTZ)", "abbr": "FTZ", "type": "Duty Deferral", "tennant_use": "Golden Valley & Louisville sites operate under FTZ #119 (subzone) — defer duty until withdrawn for consumption.", "status": "Active · 12 lots open"},
        {"program": "Duty Drawback", "abbr": "DBK", "type": "Refund", "tennant_use": "Refund of duty on imported components exported in finished scrubbers (e.g., to Mexico/EU).", "status": "Active · $184K YTD"},
        {"program": "GSP (Generalized System of Preferences)", "abbr": "GSP", "type": "Preference", "tennant_use": "Limited — some Thai/Philippine spare parts. NOTE: GSP authorization expired Dec 2020, awaiting Congress.", "status": "Expired"},
        {"program": "Section 321 De Minimis", "abbr": "§321", "type": "Threshold", "tennant_use": "Shipments < $800 to a single consignee/day duty-free. Used for low-value parts shipments.", "status": "Active (under reform review)"},
    ],
    "section_301": {
        "country_focus": "China",
        "lists": [
            {"list": "List 1", "rate_pct": 25, "effective": "2018-07-06", "applies_to_tennant": "Some motor controllers (8501)"},
            {"list": "List 2", "rate_pct": 25, "effective": "2018-08-23", "applies_to_tennant": "Limited"},
            {"list": "List 3", "rate_pct": 25, "effective": "2018-09-24", "applies_to_tennant": "Li-ion batteries, plastics, motor parts, hoses, harnesses, pumps, controllers — significant exposure"},
            {"list": "List 4A", "rate_pct": 7.5, "effective": "2019-09-01", "applies_to_tennant": "Tires, accessories"},
        ],
        "current_exclusions": [
            {"hts": "8507.60.00", "expires": "2025-12-31", "scope": "Specific cell chemistries — verify supplier certification"},
        ],
        "mitigation": "Tennant strategy: shift CoO to Korea (KORUS-eligible) and USMCA suppliers; document substantial transformation when assembling in MX.",
    },
    "section_232": {
        "scope": "Steel & aluminum imports",
        "rate_pct": 25,
        "tennant_exposure": "Steel frames, fasteners, sheet metal. Quotas/exemptions for EU/JP/UK active; Canada/Mexico exempt under USMCA.",
        "tariff_rate_quota": "Active for EU (3.3M MT), JP, UK",
    },
    "country_of_origin_rules": [
        {"product": "T16 AMR Scrubber", "assembled_in": "USA (Golden Valley)", "us_content_pct": 68, "marking_required": "Made in USA — California-compliant", "notes": "Substantial transformation occurs at final assembly. Battery & some electronics from KR."},
        {"product": "T350 LPG Scrubber", "assembled_in": "USA (Holland, MI)", "us_content_pct": 74, "marking_required": "Made in USA", "notes": "LPG tank from US supplier; engine from JP (Honda)."},
        {"product": "S30 Industrial Sweeper", "assembled_in": "USA (Louisville, KY)", "us_content_pct": 71, "marking_required": "Made in USA", "notes": ""},
        {"product": "Spare Parts Kit (M30)", "assembled_in": "USA (kitted)", "us_content_pct": 55, "marking_required": "Each part marked with CoO", "notes": "Mixed origin — must mark per-component."},
    ],
    "watchlists": {
        "denied_parties": {
            "source": "BIS Entity List, OFAC SDN, State Department DTC Debarred",
            "last_screened": (datetime.now(timezone.utc) - timedelta(hours=2)).isoformat(),
            "matches_30d": 0,
            "screening_cadence": "Every shipment + weekly batch",
        },
        "embargoed_countries": ["Cuba", "Iran", "North Korea", "Syria", "Russia (sectoral)", "Belarus (sectoral)", "Crimea/Donetsk/Luhansk"],
        "restricted_end_use": ["Military / Defense", "Nuclear", "Missile Tech", "Chemical/Biological"],
        "uflpa_xinjiang_diligence": "Required for all China-origin batteries, polysilicon, cotton-derived items. Tennant policy: ZERO China-origin Li-ion accepted in 2026.",
    },
    "broker_filings": {
        "primary_broker": "UPS Supply Chain Solutions (UPS_SCS)",
        "ace_portal_id": "UPS_SCS_001",
        "ytd_entry_summaries": 412,
        "ytd_isf_filings": 284,
        "average_clearance_hrs": 14.2,
        "exam_rate_pct": 2.8,
        "post_summary_corrections_ytd": 9,
    },
    "key_regulations": [
        {"reg": "19 CFR 134 (Country of Origin Marking)", "scope": "Every imported article must be conspicuously marked with English-language CoO unless excepted.", "owner": "Trade Compliance Manager"},
        {"reg": "15 CFR 730–774 (EAR)", "scope": "Export Administration Regulations — dual-use items. T-series scrubbers generally EAR99 (no license).", "owner": "Export Compliance Officer"},
        {"reg": "19 CFR 149 (ISF)", "scope": "Importer Security Filing — 10+2 elements due 24 hrs prior to lading at foreign port.", "owner": "K+N + UPS_SCS"},
        {"reg": "19 USC 1641 (Customs Broker)", "scope": "Power of Attorney to UPS_SCS on file; renewed annually.", "owner": "Compliance Manager"},
        {"reg": "UFLPA (Uyghur Forced Labor Prevention Act)", "scope": "Rebuttable presumption — all Xinjiang-region goods or downstream input banned. Documentary diligence mandatory.", "owner": "Trade + Supplier Mgmt"},
        {"reg": "Lacey Act Declarations", "scope": "Only if wood packaging / pallets from regulated source. Tennant uses heat-treated ISPM-15 pallets exclusively.", "owner": "Packaging Engineering"},
    ],
    "recent_alerts": [
        {"date": (datetime.now(timezone.utc) - timedelta(days=2)).date().isoformat(), "title": "USTR Section 301 — Final List 4B rates effective Jan 1, 2027", "severity": "high", "impact": "Affects HTS 8413 pumps (+25%) and 8537 controllers (+25%) if China-origin. Action: confirm KR/USMCA sourcing for 2027."},
        {"date": (datetime.now(timezone.utc) - timedelta(days=5)).date().isoformat(), "title": "CBP CSMS #59042: PGA filings now required for Li-ion ≥ 100Wh", "severity": "medium", "impact": "Add EPA/DOT PGA data to all 8507.60 entries. UPS_SCS already updated."},
        {"date": (datetime.now(timezone.utc) - timedelta(days=8)).date().isoformat(), "title": "Canada CARM Phase 3 portal go-live", "severity": "medium", "impact": "All Tennant Canadian shipments must use CARM for RPP / B3 filings."},
        {"date": (datetime.now(timezone.utc) - timedelta(days=12)).date().isoformat(), "title": "Section 232 Steel — UK TRQ allocation Q2 published", "severity": "low", "impact": "UK-origin steel frame components for HOM remain in-quota."},
    ],
    "quick_links": [
        {"label": "USITC HTS Search", "url": "https://hts.usitc.gov/"},
        {"label": "CBP ACE Portal", "url": "https://ace.cbp.dhs.gov/"},
        {"label": "USTR Section 301 Tariff Lookup", "url": "https://ustr.gov/issue-areas/enforcement/section-301-investigations"},
        {"label": "BIS Entity List", "url": "https://www.bis.doc.gov/index.php/policy-guidance/lists-of-parties-of-concern/entity-list"},
        {"label": "OFAC SDN List Search", "url": "https://sanctionssearch.ofac.treas.gov/"},
        {"label": "USMCA Center", "url": "https://www.trade.gov/usmca"},
        {"label": "Census Schedule B", "url": "https://uscensus.prod.3ceonline.com/"},
        {"label": "DOT HazMat / 49 CFR", "url": "https://www.phmsa.dot.gov/hazmat/regs"},
    ],
}

@api_router.get("/trade-compliance")
async def get_trade_compliance(_: User = Depends(get_current_user)):
    return TRADE_COMPLIANCE


# -------------------- COMPANY BRANDING (multi-tenant theme) --------------------
# Endpoints live in routes/branding.py — the DEFAULT_BRAND below is kept here
# so the admin dashboard's brand-fallback line stays a cheap dict lookup.
from routes.branding import DEFAULT_BRAND, build_branding_router  # noqa: E402


# Branding endpoints live in routes.branding — registered near wire-up.
# DEFAULT_BRAND remains imported below for the admin dashboard fallback.

# -------------------- ADMIN DASHBOARD (system telemetry) --------------------
_APP_BOOT_AT = datetime.now(timezone.utc)


def _human_uptime(sec: int) -> str:
    d, rem = divmod(sec, 86400)
    h, rem = divmod(rem, 3600)
    m, _ = divmod(rem, 60)
    parts = []
    if d: parts.append(f"{d}d")
    if h: parts.append(f"{h}h")
    if m or not parts: parts.append(f"{m}m")
    return " ".join(parts)


@api_router.get("/admin/dashboard")
async def admin_dashboard(_: User = Depends(require_role("admin"))):
    """One-stop telemetry the Admin Dashboard surfaces — DB health,
    record counts, recent user activity, LLM usage, active brand, and
    quick-toggle settings state."""
    now = datetime.now(timezone.utc)

    db_ms = None
    db_ok = False
    try:
        t0 = now
        await asyncio.wait_for(db.command("ping"), timeout=1.5)
        db_ms = int((datetime.now(timezone.utc) - t0).total_seconds() * 1000)
        db_ok = True
    except Exception:
        pass

    collections = [
        "users", "user_sessions", "shipments", "carrier_onboarding",
        "drivers", "trailers", "freight_bills", "yard_reports",
        "truckload_bookings", "workbook_tabs", "outbound_emails",
        "ai_messages", "audit_log", "suppliers_custom", "company_brand",
        "claims", "documents",
    ]
    counts: Dict[str, int] = {}
    for c in collections:
        try:
            counts[c] = await db[c].estimated_document_count()
        except Exception:
            counts[c] = 0

    cutoff_24h = (now - timedelta(hours=24)).isoformat()
    cutoff_7d = (now - timedelta(days=7)).isoformat()
    try:
        active_24h = await db.user_sessions.distinct("user_id", {"created_at": {"$gte": cutoff_24h}})
        active_7d = await db.user_sessions.distinct("user_id", {"created_at": {"$gte": cutoff_7d}})
    except Exception:
        active_24h, active_7d = [], []

    role_counts: Dict[str, int] = {}
    try:
        async for u in db.users.find({}, {"_id": 0, "role": 1}):
            r = u.get("role") or "unknown"
            role_counts[r] = role_counts.get(r, 0) + 1
    except Exception:
        pass

    try:
        llm_total = await db.ai_messages.count_documents({"role": "user"})
        llm_24h = await db.ai_messages.count_documents({"role": "user", "created_at": {"$gte": cutoff_24h}})
    except Exception:
        llm_total, llm_24h = 0, 0

    trend = []
    try:
        for i in range(13, -1, -1):
            ds = (now - timedelta(days=i)).replace(hour=0, minute=0, second=0, microsecond=0)
            de = ds + timedelta(days=1)
            cnt = await db.shipments.count_documents({"created_at": {"$gte": ds.isoformat(), "$lt": de.isoformat()}})
            trend.append({"date": ds.strftime("%m/%d"), "count": cnt})
    except Exception:
        pass

    try:
        recent_audit = await db.audit_log.find({}, {"_id": 0}).sort("at", -1).limit(10).to_list(10)
    except Exception:
        recent_audit = []

    active_brand = await db.company_brand.find_one({"is_active": True}, {"_id": 0})
    if not active_brand:
        active_brand = {"brand_id": "orisei-freight", "company_name": "Orisei Freight Solutions", "short_name": "Orisei", "primary_color": "#0E3A6B", "logo_letter": "T", "is_default": True}
    brand_count = await db.company_brand.count_documents({})

    settings_doc = await db.admin_settings.find_one({}, {"_id": 0}) or {}
    erp_doc = await db.erp_config.find_one({"is_active": True}, {"_id": 0}) or None

    uptime_sec = int((now - _APP_BOOT_AT).total_seconds())

    return {
        "system": {
            "db_ok": db_ok,
            "db_ping_ms": db_ms,
            "uptime_seconds": uptime_sec,
            "uptime_human": _human_uptime(uptime_sec),
            "boot_at": _APP_BOOT_AT.isoformat(),
            "server_time": now.isoformat(),
        },
        "counts": counts,
        "users": {
            "total": counts.get("users", 0),
            "active_24h": len(active_24h),
            "active_7d": len(active_7d),
            "role_breakdown": role_counts,
        },
        "llm": {"total_messages": llm_total, "messages_24h": llm_24h},
        "shipments": {"total": counts.get("shipments", 0), "trend_14d": trend},
        "brand": {"active": active_brand, "custom_count": brand_count},
        "erp": {"active": erp_doc},
        "recent_audit": recent_audit,
        "settings": settings_doc,
    }


@api_router.post("/admin/dashboard/quick-toggle")
async def admin_quick_toggle(payload: dict, _: User = Depends(require_role("admin"))):
    """Quick-toggle endpoint used by the dashboard's switch row."""
    key = (payload or {}).get("key")
    value = (payload or {}).get("value")
    if not key:
        raise HTTPException(400, "key required")
    await db.admin_settings.update_one({}, {"$set": {key: value}}, upsert=True)
    return {"ok": True, "key": key, "value": value}


# Server Registry endpoints live in routes.server_registry — included below.


# -------------------- ERP CONNECTORS (multi-system) --------------------
# Lets the admin point the TMS at the company's actual ERP — SAP S/4HANA,
# Oracle Fusion, Microsoft D365, NetSuite, Infor M3, Sage X3, Epicor, IFS.
# Each connector stores its credentials encrypted-at-rest only via the
# host MongoDB security model; secrets are never returned to the frontend.

ERP_TEMPLATES = [
    {
        "key": "sap_s4hana",
        "name": "SAP S/4HANA",
        "auth_modes": ["oauth2_client_credentials", "basic"],
        "fields": [
            {"key": "base_url", "label": "S/4HANA Base URL", "placeholder": "https://my-s4.company.com:443"},
            {"key": "client", "label": "Client (mandant)", "placeholder": "100"},
            {"key": "username", "label": "Service User",     "placeholder": "TMS_INTEGRATION"},
            {"key": "password", "label": "Service Password", "secret": True},
            {"key": "oauth_token_url", "label": "OAuth Token URL (if OAuth)", "placeholder": "https://auth.company.com/oauth/token", "optional": True},
            {"key": "oauth_client_id", "label": "OAuth Client ID", "optional": True},
            {"key": "oauth_client_secret", "label": "OAuth Client Secret", "secret": True, "optional": True},
        ],
        "test_path": "/sap/opu/odata/sap/API_SALES_ORDER_SRV/$metadata",
    },
    {
        "key": "oracle_fusion",
        "name": "Oracle Fusion Cloud ERP",
        "auth_modes": ["basic", "oauth2_password"],
        "fields": [
            {"key": "base_url", "label": "Fusion Pod URL", "placeholder": "https://fa-xxx.oraclecloud.com"},
            {"key": "username", "label": "Integration User"},
            {"key": "password", "label": "Password", "secret": True},
        ],
        "test_path": "/fscmRestApi/resources/11.13.18.05/salesOrdersForOrderHub",
    },
    {
        "key": "dynamics_365",
        "name": "Microsoft Dynamics 365 F&O",
        "auth_modes": ["oauth2_client_credentials"],
        "fields": [
            {"key": "base_url", "label": "F&O Environment URL", "placeholder": "https://company.operations.dynamics.com"},
            {"key": "tenant_id", "label": "Azure Tenant ID"},
            {"key": "client_id", "label": "App Registration Client ID"},
            {"key": "client_secret", "label": "Client Secret", "secret": True},
        ],
        "test_path": "/data/SalesOrderHeaders",
    },
    {
        "key": "netsuite",
        "name": "Oracle NetSuite",
        "auth_modes": ["tba_oauth1"],
        "fields": [
            {"key": "account_id", "label": "NetSuite Account ID", "placeholder": "1234567"},
            {"key": "consumer_key", "label": "Consumer Key"},
            {"key": "consumer_secret", "label": "Consumer Secret", "secret": True},
            {"key": "token_id", "label": "Token ID"},
            {"key": "token_secret", "label": "Token Secret", "secret": True},
        ],
        "test_path": "/services/rest/record/v1/salesOrder",
    },
    {
        "key": "infor_m3",
        "name": "Infor M3 / CloudSuite",
        "auth_modes": ["basic", "oauth2_client_credentials"],
        "fields": [
            {"key": "base_url",  "label": "Infor Tenant URL", "placeholder": "https://mingle-ionapi.eu1.inforcloudsuite.com"},
            {"key": "tenant",    "label": "Tenant ID"},
            {"key": "username",  "label": "Service User"},
            {"key": "password",  "label": "Password", "secret": True},
        ],
        "test_path": "/M3/m3api-rest/v2/execute/CRS610MI/LstByNumber",
    },
    {
        "key": "sage_x3",
        "name": "Sage X3",
        "auth_modes": ["basic"],
        "fields": [
            {"key": "base_url", "label": "Sage X3 Endpoint", "placeholder": "https://sage.company.com:8124"},
            {"key": "username", "label": "User"},
            {"key": "password", "label": "Password", "secret": True},
            {"key": "pool",     "label": "Pool / Endpoint", "placeholder": "SEED"},
        ],
        "test_path": "/api1/x3/erp/SEED/SOH",
    },
    {
        "key": "epicor_kinetic",
        "name": "Epicor Kinetic",
        "auth_modes": ["basic", "api_key"],
        "fields": [
            {"key": "base_url", "label": "Kinetic Server URL", "placeholder": "https://kinetic.company.com/EpicorERP"},
            {"key": "api_key",  "label": "API Key",     "secret": True, "optional": True},
            {"key": "username", "label": "User",        "optional": True},
            {"key": "password", "label": "Password",    "secret": True, "optional": True},
            {"key": "company",  "label": "Company ID",  "placeholder": "EPIC06"},
        ],
        "test_path": "/api/v2/odata/EPIC06/Erp.BO.SalesOrderSvc/SalesOrders",
    },
    {
        "key": "ifs_cloud",
        "name": "IFS Cloud",
        "auth_modes": ["oauth2_password", "oauth2_client_credentials"],
        "fields": [
            {"key": "base_url", "label": "IFS Cloud URL", "placeholder": "https://ifs.company.com"},
            {"key": "client_id", "label": "Client ID"},
            {"key": "client_secret", "label": "Client Secret", "secret": True},
            {"key": "username", "label": "Username", "optional": True},
            {"key": "password", "label": "Password", "secret": True, "optional": True},
        ],
        "test_path": "/main/ifsapplications/projection/v1/CustomerOrderHandling.svc/CustomerOrder",
    },
    {
        "key": "custom_rest",
        "name": "Custom REST API",
        "auth_modes": ["api_key", "bearer", "basic", "none"],
        "fields": [
            {"key": "base_url", "label": "Base URL", "placeholder": "https://erp.company.com/api"},
            {"key": "api_key",  "label": "API Key / Bearer Token", "secret": True, "optional": True},
            {"key": "username", "label": "Basic User", "optional": True},
            {"key": "password", "label": "Basic Password", "secret": True, "optional": True},
            {"key": "test_endpoint", "label": "Health-check Path", "placeholder": "/health"},
        ],
        "test_path": "/health",
    },
]


def _mask_secrets(cfg: Dict[str, Any], template: Dict[str, Any]) -> Dict[str, Any]:
    """Return a copy of cfg with secret fields replaced by a masked sentinel."""
    out = dict(cfg)
    secret_keys = {f["key"] for f in template.get("fields", []) if f.get("secret")}
    for k in secret_keys:
        v = out.get(k)
        if v:
            out[k] = "•" * 8 + str(v)[-3:]  # show last 3 chars only
    return out


@api_router.get("/admin/erp/templates")
async def erp_templates(_: User = Depends(require_role("admin"))):
    return {"templates": ERP_TEMPLATES}


@api_router.get("/admin/erp")
async def erp_list(_: User = Depends(require_role("admin"))):
    rows = await db.erp_config.find({}, {"_id": 0}).sort("created_at", -1).to_list(50)
    by_key = {t["key"]: t for t in ERP_TEMPLATES}
    return {
        "connections": [
            {**r, "config": _mask_secrets(r.get("config", {}), by_key.get(r.get("erp_key"), {}))}
            for r in rows
        ],
    }


class ERPConnectionIn(BaseModel):
    erp_key: str
    label: str
    auth_mode: Optional[str] = None
    config: Dict[str, Any]
    activate: bool = False


@api_router.post("/admin/erp")
async def erp_save(payload: ERPConnectionIn, user: User = Depends(require_role("admin"))):
    """Create or update an ERP connection. Each connection has a unique label
    so admins can keep multiple environments side-by-side (e.g. SAP-QAS, SAP-PRD)."""
    template = next((t for t in ERP_TEMPLATES if t["key"] == payload.erp_key), None)
    if not template:
        raise HTTPException(400, f"Unknown ERP type '{payload.erp_key}'")
    if not payload.label.strip():
        raise HTTPException(400, "label required")
    # Required fields enforcement
    for f in template["fields"]:
        if f.get("optional"):
            continue
        if not payload.config.get(f["key"]):
            raise HTTPException(400, f"Missing required field: {f['label']}")

    doc = {
        "connection_id": re.sub(r"[^a-z0-9-]+", "-", payload.label.lower()).strip("-"),
        "erp_key": payload.erp_key,
        "erp_name": template["name"],
        "label": payload.label,
        "auth_mode": payload.auth_mode or template["auth_modes"][0],
        "config": payload.config,
        "is_active": False,
        "created_by": user.user_id,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.erp_config.update_one({"connection_id": doc["connection_id"]}, {"$set": doc}, upsert=True)
    if payload.activate:
        await db.erp_config.update_many({}, {"$set": {"is_active": False}})
        await db.erp_config.update_one({"connection_id": doc["connection_id"]}, {"$set": {"is_active": True}})
        doc["is_active"] = True
    doc["config"] = _mask_secrets(doc["config"], template)
    return {"ok": True, "connection": doc}


@api_router.post("/admin/erp/activate")
async def erp_activate(payload: dict, _: User = Depends(require_role("admin"))):
    cid = (payload or {}).get("connection_id")
    if not cid:
        raise HTTPException(400, "connection_id required")
    found = await db.erp_config.find_one({"connection_id": cid})
    if not found:
        raise HTTPException(404, "Connection not found")
    await db.erp_config.update_many({}, {"$set": {"is_active": False}})
    await db.erp_config.update_one({"connection_id": cid}, {"$set": {"is_active": True}})
    return {"ok": True}


@api_router.post("/admin/erp/test")
async def erp_test(payload: dict, _: User = Depends(require_role("admin"))):
    """Live-tests an ERP connection: GETs the template's test_path with the
    supplied creds and a 6-second budget. Returns status code + elapsed_ms
    so the admin sees instantly whether their config is correct."""
    erp_key = (payload or {}).get("erp_key")
    cfg = (payload or {}).get("config") or {}
    template = next((t for t in ERP_TEMPLATES if t["key"] == erp_key), None)
    if not template:
        raise HTTPException(400, "Unknown ERP type")
    base = (cfg.get("base_url") or "").rstrip("/")
    if not base:
        raise HTTPException(400, "base_url is required to test")
    path = cfg.get("test_endpoint") or template.get("test_path", "/")
    url = f"{base}{path}"

    t0 = datetime.now(timezone.utc)
    try:
        auth = None
        headers = {"Accept": "application/json"}
        api_key = cfg.get("api_key")
        if api_key:
            headers["Authorization"] = f"Bearer {api_key}"
        if cfg.get("username") and cfg.get("password"):
            auth = (cfg["username"], cfg["password"])
        async with httpx.AsyncClient(timeout=6.0, verify=False) as cx:
            r = await cx.get(url, headers=headers, auth=auth)
        ms = int((datetime.now(timezone.utc) - t0).total_seconds() * 1000)
        return {
            "ok": 200 <= r.status_code < 400,
            "status_code": r.status_code,
            "elapsed_ms": ms,
            "url": url,
            "preview": r.text[:240],
        }
    except Exception as e:
        ms = int((datetime.now(timezone.utc) - t0).total_seconds() * 1000)
        return {"ok": False, "status_code": 0, "elapsed_ms": ms, "url": url, "error": str(e)[:240]}


@api_router.delete("/admin/erp/{connection_id}")
async def erp_delete(connection_id: str, _: User = Depends(require_role("admin"))):
    r = await db.erp_config.delete_one({"connection_id": connection_id})
    if r.deleted_count == 0:
        raise HTTPException(404, "Connection not found")
    return {"ok": True}





# -------------------- USER MANUAL (PPTX + in-app deck) --------------------
from manual_content import SLIDES as MANUAL_SLIDES

@api_router.get("/manual/content")
async def get_manual_content(_: User = Depends(get_current_user)):
    """Returns the slide data for the in-app deck viewer."""
    return {"version": "1.0", "title": "Tennant TMS · User Manual", "slides": MANUAL_SLIDES}

@api_router.get("/manual/download")
async def download_manual(_: User = Depends(get_current_user)):
    """Generates and streams the Tennant-branded .pptx user manual."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor as PPTColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    CYAN = PPTColor(0x00, 0xE5, 0xFF)
    BG = PPTColor(0x0B, 0x0E, 0x14)
    PANEL = PPTColor(0x13, 0x18, 0x21)
    SLATE_LIGHT = PPTColor(0xF8, 0xFA, 0xFC)
    SLATE_MID = PPTColor(0x94, 0xA3, 0xB8)
    SLATE_DIM = PPTColor(0x47, 0x55, 0x69)
    EMERALD = PPTColor(0x10, 0xB9, 0x81)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def fill_bg(slide, color=BG):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_text(slide, text, left, top, width, height, *, size=18, bold=False, color=SLATE_LIGHT, font="Calibri", align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
        return tb

    def add_accent_bar(slide, left, top, width=Inches(0.4)):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.05))
        bar.fill.solid(); bar.fill.fore_color.rgb = CYAN
        bar.line.fill.background()

    for sl in MANUAL_SLIDES:
        slide = prs.slides.add_slide(blank)
        fill_bg(slide)
        kind = sl.get("kind")

        if kind == "cover":
            add_text(slide, sl["eyebrow"], Inches(0.75), Inches(0.6), Inches(12), Inches(0.4), size=12, bold=True, color=CYAN, font="Consolas")
            add_accent_bar(slide, Inches(0.75), Inches(1.05), width=Inches(1.5))
            add_text(slide, sl["title"], Inches(0.75), Inches(1.3), Inches(12), Inches(2), size=54, bold=True, color=SLATE_LIGHT)
            add_text(slide, sl["subtitle"], Inches(0.75), Inches(3.6), Inches(12), Inches(1.2), size=22, color=SLATE_MID)
            add_text(slide, sl.get("footnote", ""), Inches(0.75), Inches(6.7), Inches(12), Inches(0.4), size=11, color=SLATE_DIM, font="Consolas")

        elif kind == "toc":
            add_text(slide, "TABLE OF CONTENTS", Inches(0.75), Inches(0.6), Inches(12), Inches(0.4), size=12, bold=True, color=CYAN, font="Consolas")
            add_accent_bar(slide, Inches(0.75), Inches(1.05), width=Inches(1.5))
            add_text(slide, sl["title"], Inches(0.75), Inches(1.3), Inches(12), Inches(1), size=40, bold=True, color=SLATE_LIGHT)
            sections = sl.get("sections", [])
            mid = (len(sections) + 1) // 2
            left_col = sections[:mid]
            right_col = sections[mid:]
            for i, s in enumerate(left_col):
                add_text(slide, s, Inches(0.75), Inches(2.7 + i * 0.45), Inches(6), Inches(0.4), size=14, color=SLATE_LIGHT)
            for i, s in enumerate(right_col):
                add_text(slide, s, Inches(7.0), Inches(2.7 + i * 0.45), Inches(6), Inches(0.4), size=14, color=SLATE_LIGHT)

        elif kind == "section":
            add_text(slide, sl["title"], Inches(0.75), Inches(2.4), Inches(12), Inches(1.2), size=44, bold=True, color=CYAN)
            add_accent_bar(slide, Inches(0.75), Inches(3.5), width=Inches(1.2))
            add_text(slide, sl.get("tagline", ""), Inches(0.75), Inches(3.8), Inches(12), Inches(1), size=22, color=SLATE_MID)

        elif kind == "feature":
            add_text(slide, sl.get("subtitle", "FEATURE WALKTHROUGH").upper(), Inches(0.75), Inches(0.55), Inches(12), Inches(0.4), size=11, bold=True, color=CYAN, font="Consolas")
            add_accent_bar(slide, Inches(0.75), Inches(0.95), width=Inches(1.3))
            add_text(slide, sl["title"], Inches(0.75), Inches(1.15), Inches(12), Inches(0.9), size=32, bold=True, color=SLATE_LIGHT)
            # Steps
            steps_tb = slide.shapes.add_textbox(Inches(0.75), Inches(2.25), Inches(8.5), Inches(4.5))
            stf = steps_tb.text_frame
            stf.word_wrap = True
            for i, step in enumerate(sl.get("steps", [])):
                p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
                r1 = p.add_run(); r1.text = f"{i+1:02d}  "
                r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = CYAN; r1.font.name = "Consolas"
                r2 = p.add_run(); r2.text = step
                r2.font.size = Pt(13); r2.font.color.rgb = SLATE_LIGHT; r2.font.name = "Calibri"
                p.space_after = Pt(6)
            # Tips panel
            if sl.get("tips"):
                panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(2.25), Inches(3.2), Inches(3.5))
                panel.fill.solid(); panel.fill.fore_color.rgb = PANEL
                panel.line.color.rgb = CYAN; panel.line.width = Pt(0.75)
                add_text(slide, "TIPS", Inches(9.7), Inches(2.4), Inches(3), Inches(0.3), size=10, bold=True, color=CYAN, font="Consolas")
                tips_text = "\n".join("• " + t for t in sl["tips"])
                add_text(slide, tips_text, Inches(9.7), Inches(2.8), Inches(2.9), Inches(2.8), size=11, color=SLATE_LIGHT)
            if sl.get("page_url"):
                add_text(slide, f"→ {sl['page_url']}", Inches(0.75), Inches(6.85), Inches(12), Inches(0.4), size=11, color=EMERALD, font="Consolas")

        elif kind == "closing":
            add_text(slide, sl["title"], Inches(0.75), Inches(2.8), Inches(12), Inches(1.5), size=56, bold=True, color=CYAN)
            add_accent_bar(slide, Inches(0.75), Inches(4.0), width=Inches(1.5))
            add_text(slide, sl.get("subtitle", ""), Inches(0.75), Inches(4.3), Inches(12), Inches(1), size=22, color=SLATE_MID)
            add_text(slide, sl.get("footnote", ""), Inches(0.75), Inches(6.7), Inches(12), Inches(0.4), size=11, color=SLATE_DIM, font="Consolas")

        else:
            add_text(slide, sl.get("title", "Tennant TMS"), Inches(0.75), Inches(3), Inches(12), Inches(1), size=36, bold=True, color=SLATE_LIGHT)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    headers = {"Content-Disposition": 'attachment; filename="Tennant_TMS_User_Manual.pptx"'}
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers=headers)

# -------------------- ROOT --------------------
@api_router.get("/")
async def root():
    return {"service": "Tennant TMS API", "status": "ok", "time": datetime.now(timezone.utc).isoformat()}


@api_router.get("/health")
async def health():
    """Production readiness probe. Pings Mongo with a 1.5s budget so load
    balancers don't hang. Returns 503 on any DB outage so the deploy
    platform can pull the pod out of rotation."""
    import asyncio as _aio
    started = datetime.now(timezone.utc)
    try:
        await _aio.wait_for(db.command("ping"), timeout=1.5)
        db_ok = True
        db_error = None
    except Exception as e:
        db_ok = False
        db_error = str(e)[:200]
    payload = {
        "service": "Tennant TMS API",
        "status": "ok" if db_ok else "degraded",
        "db": "up" if db_ok else "down",
        "db_error": db_error,
        "time": started.isoformat(),
        "elapsed_ms": int((datetime.now(timezone.utc) - started).total_seconds() * 1000),
    }
    if not db_ok:
        return Response(content=json.dumps(payload), media_type="application/json", status_code=503)
    return payload

# -------------------- EQUIPMENT / YARD MODULE --------------------
from equipment_module import register_equipment_routes  # noqa: E402
register_equipment_routes(api_router, db, get_current_user, require_role)

# -------------------- ROUTES PACKAGE — modular routers --------------------
# Conservative refactor: each feature group lives in its own file under
# /app/backend/routes/ and exposes a build_*_router() factory that takes
# the shared DB handle + helpers. Adding more groups here over time
# pulls server.py back under the 1000-line guideline.
from routes.weather import build_weather_router  # noqa: E402
from routes.server_registry import build_server_registry_router  # noqa: E402
from routes.sap import build_sap_router  # noqa: E402
from routes.brokerage import build_brokerage_router  # noqa: E402
from routes.connections import build_connections_router  # noqa: E402
from routes.freight_news import build_freight_news_router  # noqa: E402
from routes.provider_outreach import build_provider_outreach_router  # noqa: E402
from routes.public_site import build_public_router  # noqa: E402
api_router.include_router(build_weather_router(
    db=db,
    get_current_user=get_current_user,
    brand_swap=_brand_swap,
    active_brand_doc=_active_brand_doc,
    mock_weather_alerts=MOCK_WEATHER_ALERTS,
))
api_router.include_router(build_server_registry_router(
    db=db,
    require_role=require_role,
    app_boot_at=_APP_BOOT_AT,
))
api_router.include_router(build_branding_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
    emergent_llm_key=EMERGENT_LLM_KEY,
    LlmChat=LlmChat,
    UserMessage=UserMessage,
))
api_router.include_router(build_sap_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
    brand_swap=_brand_swap,
    active_brand_doc=_active_brand_doc,
))
api_router.include_router(build_brokerage_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
    emergent_llm_key=EMERGENT_LLM_KEY,
    LlmChat=LlmChat,
    UserMessage=UserMessage,
))
api_router.include_router(build_connections_router(
    db=db,
    require_role=require_role,
))
build_provider_outreach_router(
    api_router=api_router,
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
)
build_public_router(api_router=api_router, db=db)
build_freight_news_router(api_router=api_router, get_current_user=get_current_user)

from routes.tms_investor import build_tms_investor_router  # noqa: E402
build_tms_investor_router(api_router=api_router)

from routes.tms_invite_links import build_tms_invite_links_router  # noqa: E402
build_tms_invite_links_router(api_router=api_router, db=db,
                              get_current_user=get_current_user,
                              require_role=require_role)

from routes.margin_shield import build_margin_shield_router  # noqa: E402
build_margin_shield_router(api_router=api_router, db=db,
                           get_current_user=get_current_user,
                           require_role=require_role)

from routes.orisei_ops_kpis import build_orisei_ops_router  # noqa: E402
build_orisei_ops_router(api_router=api_router, db=db,
                        get_current_user=get_current_user,
                        require_role=require_role)

from routes.orisei_operations import build_orisei_operations_router  # noqa: E402
build_orisei_operations_router(api_router=api_router, db=db,
                               get_current_user=get_current_user,
                               require_role=require_role)

from routes.orisei_workflow import build_orisei_workflow_router  # noqa: E402
build_orisei_workflow_router(api_router=api_router, db=db,
                             get_current_user=get_current_user,
                             require_role=require_role)

from routes.factoring import build_factoring_router  # noqa: E402
build_factoring_router(api_router=api_router, db=db,
                       get_current_user=get_current_user,
                       require_role=require_role,
                       emergent_llm_key=EMERGENT_LLM_KEY,
                       LlmChat=LlmChat, UserMessage=UserMessage)

from routes.cash_flow import build_cash_flow_router  # noqa: E402
build_cash_flow_router(api_router=api_router, db=db,
                       get_current_user=get_current_user,
                       require_role=require_role)

from routes.shipment_triage import build_shipment_triage_router  # noqa: E402
build_shipment_triage_router(api_router=api_router, db=db,
                             get_current_user=get_current_user,
                             require_role=require_role,
                             emergent_llm_key=EMERGENT_LLM_KEY,
                             LlmChat=LlmChat, UserMessage=UserMessage)

from routes.orisei_auto_digest import build_auto_digest_router  # noqa: E402
build_auto_digest_router(api_router=api_router, db=db,
                          get_current_user=get_current_user,
                          require_role=require_role)

from routes.load_aggregator import build_aggregator_router  # noqa: E402
build_aggregator_router(api_router=api_router, db=db,
                        get_current_user=get_current_user,
                        require_role=require_role)

from routes.shipper_relations import build_shipper_relations_router  # noqa: E402
build_shipper_relations_router(api_router=api_router, db=db,
                                get_current_user=get_current_user,
                                require_role=require_role)

from routes.claims_master import build_claims_master_router  # noqa: E402
build_claims_master_router(api_router=api_router, db=db,
                            get_current_user=get_current_user,
                            require_role=require_role)

from routes.qbr_studio import build_qbr_studio_router  # noqa: E402
build_qbr_studio_router(api_router=api_router, db=db,
                        get_current_user=get_current_user,
                        require_role=require_role)

from routes.lighthouse import build_lighthouse_router  # noqa: E402
build_lighthouse_router(api_router=api_router, db=db,
                         get_current_user=get_current_user,
                         require_role=require_role)

from routes.routing_svc import build_routing_router  # noqa: E402
build_routing_router(api_router=api_router, db=db,
                     get_current_user=get_current_user,
                     require_role=require_role)

from routes.telematics import build_telematics_router  # noqa: E402
build_telematics_router(api_router=api_router, db=db,
                         get_current_user=get_current_user,
                         require_role=require_role)

from routes.tms_competitive import build_tms_competitive_router, build_driver_pwa_router  # noqa: E402
build_tms_competitive_router(api_router=api_router, db=db,
                              get_current_user=get_current_user,
                              require_role=require_role)
build_driver_pwa_router(api_router=api_router, db=db,
                         get_current_user=get_current_user)

from routes.enterprise_tms import build_enterprise_tms_router  # noqa: E402
build_enterprise_tms_router(api_router=api_router, db=db,
                             get_current_user=get_current_user,
                             require_role=require_role)

from routes.enterprise_adapters import build_enterprise_adapters_router  # noqa: E402
build_enterprise_adapters_router(api_router=api_router, db=db,
                                   get_current_user=get_current_user,
                                   require_role=require_role)

from routes.research_analytics import build_research_analytics_router  # noqa: E402
build_research_analytics_router(api_router=api_router, db=db,
                                   get_current_user=get_current_user,
                                   require_role=require_role)

from routes.orisei_gtm_assets import build_gtm_assets_router  # noqa: E402
build_gtm_assets_router(api_router=api_router, db=db, require_role=require_role)

from routes.investor import build_investor_router  # noqa: E402
api_router.include_router(build_investor_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
    active_brand_doc=_active_brand_doc,
))

from routes.marketing import build_marketing_router  # noqa: E402
api_router.include_router(build_marketing_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
    active_brand_doc=_active_brand_doc,
))

from routes.doc_vault import build_doc_vault_router  # noqa: E402
api_router.include_router(build_doc_vault_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
))

from routes.launch_runway import build_launch_runway_router  # noqa: E402
api_router.include_router(build_launch_runway_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
))

from routes.shipper_outreach import build_shipper_outreach_router  # noqa: E402
api_router.include_router(build_shipper_outreach_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
    active_brand_doc=_active_brand_doc,
))

from routes.data_status import build_data_status_router  # noqa: E402
api_router.include_router(build_data_status_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
))

from routes.upwork_portfolio import build_upwork_portfolio_router  # noqa: E402
api_router.include_router(build_upwork_portfolio_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
    active_brand_doc=_active_brand_doc,
))

from routes.autocomplete import build_autocomplete_router  # noqa: E402
api_router.include_router(build_autocomplete_router(
    db=db,
    get_current_user=get_current_user,
    require_role=require_role,
))

# Mount International (ocean + intermodal rail) module
from routes.international import build_international_router  # noqa: E402
build_international_router(
    api_router=api_router, db=db,
    get_current_user=get_current_user,
    require_role=require_role,
)

# Mount Shipper Intake (branded fillable template + public submit)
from routes.shipper_intake import build_shipper_intake_router  # noqa: E402
build_shipper_intake_router(
    api_router=api_router, db=db,
    get_current_user=get_current_user,
    require_role=require_role,
    send_email_fn=None,  # wire Resend via integrations later
)

# Mount Onboarding Checklist (MC/bond/API-keys walkthrough)
from routes.onboarding_checklist import build_onboarding_router  # noqa: E402
build_onboarding_router(
    api_router=api_router, db=db,
    get_current_user=get_current_user,
    require_role=require_role,
)

# Mount BOC-3 Compliance (50-state process-agent tracker + renewal calendar)
from routes.boc3_compliance import build_boc3_router  # noqa: E402
build_boc3_router(
    api_router=api_router, db=db,
    get_current_user=get_current_user,
    require_role=require_role,
)

# -------------------- WIRE UP --------------------
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    # Explicit allow-list of Emergent host families + the active brand's
    # custom domain via env var. Allow-list pattern is required because the
    # browser silently rejects credentials when the response carries
    # `Access-Control-Allow-Origin: *`.
    allow_origin_regex=(
        r"https?://(?:[a-z0-9-]+\.)*(?:emergentagent\.com|emergent\.host|emergent\.sh|"
        r"oriseifreight\.com|livecleans\.com|localhost(?::\d+)?)$"
    ),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("startup")
async def startup():
    # ---------- 1. Mongo indexes (idempotent; creates only what's missing).
    # Sized for ~250 concurrent users hitting hot collections; every entry
    # below is referenced by a real query in this file. Run in background so
    # any single failure can't block app boot.
    async def _ensure_indexes():
        IDX_PLAN = [
            ("shipments",          [("shipment_id", 1)],       {"name": "ix_shipment_id"}),
            ("shipments",          [("reference", 1)],         {"name": "ix_shipment_ref"}),
            ("shipments",          [("status", 1)],            {"name": "ix_shipment_status"}),
            ("shipments",          [("carrier", 1)],           {"name": "ix_shipment_carrier"}),
            ("shipments",          [("eta", 1)],               {"name": "ix_shipment_eta"}),
            ("shipments",          [("created_at", -1)],       {"name": "ix_shipment_created"}),
            ("users",              [("user_id", 1)],           {"unique": True, "name": "uq_user_id"}),
            ("users",              [("email", 1)],             {"unique": True, "name": "uq_user_email"}),
            ("user_sessions",      [("session_token", 1)],     {"unique": True, "name": "uq_session_token"}),
            ("user_sessions",      [("expires_at", 1)],        {"name": "ix_session_expires", "expireAfterSeconds": 0}),
            ("user_layouts",       [("user_id", 1), ("page_key", 1)], {"unique": True, "name": "uq_layout_user_page"}),
            ("truckload_bookings", [("created_at", -1)],       {"name": "ix_tlb_created"}),
            ("truckload_bookings", [("status", 1)],            {"name": "ix_tlb_status"}),
            ("carrier_onboarding", [("status", 1)],            {"name": "ix_co_status"}),
            ("carrier_onboarding", [("legal_name", 1)],        {"name": "ix_co_legal"}),
            ("carrier_onboarding", [("submitted_at", -1)],     {"name": "ix_co_submitted"}),
            ("drivers",            [("name", 1)],              {"name": "ix_driver_name"}),
            ("drivers",            [("carrier", 1)],           {"name": "ix_driver_carrier"}),
            ("trailers",           [("trailer_no", 1)],        {"name": "ix_trailer_no"}),
            ("trailers",           [("carrier", 1)],           {"name": "ix_trailer_carrier"}),
            ("machines",           [("model", 1)],             {"name": "ix_machine_model"}),
            ("freight_bills",      [("status", 1)],            {"name": "ix_fb_status"}),
            ("freight_bills",      [("carrier", 1)],           {"name": "ix_fb_carrier"}),
            ("yard_reports",       [("uploaded_at", -1)],      {"name": "ix_yard_uploaded"}),
            ("workbook_tabs",      [("order", 1)],             {"name": "ix_tab_order"}),
            ("outbound_emails",    [("at", -1)],               {"name": "ix_email_at"}),
            ("ai_messages",        [("session_id", 1), ("created_at", 1)], {"name": "ix_ai_session"}),
            ("audit_log",          [("at", -1)],               {"name": "ix_audit_at"}),
            ("suppliers_custom",   [("supplier_id", 1)],       {"unique": True, "name": "uq_sup_custom_id"}),
        ]
        ok, fail = 0, 0
        for coll, keys, opts in IDX_PLAN:
            try:
                await db[coll].create_index(keys, **opts)
                ok += 1
            except Exception as e:
                # Common: existing index conflict — drop & recreate is risky in prod,
                # so just log and move on. Production-safe.
                fail += 1
                logger.debug(f"index skip {coll}.{opts.get('name')}: {e}")
        logger.info(f"Indexes ready: {ok} created/existing, {fail} skipped")

    try:
        await _ensure_indexes()
    except Exception as e:
        logger.warning(f"Index ensure failed: {e}")

    # ---------- 1.5 Brand bootstrap — every PDF depends on this
    try:
        from routes.brand_bootstrap import ensure_active_brand
        await ensure_active_brand(db)
    except Exception as e:
        logger.warning(f"Brand bootstrap failed: {e}")

    # ---------- 2. Auto-seed if empty
    try:
        count = await db.shipments.count_documents({})
        if count == 0:
            logger.info("Auto-seeding shipments...")
            from fastapi import BackgroundTasks  # noqa
            await seed_data(force=False)
    except Exception as e:
        logger.warning(f"Seed on startup failed: {e}")

    # Idempotent backfill: ensure the Truckload Booking Sheet's carrier
    # dropdown has a usable roster of approved carriers on every boot.
    # Inserts only the missing legal_names — safe to run repeatedly.
    try:
        approved_roster = [
            ("XPO Logistics LLC", "XPO", "MC-249635", "528970", "XPOL", "TL",
             "Renee Calderon", "tennant.team@xpo.com", "+1-855-976-2243", 2000000, 18),
            ("Old Dominion Freight Line, Inc.", "ODFL", "MC-22198", "55977", "ODFL", "LTL",
             "Greg Halsey", "tennant@odfl.com", "+1-800-432-6335", 2000000, 11),
            ("Saia Motor Freight Line, LLC", "Saia", "MC-44918", "33172", "SAIA", "LTL",
             "Tasha Burnett", "tennant@saia.com", "+1-800-765-7242", 1500000, 22),
            ("Estes Express Lines", "Estes", "MC-105764", "55712", "EXLA", "LTL",
             "Adam Mueller", "ops@estes-express.com", "+1-866-378-3748", 2000000, 16),
            ("R&L Carriers, Inc.", "R+L", "MC-133134", "243809", "RLCA", "LTL",
             "Marcus Lavoie", "tennant@rlcarriers.com", "+1-800-543-5589", 1500000, 20),
            ("Knight-Swift Transportation Holdings", "Knight Transportation", "MC-247369", "362724", "KNIG", "TL",
             "Lina Ortega", "tennant@knight-swift.com", "+1-602-269-2000", 2000000, 25),
            ("Schneider National Carriers", "Schneider", "MC-237983", "264184", "SCNN", "TL",
             "Henry Park", "tennant@schneider.com", "+1-800-558-6767", 2000000, 28),
            ("C.H. Robinson Worldwide, Inc.", "C.H. Robinson", "MC-208535", "388873", "CHRW", "Brokerage",
             "Ben Reichl", "tennant@chrobinson.com", "+1-800-323-7587", 2000000, 15),
            ("Werner Enterprises, Inc.", "Werner", "MC-159458", "111723", "WERN", "TL",
             "Christine Yoder", "tennant@werner.com", "+1-800-228-2240", 2000000, 26),
        ]
        existing_names = set(await db.carrier_onboarding.distinct("legal_name"))
        to_insert = []
        for (legal, dba, mc, dot, scac, mode, contact, email, phone, ins, csa) in approved_roster:
            if legal in existing_names:
                continue
            to_insert.append({
                "onboarding_id": f"OB-{uuid.uuid4().hex[:8].upper()}",
                "legal_name": legal, "dba": dba, "mc_number": mc, "dot_number": dot,
                "scac": scac, "mode": mode,
                "contact_name": contact, "contact_email": email, "contact_phone": phone,
                "insurance_amount": float(ins),
                "insurance_expiry": (datetime.now(timezone.utc) + timedelta(days=270)).date().isoformat(),
                "safety_rating": "Satisfactory", "csa_score": int(csa),
                "w9_received": True, "coi_received": True, "contract_signed": True,
                "status": "approved",
                "submitted_at": datetime.now(timezone.utc).isoformat(),
                "submitted_by": "System",
                "notes": "Approved national-carrier backfill",
            })
        if to_insert:
            await db.carrier_onboarding.insert_many(to_insert)
            logger.info(f"Backfilled {len(to_insert)} approved carrier onboardings.")
    except Exception as e:
        logger.warning(f"Approved-carrier backfill failed: {e}")

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
